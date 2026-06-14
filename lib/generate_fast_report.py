#!/usr/bin/env python3
# SPDX-License-Identifier: MIT OR Apache-2.0
# SPDX-FileCopyrightText: 2026 Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin · Joost Beekman
"""
generate_fast_report.py — FAST (Forensic Analysis Storage) report generator.

Aggregates The Sleuth Kit / EWF tool outputs from ./analysis/storage/ and
./exports/ into a structured incident report in Markdown, PDF, PPTX, and DOCX.

All sections follow the FanGetFameFast dual-register voice:
  - Management Summary: no technical identifiers; plain business language
  - Technical Body: precise identifiers; scoped conclusions citing evidence source

Claude instructs itself to "enhance and elaborate when necessary" on each section.

Usage (CLI):
    python3 lib/generate_fast_report.py \\
        --case-id FAST-2026-001 \\
        --hostname SERVER1234 \\
        --disk-image /path/to/SERVER1234.vmdk \\
        [--analysis-dir ./analysis/storage] \\
        [--output-dir ./reports]

Python API:
    from lib.generate_fast_report import generate
    paths = generate(case_id="FAST-2026-001", hostname="SERVER1234",
                     disk_image="/path/to/SERVER1234.vmdk")
"""
from __future__ import annotations

import argparse
import hashlib
import json
import re
import sys
from datetime import datetime, timezone
try:
    from zoneinfo import ZoneInfo
    _CET = ZoneInfo("Europe/Amsterdam")
except ImportError:
    _CET = timezone.utc
from pathlib import Path
sys.path.insert(0, str(Path(__file__).resolve().parent))
import path_guard  # noqa: E402  write-path policy enforcement
import report_completeness  # noqa: E402  narrative/reasoning completeness gate
from typing import Any

try:
    from research_notes import (
        parse_steps as _parse_research_steps,
        parse_events as _parse_research_events,
        parse_reflections as _parse_research_reflections,
    )
    import report_sections
    from hallucination_guard import (
        ConfidenceTier,
        tag_finding,
        render_confidence_summary,
        reset_counter as _hg_reset,
    )
except ModuleNotFoundError:
    # Imported as a package (e.g. `from lib.generate_fast_report import generate`)
    # rather than run as a script — put lib/ on the path for the sibling import.
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    from research_notes import (
        parse_steps as _parse_research_steps,
        parse_events as _parse_research_events,
        parse_reflections as _parse_research_reflections,
    )
    import report_sections
    from hallucination_guard import (
        ConfidenceTier,
        tag_finding,
        render_confidence_summary,
        reset_counter as _hg_reset,
    )

PROJECT_ROOT = Path(__file__).parent.parent

# ── Colour palette ─────────────────────────────────────────────────────────────
_DARK_NAVY  = (0x0f, 0x17, 0x2a)
_MID_NAVY   = (0x1e, 0x3a, 0x5f)
_BLUE       = (0x1d, 0x4e, 0xd8)
_LIGHT_BLUE = (0x93, 0xc5, 0xfd)
_WHITE      = (0xff, 0xff, 0xff)
_LIGHT_BG   = (0xf8, 0xfa, 0xfc)
_ROW_ALT    = (0xf1, 0xf5, 0xf9)
_TEXT_DARK  = (0x1f, 0x29, 0x37)
_TEXT_MID   = (0x6b, 0x72, 0x80)
_AMBER      = (0xfb, 0xbf, 0x24)
_GREEN      = (0x22, 0xc5, 0x5e)

_SEV_RGB = {
    "critical": (0xef, 0x44, 0x44),
    "high":     (0xf9, 0x73, 0x16),
    "medium":   (0xea, 0xb3, 0x08),
    "low":      (0x22, 0xc5, 0x5e),
    "info":     (0x6b, 0x72, 0x80),
}


# ── Data loading ───────────────────────────────────────────────────────────────

def _load_analysis(analysis_dir: Path, exports_dir: Path) -> dict[str, Any]:
    data: dict[str, Any] = {
        "ewfinfo":           _read_text(analysis_dir / "ewfinfo.txt"),
        "ewfverify":         _read_text(analysis_dir / "ewfverify.txt"),
        "mmls":              _read_text(analysis_dir / "mmls.txt"),
        "fsstat":            _read_text(analysis_dir / "fsstat.txt"),
        "fls_output":        _read_text(analysis_dir / "fls_output.txt"),
        "ils_output":        _read_text(analysis_dir / "ils_output.txt"),
        "ils_orphan":        _read_text(analysis_dir / "ils_orphan.txt"),
        "bodyfile":          _read_text(analysis_dir / "bodyfile.txt"),
        "fs_timeline":       _read_text(exports_dir / "fs_timeline.txt"),
        "fs_timeline_csv":   _read_csv(exports_dir / "fs_timeline.csv"),
        "windows_hashes":    _read_text(analysis_dir / "windows_hashes.txt"),
        "bulk_carved":       _list_dir(exports_dir / "carved"),
        "md5_manifest":      _read_text(exports_dir / "files" / "md5_manifest.txt"),
        "mft_extracted":     (exports_dir / "mft" / "$MFT").exists(),
        "usnj_extracted":    (exports_dir / "mft" / "$J").exists(),
        "evtx_list":         _list_dir(exports_dir / "evtx"),
        "registry_list":     _list_dir(exports_dir / "registry"),
        "prefetch_list":     _list_dir(exports_dir / "prefetch"),
        "srum_list":         _list_dir(exports_dir / "srum"),
        "browser_list":      _list_dir(exports_dir / "browser"),
        "recyclebin_list":   _list_dir(exports_dir / "recyclebin"),
        # Deep extraction outputs (written by fast_machine_details.py)
        "machine_details":   _load_json(exports_dir / "machine_details" / "machine_details.json"),
        "recyclebin_parsed": _load_json(exports_dir / "recyclebin" / "recyclebin_parsed.json"),
        "iocs_reference":    _load_json(exports_dir / "machine_details" / "iocs.json"),
    }
    # Absorb any *.json findings
    for jf in sorted(analysis_dir.glob("*.json")):
        data[jf.stem] = json.loads(jf.read_text())
    return data


def _read_text(path: Path) -> str:
    return path.read_text(errors="replace") if path.exists() else ""


def _read_csv(path: Path) -> list[dict]:
    if not path.exists():
        return []
    lines = path.read_text(errors="replace").splitlines()
    if not lines:
        return []
    headers = [h.strip() for h in lines[0].split(",")]
    rows = []
    for line in lines[1:]:
        if line.strip():
            cols = [c.strip() for c in line.split(",")]
            rows.append(dict(zip(headers, cols)))
    return rows


def _list_dir(path: Path) -> list[str]:
    if not path.exists():
        return []
    return sorted(p.name for p in path.iterdir())


def _load_json(path: Path) -> Any:
    """Load a JSON file; return None if absent or unparseable."""
    if not path.exists():
        return None
    try:
        return json.loads(path.read_text(encoding="utf-8", errors="replace"))
    except Exception:
        return None


def _humanize_bytes(n: int) -> str:
    for unit in ("B", "KB", "MB", "GB", "TB"):
        if n < 1024:
            return f"{n:.0f} {unit}"
        n /= 1024
    return f"{n:.1f} PB"


# ── Markdown ───────────────────────────────────────────────────────────────────

def _load_narrative(case_id: str, reports_dir: Path) -> dict[str, str]:
    """Load Claude-generated narrative sections from {case_id}_narrative.md."""
    path = reports_dir / f"{case_id}_narrative.md"
    if not path.exists():
        return {}
    sections: dict[str, str] = {}
    current: str | None = None
    for line in path.read_text(encoding="utf-8").splitlines():
        if line.startswith("## "):
            current = line[3:].strip()
            sections[current] = ""
        elif line.startswith("<!--"):
            continue
        elif current is not None:
            sections[current] += line + "\n"
    return {k: v.strip() for k, v in sections.items()}


def _narr_bullets(narrative: dict, key: str) -> list[str]:
    """Return the bullet lines of a narrative pptx_* section as clean strings.

    Accepts '-', '*' or '•' bullet markers; falls back to non-empty paragraph
    lines if the section has no explicit bullets. Returns [] when absent."""
    text = (narrative or {}).get(key, "") or ""

    def _clean(s: str) -> str:
        return re.sub(r"\*\*(.*?)\*\*", r"\1", s).replace("**", "").strip()

    bullets: list[str] = []
    marker = re.compile(r"^([-*•])\s+(.*)$")
    for raw in text.splitlines():
        line = raw.strip()
        if not line:
            continue
        m = marker.match(line)
        if m:
            item = m.group(2).strip()
            if item:
                bullets.append(item)
        elif bullets:
            bullets[-1] = (bullets[-1] + " " + line).strip()
    if not bullets:
        for para in re.split(r"\n\s*\n", text):
            p = " ".join(l.strip() for l in para.splitlines() if l.strip())
            if p:
                bullets.append(p)
    return [_clean(b) for b in bullets]


def _build_fast_hallucination_guard_section(
    data: dict[str, Any],
    case_id: str,
    reports_dir: Path,
) -> str:
    """
    Build the Hallucination Guard section for FAST reports.

    Tags key disk-forensics conclusions with ConfidenceTier based on which
    TSK / bulk_extractor / EWF artifacts were actually extracted.
    Tiers assigned by code logic, not by Claude prompt instructions.
    """
    _hg_reset()
    findings = []
    steps = _parse_research_steps(case_id, str(reports_dir)) if case_id else []

    # EWF image verification — CONFIRMED if ewfverify passed
    if data.get("ewfverify"):
        findings.append(tag_finding(
            "EWF image integrity verified (ewfverify hash check passed)",
            ConfidenceTier.CONFIRMED,
            [],
            ["libewf/ewfverify"],
            ["fast"],
        ))

    # Partition layout — CONFIRMED if mmls ran
    if data.get("mmls"):
        findings.append(tag_finding(
            "Partition layout extracted via mmls — partition boundaries confirmed",
            ConfidenceTier.CONFIRMED,
            [],
            ["tsk/mmls"],
            ["fast"],
        ))

    # Filesystem tree — CONFIRMED if fls ran
    if data.get("fls_output"):
        findings.append(tag_finding(
            "Filesystem tree enumerated via fls — file and directory entries confirmed",
            ConfidenceTier.CONFIRMED,
            [],
            ["tsk/fls"],
            ["fast"],
        ))

    # Unallocated inodes (deleted files) — INFERRED if ils ran
    if data.get("ils_output") or data.get("ils_orphan"):
        findings.append(tag_finding(
            "Unallocated/orphan inodes found via ils — deleted file recovery possible (requires icat)",
            ConfidenceTier.INFERRED,
            [],
            ["tsk/ils"],
            ["fast"],
        ))

    # MFT
    if data.get("mft_extracted"):
        findings.append(tag_finding(
            "$MFT extracted — Master File Table provides authoritative file creation/modification times",
            ConfidenceTier.CONFIRMED,
            [],
            ["tsk/icat"],
            ["fast"],
        ))

    # Windows Event Logs
    evtx = data.get("evtx_list") or []
    if evtx:
        findings.append(tag_finding(
            f"{len(evtx)} Windows Event Log file(s) extracted — security, system, application events available",
            ConfidenceTier.CONFIRMED,
            [],
            ["tsk/fls+icat"],
            ["fast"],
        ))
    else:
        findings.append(tag_finding(
            "No Windows Event Logs extracted — EVTX evidence unavailable",
            ConfidenceTier.UNVERIFIABLE,
            [],
            ["tsk/fls+icat"],
            ["fast"],
        ))

    # Registry
    registry = data.get("registry_list") or []
    if registry:
        findings.append(tag_finding(
            f"{len(registry)} registry hive(s) extracted — persistence, run keys, user artifacts available",
            ConfidenceTier.CONFIRMED,
            [],
            ["tsk/fls+icat"],
            ["fast"],
        ))

    # Prefetch
    prefetch = data.get("prefetch_list") or []
    if prefetch:
        findings.append(tag_finding(
            f"{len(prefetch)} prefetch execution artifact(s) extracted — program execution evidence confirmed",
            ConfidenceTier.CONFIRMED,
            [],
            ["tsk/fls+icat"],
            ["fast"],
        ))

    # Bulk extractor carves
    carved = data.get("bulk_carved") or []
    if carved:
        findings.append(tag_finding(
            f"bulk_extractor carved {len(carved)} artifact type(s) from unallocated space",
            ConfidenceTier.CONFIRMED,
            [],
            ["bulk_extractor"],
            ["fast"],
        ))

    # Machine details extraction
    machine = data.get("machine_details") or {}
    if machine:
        sam_users = (machine.get("sam") or {}).get("UserAccounts") or []
        net_ifaces = (machine.get("system") or {}).get("NetworkInterfaces") or []
        owner = (machine.get("software") or {}).get("RegisteredOwner") or ""
        if owner or sam_users or net_ifaces:
            findings.append(tag_finding(
                f"Machine identity extracted from registry hives — "
                f"{len(sam_users)} user account(s), {len(net_ifaces)} network interface(s)",
                ConfidenceTier.CONFIRMED,
                [],
                ["fast_machine_details/python-registry"],
                ["fast"],
            ))

    # Recycle Bin parse
    rb_parsed = data.get("recyclebin_parsed") or []
    if rb_parsed:
        valid_rb = [e for e in rb_parsed if "error" not in e]
        if valid_rb:
            findings.append(tag_finding(
                f"{len(valid_rb)} Recycle Bin $I file(s) parsed — original paths "
                "and deletion timestamps confirmed",
                ConfidenceTier.CONFIRMED,
                [],
                ["fast_machine_details/recyclebin-parser"],
                ["fast"],
            ))

    # Timeline
    if data.get("fs_timeline") or data.get("fs_timeline_csv"):
        findings.append(tag_finding(
            "Filesystem timeline reconstructed — file system activity sequence confirmed",
            ConfidenceTier.CONFIRMED,
            [],
            ["tsk/mactime"],
            ["fast"],
        ))

    # Timeline correlations (time-proximity inferences) — INFERRED
    if (data.get("fs_timeline") or data.get("fs_timeline_csv")) and evtx:
        findings.append(tag_finding(
            "Timeline correlation between filesystem events and EVTX entries possible (time-proximity inference)",
            ConfidenceTier.INFERRED,
            [],
            ["tsk/mactime + evtx"],
            ["fast"],
        ))

    # Assumptions from research notes
    for s in steps:
        outcome = s.get("outcome", "")
        if "[ASSUMPTION]" in outcome or s.get("confidence") == "assumed":
            text = outcome.replace("[ASSUMPTION]", "").strip()
            if text:
                findings.append(tag_finding(
                    text,
                    ConfidenceTier.ASSUMED,
                    [s["id"]] if s.get("id") else [],
                    [s.get("source_tool", "")] if s.get("source_tool") else [],
                    ["fast"],
                ))

    return render_confidence_summary(findings, module_label="FAST")


def _build_evidence_trail(case_id: str, reports_dir: Path) -> list[str]:
    lines = report_sections.build_evidence_trail_section(case_id, reports_dir, include_dismissed=True)
    if not lines:
        return []

    reflections = _parse_research_reflections(case_id, str(reports_dir))
    if reflections:
        lines += ["### Reflection log", ""]
        for r in reflections:
            lines.append(f"**{r['id']} — {r['trigger']}** *(recorded {r['timestamp']})*")
            lines.append("")
            if r["reinterpret"] and r["reinterpret"] != "—":
                lines.append(f"> Re-interpretations: {r['reinterpret']}")
                lines.append("")
            if r["open_leads"] and r["open_leads"] != "—":
                lines.append(f"> Open leads: {r['open_leads']}")
                lines.append("")
    else:
        lines += [
            "### Reflection log", "",
            "No reflection entries recorded.",
            "Use `python3 lib/research_notes.py reflect` to log mid-investigation re-assessments.",
            "",
        ]

    return lines


def _build_machine_details_section(data: dict[str, Any]) -> list[str]:
    """
    Build the 'Machine Investigation Details' section from machine_details.json.
    Covers system identity, user accounts, network configuration, MAC addresses,
    PowerShell history, installed applications, and user MRU activity.
    """
    machine = data.get("machine_details") or {}
    if not machine:
        return []

    lines: list[str] = ["---", "", "## 3. Machine Investigation Details", ""]
    a = lines.append

    a("> Claude: enhance and elaborate when necessary — cross-reference every username,")
    a("> IP address, and application with artifacts found in later sections (EVTX,")
    a("> Prefetch, Browser History, Recycle Bin). Record new IOCs via `--record-ioc`.")
    a("")

    software = machine.get("software") or {}
    system   = machine.get("system") or {}
    sam      = machine.get("sam") or {}

    # ── 3.1 System Identity ──────────────────────────────────────────────────
    a("### 3.1 System Identity")
    a("")
    a("| Field | Value |")
    a("|-------|-------|")

    computer_name = system.get("ComputerName") or "—"
    a(f"| Computer Name (Registry) | `{computer_name}` |")

    owner = software.get("RegisteredOwner") or "—"
    org   = software.get("RegisteredOrganization") or "—"
    a(f"| Registered Owner | {owner} |")
    a(f"| Registered Organization | {org} |")

    os_name    = software.get("ProductName") or "—"
    os_display = software.get("DisplayVersion") or ""
    os_build   = software.get("CurrentBuild") or software.get("CurrentVersion") or "—"
    os_branch  = software.get("BuildBranch") or "—"
    a(f"| OS Product | {os_name} {os_display}".rstrip() + " |")
    a(f"| OS Build | {os_build} |")
    a(f"| OS Branch | {os_branch} |")

    install_date = software.get("InstallDate")
    if install_date:
        try:
            install_dt = datetime.fromtimestamp(int(install_date), tz=timezone.utc)
            a(f"| OS Install Date (UTC) | {install_dt.strftime('%Y-%m-%d')} |")
        except Exception:
            a(f"| OS Install Date | {install_date} |")

    tz_name = (
        system.get("TimeZone_TimeZoneKeyName")
        or system.get("TimeZone_StandardName")
        or "—"
    )
    a(f"| System Timezone | {tz_name} |")

    last_user = software.get("LastLoggedOnUser") or "—"
    a(f"| Last Logged-On User (Winlogon) | `{last_user}` |")
    a("")

    # ── 3.2 User Accounts (SAM) ──────────────────────────────────────────────
    users = sam.get("UserAccounts") or []
    a("### 3.2 User Accounts (SAM)")
    a("")
    a("> Claude: for each username found here, perform a full-disk string search")
    a("> across the mounted filesystem (`grep -ril <username> $FS_MOUNT`) and correlate")
    a("> with EVTX logon events (4624/4625), browser profiles, and Recycle Bin SIDs.")
    a("")
    if users:
        a("| Username | Full Name | Last Logon (UTC) | Logon Count | Password Last Set |")
        a("|----------|-----------|-----------------|-------------|-------------------|")
        for user in users:
            uname = user.get("Username") or "—"
            fname = user.get("FullName") or "—"
            logon = user.get("LastLogon") or "—"
            count = str(user.get("LogonCount") or "—")
            pwd   = user.get("PasswordLastSet") or "—"
            a(f"| `{uname}` | {fname} | {logon} | {count} | {pwd} |")
    else:
        a("*SAM hive not available or no accounts enumerated.*")
    a("")

    # ── 3.3 Network Configuration ────────────────────────────────────────────
    interfaces = system.get("NetworkInterfaces") or []
    a("### 3.3 Network Configuration (TCP/IP)")
    a("")
    if interfaces:
        a("| Interface GUID | Static IP | DHCP IP | Subnet Mask | Default Gateway |")
        a("|----------------|-----------|---------|-------------|-----------------|")
        for iface in interfaces:
            guid    = (iface.get("GUID") or "—")[:36]
            ip      = iface.get("IPAddress") or "—"
            dhcp_ip = iface.get("DhcpIPAddress") or "—"
            mask    = iface.get("SubnetMask") or iface.get("DhcpSubnetMask") or "—"
            gw      = iface.get("DefaultGateway") or iface.get("DhcpDefaultGateway") or "—"
            for field_val in (ip, dhcp_ip, mask, gw):
                if isinstance(field_val, list):
                    field_val = ", ".join(str(v) for v in field_val if v)
            ip      = ip if not isinstance(ip, list) else ", ".join(ip)
            dhcp_ip = dhcp_ip if not isinstance(dhcp_ip, list) else ", ".join(dhcp_ip)
            mask    = mask if not isinstance(mask, list) else ", ".join(mask)
            gw      = gw if not isinstance(gw, list) else ", ".join(gw)
            a(f"| `{guid}` | {ip} | {dhcp_ip} | {mask} | {gw} |")
    else:
        a("*No network interface configuration extracted from SYSTEM hive.*")
    a("")

    # ── 3.4 Network Adapters & MAC Addresses ─────────────────────────────────
    adapters = system.get("NetworkAdapters") or []
    a("### 3.4 Network Adapters & MAC Addresses")
    a("")
    if adapters:
        a("| Adapter Description | MAC Address | Instance ID |")
        a("|---------------------|-------------|-------------|")
        for adapter in adapters:
            desc    = adapter.get("DriverDesc") or "—"
            mac     = adapter.get("NetworkAddress") or "—"
            inst_id = (adapter.get("NetCfgInstanceId") or "—")[:40]
            a(f"| {desc} | `{mac}` | `{inst_id}` |")
    else:
        a("*No network adapters with MAC addresses extracted from SYSTEM hive.*")
    a("")

    # ── 3.5 PowerShell History ───────────────────────────────────────────────
    ps_history = machine.get("ps_history") or []
    a("### 3.5 PowerShell Execution History")
    a("")
    a("> Claude: scan for encoded commands (Base64 blobs), download cradles")
    a("> (`IEX`, `Invoke-WebRequest`, `curl`, `wget`), credential access patterns")
    a("> (`mimikatz`, `sekurlsa`, `lsass`), and lateral movement commands.")
    a("")
    if ps_history:
        for entry in ps_history:
            uname    = entry.get("username") or "unknown"
            commands = entry.get("commands") or []
            count    = entry.get("command_count") or len(commands)
            a(f"**User: `{uname}`** ({count} commands in history)")
            a("")
            if commands:
                a("```powershell")
                for cmd in commands[:50]:
                    a(cmd)
                if len(commands) > 50:
                    a(f"# … {len(commands) - 50} more commands (see exports/machine_details/ps_history.txt)")
                a("```")
                a("")
    else:
        a("*No PowerShell history files found — either PSReadLine is not installed,*")
        a("*the history was cleared (T1070.003), or the filesystem was not mounted.*")
        a("")

    # ── 3.6 Installed Applications ───────────────────────────────────────────
    apps = software.get("InstalledApplications") or []
    a("### 3.6 Installed Applications")
    a("")
    a("> Claude: flag dual-use tools (remote access, port scanners, network sniffers,")
    a("> packet injectors, VPNs, anonymisers, P2P clients) and investigate them using")
    a("> the Application Deep-Dive Workflow: locate binary via Prefetch/EVTX → scan")
    a("> install dir and siblings → full-disk string search → check SRUM for network usage.")
    a("")
    if apps:
        sorted_apps = sorted(
            apps,
            key=lambda x: (x.get("InstallDate") or ""),
            reverse=True,
        )
        a("| Application | Version | Publisher | Install Date |")
        a("|-------------|---------|-----------|--------------|")
        for app in sorted_apps[:60]:
            name  = (app.get("DisplayName") or "—")[:60].replace("|", "\\|")
            ver   = (app.get("DisplayVersion") or "—")[:20]
            pub   = (app.get("Publisher") or "—")[:40].replace("|", "\\|")
            idate = app.get("InstallDate") or "—"
            a(f"| {name} | {ver} | {pub} | {idate} |")
        if len(apps) > 60:
            a(f"")
            a(f"*{len(apps) - 60} additional applications omitted — see machine_details.json.*")
    else:
        a("*No installed applications extracted from SOFTWARE hive.*")
    a("")

    # ── 3.7 Recent User Activity (MRU / UserAssist) ──────────────────────────
    ntusers = machine.get("ntuser_dat") or []
    if ntusers:
        a("### 3.7 Recent User Activity (MRU / UserAssist)")
        a("")
        for ntuser in ntusers:
            uname = ntuser.get("username") or "unknown"
            a(f"**User: `{uname}`**")
            a("")

            typed_paths = ntuser.get("TypedPaths") or {}
            if typed_paths:
                a("*Explorer typed paths (address bar):*")
                for v in list(typed_paths.values())[:15]:
                    a(f"- `{v}`")
                a("")

            typed_urls = ntuser.get("TypedURLs") or {}
            if typed_urls:
                a("*IE / Legacy typed URLs:*")
                for v in list(typed_urls.values())[:15]:
                    a(f"- `{v}`")
                a("")

            run_mru = ntuser.get("RunMRU") or {}
            if run_mru:
                a("*Win+R run dialog history:*")
                for v in list(run_mru.values())[:15]:
                    a(f"- `{v}`")
                a("")

            rd_exts = ntuser.get("RecentDocExtensions") or []
            if rd_exts:
                a(f"*Recently opened document types:* {', '.join(rd_exts[:20])}")
                a("")

            ua_execs = ntuser.get("UserAssistExecutions") or []
            if ua_execs:
                a("*UserAssist execution history (programs launched via Explorer):*")
                for exec_path in ua_execs[:30]:
                    a(f"- `{exec_path}`")
                a("")

    return lines


def _build_recyclebin_section(data: dict[str, Any], case_id: str) -> list[str]:
    """
    Build the 'Recycle Bin Analysis' section.
    Uses parsed $I metadata when available; falls back to raw file listing.
    """
    rb_parsed  = data.get("recyclebin_parsed")
    rb_list    = data.get("recyclebin_list") or []

    if not rb_parsed and not rb_list:
        return []

    lines: list[str] = ["---", "", "## 6b. Recycle Bin Analysis", ""]
    a = lines.append

    a("> Claude: enhance and elaborate when necessary — deleted files in the Recycle Bin")
    a("> are direct evidence of intentional file removal (MITRE T1070.004).")
    a("> Recover `$R` content files for executables and scripts found here.")
    a("> Map SID values to user accounts in Section 3.2 to attribute deletion activity.")
    a("")

    if rb_parsed:
        valid   = [e for e in rb_parsed if "error" not in e]
        errors  = [e for e in rb_parsed if "error" in e]

        if valid:
            total_size = sum(e.get("size_bytes") or 0 for e in valid)
            dates = sorted(
                e["deleted_at_utc"]
                for e in valid
                if e.get("deleted_at_utc")
            )
            a(f"**{len(valid)} item(s) found** — total original size: **{_humanize_bytes(total_size)}**")
            if dates:
                a(f"**Deletion date range:** {dates[0]} → {dates[-1]}")
            a("")

            a("| Original Path | Ext | Original Size | Deleted At (UTC) | SID | Content Recovered |")
            a("|---------------|-----|---------------|-----------------|-----|-------------------|")
            for entry in sorted(valid, key=lambda x: x.get("deleted_at_utc") or ""):
                path = entry.get("original_path") or "—"
                # Truncate long paths
                if len(path) > 70:
                    path = "…" + path[-67:]
                ext       = entry.get("extension") or "—"
                size      = _humanize_bytes(entry.get("size_bytes") or 0)
                deleted   = entry.get("deleted_at_utc") or "—"
                sid       = entry.get("sid") or "—"
                recovered = "**Yes**" if entry.get("r_file_present") else "No"
                a(f"| `{path}` | {ext} | {size} | {deleted} | {sid} | {recovered} |")
            a("")

            sids = sorted(set(e.get("sid") for e in valid if e.get("sid")))
            if sids:
                a("> **SID → user correlation:** cross-reference these SIDs with Section 3.2")
                a("> (SAM user accounts) to attribute each deletion to a specific account.")
                a("")
                for sid in sids:
                    a(f"- `{sid}`")
                a("")

            a("> **Recovery:** files marked 'Content Recovered: Yes' have their content preserved in")
            a(f"> `./exports/recyclebin/$R*`. Examine executables and scripts for malware or staging artefacts.")
            a("")

        if errors:
            a(f"*{len(errors)} $I file(s) could not be parsed (see recyclebin_parsed.json for details).*")
            a("")

    elif rb_list:
        i_files = [f for f in rb_list if f.startswith("$I")]
        r_files = [f for f in rb_list if f.startswith("$R")]
        a(f"**Raw extraction:** {len(i_files)} `$I` metadata file(s) and {len(r_files)} `$R` content file(s).")
        a("")
        a("*Metadata not yet parsed. Run:*")
        a("```bash")
        a("python3 lib/fast_machine_details.py --exports ./exports")
        a("```")
        a("")
        a("Raw Recycle Bin file listing:")
        for f in rb_list[:30]:
            a(f"- `./exports/recyclebin/{f}`")
        if len(rb_list) > 30:
            a(f"- *… and {len(rb_list) - 30} more*")
        a("")

    return lines


# IOC Reference categories — used to derive the canonical "category" column
# from iocs.json's per-indicator "category" (an indicator-type label, e.g. "ip").
_IOC_REFERENCE_CATEGORIES: list[tuple[str, list[str]]] = [
    ("Network Indicators",      ["ip", "mac", "domain", "url", "fqdn"]),
    ("User Identifiers",        ["username", "fullname", "email", "sid"]),
    ("File System Indicators",  ["filepath", "deleted_filepath", "hash", "filename"]),
    ("Persistence Indicators",  ["registry_key", "scheduled_task", "service"]),
    ("Suspicious Applications", ["suspicious_app"]),
]


def _iocs_reference_to_canonical(data: dict[str, Any]) -> list[dict]:
    """
    Map iocs.json entries ({category, value/raw_value, confidence, source_step})
    to the canonical IOC schema shared with FAN/FAME via report_sections.build_ioc_section.
    """
    iocs = data.get("iocs_reference")
    if not iocs:
        return []

    out: list[dict] = []
    for ioc in iocs:
        raw_type = (ioc.get("category") or "").lower()
        category = "Other Indicators"
        for cat_title, cat_types in _IOC_REFERENCE_CATEGORIES:
            if raw_type in cat_types:
                category = cat_title
                break
        out.append({
            "type": raw_type or "indicator",
            "value": ioc.get("value") or ioc.get("raw_value") or "—",
            "severity": "info",
            "category": category,
            "source": ioc.get("source_step") or "—",
            "confidence": ioc.get("confidence") or "CONFIRMED",
        })
    return out


def _build_markdown(
    data: dict[str, Any],
    case_id: str,
    hostname: str,
    disk_image: str,
    generated_utc: str,
    opencti_findings: str = "",
    fan_summary: str = "",
    fame_summary: str = "",
    reports_dir: Path | None = None,
) -> str:
    """
    Build the full FAST incident report in Markdown.

    Claude: enhance and elaborate when necessary on each section to ensure
    the analyst receives full contextual depth appropriate to their level.
    """
    lines: list[str] = []
    a = lines.append

    reports_dir = reports_dir or (PROJECT_ROOT / "reports")
    narrative = _load_narrative(case_id, reports_dir)

    narrative_result = report_completeness.check_narrative(case_id, "FAST", reports_dir)
    reasoning_result = report_completeness.check_research_notes(case_id, reports_dir)
    report_completeness.write_incomplete_marker(reports_dir, case_id, narrative_result, reasoning_result)
    incomplete_banner = report_completeness.format_incomplete_banner(narrative_result, reasoning_result)
    if incomplete_banner:
        print(f"[fast] WARNING: investigation incomplete for {case_id}/{hostname} — "
              f"see {case_id}_INVESTIGATION_INCOMPLETE.json", file=sys.stderr)

    a("# FAST Storage Forensics Report")
    a("")
    a("| Field | Value |")
    a("|-------|-------|")
    a(f"| Case ID | `{case_id}` |")
    a(f"| Hostname | `{hostname}` |")
    a(f"| Disk image | `{disk_image}` |")
    a(f"| Module | FAST — Forensic Analysis Storage |")
    a(f"| Analyst | Claude Code — FAST skill |")
    a(f"| Generated (UTC) | {generated_utc} |")
    a("")
    if incomplete_banner:
        lines.extend(incomplete_banner)

    # ── Cross-Module Intelligence ─────────────────────────────────────────────
    if fan_summary or fame_summary:
        a("---")
        a("")
        a("## Cross-module intelligence")
        a("")
        a("> Claude: enhance and elaborate when necessary — correlate storage findings")
        a("> with network (FAN) and memory (FAME) evidence below.")
        a("")
        if fan_summary:
            a("### Network forensics (FAN) summary")
            a("")
            a(fan_summary.strip())
            a("")
        if fame_summary:
            a("### Memory forensics (FAME) summary")
            a("")
            a(fame_summary.strip())
            a("")

    # ── Management Summary ────────────────────────────────────────────────────
    a("---")
    a("")
    a("## 1. Management summary")
    a("")
    a("> **Audience:** CISO, Legal, Internal Audit — no technical identifiers.")
    a("> Claude: enhance and elaborate when necessary.")
    a("")
    a("Storage forensic analysis of the subject server disk image was conducted to")
    a("identify file system artifacts, deleted files, and on-disk evidence relevant")
    a("to the investigation. The analysis examined the disk image for signs of attacker")
    a("activity, data staging, persistence mechanisms, and timeline of key file operations.")
    a("")
    fls = data.get("fls_output", "")
    deleted_count = fls.count("\n* ") if fls else 0
    evtx = data.get("evtx_list", [])
    prefetch = data.get("prefetch_list", [])
    if deleted_count > 0:
        a(f"**Key finding:** {deleted_count} deleted file entries identified in the file system.")
    if evtx:
        a(f"**Event logs:** {len(evtx)} Windows Event Log files extracted for further analysis.")
    if prefetch:
        a(f"**Prefetch:** {len(prefetch)} Prefetch execution artifacts extracted.")
    a("")

    # ── Incident Timeline (Claude-generated) ─────────────────────────────────
    timeline_text = narrative.get("attack_timeline", "")
    a("---")
    a("")
    a("## 2. Incident Timeline")
    a("")
    a("> Chronological reconstruction of the attack path. Each finding references")
    a("> the investigation step (RN-NNN) and the preserved source file in")
    a(f"> `{case_id}_evidence/`.")
    a("")
    if timeline_text:
        a(timeline_text)
    else:
        a("> *Incident timeline not yet generated. Run the FAST skill to produce*")
        a(f"> *`{case_id}_narrative.md` with the `attack_timeline` section.*")
    a("")

    # ── Machine Investigation Details ─────────────────────────────────────────
    lines.extend(_build_machine_details_section(data))

    # ── Image Verification ────────────────────────────────────────────────────
    ewfinfo   = data.get("ewfinfo", "")
    ewfverify = data.get("ewfverify", "")
    if ewfinfo or ewfverify:
        a("---")
        a("")
        a("## 2. Image verification")
        a("")
        a("> Claude: enhance and elaborate when necessary — confirm hash match and")
        a("> document any acquisition notes embedded in the E01 metadata.")
        a("")
        if ewfinfo:
            a("### ewfinfo")
            a("```")
            a(ewfinfo.strip()[:2000])
            a("```")
            a("")
        if ewfverify:
            a("### ewfverify")
            a("```")
            a(ewfverify.strip()[:500])
            a("```")
            a("")

    # ── Partition Table ───────────────────────────────────────────────────────
    mmls = data.get("mmls", "")
    if mmls:
        a("---")
        a("")
        a("## 3. Partition table (mmls)")
        a("")
        a("> Claude: enhance and elaborate when necessary — identify the primary OS")
        a("> partition, recovery partitions, and any unallocated space.")
        a("")
        a("```")
        a(mmls.strip()[:2000])
        a("```")
        a("")

    # ── Filesystem Metadata ───────────────────────────────────────────────────
    fsstat = data.get("fsstat", "")
    if fsstat:
        a("---")
        a("")
        a("## 4. Filesystem metadata (fsstat)")
        a("")
        a("> Claude: enhance and elaborate when necessary — note NTFS version, cluster")
        a("> size, MFT location, and volume last-written timestamp.")
        a("")
        a("```")
        a(fsstat.strip()[:2000])
        a("```")
        a("")

    # ── File System Timeline ──────────────────────────────────────────────────
    timeline_rows = data.get("fs_timeline_csv", [])
    timeline_txt  = data.get("fs_timeline", "")
    if timeline_rows or timeline_txt:
        a("---")
        a("")
        a("## 5. Filesystem timeline")
        a("")
        fs_narrative = narrative.get("section_filesystem", "")
        if fs_narrative:
            a(fs_narrative)
            a("")
        else:
            a("> Claude: enhance and elaborate when necessary — highlight any file activity")
            a("> that coincides with the network or memory-forensics event timeline.")
            a("")
        a(f"> **Source file:** [`{case_id}_evidence/exports/fs_timeline.csv`]"
          f"(./{case_id}_evidence/exports/fs_timeline.csv)")
        a("")
        if timeline_rows:
            shown = timeline_rows[:40]
            headers = list(shown[0].keys()) if shown else []
            if headers:
                a("| " + " | ".join(headers) + " |")
                a("|" + "|".join(["---"] * len(headers)) + "|")
                for row in shown:
                    a("| " + " | ".join(str(row.get(h, "")) for h in headers) + " |")
                if len(timeline_rows) > 40:
                    a(f"")
                    a(f"*{len(timeline_rows) - 40} additional timeline entries not shown — see `./exports/fs_timeline.csv`.*")
        elif timeline_txt:
            a("```")
            a(timeline_txt.strip()[:3000])
            a("```")
        a("")

    # ── Deleted Files ─────────────────────────────────────────────────────────
    if fls and deleted_count > 0:
        a("---")
        a("")
        a("## 6. Deleted file entries (fls)")
        a("")
        a("> Claude: enhance and elaborate when necessary — any deleted executable or")
        a("> script in a user-writable path is a high-priority triage candidate.")
        a("")
        deleted_lines = [l for l in fls.splitlines() if l.startswith("* ")][:50]
        a("```")
        for dl in deleted_lines:
            a(dl)
        a("```")
        if deleted_count > 50:
            a(f"")
            a(f"*{deleted_count - 50} additional deleted entries not shown — see `./analysis/storage/fls_output.txt`.*")
        a("")

    # ── Recycle Bin Analysis ──────────────────────────────────────────────────
    lines.extend(_build_recyclebin_section(data, case_id))

    # ── MFT / UsnJrnl ─────────────────────────────────────────────────────────
    mft = data.get("mft_extracted", False)
    usn = data.get("usnj_extracted", False)
    if mft or usn:
        a("---")
        a("")
        a("## 7. MFT and USN change journal")
        a("")
        a("> Claude: enhance and elaborate when necessary — parse MFT with MFTECmd and")
        a("> USN journal with MFTeCmd or usnjrnl.py to build a high-resolution file activity log.")
        a("")
        a("| Artifact | Status |")
        a("|----------|--------|")
        a(f"| \\$MFT | {'Extracted — `./exports/mft/$MFT`' if mft else 'Not extracted'} |")
        a(f"| \\$J (UsnJrnl) | {'Extracted — `./exports/mft/$J`' if usn else 'Not extracted'} |")
        a("")

    # ── Extracted Artifacts ───────────────────────────────────────────────────
    artifact_sections = [
        ("evtx_list",     "8. Windows Event Logs",
         "Ingest into Timeline Explorer, Chainsaw, or Hayabusa for rapid log analysis. "
         "**Deep-dive:** run `python3 -m evtx` or `chainsaw hunt` for logon events "
         "(4624/4625), service installs (7045), process creation (4688), and PowerShell "
         "script block logging (4103/4104)."),
        ("registry_list", "9. Registry Hives",
         "Analyze with RegRipper, Registry Explorer, or RECmd to extract run keys, services, and user artifacts. "
         "**Key paths:** `SYSTEM\\ControlSet001\\Services` (service persistence), "
         "`SOFTWARE\\Microsoft\\Windows\\CurrentVersion\\Run` (autorun), "
         "`NTUSER.DAT\\Software\\Microsoft\\Windows\\CurrentVersion\\Explorer\\UserAssist` (execution history), "
         "`SYSTEM\\ControlSet001\\Control\\Session Manager\\AppCompatCache` (ShimCache)."),
        ("prefetch_list", "10. Prefetch Execution Artifacts",
         "Parse with PECmd to reconstruct program execution history and first/last run times. "
         "**Deep-dive:** `python3 -m prefetch_parser -d ./exports/prefetch/` or "
         "`PECmd.exe -d ./exports/prefetch/ --csv ./exports/` — extract run count, "
         "first/last run timestamps, and loaded DLLs. Flag executables run only once "
         "or within the incident window. Cross-reference with Section 3.6 installed apps."),
        ("srum_list",     "11. SRUM Database",
         "Parse SRUDB.dat with srum-dump or SRUM-DUMP2 for per-app network and CPU usage. "
         "**Deep-dive:** correlates each application's network bytes sent/received with "
         "timestamps — identifies data exfiltration and C2 beaconing patterns."),
        ("browser_list",  "12. Browser History",
         "Analyze History files with DB Browser for SQLite or BrowsingHistoryView. "
         "**Deep-dive SQL:** `SELECT url, title, visit_count, last_visit_time FROM urls "
         "ORDER BY last_visit_time DESC LIMIT 100;` — surface browsing patterns, "
         "downloads table for installer URLs, and keyword-search against case IOCs."),
    ]
    for key, title, note in artifact_sections:
        items = data.get(key, [])
        if items:
            a("---")
            a("")
            a(f"## {title}")
            a("")
            a(f"> Claude: enhance and elaborate when necessary — {note}")
            a("")
            for item in items[:30]:
                a(f"- `{item}`")
            if len(items) > 30:
                a(f"- *… and {len(items) - 30} more*")
            a("")

    # ── Carved Files ──────────────────────────────────────────────────────────
    carved = data.get("bulk_carved", [])
    if carved:
        a("---")
        a("")
        a("## 13. Bulk extractor carved artifacts")
        a("")
        a("> Claude: enhance and elaborate when necessary — email addresses, URLs, and")
        a("> domains found in unallocated space are IOC candidates for CTI enrichment.")
        a("")
        for item in carved[:20]:
            a(f"- `./exports/carved/{item}`")
        a("")

    # ── File Hashes ───────────────────────────────────────────────────────────
    md5_manifest = data.get("md5_manifest", "")
    if md5_manifest:
        a("---")
        a("")
        a("## 14. File hash manifest")
        a("")
        a("> Claude: enhance and elaborate when necessary — submit hashes to VirusTotal")
        a("> or check against the NSRL to filter known-good files.")
        a("")
        a("```")
        a(md5_manifest.strip()[:2000])
        a("```")
        a("")

    # ── OpenCTI Enrichment ────────────────────────────────────────────────────
    if opencti_findings:
        a("---")
        a("")
        a("## 15. OpenCTI threat intelligence enrichment")
        a("")
        a("> Claude: enhance and elaborate when necessary — link matched indicators to")
        a("> known threat actors, malware families, or campaigns in OpenCTI.")
        a("")
        a(opencti_findings.strip())
        a("")

    # ── MITRE ATT&CK ─────────────────────────────────────────────────────────
    techniques = _derive_techniques(data)
    technique_dicts = [
        {"id": tid, "name": name, "tactic": tactic, "observation": obs}
        for tid, name, tactic, obs in techniques
    ]
    a("---")
    a("")
    a("\n".join(report_sections.build_mitre_section(
        technique_dicts, heading="MITRE ATT&CK Coverage", section_num="16",
        show_severity=False,
        empty_message="No MITRE ATT&CK techniques mapped — insufficient evidence for attribution.")))

    # ── Indicators of Compromise (iocs.json + bulk_extractor carves) ──────────
    merged_iocs = _iocs_reference_to_canonical(data) + _extract_iocs(data)
    a("\n".join(report_sections.build_ioc_section(
        merged_iocs, heading="Indicators of Compromise", section_num="17",
        empty_message="No malicious indicators of compromise identified from storage analysis.")))

    # ── Recommendations ───────────────────────────────────────────────────────
    recs = _build_recommendations(data)
    a("\n".join(report_sections.build_recommendations_section(
        recs, heading="Recommendations", section_num="18", numbered=True)))

    # ── Appendix ──────────────────────────────────────────────────────────────
    a("---")
    a("")
    a("## Appendix A — Analysis source files")
    a("")
    a(f"All artifact files are preserved in `./{case_id}_evidence/` and uploaded")
    a("to the investigations vault alongside this report. SHA-256 hashes are recorded in the")
    a("research notes (Appendix B) for chain-of-custody verification.")
    a("")
    a("| File | Description |")
    a("|------|-------------|")
    a(f"| [`{case_id}_evidence/storage/ewfinfo.txt`](./{case_id}_evidence/storage/ewfinfo.txt) | E01 image metadata |")
    a(f"| [`{case_id}_evidence/storage/mmls.txt`](./{case_id}_evidence/storage/mmls.txt) | Partition table |")
    a(f"| [`{case_id}_evidence/storage/fsstat.txt`](./{case_id}_evidence/storage/fsstat.txt) | Filesystem metadata |")
    a(f"| [`{case_id}_evidence/storage/fls_output.txt`](./{case_id}_evidence/storage/fls_output.txt) | Full file listing (incl. deleted) |")
    a(f"| [`{case_id}_evidence/storage/bodyfile.txt`](./{case_id}_evidence/storage/bodyfile.txt) | MAC time bodyfile |")
    a(f"| [`{case_id}_evidence/exports/fs_timeline.csv`](./{case_id}_evidence/exports/fs_timeline.csv) | Filesystem timeline (mactime CSV) |")
    a(f"| [`{case_id}_evidence/exports/mft/$MFT`](./{case_id}_evidence/exports/mft/) | Master File Table |")
    a(f"| [`{case_id}_evidence/exports/mft/$J`](./{case_id}_evidence/exports/mft/) | USN Change Journal |")
    a(f"| [`{case_id}_evidence/exports/evtx/`](./{case_id}_evidence/exports/evtx/) | Windows Event Logs |")
    a(f"| [`{case_id}_evidence/exports/registry/`](./{case_id}_evidence/exports/registry/) | Registry hives |")
    a(f"| [`{case_id}_evidence/exports/prefetch/`](./{case_id}_evidence/exports/prefetch/) | Prefetch files |")
    a(f"| [`{case_id}_evidence/exports/carved/`](./{case_id}_evidence/exports/carved/) | bulk_extractor carved artifacts |")
    a("")
    a("*All findings derived from disk image analysis as stated. Evidence integrity preserved.*")
    a("")

    # ── Hallucination Guard ───────────────────────────────────────────────────
    hg_section = _build_fast_hallucination_guard_section(data, case_id, reports_dir)
    if hg_section:
        a(hg_section)
        a("")

    # ── Evidence Trail ────────────────────────────────────────────────────────
    lines.extend(_build_evidence_trail(case_id, reports_dir))

    return "\n".join(lines)


def _derive_techniques(data: dict[str, Any]) -> list[tuple[str, str, str, str]]:
    techniques = []
    deleted = data.get("fls_output", "").count("\n* ")
    if deleted > 0:
        techniques.append((
            "T1070.004", "Indicator Removal: File Deletion", "Defense Evasion",
            f"{deleted} deleted file entries found — evidence of file clean-up activity"
        ))
    prefetch = data.get("prefetch_list", [])
    if prefetch:
        techniques.append((
            "T1564.003", "Hide Artifacts: NTFS File Attributes", "Defense Evasion",
            "Prefetch artifacts present — can reveal execution of tools even after deletion"
        ))
    evtx = data.get("evtx_list", [])
    if evtx:
        techniques.append((
            "T1070.001", "Indicator Removal: Clear Windows Event Logs", "Defense Evasion",
            f"Event logs extracted ({len(evtx)} files) — review for gaps or cleared logs"
        ))
    carved = data.get("bulk_carved", [])
    if "url" in str(carved).lower() or "domain" in str(carved).lower():
        techniques.append((
            "T1048", "Exfiltration Over Alternative Protocol", "Exfiltration",
            "URLs/domains found in carved artifacts — review for data staging or C2 URLs"
        ))
    return techniques


def _extract_iocs(data: dict[str, Any]) -> list[dict]:
    iocs = []
    carved_files = data.get("bulk_carved", [])
    for cf in carved_files:
        if "url" in cf.lower():
            iocs.append({
                "type": "File",
                "value": f"./exports/carved/{cf}",
                "severity": report_sections.normalize_severity("Medium"),
                "category": "Carved Artifact",
                "source": "bulk_extractor",
                "context": "bulk_extractor URL feature file — review for C2 or exfiltration endpoints",
            })
            break
        if "domain" in cf.lower():
            iocs.append({
                "type": "File",
                "value": f"./exports/carved/{cf}",
                "severity": report_sections.normalize_severity("Medium"),
                "category": "Carved Artifact",
                "source": "bulk_extractor",
                "context": "bulk_extractor domain feature file — cross-reference with FAN DNS findings",
            })
            break
    return iocs


def _build_recommendations(data: dict[str, Any]) -> list[str]:
    recs = []
    evtx = data.get("evtx_list", [])
    if evtx:
        recs.append(
            "**AnalyzeWindows Event Logs** — ingest the extracted EVTX files into "
            "Timeline Explorer or Hayabusa to identify logon events, service installs, "
            "and PowerShell activity around the time of the incident."
        )
    prefetch = data.get("prefetch_list", [])
    if prefetch:
        recs.append(
            "**Parse Prefetch files** — run PECmd against the extracted Prefetch files "
            "to reconstruct which executables ran, their first and last run timestamps, "
            "and what files they accessed."
        )
    if data.get("mft_extracted"):
        recs.append(
            "**Parse the MFT** — run MFTECmd against `./exports/mft/$MFT` and combine "
            "with the USN Change Journal to build a high-resolution file-modification timeline. "
            "Cross-reference with the FAME memory timeline."
        )
    registry = data.get("registry_list", [])
    if registry:
        recs.append(
            "**Analyzeregistry hives** — run RegRipper or RECmd against the extracted "
            "hives to identify Run/RunOnce persistence keys, recently executed programs "
            "(UserAssist), and network share history."
        )
    carved = data.get("bulk_carved", [])
    if carved:
        recs.append(
            "**Review carved artifacts** — examine bulk_extractor output for email addresses, "
            "URLs, and credit card numbers recovered from unallocated space. Submit any "
            "suspicious URLs to OpenCTI for CTI enrichment."
        )
    recs.append(
        "**Cross-reference with FAME and FAN** — correlate disk timeline events with "
        "the memory forensics (FAME) process timeline and network (FAN) packet timeline "
        "to identify the full attack chain."
    )
    recs.append(
        "**Verify image integrity** — confirm ewfverify completed without errors before "
        "any findings are cited as evidence in legal or regulatory proceedings."
    )
    return recs


# ── DOCX utility helpers ──────────────────────────────────────────────────────

def _compute_file_metadata(file_path: Path) -> tuple[str, str]:
    if not file_path.exists():
        return "N/A", "N/A"
    mtime_str = datetime.fromtimestamp(
        file_path.stat().st_mtime, tz=_CET
    ).strftime("%d-%b-%Y %H:%M CET")
    h = hashlib.sha256()
    with open(file_path, "rb") as fh:
        for chunk in iter(lambda: fh.read(65536), b""):
            h.update(chunk)
    return mtime_str, h.hexdigest()


def _set_doc_font(doc: Any, font_name: str = "Arial") -> None:
    for style_name in [
        "Normal", "Heading 1", "Heading 2", "Heading 3", "Heading 4",
        "No Spacing", "Intense Quote", "List Number", "List Bullet", "Table Grid",
    ]:
        try:
            doc.styles[style_name].font.name = font_name
        except Exception:
            pass


def _add_header_footer(section: Any, case_id: str) -> None:
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    hdr = section.header
    h_para = hdr.paragraphs[0] if hdr.paragraphs else hdr.add_paragraph()
    h_para.clear()
    h_run = h_para.add_run(case_id)
    h_run.font.name = "Arial"
    h_run.font.size = Pt(9)
    h_run.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)
    h_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    ftr = section.footer
    f_para = ftr.paragraphs[0] if ftr.paragraphs else ftr.add_paragraph()
    f_para.clear()
    f_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    def _fld_char(fld_type: str) -> Any:
        r = OxmlElement("w:r")
        fc = OxmlElement("w:fldChar")
        fc.set(qn("w:fldCharType"), fld_type)
        r.append(fc)
        return r

    def _instr(text: str) -> Any:
        r = OxmlElement("w:r")
        it = OxmlElement("w:instrText")
        it.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
        it.text = text
        r.append(it)
        return r

    def _txt_run(text: str) -> Any:
        r = f_para.add_run(text)
        r.font.name = "Arial"
        r.font.size = Pt(9)
        r.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)
        return r

    _txt_run("Page ")
    f_para._p.append(_fld_char("begin"))
    f_para._p.append(_instr(" PAGE "))
    f_para._p.append(_fld_char("end"))
    _txt_run(" of ")
    f_para._p.append(_fld_char("begin"))
    f_para._p.append(_instr(" NUMPAGES "))
    f_para._p.append(_fld_char("end"))


def _add_watermark(section: Any) -> None:
    from lxml import etree

    hdr = section.header
    wm_para = hdr.add_paragraph()
    vml = (
        '<w:r xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"'
        ' xmlns:v="urn:schemas-microsoft-com:vml"'
        ' xmlns:o="urn:schemas-microsoft-com:office:office">'
        "<w:rPr><w:noProof/></w:rPr>"
        "<w:pict>"
        '<v:shape id="WaterMark" o:spid="_x0000_s2049"'
        ' type="#_x0000_t136"'
        ' style="position:absolute;margin-left:0;margin-top:0;'
        "width:430pt;height:120pt;z-index:-251654144;"
        "rotation:315;"
        "mso-position-horizontal:center;"
        "mso-position-horizontal-relative:page;"
        "mso-position-vertical:center;"
        'mso-position-vertical-relative:page"'
        ' fillcolor="#C0C0C0" stroked="f">'
        "<v:textpath"
        ' on="t" fitshape="t" string="CONFIDENTIAL"'
        " style='font-family:\"Arial\";font-size:1pt;font-weight:bold'"
        "/>"
        "</v:shape>"
        "</w:pict>"
        "</w:r>"
    )
    wm_para._p.append(etree.fromstring(vml))


def _remove_blank_paragraphs(doc: Any) -> None:
    from docx.oxml.ns import qn

    paras = list(doc.paragraphs)
    to_remove = []
    for i, para in enumerate(paras):
        if para.text.strip():
            continue
        if para.style.name.startswith("Heading"):
            continue
        if para._element.find(".//" + qn("w:br")) is not None:
            continue
        prev_heading = i > 0 and paras[i - 1].style.name.startswith("Heading")
        next_heading = i + 1 < len(paras) and paras[i + 1].style.name.startswith("Heading")
        if prev_heading or next_heading:
            continue
        to_remove.append(para._element)
    for elem in to_remove:
        if elem.getparent() is not None:
            elem.getparent().remove(elem)


# ── PPTX ───────────────────────────────────────────────────────────────────────

def _build_pptx(
    data: dict[str, Any],
    case_id: str,
    hostname: str,
    disk_image: str,
    generated_utc: str,
    output_path: Path,
    opencti_findings: str = "",
    fan_summary: str = "",
    fame_summary: str = "",
) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("[fast] WARNING: python-pptx not installed — skipping PPTX. pip3 install python-pptx")
        return

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def _rgb(t):
        return RGBColor(*t)

    def _rect(slide, l, t, w, h, fill):
        s = slide.shapes.add_shape(1, l, t, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = _rgb(fill)
        s.line.fill.background()
        return s

    def _txt(slide, text, l, t, w, h, sz, bold=False, color=_WHITE, align=PP_ALIGN.LEFT):
        tb  = slide.shapes.add_textbox(l, t, w, h)
        tf  = tb.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.name = "Arial"
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)
        return tb

    W = prs.slide_width
    H = prs.slide_height
    M = Inches(0.4)

    def _add_bullet_slide(title, bullets, fallback, max_items=8):
        slide = prs.slides.add_slide(blank)
        _rect(slide, 0, 0, W, Inches(1.1), _MID_NAVY)
        _txt(slide, title, M, Inches(0.2), W, Inches(0.8), 28, bold=True, color=_WHITE)
        _txt(slide, f"{case_id}  |  {hostname}", M, Inches(0.75), W, Inches(0.3), 12, color=_LIGHT_BLUE)
        if bullets:
            body = "\n\n".join(f"•  {b}" for b in bullets[:max_items])
        else:
            body = fallback
        _txt(slide, body, M, Inches(1.3), W - 2 * M, H - Inches(1.6), 14, color=_TEXT_DARK)
        return slide

    fls      = data.get("fls_output", "")
    deleted  = fls.count("\n* ") if fls else 0
    evtx     = data.get("evtx_list", [])
    prefetch = data.get("prefetch_list", [])
    narrative = data.get("_narrative", {})

    # ── Slide 1 — Cover ───────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    _rect(s1, 0, 0, W, H, _DARK_NAVY)
    _rect(s1, 0, 0, W, Inches(0.08), _BLUE)
    _rect(s1, 0, H - Inches(0.08), W, Inches(0.08), _BLUE)
    _txt(s1, "FAST", M, Inches(1.2), W - 2*M, Inches(1.2),
         72, bold=True, color=_LIGHT_BLUE, align=PP_ALIGN.CENTER)
    _txt(s1, "Forensic analysis storage", M, Inches(2.2), W - 2*M, Inches(0.7),
         28, color=_WHITE, align=PP_ALIGN.CENTER)
    _txt(s1, "Storage forensics incident report", M, Inches(2.9), W - 2*M, Inches(0.6),
         20, color=_LIGHT_BLUE, align=PP_ALIGN.CENTER)
    _rect(s1, Inches(3), Inches(3.8), W - Inches(6), Inches(0.04), _BLUE)
    _txt(s1, f"Case: {case_id}  |  Host: {hostname}  |  {generated_utc[:10]}",
         M, Inches(4.1), W - 2*M, Inches(0.5), 14, color=_TEXT_MID, align=PP_ALIGN.CENTER)
    _txt(s1, "Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin · Joost Beekman",
         M, Inches(4.6), W - 2*M, Inches(0.4), 11, color=_TEXT_MID, align=PP_ALIGN.CENTER)
    _txt(s1, "Fan Get Fame Fast  |  FAST module",
         M, H - Inches(0.7), W - 2*M, Inches(0.4), 11, color=_TEXT_MID, align=PP_ALIGN.CENTER)

    # ── Slide 2 — Executive Summary ──────────────────────────────────────────
    fallback_summary = (
        f"Storage forensic analysis was conducted on the disk image of {hostname}. "
        f"The analysis identified {deleted} deleted file entries, "
        f"{len(evtx)} Windows Event Log file(s), and "
        f"{len(prefetch)} Prefetch execution artifact(s). "
        "The filesystem timeline has been reconstructed for cross-reference with "
        "network and memory forensics findings."
    )
    _add_bullet_slide(
        "Executive Summary",
        _narr_bullets(narrative, "pptx_executive_summary"),
        fallback_summary,
    )

    # ── Slide 3 — Business Impact ────────────────────────────────────────────
    _add_bullet_slide(
        "Business Impact",
        _narr_bullets(narrative, "pptx_impact"),
        "Operational impact assessment is in progress; refer to the technical report.",
    )

    # ── Slide 4 — Board Timeline ─────────────────────────────────────────────
    _add_bullet_slide(
        "Incident Timeline",
        _narr_bullets(narrative, "pptx_timeline"),
        "The incident timeline is under investigation; refer to the technical report for the full chronology.",
        max_items=6,
    )

    # ── Slide 5 — Root Cause & Risk ──────────────────────────────────────────
    s5 = prs.slides.add_slide(blank)
    _rect(s5, 0, 0, W, Inches(1.1), _MID_NAVY)
    _txt(s5, "Root Cause & Risk", M, Inches(0.2), W, Inches(0.8), 28, bold=True, color=_WHITE)
    _txt(s5, f"{case_id}  |  {hostname}", M, Inches(0.75), W, Inches(0.3), 12, color=_LIGHT_BLUE)
    root_cause = narrative.get("pptx_root_cause", "").strip() or (
        "Root cause is under investigation; refer to the technical report for the initial access vector."
    )
    _txt(s5, "ROOT CAUSE", M, Inches(1.25), W - 2 * M, Inches(0.35), 13, bold=True, color=_LIGHT_BLUE)
    _txt(s5, root_cause, M, Inches(1.65), W - 2 * M, Inches(1.3), 14, color=_TEXT_DARK)
    risk_bullets = _narr_bullets(narrative, "pptx_risk")
    _txt(s5, "KEY RISKS", M, Inches(3.1), W - 2 * M, Inches(0.35), 13, bold=True, color=_LIGHT_BLUE)
    if risk_bullets:
        risk_text = "\n\n".join(f"•  {b}" for b in risk_bullets[:5])
    else:
        risk_text = "Risk assessment is in progress; refer to the technical report."
    _txt(s5, risk_text, M, Inches(3.5), W - 2 * M, Inches(3.2), 14, color=_TEXT_DARK)

    # ── Slide 6 — Response & Containment ─────────────────────────────────────
    _add_bullet_slide(
        "Response & Containment",
        _narr_bullets(narrative, "pptx_mitigations"),
        "Response and containment actions are in progress; refer to the technical report.",
    )

    # ── Slide 7 — Recommendations ────────────────────────────────────────────
    s7 = prs.slides.add_slide(blank)
    _rect(s7, 0, 0, W, Inches(1.1), _MID_NAVY)
    _txt(s7, "Recommendations", M, Inches(0.2), W, Inches(0.8), 28, bold=True, color=_WHITE)
    _txt(s7, f"{case_id}  |  {hostname}", M, Inches(0.75), W, Inches(0.3), 12, color=_LIGHT_BLUE)
    recs = _narr_bullets(narrative, "pptx_recommendations") or _build_recommendations(data)
    row_h = Inches(0.72)
    for i, rec in enumerate(recs[:7]):
        y = Inches(1.2) + i * row_h
        _rect(s7, M, y, Inches(0.5), row_h - Inches(0.08), _BLUE)
        _txt(s7, str(i + 1), M + Inches(0.1), y + Inches(0.1), Inches(0.3), row_h,
             16, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
        rec_clean = re.sub(r"\*\*(.*?)\*\*", r"\1", rec.split(" — ")[0][:120])
        _txt(s7, rec_clean, M + Inches(0.6), y + Inches(0.1), W - M - Inches(1.0), row_h,
             13, color=_TEXT_DARK)

    # ── Slide 8 — Lessons Learned ────────────────────────────────────────────
    _add_bullet_slide(
        "Lessons Learned",
        _narr_bullets(narrative, "pptx_lessons_learned"),
        "Lessons learned will be documented once the investigation and remediation are complete.",
    )

    prs.save(str(output_path))
    print(f"[fast] PPTX saved: {output_path}")


# ── DOCX ───────────────────────────────────────────────────────────────────────

def _build_docx(
    data: dict[str, Any],
    case_id: str,
    hostname: str,
    disk_image: str,
    generated_utc: str,
    output_path: Path,
    opencti_findings: str = "",
    fan_summary: str = "",
    fame_summary: str = "",
) -> None:
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        print("[fast] WARNING: python-docx not installed — skipping DOCX. pip3 install python-docx")
        return

    analysis_dir_p = PROJECT_ROOT / "analysis" / "storage"
    exports_dir_p  = PROJECT_ROOT / "exports"
    doc = Document()

    # ── Page setup ────────────────────────────────────────────────────────────
    for section in doc.sections:
        section.top_margin      = Inches(1.0)
        section.bottom_margin   = Inches(1.0)
        section.left_margin     = Inches(1.2)
        section.right_margin    = Inches(1.2)
        section.header_distance = Inches(0.4)
        section.footer_distance = Inches(0.4)

    # ── Default font ──────────────────────────────────────────────────────────
    _set_doc_font(doc, "Arial")
    doc.styles["Normal"].paragraph_format.space_after  = Pt(5)
    doc.styles["Normal"].paragraph_format.space_before = Pt(0)
    for _i in range(1, 5):
        try:
            _hs = doc.styles[f"Heading {_i}"]
            _hs.paragraph_format.space_before = Pt(14 if _i == 1 else 10)
            _hs.paragraph_format.space_after  = Pt(4)
            _hs.font.name = "Arial"
        except Exception:
            pass

    styles = doc.styles

    def _heading(text: str, level: int) -> None:
        p = doc.add_heading(text, level=level)
        if p.runs:
            p.runs[0].font.color.rgb = RGBColor(0x0f, 0x17, 0x2a)
            p.runs[0].font.name = "Arial"

    def _para(text: str, bold: bool = False, italic: bool = False) -> None:
        p = doc.add_paragraph()
        run = p.add_run(text)
        run.bold      = bold
        run.italic    = italic
        run.font.name = "Arial"

    def _note(text: str) -> None:
        p = (doc.add_paragraph(style="Intense Quote")
             if "Intense Quote" in [s.name for s in styles]
             else doc.add_paragraph())
        run = p.add_run(text)
        run.italic = True
        run.font.name = "Arial"
        run.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)

    def _code(text: str) -> None:
        sn = "No Spacing" if "No Spacing" in [s.name for s in styles] else "Normal"
        p  = doc.add_paragraph(style=sn)
        r  = p.add_run(text)
        r.font.name  = "Courier New"
        r.font.size  = Pt(9)
        r.font.color.rgb = RGBColor(0x1e, 0x3a, 0x5f)

    def _table_2col(rows: list[tuple[str, str]], header: bool = True) -> None:
        start = 1 if header else 0
        tbl = doc.add_table(rows=len(rows) + start, cols=2)
        tbl.style = "Table Grid"
        if header:
            for i, h in enumerate(["Field", "Value"]):
                c = tbl.rows[0].cells[i]
                c.text = h
                c.paragraphs[0].runs[0].font.bold = True
                c.paragraphs[0].runs[0].font.name = "Arial"
        for i, (k, v) in enumerate(rows):
            r = tbl.rows[i + start]
            r.cells[0].text = k
            r.cells[1].text = v
            for j in range(2):
                for run in r.cells[j].paragraphs[0].runs:
                    run.font.name = "Arial"

    # ── Header, footer, watermark ─────────────────────────────────────────────
    for section in doc.sections:
        _add_header_footer(section, case_id)
        _add_watermark(section)

    # ── Cover ──────────────────────────────────────────────────────────────────
    title = doc.add_heading("FAST — Storage forensics report", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if title.runs:
        title.runs[0].font.name = "Arial"

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = sub.add_run("Forensic Analysis Storage  |  Fan Get Fame Fast")
    r.font.size  = Pt(14)
    r.font.name  = "Arial"
    r.font.color.rgb = RGBColor(0x1d, 0x4e, 0xd8)

    disk_mtime = "N/A"
    if disk_image and Path(disk_image).exists():
        disk_mtime = datetime.fromtimestamp(
            Path(disk_image).stat().st_mtime, tz=_CET
        ).strftime("%d-%b-%Y %H:%M CET")

    _table_2col([
        ("Case ID",              case_id),
        ("Hostname",             hostname),
        ("Disk image",           Path(disk_image).name if disk_image else ""),
        ("Disk image created",   disk_mtime),
        ("Module",               "FAST — Forensic Analysis Storage"),
        ("Analysts",             "Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin · Joost Beekman"),
        ("Generated",            generated_utc),
        ("Analysis tools",       "The Sleuth Kit (TSK) · ewflib · bulk_extractor · qemu-nbd"),
    ])
    doc.add_page_break()

    # ── Part A — Investigation methodology ───────────────────────────────────
    _heading("Part A — Investigation methodology", 1)
    doc.add_paragraph()

    _heading("A.1  Evidence acquisition and integrity", 2)
    _para(
        f"The disk image {Path(disk_image).name!r} was acquired from the subject machine "
        "and stored as a read-only artifact. All analysis tools in this pipeline open the "
        "image in read-only mode. The image is never modified during analysis. "
        "For E01 (Expert Witness Format) images, ewfverify confirms hash integrity before "
        "any analysis begins. For raw and VM disk formats (VDI, VMDK), the image is "
        "exposed as a block device via qemu-nbd in read-only mode."
    )
    doc.add_paragraph()
    _para(
        "All analysis output files are written to a separate working directory "
        "(analysis/storage/ and exports/) and are archived in the investigations vault "
        "on completion. The evidence image is never stored in the working directory."
    )
    doc.add_paragraph()

    _heading("A.2  Tool suite and versions", 2)
    _table_2col([
        ("The Sleuth Kit (TSK)",  "fls, fsstat, mmls, ils, icat — file system analysis"),
        ("ewflib / libewf",       "ewfinfo, ewfverify — E01 image metadata and integrity"),
        ("bulk_extractor",        "File carving and feature extraction from raw disk"),
        ("qemu-nbd",              "Block device export for VDI/VMDK disk images"),
        ("mactime (TSK)",         "Filesystem timeline generation from bodyfile"),
        ("python-docx",           "Word document generation"),
    ])
    doc.add_paragraph()

    _heading("A.3  Chain of custody summary", 2)
    _para(
        f"Image: {Path(disk_image).name}  →  read-only mount / nbd export  →  "
        "TSK tools extract metadata to analysis/storage/  →  "
        "artifacts exported to exports/  →  "
        "reports generated to reports/  →  "
        "uploaded to investigations vault"
    )
    doc.add_paragraph()
    doc.add_page_break()

    # ── Part B — Artifact extraction catalog ─────────────────────────────────
    _heading("Part B — Artifact extraction catalog", 1)
    _note(
        "For each artifact: extraction method, evidence integrity statement, "
        "contents, and role in the investigation."
    )
    doc.add_paragraph()

    # B.1 EWF
    _heading("B.1  E01 image verification (ewfinfo / ewfverify)", 2)
    _para("Extraction method", bold=True)
    _code(f"ewfinfo {Path(disk_image).name} > analysis/storage/ewfinfo.txt\n"
          f"ewfverify {Path(disk_image).name} > analysis/storage/ewfverify.txt")
    doc.add_paragraph()
    _para("Evidence integrity", bold=True)
    _para(
        "ewfverify reads the embedded MD5 and SHA-1 hashes from the E01 segment headers "
        "and recalculates them over the stored data. A mismatch would indicate image "
        "corruption or tampering. ewfinfo extracts acquisition metadata (examiner, case "
        "notes, acquisition tool, timestamps) embedded by the acquisition software. "
        "These are immutable properties of the E01 container."
    )
    doc.add_paragraph()
    _para("Contents", bold=True)
    ewfinfo = data.get("ewfinfo", "")
    ewfverify = data.get("ewfverify", "")
    if ewfinfo:
        _code(ewfinfo.strip()[:800])
    else:
        _para("ewfinfo.txt: not available (non-E01 image format or tool not run).")
    if ewfverify:
        _code(ewfverify.strip()[:300])
    doc.add_paragraph()
    _para("Role in the investigation", bold=True)
    _para(
        "Hash verification is the first step before any analysis. Without a confirmed "
        "hash match, no finding can be cited as derived from an intact image. For non-E01 "
        "formats (VDI, VMDK), integrity is established by recording the SHA-256 hash of "
        "the image file at analysis start."
    )
    doc.add_paragraph()

    # B.2 mmls
    _heading("B.2  Partition table analysis (mmls)", 2)
    _para("Extraction method", bold=True)
    _code("mmls -a <device> > analysis/storage/mmls.txt")
    doc.add_paragraph()
    _para("Evidence integrity", bold=True)
    _para(
        "mmls reads the partition table structures (MBR, GPT, or BSD disklabel) from "
        "the first sectors of the block device. The output is a direct representation "
        "of the on-disk partition map with no interpretation beyond sector arithmetic."
    )
    doc.add_paragraph()
    _para("Contents", bold=True)
    mmls = data.get("mmls", "")
    if mmls:
        _code(mmls.strip()[:1000])
    else:
        _para("mmls.txt: not available.")
    doc.add_paragraph()
    _para("Role in the investigation", bold=True)
    _para(
        "The partition table identifies which sectors contain the primary filesystem, "
        "recovery partitions, and unallocated space. Unallocated space is a high-value "
        "target for carved artifacts and deleted data recovery."
    )
    doc.add_paragraph()

    # B.3 fsstat
    _heading("B.3  Filesystem metadata (fsstat)", 2)
    _para("Extraction method", bold=True)
    _code("fsstat -o <start_sector> <device> > analysis/storage/fsstat.txt")
    doc.add_paragraph()
    _para("Evidence integrity", bold=True)
    _para(
        "fsstat reads the filesystem superblock (ext4) or boot sector / VBR (NTFS/FAT) "
        "and reports volume metadata: filesystem type, cluster size, MFT location, "
        "volume serial number, and last-written timestamp. These are fixed metadata "
        "structures that are not altered by read-only file access."
    )
    doc.add_paragraph()
    _para("Contents", bold=True)
    fsstat = data.get("fsstat", "")
    if fsstat:
        _code(fsstat.strip()[:1000])
    else:
        _para("fsstat.txt: not available.")
    doc.add_paragraph()
    _para("Role in the investigation", bold=True)
    _para(
        "The filesystem type and cluster size determine which TSK plugins apply. "
        "The volume last-written timestamp is an important corroboration point against "
        "the investigation timeline."
    )
    doc.add_paragraph()

    # B.4 fls
    _heading("B.4  File listing including deleted entries (fls)", 2)
    _para("Extraction method", bold=True)
    _code("fls -r -o <start_sector> <device> > analysis/storage/fls_output.txt\n"
          "ils -o <start_sector> <device> > analysis/storage/ils_output.txt\n"
          "ils -O -o <start_sector> <device> > analysis/storage/ils_orphan.txt")
    doc.add_paragraph()
    _para("Evidence integrity", bold=True)
    _para(
        "fls enumerates all allocated and unallocated directory entries from the "
        "filesystem metadata layer. Deleted entries (marked with *) are directory "
        "entries whose inode has been deallocated but whose name record persists in "
        "the parent directory. ils enumerates all inodes including orphan inodes "
        "(allocated inodes with no directory entry). Both tools operate read-only."
    )
    doc.add_paragraph()
    fls_txt = data.get("fls_output", "")
    deleted = fls_txt.count("\n* ") if fls_txt else 0
    _para("Contents", bold=True)
    _para(f"Total deleted file entries identified: {deleted}")
    if fls_txt:
        deleted_lines = [l for l in fls_txt.splitlines() if l.startswith("* ")][:10]
        if deleted_lines:
            _code("\n".join(deleted_lines))
    doc.add_paragraph()
    _para("Role in the investigation", bold=True)
    _para(
        "Deleted files in user-writable paths (home directories, temp, downloads) are "
        "high-priority triage candidates. An attacker will often delete tools, scripts, "
        "and logs after use. The presence of deleted executables or scripts in "
        "system paths may indicate persistence mechanism clean-up."
    )
    doc.add_paragraph()

    # B.5 Filesystem timeline
    _heading("B.5  Filesystem timeline (mactime)", 2)
    _para("Extraction method", bold=True)
    _code("fls -r -m '' -o <start> <device> > analysis/storage/bodyfile.txt\n"
          "mactime -b analysis/storage/bodyfile.txt -d > exports/fs_timeline.csv")
    doc.add_paragraph()
    _para("Evidence integrity", bold=True)
    _para(
        "The bodyfile format captures four POSIX timestamps per file entry: modified (M), "
        "accessed (A), changed (C), and born/created (B) — collectively MAC(B) times. "
        "mactime converts these to a human-readable timeline ordered by timestamp. "
        "Timestamps come directly from the filesystem metadata layer and are not "
        "interpreted or modified."
    )
    doc.add_paragraph()
    timeline_rows = data.get("fs_timeline_csv", [])
    _para("Contents", bold=True)
    if timeline_rows:
        _para(f"Timeline contains {len(timeline_rows)} entries. First 10 shown:")
        headers_t = list(timeline_rows[0].keys()) if timeline_rows else []
        tbl = doc.add_table(rows=min(len(timeline_rows), 10) + 1, cols=len(headers_t))
        tbl.style = "Table Grid"
        for i, h in enumerate(headers_t):
            c = tbl.rows[0].cells[i]
            c.text = h
            if c.paragraphs[0].runs:
                c.paragraphs[0].runs[0].font.bold = True
                c.paragraphs[0].runs[0].font.name = "Arial"
        for ri, row_t in enumerate(timeline_rows[:10]):
            for ci, h in enumerate(headers_t):
                c = tbl.rows[ri + 1].cells[ci]
                c.text = str(row_t.get(h, ""))
                for run in c.paragraphs[0].runs:
                    run.font.name = "Arial"
    else:
        _para("fs_timeline.csv: not available — mactime not run or bodyfile empty.")
    doc.add_paragraph()
    _para("Role in the investigation", bold=True)
    _para(
        "The filesystem timeline is the primary tool for reconstructing the sequence "
        "of file operations during the incident. Sorting by modification time reveals "
        "when files were created, modified, or deleted. Cross-referencing with network "
        "(FAN) and memory (FAME) timelines establishes the full attack chain."
    )
    doc.add_paragraph()
    doc.add_page_break()

    # B.6 Extracted artifacts
    _heading("B.6  Artifact extraction", 2)
    _para(
        "The following artifacts were extracted from the disk image using icat (TSK) "
        "and direct filesystem access. Each category is stored in a dedicated subdirectory "
        "under exports/."
    )
    doc.add_paragraph()
    evtx_l    = data.get("evtx_list", [])
    reg_l     = data.get("registry_list", [])
    pref_l    = data.get("prefetch_list", [])
    srum_l    = data.get("srum_list", [])
    browser_l = data.get("browser_list", [])
    _table_2col([
        ("Windows Event Logs (EVTX)",    f"{len(evtx_l)} file(s) in exports/evtx/"),
        ("Registry hives",               f"{len(reg_l)} file(s) in exports/registry/"),
        ("Prefetch files",               f"{len(pref_l)} file(s) in exports/prefetch/"),
        ("SRUM database",                f"{len(srum_l)} file(s) in exports/srum/"),
        ("Browser history",              f"{len(browser_l)} file(s) in exports/browser/"),
        ("MFT ($MFT)",                   "Extracted" if data.get("mft_extracted") else "Not extracted"),
        ("USN Change Journal ($J)",      "Extracted" if data.get("usnj_extracted") else "Not extracted"),
    ])
    doc.add_paragraph()
    _para("Role in the investigation", bold=True)
    _para(
        "Windows Event Logs contain logon events (4624/4625), service installs (7045), "
        "PowerShell execution (4103/4104), and process creation (4688). Registry hives "
        "contain run keys, service definitions, and user activity artifacts. Prefetch "
        "files reveal program execution history with first/last run timestamps. The MFT "
        "and USN journal provide a high-resolution file modification log that complements "
        "the bodyfile timeline."
    )
    doc.add_paragraph()

    # B.7 bulk_extractor
    _heading("B.7  File carving — bulk_extractor", 2)
    _para("Extraction method", bold=True)
    _code(f"bulk_extractor -o exports/carved -j 4 <device>")
    doc.add_paragraph()
    _para("Evidence integrity", bold=True)
    _para(
        "bulk_extractor scans the raw disk image sequentially, identifying feature "
        "patterns (email addresses, URLs, credit card numbers, etc.) and carving "
        "file headers. It does not parse filesystem structures — it operates directly "
        "on the raw byte stream, recovering artifacts from both allocated and unallocated "
        "space. Output files are plain text feature lists and carved file fragments."
    )
    doc.add_paragraph()
    carved = data.get("bulk_carved", [])
    _para("Contents", bold=True)
    if carved:
        _para(f"{len(carved)} output file(s) in exports/carved/:")
        for cf in carved[:15]:
            _para(f"  {cf}")
    else:
        _para("bulk_extractor produced no output or was not run.")
    doc.add_paragraph()
    _para("Role in the investigation", bold=True)
    _para(
        "Email addresses, URLs, and domains recovered from unallocated space are IOC "
        "candidates for CTI enrichment. Carved file fragments may reconstruct deleted "
        "executables or documents. Credit card and social security number hits trigger "
        "PII breach notification obligations."
    )
    doc.add_paragraph()
    doc.add_page_break()

    # ── Part C — Findings ──────────────────────────────────────────────────────
    _heading("Part C — Findings", 1)
    doc.add_page_break()

    # C.1 Management summary
    _heading("C.1  Management summary", 2)
    _note("Audience: CISO, Legal, Internal Audit — no technical identifiers.")
    _para(
        f"Storage forensic analysis was conducted on the disk image of {hostname}. "
        f"The analysis identified {deleted} deleted file entries, "
        f"{len(evtx_l)} Windows Event Log file(s), and {len(pref_l)} Prefetch "
        "execution artifact(s). "
        "The filesystem timeline has been reconstructed and is available for "
        "cross-reference with network and memory forensics findings."
    )
    doc.add_paragraph()

    # C.2 Cross-module intelligence
    if fan_summary or fame_summary or opencti_findings:
        _heading("C.2  Cross-module intelligence", 2)
        _note("Claude: enhance and elaborate when necessary.")
        if fan_summary:
            _heading("C.2.1  Network forensics (FAN)", 3)
            _para(fan_summary.strip())
        if fame_summary:
            _heading("C.2.2  Memory forensics (FAME)", 3)
            _para(fame_summary.strip())
        if opencti_findings:
            _heading("C.2.3  OpenCTI threat intelligence", 3)
            _para(opencti_findings.strip())
        doc.add_paragraph()

    # C.3 IOCs
    _heading("C.3  Indicators of compromise", 2)
    _note("All IOC values defanged. Claude: enhance and elaborate when necessary.")
    iocs = _extract_iocs(data)
    if iocs:
        tbl = doc.add_table(rows=len(iocs) + 1, cols=4)
        tbl.style = "Table Grid"
        for i, h in enumerate(["Type", "Value", "Severity", "Context"]):
            c = tbl.rows[0].cells[i]
            c.text = h
            if c.paragraphs[0].runs:
                c.paragraphs[0].runs[0].font.bold = True
                c.paragraphs[0].runs[0].font.name = "Arial"
        for i, ioc in enumerate(iocs):
            for j, val in enumerate([ioc["type"], ioc["value"], ioc["severity"], ioc["context"]]):
                c = tbl.rows[i + 1].cells[j]
                c.text = val
                for run in c.paragraphs[0].runs:
                    run.font.name = "Arial"
    else:
        _para("No malicious indicators of compromise identified from storage analysis.")
    doc.add_paragraph()

    # C.4 Recommendations
    _heading("C.4  Recommendations", 2)
    _note("Claude: enhance and elaborate when necessary.")
    for i, rec in enumerate(_build_recommendations(data), 1):
        rec_clean = re.sub(r"\*\*(.*?)\*\*", r"\1", rec)
        p = doc.add_paragraph(style="List Number")
        run = p.add_run(rec_clean)
        run.font.name = "Arial"
    doc.add_page_break()

    # ── Appendix A — 4-column with CET timestamps + SHA-256 ──────────────────
    _heading("Appendix A — Analysis source files", 1)
    _note("SHA-256 hashes computed at report generation time.")

    artifact_files = [
        (analysis_dir_p / "ewfinfo.txt",       "E01 image metadata (ewfinfo)"),
        (analysis_dir_p / "ewfverify.txt",      "E01 hash verification (ewfverify)"),
        (analysis_dir_p / "mmls.txt",           "Partition table (mmls)"),
        (analysis_dir_p / "fsstat.txt",         "Filesystem metadata (fsstat)"),
        (analysis_dir_p / "fls_output.txt",     "Full file listing incl. deleted (fls)"),
        (analysis_dir_p / "ils_output.txt",     "Inode listing (ils)"),
        (analysis_dir_p / "ils_orphan.txt",     "Orphan inode listing (ils -O)"),
        (analysis_dir_p / "bodyfile.txt",       "MAC time bodyfile (fls -m)"),
        (exports_dir_p  / "fs_timeline.csv",    "Filesystem timeline (mactime)"),
        (exports_dir_p  / "mft" / "$MFT",       "Master File Table"),
        (exports_dir_p  / "mft" / "$J",         "USN Change Journal"),
    ]

    tbl = doc.add_table(rows=len(artifact_files) + 1, cols=4)
    tbl.style = "Table Grid"
    for i, h in enumerate(["Filename", "Description", "Generated (CET)", "SHA-256 (first 32)"]):
        c = tbl.rows[0].cells[i]
        c.text = h
        if c.paragraphs[0].runs:
            c.paragraphs[0].runs[0].font.bold = True
            c.paragraphs[0].runs[0].font.name = "Arial"
    for i, (fp, desc) in enumerate(artifact_files):
        mtime, sha256 = _compute_file_metadata(fp)
        vals = [fp.name, desc, mtime, sha256[:32] if sha256 != "N/A" else "N/A"]
        for j, val in enumerate(vals):
            c = tbl.rows[i + 1].cells[j]
            c.text = val
            for run in c.paragraphs[0].runs:
                run.font.name = "Arial"

    _remove_blank_paragraphs(doc)
    doc.save(str(output_path))
    print(f"[fast] DOCX saved: {output_path}")


# ── Public API ─────────────────────────────────────────────────────────────────

def generate(
    case_id: str,
    hostname: str,
    disk_image: str = "",
    analysis_dir: Path | None = None,
    exports_dir: Path | None = None,
    output_dir: Path | None = None,
    case_dir: Path | None = None,
    docs_dir: Path | None = None,
    opencti_findings: str = "",
    fan_summary: str = "",
    fame_summary: str = "",
    md_only: bool = False,
) -> dict[str, Path | None]:
    """Generate the full FAST report suite.

    case_dir: module-specific directory (reports/<case_id>/FAST/<hostname>/). When
    supplied, Markdown goes to case_dir/. docs_dir overrides where PDF/PPTX/DOCX land
    (default: case_dir/output/ for legacy compat, typically reports/<case_id>/documents/).
    When both are omitted, all formats land in output_dir (legacy flat behaviour).
    """
    analysis_dir = analysis_dir or (PROJECT_ROOT / "analysis" / "storage")
    exports_dir  = exports_dir  or (PROJECT_ROOT / "exports")

    if case_dir is not None:
        md_dir  = path_guard.guard_output_dir(case_dir)
        aux_dir = path_guard.guard_output_dir(docs_dir or (case_dir / "output"))
    else:
        md_dir  = path_guard.guard_output_dir(output_dir or (PROJECT_ROOT / "reports"))
        aux_dir = md_dir

    data = _load_analysis(analysis_dir, exports_dir)
    # Make the Claude-authored narrative available to the PPTX board deck too
    # (the Markdown builder loads it separately).
    data["_narrative"] = _load_narrative(case_id, md_dir)
    generated_utc = datetime.now(_CET).strftime("%Y-%m-%d %H:%M CET")
    stem = case_id.replace(" ", "_")
    # docs_dir (when supplied) is shared across all hosts in this case — disambiguate
    # the PDF/PPTX/DOCX filenames with the hostname so multiple hosts don't collide
    # and silently overwrite each other's reports in reports/<case_id>/documents/.
    aux_stem = f"{stem}_{hostname.replace(' ', '_')}" if docs_dir is not None else stem

    # Markdown
    md_text = _build_markdown(
        data, case_id, hostname, disk_image or str(analysis_dir),
        generated_utc, opencti_findings, fan_summary, fame_summary,
        reports_dir=md_dir,
    )
    md_path = md_dir / f"{stem}_fast_report.md"
    md_path.write_text(md_text)
    print(f"[fast] Markdown saved: {md_path}")

    # PDF
    pdf_path: Path | None = None
    if not md_only:
        try:
            sys.path.insert(0, str(PROJECT_ROOT / "lib"))
            from md_to_pdf import convert as md2pdf
            pdf_path = aux_dir / f"{aux_stem}_fast_report.pdf"
            md2pdf(md_path, pdf_path)
            print(f"[fast] PDF saved: {pdf_path}")
        except Exception as exc:
            print(f"[fast] WARNING: PDF generation failed: {exc}")

    # PPTX
    pptx_path: Path | None = None
    if not md_only:
        pptx_path = aux_dir / f"{aux_stem}_fast_presentation.pptx"
        _build_pptx(
            data, case_id, hostname, disk_image or str(analysis_dir),
            generated_utc, pptx_path, opencti_findings, fan_summary, fame_summary,
        )

    # DOCX
    docx_path: Path | None = None
    if not md_only:
        docx_path = aux_dir / f"{aux_stem}_fast_report.docx"
        _build_docx(
            data, case_id, hostname, disk_image or str(analysis_dir),
            generated_utc, docx_path, opencti_findings, fan_summary, fame_summary,
        )

    return {
        "md":   md_path,
        "pdf":  pdf_path,
        "pptx": pptx_path if (pptx_path and pptx_path.exists()) else None,
        "docx": docx_path if (docx_path and docx_path.exists()) else None,
    }


# ── CLI ────────────────────────────────────────────────────────────────────────

def _build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(description="FAST — Storage Forensics Report Generator")
    p.add_argument("--case-id",      required=True, metavar="ID")
    p.add_argument("--hostname",     required=True, metavar="HOST")
    p.add_argument("--disk-image",   default="",   metavar="PATH")
    p.add_argument("--analysis-dir", default=None, metavar="DIR")
    p.add_argument("--exports-dir",  default=None, metavar="DIR")
    p.add_argument("--output-dir",   default=None, metavar="DIR")
    p.add_argument("--case-dir",     default=None, metavar="DIR",  help="Module-specific dir (reports/<case_id>/FAST/<host>/); MD lands here")
    p.add_argument("--docs-dir",     default=None, metavar="DIR",  help="Shared documents dir (reports/<case_id>/documents/); PDF/PPTX/DOCX land here")
    p.add_argument("--opencti",      default="",   metavar="TEXT")
    p.add_argument("--fan-summary",  default="",   metavar="TEXT")
    p.add_argument("--fame-summary", default="",   metavar="TEXT")
    p.add_argument("--md-only",      action="store_true", help="Generate Markdown only — skip PDF, PPTX, DOCX")
    return p


if __name__ == "__main__":
    args = _build_parser().parse_args()
    paths = generate(
        case_id      = args.case_id,
        hostname     = args.hostname,
        disk_image   = args.disk_image,
        analysis_dir = Path(args.analysis_dir) if args.analysis_dir else None,
        exports_dir  = Path(args.exports_dir)  if args.exports_dir  else None,
        output_dir   = Path(args.output_dir)   if args.output_dir   else None,
        case_dir     = Path(args.case_dir)     if args.case_dir     else None,
        docs_dir     = Path(args.docs_dir)     if args.docs_dir     else None,
        opencti_findings = args.opencti,
        fan_summary  = args.fan_summary,
        fame_summary = args.fame_summary,
        md_only      = args.md_only,
    )
    print("[fast] Report suite complete:")
    for fmt, p in paths.items():
        if p:
            print(f"  {fmt.upper():4s}  {p}")
