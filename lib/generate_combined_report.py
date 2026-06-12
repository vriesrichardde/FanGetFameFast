#!/usr/bin/env python3
# SPDX-License-Identifier: MIT OR Apache-2.0
# SPDX-FileCopyrightText: 2026 Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin · Joost Beekman
"""
generate_combined_report.py — Unified FAN + FAME + FAST campaign report generator.

DEPRECATED for interactive campaign reports. The per-case campaign report
(`<case_id>_campaign_report.*`) is, per CLAUDE.md, a "carefully constructed
synthesis" — Claude hand-authors it against `docs/campaign_report_template.md`
and renders it via `lib/render_campaign_report.render()`. This module's
`generate()` produces a structurally generic report (its Cross-Domain
Correlation section just embeds `<case_id>_correlation.md` verbatim, and its
Hallucination Guard is a sparse per-module summary, not a curated
finding-by-finding list) and is retained only as an automated fallback for
`--md-only`/headless batch runs or very-low-evidence cases.

The table-merge helpers below (`_merge_mitre`, `_mitre_table_md`,
`_merge_iocs`, `_ioc_tables_md`, `_merge_recommendations`,
`_extract_summary`, `_extract_mitre`, `_extract_iocs`,
`_extract_recommendations`) remain importable utilities for hand-authoring
campaign reports — use them to pre-populate the unified MITRE/IOC/
recommendations tables in Sections 4, 5, and the recommendations section of
the template before refining by hand.

When multiple module reports exist for the same case ID, this module merges them
into a single campaign report (Markdown, PDF, PPTX, DOCX): a management summary,
cross-domain correlation, deduplicated MITRE ATT&CK coverage, deduplicated
indicators of compromise, per-module technical summaries (linking to the full
module reports rather than re-embedding them), unified recommendations, and a
combined Hallucination Guard confidence assessment.

Usage (CLI):
    python3 lib/generate_combined_report.py \\
        --case-id CASE-2026-001 \\
        --hostname SERVER1234 \\
        [--reports-dir ./reports] \\
        [--output-dir ./reports]

Python API:
    from lib.generate_combined_report import generate
    paths = generate(case_id="CASE-2026-001", hostname="SERVER1234")
"""
from __future__ import annotations

import argparse
import os
import re
import sys
from datetime import datetime, timezone
try:
    from zoneinfo import ZoneInfo
    _CET = ZoneInfo("Europe/Amsterdam")
except ImportError:
    _CET = timezone.utc
from pathlib import Path
sys.path.insert(0, str(Path(__file__).resolve().parent))
import path_guard  # noqa: E402  write-path policy enforcement

PROJECT_ROOT = Path(__file__).parent.parent

try:
    from hallucination_guard import (
        ConfidenceTier,
        tag_finding,
        render_confidence_summary,
        reset_counter as _hg_reset,
    )
except ModuleNotFoundError:
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    from hallucination_guard import (
        ConfidenceTier,
        tag_finding,
        render_confidence_summary,
        reset_counter as _hg_reset,
    )

_DARK_NAVY  = (0x0f, 0x17, 0x2a)
_MID_NAVY   = (0x1e, 0x3a, 0x5f)
_BLUE       = (0x1d, 0x4e, 0xd8)
_LIGHT_BLUE = (0x93, 0xc5, 0xfd)
_WHITE      = (0xff, 0xff, 0xff)
_LIGHT_BG   = (0xf8, 0xfa, 0xfc)
_ROW_ALT    = (0xf1, 0xf5, 0xf9)
_TEXT_DARK  = (0x1f, 0x29, 0x37)
_TEXT_MID   = (0x6b, 0x72, 0x80)
_AMBER      = (0xfb, 0xbf, 0x24)
_GREEN      = (0x22, 0xc5, 0x5e)

_MODULE_LABELS = {
    "FAN":  "FAN (Network)",
    "FAME": "FAME (Memory)",
    "FAST": "FAST (Storage)",
}
_MODULE_SECTION_TITLES = {
    "FAN":  "FAN — Network Forensics",
    "FAME": "FAME — Memory Forensics",
    "FAST": "FAST — Storage Forensics",
}


# ── Source discovery ───────────────────────────────────────────────────────────

_RESERVED_DIRS = {"FAN", "FAME", "FAST", "documents", "raw", "output", "final"}


def _find_flat(search_dirs: list[Path], filenames: list[str]) -> tuple[str, Path] | None:
    """Search for an exact filename in a flat (legacy) layout."""
    for d in search_dirs:
        for fn in filenames:
            candidate = d / fn
            if candidate.is_file():
                return candidate.read_text(errors="replace"), candidate
    return None


def _find_module_report(case_root: Path, module_dirname: str, glob_patterns: list[str]) -> tuple[str, Path] | None:
    """Search case_root/<module_dirname>/<investigation>/ for a file matching
    any of the given glob patterns (e.g. "*_fast_report.md"), regardless of
    whether the filename is prefixed with the case ID or the investigation
    name."""
    mod_dir = case_root / module_dirname
    if not mod_dir.is_dir():
        return None
    for inv_dir in sorted(mod_dir.iterdir()):
        if not inv_dir.is_dir():
            continue
        for pattern in glob_patterns:
            matches = sorted(inv_dir.glob(pattern))
            if matches:
                return matches[0].read_text(errors="replace"), matches[0]
    return None


def _find_legacy_2level(case_root: Path, suffixes: list[str]) -> tuple[str, Path] | None:
    """Search case_root/<sub>/ (excluding module/reserved dirs) for files
    ending in one of the given suffixes — pre-3-level-layout investigations."""
    for sub in sorted(case_root.iterdir()):
        if not sub.is_dir() or sub.name in _RESERVED_DIRS:
            continue
        for suf in suffixes:
            for candidate in sorted(sub.glob(f"*{suf}")):
                return candidate.read_text(errors="replace"), candidate
    return None


def _find_via_narrative(search_dirs: list[Path], stem: str) -> tuple[str, Path] | None:
    """Legacy: a fan_narrative file records the PCAP stem in its source comment,
    used to locate the actual FAN report when it was triggered cross-module and
    named after the PCAP rather than the case ID."""
    for d in search_dirs:
        narrative = d / f"{stem}_fan_narrative.md"
        if not narrative.is_file():
            continue
        text = narrative.read_text(errors="replace")
        m = re.search(r"<!-- source: .*?exports/([A-Za-z0-9_\-]+)", text)
        if m:
            pcap_stem = m.group(1)
            found = _find_flat(search_dirs, [f"{pcap_stem}_incident_report.md", f"{pcap_stem}_fan_report.md"])
            if found:
                return found
        return text, narrative
    return None


def _discover_sources(reports_dir: Path, case_id: str) -> dict[str, object]:
    """
    Find existing module reports and the cross-module correlation report for
    this case. Returns a dict with any of the following keys present:
    fan_md/fan_path, fame_md/fame_path, fast_md/fast_path,
    correlation_md/correlation_path.

    Supports both the legacy flat layout (all files in reports_dir) and the
    3-level case layout (reports_dir/<case_id>/{FAN,FAME,FAST}/<investigation>/).
    """
    stem = case_id.replace(" ", "_")
    sources: dict[str, object] = {}
    case_root = reports_dir / stem
    search_dirs = [case_root, reports_dir] if case_root.is_dir() else [reports_dir]

    # FAN
    found = _find_flat(search_dirs, [f"{stem}_incident_report.md", f"{stem}_fan_report.md"])
    if not found and case_root.is_dir():
        found = _find_module_report(case_root, "FAN", ["*_incident_report.md", "*_fan_report.md"])
    if not found and case_root.is_dir():
        found = _find_legacy_2level(case_root, ["_incident_report.md", "_fan_report.md"])
    if not found:
        found = _find_via_narrative(search_dirs, stem)
    if found:
        sources["fan_md"], sources["fan_path"] = found

    # FAME
    found = _find_flat(search_dirs, [f"{stem}_fame_report.md"])
    if not found and case_root.is_dir():
        found = _find_module_report(case_root, "FAME", ["*_fame_report.md"])
    if found:
        sources["fame_md"], sources["fame_path"] = found

    # FAST
    found = _find_flat(search_dirs, [f"{stem}_fast_report.md"])
    if not found and case_root.is_dir():
        found = _find_module_report(case_root, "FAST", ["*_fast_report.md"])
    if found:
        sources["fast_md"], sources["fast_path"] = found

    # Cross-module correlation (computed by lib/correlate_findings.py)
    found = _find_flat(search_dirs, [f"{stem}_correlation.md"])
    if not found and case_root.is_dir():
        for module_dirname in ("FAST", "FAME", "FAN"):
            found = _find_module_report(case_root, module_dirname, ["*_correlation.md"])
            if found:
                break
    if found:
        sources["correlation_md"], sources["correlation_path"] = found

    return sources


def _strip_directives(text: str) -> str:
    """Drop leaked authoring instructions (`> Claude: ...`) and the per-module
    `> **Audience:** ...` line — the campaign report states the audience once
    at the top of each unified section."""
    lines = []
    for line in text.splitlines():
        stripped = line.strip()
        if stripped.startswith(">") and ("claude" in stripped.lower() or stripped.lower().startswith("> **audience")):
            continue
        lines.append(line)
    return "\n".join(lines).strip()


def _extract_section(md_text: str, section_marker: str, max_chars: int = 2000) -> str:
    """Extract a named section from a Markdown report.

    Stops at the next same-level (## ) or higher (# ) heading so that adjacent
    sections do not bleed into the extracted text.
    """
    lines = md_text.splitlines()
    in_section = False
    collected: list[str] = []
    for line in lines:
        if section_marker.lower() in line.lower() and line.startswith("#"):
            in_section = True
            continue
        if in_section:
            # Stop at any ## or # level heading (but not ### or deeper)
            if re.match(r"^#{1,2} ", line):
                break
            collected.append(line)
    result = _strip_directives("\n".join(collected))
    return result[:max_chars] if len(result) > max_chars else result


def _extract_summary(md_text: str) -> str:
    """Pull the Management Summary section from any module report."""
    return _extract_section(md_text, "Management Summary", max_chars=2000)


def _extract_mitre(md_text: str) -> str:
    """Pull the MITRE ATT&CK section from any module report."""
    return _extract_section(md_text, "MITRE ATT", max_chars=4000)


def _extract_iocs(md_text: str) -> str:
    """Pull the IOC section from any module report."""
    result = _extract_section(md_text, "IOC Reference", max_chars=4000)
    if not result:
        result = _extract_section(md_text, "Indicators of Compromise", max_chars=4000)
    return result


def _extract_recommendations(md_text: str) -> str:
    """Pull the Recommendations section from any module report."""
    return _extract_section(md_text, "Recommendations", max_chars=3000)


# ── Markdown table parsing & cross-module merging ───────────────────────────────

_TABLE_ROW_RE = re.compile(r"^\|(.+)\|\s*$")
_TID_RE = re.compile(r"\[(T\d[\d.]*)\]")


def _is_separator_row(cells: list[str]) -> bool:
    return all(re.fullmatch(r":?-{2,}:?", c) for c in cells)


def _parse_table_rows(md_text: str) -> list[list[str]]:
    """Parse all pipe-table data rows from a markdown fragment, skipping
    separator rows (|---|---|)."""
    rows: list[list[str]] = []
    for line in md_text.splitlines():
        line = line.strip()
        m = _TABLE_ROW_RE.match(line)
        if not m:
            continue
        cells = [c.strip() for c in m.group(1).split("|")]
        if _is_separator_row(cells):
            continue
        rows.append(cells)
    return rows


def _iter_tables(md_text: str):
    """Yield (heading, rows) for every markdown table in *md_text*, where
    *heading* is the nearest preceding ###/#### heading (or "" if the table
    sits directly under the extracted section) and *rows* is a list of dicts
    keyed by the table's lower-cased header column names. Different module
    report generators use different column layouts for the same logical
    table (e.g. FAN's MITRE table has a `Triggered By` column where FAST's has
    `Observation`) — keying by header name lets callers handle both."""
    lines = md_text.splitlines()
    current_heading = ""
    i, n = 0, len(lines)
    while i < n:
        stripped = lines[i].strip()
        heading_m = re.match(r"^#{2,4}\s+(.*)", stripped)
        if heading_m:
            current_heading = heading_m.group(1).strip()
            i += 1
            continue
        if _TABLE_ROW_RE.match(stripped):
            table_lines = []
            while i < n and _TABLE_ROW_RE.match(lines[i].strip()):
                table_lines.append(lines[i].strip())
                i += 1
            table_rows = _parse_table_rows("\n".join(table_lines))
            if table_rows:
                header = [h.strip().lower() for h in table_rows[0]]
                data = [dict(zip(header, r)) for r in table_rows[1:]]
                yield current_heading, data
            continue
        i += 1


def _strip_md(text: str) -> str:
    """Strip markdown links/emphasis/code spans for use in plain-text contexts (PPTX)."""
    text = re.sub(r"\[([^\]]*)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"\*(.*?)\*", r"\1", text)
    text = re.sub(r"`(.*?)`", r"\1", text)
    return text


def _merge_mitre(module_texts: dict[str, str]) -> list[dict]:
    """Merge per-module MITRE ATT&CK tables into one row per technique,
    deduplicated by technique ID and annotated with the module(s) that
    observed it. Tolerates the differing column layouts used by the FAN
    (Technique | Name | Tactic | Severity | Triggered By) and FAME/FAST
    (Technique | Name | Tactic | Observation) report generators."""
    merged: dict[str, dict] = {}
    order: list[str] = []
    for label, md_text in module_texts.items():
        for _heading, rows in _iter_tables(_extract_mitre(md_text)):
            for row in rows:
                tid_raw = row.get("technique", "")
                if not tid_raw:
                    continue
                m = _TID_RE.search(tid_raw)
                key = m.group(1) if m else tid_raw
                name = row.get("name", "").strip()
                tactic = row.get("tactic", "").strip()
                if "observation" in row:
                    observation = row["observation"].strip()
                else:
                    parts = [row.get("severity", "").strip(), row.get("triggered by", "").strip()]
                    observation = " — ".join(p for p in parts if p)
                entry = merged.setdefault(key, {
                    "id_md": tid_raw, "name": name, "tactic": tactic,
                    "modules": [], "observations": {},
                })
                if not entry["name"] and name:
                    entry["name"] = name
                if not entry["tactic"] and tactic:
                    entry["tactic"] = tactic
                if key not in order:
                    order.append(key)
                if label not in entry["modules"]:
                    entry["modules"].append(label)
                if observation:
                    entry["observations"][label] = observation
    rows_out: list[dict] = []
    for key in order:
        e = merged[key]
        observation = " ".join(f"**{m}:** {o}" for m, o in e["observations"].items())
        rows_out.append({
            "id_md": e["id_md"], "name": e["name"], "tactic": e["tactic"],
            "modules": " + ".join(e["modules"]), "observation": observation,
        })
    return rows_out


def _mitre_table_md(rows: list[dict]) -> str:
    if not rows:
        return ""
    lines = [
        "| Technique | Name | Tactic | Module(s) | Observation |",
        "|-----------|------|--------|-----------|-------------|",
    ]
    for r in rows:
        lines.append(f"| {r['id_md']} | {r['name']} | {r['tactic']} | {r['modules']} | {r['observation']} |")
    return "\n".join(lines)


def _merge_iocs(module_texts: dict[str, str]) -> tuple[dict[str, list[dict]], list[str]]:
    """Merge per-module IOC tables into one table per category, deduplicated
    by IOC value. Tolerates the differing layouts used by FAST (IOCs grouped
    under ### category headings, columns Value | Confidence | Source) and FAN
    (a single flat table: Severity | Type | Value | Category | Source). IOCs
    confirmed by more than one module are upgraded to
    "CONFIRMED (multi-module)"."""
    categories: dict[str, dict[str, dict]] = {}
    cat_order: list[str] = []
    for label, md_text in module_texts.items():
        for heading, rows in _iter_tables(_extract_iocs(md_text)):
            for row in rows:
                value = row.get("value", "").strip()
                if not value:
                    continue
                if heading and "confidence" in row:
                    cat = heading
                    confidence = row.get("confidence", "").strip()
                    source = row.get("source", "").strip()
                else:
                    # FAN-style flat table: Severity | Type | Value | Category | Source
                    type_ = row.get("type", "").strip().lower()
                    cat = "Network Indicators" if type_ in ("ip", "domain", "url", "fqdn") else "Other Indicators"
                    confidence = row.get("confidence", "").strip() or row.get("severity", "").strip()
                    source = " / ".join(p for p in (row.get("source", "").strip(), row.get("category", "").strip()) if p)
                cat_entries = categories.setdefault(cat, {})
                if cat not in cat_order:
                    cat_order.append(cat)
                key = value.lower()
                entry = cat_entries.setdefault(key, {"value": value, "confidence": confidence, "sources": []})
                if source:
                    entry["sources"].append(f"{label} {source}")
                modules_seen = {s.split()[0] for s in entry["sources"]}
                if len(modules_seen) > 1:
                    entry["confidence"] = "CONFIRMED (multi-module)"
    return {cat: list(entries.values()) for cat, entries in categories.items()}, cat_order


def _ioc_tables_md(merged: dict[str, list[dict]], cat_order: list[str]) -> str:
    blocks: list[str] = []
    for cat in cat_order:
        entries = merged.get(cat) or []
        if not entries:
            continue
        lines = [f"### {cat}", "", "| Value | Confidence | Source |", "|-------|------------|--------|"]
        for e in entries:
            lines.append(f"| {e['value']} | {e['confidence']} | {'; '.join(e['sources'])} |")
        blocks.append("\n".join(lines))
    return "\n\n".join(blocks)


def _merge_recommendations(module_texts: dict[str, str]) -> list[str]:
    """Flatten and deduplicate the recommendation lists from each module report."""
    recs: list[str] = []
    for md_text in module_texts.values():
        section = _extract_recommendations(md_text)
        for line in section.splitlines():
            line = line.strip()
            if line and (line[0].isdigit() or line.startswith("-")):
                cleaned = re.sub(r"^\d+\.\s*", "", re.sub(r"^-\s*", "", line))
                cleaned = re.sub(r"\*\*(.*?)\*\*", r"\1 ", cleaned)
                cleaned = re.sub(r"\s+", " ", cleaned).strip()
                if cleaned and not re.fullmatch(r"-+", cleaned) and cleaned not in recs:
                    recs.append(cleaned)
    return recs


def _demote_headings(text: str, levels: int = 1) -> str:
    """Add *levels* extra `#` to every markdown heading, so a document fragment
    can be nested under a higher-level section without colliding heading
    levels."""
    out = []
    for line in text.splitlines():
        m = re.match(r"^(#{1,5})(\s+.*)$", line)
        if m:
            out.append("#" * (len(m.group(1)) + levels) + m.group(2))
        else:
            out.append(line)
    return "\n".join(out)


def _clean_correlation_body(correlation_md: str) -> str:
    """Strip the metadata header, any trailing Hallucination Guard subsection
    (the campaign report has its own combined one), and any leaked authoring
    directives from a lib/correlate_findings.py output, then demote its
    headings so they nest under the campaign report's "Cross-Domain
    Correlation" section."""
    lines = correlation_md.splitlines()
    past_meta = False
    out: list[str] = []
    for line in lines:
        if not past_meta:
            if line.strip() == "---":
                past_meta = True
            continue
        if "Claude:" in line:
            continue
        out.append(line)
    text = "\n".join(out)
    text = re.split(r"\n##\s+Hallucination Guard", text)[0]
    text = _demote_headings(text.strip(), levels=1)
    # Drop redundant standalone "---" separators left over from the source doc;
    # the campaign report's own collapsing pass tidies the rest.
    text = re.sub(r"\n-{3,}\n", "\n\n", text)
    return text.strip()


def _final_cleanup(text: str) -> str:
    """Collapse runs of blank lines and repeated "---" separators left over
    from concatenating cleaned-up fragments of several source documents."""
    text = re.sub(r"\n{3,}", "\n\n", text)
    while True:
        new_text = re.sub(r"\n---\n\n---\n", "\n---\n", text)
        if new_text == text:
            break
        text = new_text
    return re.sub(r"\n{3,}", "\n\n", text)


# ── Markdown ───────────────────────────────────────────────────────────────────

def _build_markdown(
    sources: dict[str, object],
    case_id: str,
    hostname: str,
    generated_utc: str,
    case_dir: Path,
) -> str:
    """Build the unified campaign report."""
    module_texts: dict[str, str] = {}
    for m in ("FAN", "FAME", "FAST"):
        text = sources.get(f"{m.lower()}_md")
        if text:
            module_texts[m] = text
    modules_run = list(module_texts.keys())

    lines: list[str] = []
    a = lines.append

    a(f"# {case_id} — Campaign Forensics Report")
    a("")
    a("| Field | Value |")
    a("|-------|-------|")
    a(f"| Case ID | `{case_id}` |")
    a(f"| Hostname | `{hostname}` |")
    a(f"| Modules | {', '.join(_MODULE_LABELS[m] for m in modules_run) if modules_run else 'None detected'} |")
    a(f"| Generated (UTC) | {generated_utc} |")
    a("| Prepared by | Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin · Joost Beekman |")
    a("")

    # ── 1. Management Summary ─────────────────────────────────────────────────
    a("---")
    a("")
    a("## 1. Management Summary")
    a("")
    a("> **Audience:** CISO, Legal, Law Enforcement, Internal Audit — no technical identifiers.")
    a("")
    if modules_run:
        for m in modules_run:
            summary = _extract_summary(module_texts[m])
            if summary:
                a(f"**{_MODULE_LABELS[m]}:**")
                a("")
                a(summary)
                a("")
    else:
        a("No module reports found for this case ID. Run FAN, FAME, and/or FAST first.")
        a("")

    # ── 2. Cross-Domain Correlation ───────────────────────────────────────────
    a("---")
    a("")
    a("## 2. Cross-Domain Correlation")
    a("")
    correlation_md = sources.get("correlation_md")
    if correlation_md:
        body = _clean_correlation_body(correlation_md)
        if body:
            a(body)
            a("")
    elif len(modules_run) > 1:
        a(f"Cross-domain correlation has not yet been computed for this case. Run "
          f"`python3 lib/correlate_findings.py --case-id {case_id}` while `./analysis/` "
          f"still contains the raw artifacts to compute actual matches.")
        a("")
        a("Candidate correlation pairs based on the modules that ran:")
        a("")
        if "FAN" in module_texts and "FAME" in module_texts:
            a("- **FAN ↔ FAME:** match network connections from memory `netscan` output "
              "against PCAP flow data to link a process to observed traffic.")
        if "FAME" in module_texts and "FAST" in module_texts:
            a("- **FAME ↔ FAST:** match process image paths from memory `filescan`/`dlllist` "
              "against the `fls` file listing — a path present in memory but deleted on "
              "disk indicates persistence or clean-up.")
        if "FAN" in module_texts and "FAST" in module_texts:
            a("- **FAN ↔ FAST:** match carved URLs/domains from `bulk_extractor` against "
              "DNS queries observed in the PCAP.")
        a("")
    else:
        a("Only one module has run for this case. Cross-domain correlation requires at "
          "least two of FAN, FAME, and FAST. Run the remaining modules with the same "
          f"`--case-id {case_id}`, then run `python3 lib/correlate_findings.py --case-id {case_id}`.")
        a("")

    # ── 3. Unified MITRE ATT&CK ───────────────────────────────────────────────
    a("---")
    a("")
    a("## 3. Unified MITRE ATT&CK Coverage")
    a("")
    mitre_table = _mitre_table_md(_merge_mitre(module_texts))
    if mitre_table:
        a(mitre_table)
        a("")
    else:
        a("No MITRE ATT&CK techniques were mapped by the modules that ran.")
        a("")

    # ── 4. Indicators of Compromise ───────────────────────────────────────────
    a("---")
    a("")
    a("## 4. Indicators of Compromise")
    a("")
    ioc_merged, ioc_cat_order = _merge_iocs(module_texts)
    ioc_table = _ioc_tables_md(ioc_merged, ioc_cat_order)
    if ioc_table:
        a(ioc_table)
        a("")
    else:
        a("No indicators of compromise were extracted by the modules that ran.")
        a("")

    # ── 5. Module Reports ─────────────────────────────────────────────────────
    a("---")
    a("")
    a("## 5. Module Reports")
    a("")
    sub_n = 0
    for m in ("FAN", "FAME", "FAST"):
        if m not in module_texts:
            continue
        sub_n += 1
        a(f"### 5.{sub_n} {_MODULE_SECTION_TITLES[m]}")
        a("")
        summary = _extract_summary(module_texts[m])
        if summary:
            a(summary)
            a("")
        path = sources.get(f"{m.lower()}_path")
        if path:
            rel = os.path.relpath(path, case_dir)
            a(f"*Full technical report: [`{rel}`](./{rel})*")
            a("")
    if not modules_run:
        a("No module reports available.")
        a("")

    # ── 6. Unified Recommendations ────────────────────────────────────────────
    a("---")
    a("")
    a("## 6. Unified Recommendations")
    a("")
    recs = _merge_recommendations(module_texts)
    if recs:
        for i, rec in enumerate(recs, start=1):
            a(f"{i}. {rec}")
        a("")
    else:
        a("No recommendations were extracted from the module reports that ran.")
        a("")

    # ── 7. Confidence Assessment ──────────────────────────────────────────────
    a("---")
    a("")
    hg = _build_combined_hallucination_guard(module_texts, ioc_merged)
    if hg:
        a(hg)
        a("")

    # ── Appendix — Evidence Sources ───────────────────────────────────────────
    appendix_rows = []
    for m, label in (("fan", "FAN module report"), ("fame", "FAME module report"),
                      ("fast", "FAST module report"), ("correlation", "Cross-module correlation")):
        path = sources.get(f"{m}_path")
        if path:
            rel = os.path.relpath(path, case_dir)
            appendix_rows.append((label, rel))
    if appendix_rows:
        a("---")
        a("")
        a("## Appendix — Evidence Sources")
        a("")
        a("| Source | Path |")
        a("|--------|------|")
        for label, rel in appendix_rows:
            a(f"| {label} | [`{rel}`](./{rel}) |")
        a("")

    return _final_cleanup("\n".join(lines))


def _build_combined_hallucination_guard(
    module_texts: dict[str, str],
    ioc_merged: dict[str, list[dict]],
) -> str:
    """
    Build a cross-module Hallucination Guard summary for the campaign report:
    one meta-finding per module (CONFIRMED if it ran, UNVERIFIABLE if it
    didn't), plus a CONFIRMED finding for every IOC corroborated across two or
    more modules.
    """
    _hg_reset()
    findings = []

    module_sources = {
        "FAN":  ["fan_protocol_analyzers"],
        "FAME": ["volatility3", "memory_baseliner"],
        "FAST": ["tsk", "bulk_extractor"],
    }
    module_names = {
        "FAN":  "FAN (Network Forensics)",
        "FAME": "FAME (Memory Forensics)",
        "FAST": "FAST (Storage Forensics)",
    }
    for m in ("FAN", "FAME", "FAST"):
        if m in module_texts:
            findings.append(tag_finding(
                f"{module_names[m]} analysis completed — findings backed by module-level evidence",
                ConfidenceTier.CONFIRMED,
                [],
                module_sources[m],
                [m.lower()],
            ))
        else:
            findings.append(tag_finding(
                f"{module_names[m]} did not run — corresponding evidence absent from this case",
                ConfidenceTier.UNVERIFIABLE,
                [],
                module_sources[m],
                [m.lower()],
            ))

    # Cross-module IOC corroboration
    for cat, entries in ioc_merged.items():
        for e in entries:
            modules_seen = sorted({s.split()[0] for s in e["sources"]})
            if len(modules_seen) > 1:
                findings.append(tag_finding(
                    f"{e['value']} ({cat}) confirmed across {' + '.join(modules_seen)} — cross-domain corroboration",
                    ConfidenceTier.CONFIRMED,
                    [],
                    e["sources"],
                    [mm.lower() for mm in modules_seen],
                ))

    return render_confidence_summary(findings, module_label="Campaign report")


# ── PPTX ───────────────────────────────────────────────────────────────────────

def _build_pptx(
    sources: dict[str, object],
    case_id: str,
    hostname: str,
    generated_utc: str,
    output_path: Path,
) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("[combined] WARNING: python-pptx not installed — skipping PPTX.")
        return

    from artifact_guard import resolve_output, record_generated
    output_path, diverted = resolve_output(output_path)

    module_texts: dict[str, str] = {}
    for m in ("FAN", "FAME", "FAST"):
        text = sources.get(f"{m.lower()}_md")
        if text:
            module_texts[m] = text
    fan, fame, fast = module_texts.get("FAN", ""), module_texts.get("FAME", ""), module_texts.get("FAST", "")
    modules_run = list(module_texts.keys())

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def _rgb(t):
        return RGBColor(*t)

    def _rect(slide, l, t, w, h, fill):
        s = slide.shapes.add_shape(1, l, t, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = _rgb(fill)
        s.line.fill.background()
        return s

    def _txt(slide, text, l, t, w, h, sz, bold=False, color=_WHITE, align=PP_ALIGN.LEFT):
        tb  = slide.shapes.add_textbox(l, t, w, h)
        tf  = tb.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)[:500]
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)

    W = prs.slide_width
    H = prs.slide_height
    M = Inches(0.4)

    # Slide 1 — Cover
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, W, H, _DARK_NAVY)
    _rect(s, 0, 0, W, Inches(0.08), _BLUE)
    _rect(s, 0, H - Inches(0.08), W, Inches(0.08), _BLUE)
    _txt(s, "Fan Get Fame Fast", M, Inches(1.0), W - 2*M, Inches(1.2),
         52, bold=True, color=_LIGHT_BLUE, align=PP_ALIGN.CENTER)
    _txt(s, "Campaign forensics report", M, Inches(2.1), W - 2*M, Inches(0.7),
         24, color=_WHITE, align=PP_ALIGN.CENTER)
    _txt(s, "FAN  ·  FAME  ·  FAST", M, Inches(2.7), W - 2*M, Inches(0.5),
         18, color=_LIGHT_BLUE, align=PP_ALIGN.CENTER)
    _rect(s, Inches(3), Inches(3.6), W - Inches(6), Inches(0.04), _BLUE)
    _txt(s, f"Case: {case_id}  |  Host: {hostname}  |  {generated_utc[:10]}",
         M, Inches(3.9), W - 2*M, Inches(0.5), 14, color=_TEXT_MID, align=PP_ALIGN.CENTER)
    modules_str = "  ·  ".join(_MODULE_LABELS[m] for m in modules_run) if modules_run else "No modules run"
    _txt(s, f"Modules: {modules_str}", M, Inches(4.5), W - 2*M, Inches(0.4),
         12, color=_TEXT_MID, align=PP_ALIGN.CENTER)
    _txt(s, "Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin · Joost Beekman",
         M, Inches(5.0), W - 2*M, Inches(0.4), 11, color=_TEXT_MID, align=PP_ALIGN.CENTER)
    _txt(s, "CONFIDENTIAL — FOR AUTHORISED PERSONNEL ONLY",
         M, H - Inches(0.7), W - 2*M, Inches(0.4), 11, color=_TEXT_MID, align=PP_ALIGN.CENTER)

    # Slide 2 — Module Coverage Overview
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, W, Inches(1.1), _MID_NAVY)
    _txt(s, "Investigation scope", M, Inches(0.2), W, Inches(0.8), 28, bold=True, color=_WHITE)

    col_data = [
        ("FAN", "Network forensics", "PCAP analysis · 22 protocol detectors · IDS/YARA alerts", fan),
        ("FAME", "Memory forensics", "Volatility 3 · Memory Baseliner · Process/network/code", fame),
        ("FAST", "Storage forensics", "TSK · EWF tools · Timeline · Artifact extraction", fast),
    ]
    col_w = (W - 2*M - Inches(0.4)) // 3
    for i, (abbr, title, desc, content) in enumerate(col_data):
        cx = M + i * (col_w + Inches(0.2))
        color = _MID_NAVY if content else (0x3b, 0x44, 0x5b)
        _rect(s, cx, Inches(1.2), col_w, Inches(5.8), color)
        status = "COMPLETE" if content else "NOT RUN"
        status_color = _GREEN if content else _AMBER
        _rect(s, cx, Inches(1.2), col_w, Inches(0.35), status_color if content else _AMBER)
        _txt(s, status, cx + Inches(0.1), Inches(1.22), col_w - Inches(0.2), Inches(0.3),
             11, bold=True, color=_WHITE)
        _txt(s, abbr, cx + Inches(0.1), Inches(1.65), col_w - Inches(0.2), Inches(0.8),
             32, bold=True, color=_LIGHT_BLUE if content else _TEXT_MID)
        _txt(s, title, cx + Inches(0.1), Inches(2.45), col_w - Inches(0.2), Inches(0.5),
             13, bold=True, color=_WHITE if content else _TEXT_MID)
        _txt(s, desc, cx + Inches(0.1), Inches(3.0), col_w - Inches(0.2), Inches(2.5),
             11, color=_LIGHT_BLUE if content else _TEXT_MID)
        if content:
            summary = _strip_md(_extract_summary(content))[:250]
            _txt(s, summary, cx + Inches(0.1), Inches(5.0), col_w - Inches(0.2), Inches(1.8),
                 10, color=_TEXT_MID)

    # Slide 3 — Cross-Domain Correlation
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, W, Inches(1.1), _MID_NAVY)
    _txt(s, "Cross-domain correlation", M, Inches(0.2), W, Inches(0.8), 28, bold=True, color=_WHITE)

    correlation_md = sources.get("correlation_md")
    if correlation_md:
        # Extract match counts from the confidence table in the correlation report
        ff_n = mf_n = fd_n = 0
        for line in correlation_md.splitlines():
            if "FAN ↔ FAME" in line and "|" in line:
                m = re.search(r"\|\s*(\d+)\s*\|", line)
                if m: ff_n = int(m.group(1))
            elif "FAME ↔ FAST" in line and "|" in line:
                m = re.search(r"\|\s*(\d+)\s*\|", line)
                if m: mf_n = int(m.group(1))
            elif "FAN ↔ FAST" in line and "|" in line:
                m = re.search(r"\|\s*(\d+)\s*\|", line)
                if m: fd_n = int(m.group(1))
        total_corr = ff_n + mf_n + fd_n
        correlations = [
            ("FAN ↔ FAME",
             f"{ff_n} match(es) — process-to-network: links running processes to flagged PCAP connections"),
            ("FAME ↔ FAST",
             f"{mf_n} match(es) — process-to-disk: identifies executables deleted post-execution (T1070.004)"),
            ("FAN ↔ FAST",
             f"{fd_n} match(es) — domain-to-URL: confirms endpoints seen in both DNS traffic and carved artifacts"),
        ]
        if total_corr > 0:
            correlations.append(("Total matches",
                f"{total_corr} cross-domain linkages — see the correlation report for full detail"))
    else:
        correlations = []
        if "FAN" in module_texts and "FAME" in module_texts:
            correlations.append(("FAN ↔ FAME",
                "Match netscan process IDs to PCAP flows — links specific processes to observed network traffic"))
        if "FAME" in module_texts and "FAST" in module_texts:
            correlations.append(("FAME ↔ FAST",
                "Cross-reference process image paths in memory with deleted file entries on disk"))
        if "FAN" in module_texts and "FAST" in module_texts:
            correlations.append(("FAN ↔ FAST",
                "Match carved URLs from bulk_extractor with DNS queries in PCAP"))
        if len(modules_run) == 3:
            correlations.append(("FAN + FAME + FAST",
                f"Run python3 lib/correlate_findings.py --case-id {case_id} to compute full kill-chain correlations"))

    if correlations:
        row_h = Inches(1.1)
        for i, (pair, desc) in enumerate(correlations):
            y = Inches(1.2) + i * row_h
            _rect(s, M, y, Inches(2.8), row_h - Inches(0.1), _BLUE)
            _txt(s, pair, M + Inches(0.1), y + Inches(0.2), Inches(2.6), row_h,
                 14, bold=True, color=_WHITE)
            _txt(s, desc, M + Inches(3.0), y + Inches(0.15), W - M - Inches(3.4), row_h,
                 14, color=_TEXT_DARK)
    else:
        _txt(s, "Run all three modules (FAN, FAME, FAST) for the same case ID, then run "
             f"python3 lib/correlate_findings.py --case-id {case_id} to generate cross-domain correlation.",
             M, Inches(2.0), W - 2*M, Inches(2.0), 15, color=_TEXT_MID)

    # Slide 4 — Unified MITRE ATT&CK
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, W, Inches(1.1), _MID_NAVY)
    _txt(s, "Unified MITRE ATT&CK kill chain", M, Inches(0.2), W, Inches(0.8), 28, bold=True, color=_WHITE)

    mitre_rows = _merge_mitre(module_texts)
    if mitre_rows:
        headers = ["Technique", "Name", "Tactic", "Module(s)", "Observation"]
        col_ws  = [Inches(1.3), Inches(2.0), Inches(2.0), Inches(1.2), W - M - Inches(7.3)]
        row_h   = Inches(0.63)
        hx = M
        for h, cw in zip(headers, col_ws):
            _rect(s, hx, Inches(1.2), cw - Inches(0.05), row_h - Inches(0.04), _MID_NAVY)
            _txt(s, h, hx + Inches(0.08), Inches(1.25), cw, row_h, 12, bold=True, color=_WHITE)
            hx += cw
        for i, row in enumerate(mitre_rows[:8]):
            y = Inches(1.2) + (i+1) * row_h
            bg = _LIGHT_BG if i % 2 == 0 else _ROW_ALT
            rx = M
            values = [
                _strip_md(row["id_md"]), row["name"], row["tactic"],
                row["modules"], _strip_md(row["observation"])[:60],
            ]
            for val, cw in zip(values, col_ws):
                _rect(s, rx, y, cw - Inches(0.05), row_h - Inches(0.04), bg)
                _txt(s, val, rx + Inches(0.08), y + Inches(0.08), cw - Inches(0.13), row_h,
                     10, color=_TEXT_DARK)
                rx += cw
    else:
        _txt(s, "No MITRE ATT&CK techniques mapped across modules.",
             M, Inches(2.0), W - 2*M, Inches(1.0), 16, color=_TEXT_DARK)

    # Slide 5 — Unified Recommendations
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, W, Inches(1.1), _MID_NAVY)
    _txt(s, "Unified recommendations", M, Inches(0.2), W, Inches(0.8), 28, bold=True, color=_WHITE)
    all_recs = [r[:130] for r in _merge_recommendations(module_texts)]

    row_h = Inches(0.62)
    for i, rec in enumerate(all_recs[:8]):
        y = Inches(1.2) + i * row_h
        _rect(s, M, y, Inches(0.45), row_h - Inches(0.08), _BLUE)
        _txt(s, str(i+1), M + Inches(0.08), y + Inches(0.08), Inches(0.3), row_h,
             14, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
        _txt(s, rec, M + Inches(0.55), y + Inches(0.08), W - M - Inches(0.95), row_h,
             12, color=_TEXT_DARK)
    if not all_recs:
        _txt(s, "No recommendations were extracted from the module reports that ran.",
             M, Inches(1.4), W - 2*M, Inches(1.0), 16, color=_TEXT_DARK)

    prs.save(str(output_path))
    print(f"[combined] PPTX saved: {output_path}")
    if not diverted:
        record_generated(output_path)


# ── DOCX ───────────────────────────────────────────────────────────────────────

def _build_docx(
    sources: dict[str, object],
    case_id: str,
    hostname: str,
    generated_utc: str,
    output_path: Path,
) -> None:
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        print("[combined] WARNING: python-docx not installed — skipping DOCX.")
        return

    module_texts: dict[str, str] = {}
    for m in ("FAN", "FAME", "FAST"):
        text = sources.get(f"{m.lower()}_md")
        if text:
            module_texts[m] = text
    modules_run = list(module_texts.keys())

    doc = Document()
    for section in doc.sections:
        section.top_margin    = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin   = Inches(1.2)
        section.right_margin  = Inches(1.2)

    def _xml_safe(text: str) -> str:
        return re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", str(text))

    def _h(text, level):
        p = doc.add_heading(_xml_safe(text), level=level)
        p.runs[0].font.color.rgb = RGBColor(0x0f, 0x17, 0x2a)

    def _p(text, bold=False, italic=False):
        p = doc.add_paragraph()
        run = p.add_run(_xml_safe(text))
        run.bold   = bold
        run.italic = italic

    def _note(text):
        p = doc.add_paragraph()
        run = p.add_run(_xml_safe(text))
        run.italic = True
        run.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)

    def _md_blocks(md_text: str) -> None:
        """Render a markdown fragment as DOCX headings/paragraphs/notes."""
        for line in md_text.splitlines():
            stripped = line.strip()
            if not stripped:
                doc.add_paragraph()
            elif stripped.startswith("#### "):
                doc.add_heading(_xml_safe(stripped[5:]), 4)
            elif stripped.startswith("### "):
                doc.add_heading(_xml_safe(stripped[4:]), 3)
            elif stripped.startswith("## "):
                doc.add_heading(_xml_safe(stripped[3:]), 2)
            elif stripped.startswith("# "):
                doc.add_heading(_xml_safe(stripped[2:]), 1)
            elif stripped.startswith("|"):
                _p(stripped)
            elif stripped.startswith(">"):
                _note(stripped.lstrip("> ").strip())
            elif stripped == "---":
                pass
            else:
                clean = re.sub(r"\*\*(.*?)\*\*", r"\1", stripped)
                clean = re.sub(r"\*(.*?)\*",     r"\1", clean)
                clean = re.sub(r"`(.*?)`",       r"\1", clean)
                if clean:
                    _p(clean)

    # Cover
    doc.add_paragraph()
    t = doc.add_heading(f"{case_id} — Campaign Forensics Report", 0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = sub.add_run("FAN · FAME · FAST  |  Integrated Investigation Report")
    r.font.size = Pt(14)
    r.font.color.rgb = RGBColor(0x1d, 0x4e, 0xd8)
    doc.add_paragraph()

    tbl = doc.add_table(rows=5, cols=2)
    tbl.style = "Table Grid"
    rows_data = [
        ("Case ID", case_id),
        ("Hostname", hostname),
        ("Modules", ", ".join(_MODULE_LABELS[m] for m in modules_run) if modules_run else "None"),
        ("Analyst", "Claude Code — Campaign Report"),
        ("Generated UTC", generated_utc),
    ]
    for i, (k, v) in enumerate(rows_data):
        tbl.rows[i].cells[0].text = k
        tbl.rows[i].cells[1].text = v

    doc.add_paragraph()
    conf = doc.add_paragraph("CONFIDENTIAL — FOR AUTHORISED PERSONNEL ONLY")
    conf.alignment = WD_ALIGN_PARAGRAPH.CENTER
    conf.runs[0].font.bold = True
    conf.runs[0].font.color.rgb = RGBColor(0xef, 0x44, 0x44)
    doc.add_page_break()

    # 1. Management Summary
    _h("1. Management Summary", 1)
    _note("Audience: CISO, Legal, Law Enforcement, Internal Audit — no technical identifiers.")
    if modules_run:
        for m in modules_run:
            summary = _extract_summary(module_texts[m])
            if summary:
                _h(_MODULE_LABELS[m], 2)
                _p(summary.strip())
                doc.add_paragraph()
    else:
        _p("No module reports found for this case ID. Run FAN, FAME, and/or FAST first.")
    doc.add_page_break()

    # 2. Cross-Domain Correlation
    _h("2. Cross-Domain Correlation", 1)
    correlation_md = sources.get("correlation_md")
    if correlation_md:
        _md_blocks(_clean_correlation_body(correlation_md))
    elif len(modules_run) > 1:
        _p(
            f"Cross-domain correlation has not yet been computed for this case. Run "
            f"python3 lib/correlate_findings.py --case-id {case_id} while ./analysis/ "
            f"still contains the raw artifacts to compute actual matches."
        )
        if "FAN" in module_texts and "FAME" in module_texts:
            p = doc.add_paragraph(style="List Bullet")
            p.add_run("FAN ↔ FAME: match memory netscan output to PCAP flows — links "
                      "specific processes to observed network connections.")
        if "FAME" in module_texts and "FAST" in module_texts:
            p = doc.add_paragraph(style="List Bullet")
            p.add_run("FAME ↔ FAST: cross-reference process image paths from memory "
                      "filescan with deleted entries in fls output.")
        if "FAN" in module_texts and "FAST" in module_texts:
            p = doc.add_paragraph(style="List Bullet")
            p.add_run("FAN ↔ FAST: match carved URLs/domains from bulk_extractor with "
                      "DNS queries in the PCAP.")
    else:
        _p("Only one module has run for this case. Cross-domain correlation requires "
           "at least two of FAN, FAME, and FAST.")
    doc.add_paragraph()
    doc.add_page_break()

    # 3. Unified MITRE ATT&CK
    _h("3. Unified MITRE ATT&CK Coverage", 1)
    mitre_rows = _merge_mitre(module_texts)
    if mitre_rows:
        table = doc.add_table(rows=1, cols=5)
        table.style = "Table Grid"
        hdr = table.rows[0].cells
        for cell, text in zip(hdr, ["Technique", "Name", "Tactic", "Module(s)", "Observation"]):
            cell.text = text
        for row in mitre_rows:
            cells = table.add_row().cells
            for cell, text in zip(cells, [
                _strip_md(row["id_md"]), row["name"], row["tactic"],
                row["modules"], _strip_md(row["observation"]),
            ]):
                cell.text = text
    else:
        _p("No MITRE ATT&CK techniques were mapped by the modules that ran.")
    doc.add_page_break()

    # 4. Indicators of Compromise
    _h("4. Indicators of Compromise", 1)
    ioc_merged, ioc_cat_order = _merge_iocs(module_texts)
    any_iocs = False
    for cat in ioc_cat_order:
        entries = ioc_merged.get(cat) or []
        if not entries:
            continue
        any_iocs = True
        _h(cat, 2)
        table = doc.add_table(rows=1, cols=3)
        table.style = "Table Grid"
        hdr = table.rows[0].cells
        for cell, text in zip(hdr, ["Value", "Confidence", "Source"]):
            cell.text = text
        for e in entries:
            cells = table.add_row().cells
            cells[0].text = _strip_md(e["value"])
            cells[1].text = e["confidence"]
            cells[2].text = "; ".join(e["sources"])
        doc.add_paragraph()
    if not any_iocs:
        _p("No indicators of compromise were extracted by the modules that ran.")
    doc.add_page_break()

    # 5. Module Reports
    _h("5. Module Reports", 1)
    for m in ("FAN", "FAME", "FAST"):
        if m not in module_texts:
            continue
        _h(_MODULE_SECTION_TITLES[m], 2)
        summary = _extract_summary(module_texts[m])
        if summary:
            _p(summary.strip())
        path = sources.get(f"{m.lower()}_path")
        if path:
            _note(f"Full technical report: {Path(path).name}")
        doc.add_paragraph()
    if not modules_run:
        _p("No module reports available.")
    doc.add_page_break()

    # 6. Unified Recommendations
    _h("6. Unified Recommendations", 1)
    all_recs = _merge_recommendations(module_texts)
    if all_recs:
        for rec in all_recs:
            p = doc.add_paragraph(style="List Number")
            p.add_run(rec)
    else:
        _p("No recommendations were extracted from the module reports that ran.")
    doc.add_page_break()

    # 7. Confidence Assessment
    hg = _build_combined_hallucination_guard(module_texts, ioc_merged)
    if hg:
        _h("7. Confidence Assessment", 1)
        _md_blocks(hg)

    doc.save(str(output_path))
    print(f"[combined] DOCX saved: {output_path}")


# ── Public API ─────────────────────────────────────────────────────────────────

def generate(
    case_id: str,
    hostname: str,
    reports_dir: Path | None = None,
    output_dir: Path | None = None,
    case_dir: Path | None = None,
    docs_dir: Path | None = None,
) -> dict[str, Path | None]:
    """Generate the unified campaign report.

    case_dir: case root directory (reports/<case_id>/). When supplied, the campaign
    Markdown lands in case_dir/ and PDF/PPTX/DOCX go to docs_dir (default:
    case_dir/documents/). Module reports are auto-discovered from the FAN/FAME/FAST
    subdirectories. When omitted, falls back to output_dir / reports_dir (legacy).
    """
    reports_dir = reports_dir or (PROJECT_ROOT / "reports")

    if case_dir is not None:
        md_dir  = path_guard.guard_output_dir(case_dir)
        doc_dir = path_guard.guard_output_dir(docs_dir or (case_dir / "documents"))
        discover_dir = reports_dir
    else:
        md_dir  = path_guard.guard_output_dir(output_dir or reports_dir)
        doc_dir = md_dir
        discover_dir = reports_dir

    sources = _discover_sources(discover_dir, case_id)
    generated_utc = datetime.now(_CET).strftime("%Y-%m-%d %H:%M CET")
    stem = case_id.replace(" ", "_")

    md_text = _build_markdown(sources, case_id, hostname, generated_utc, md_dir)
    md_path = md_dir / f"{stem}_campaign_report.md"
    md_path.write_text(md_text)
    print(f"[combined] Markdown saved: {md_path}")

    pdf_path: Path | None = None
    try:
        sys.path.insert(0, str(PROJECT_ROOT / "lib"))
        from md_to_pdf import convert as md2pdf
        pdf_path = doc_dir / f"{stem}_campaign_report.pdf"
        md2pdf(md_path, pdf_path)
        print(f"[combined] PDF saved: {pdf_path}")
    except Exception as exc:
        print(f"[combined] WARNING: PDF generation failed: {exc}")

    pptx_path = doc_dir / f"{stem}_campaign_presentation.pptx"
    _build_pptx(sources, case_id, hostname, generated_utc, pptx_path)

    docx_path = doc_dir / f"{stem}_campaign_report.docx"
    _build_docx(sources, case_id, hostname, generated_utc, docx_path)

    print(f"[combined] Campaign report suite complete:")
    print(f"  MD    {md_path}")
    print(f"  PDF   {pdf_path}")
    print(f"  PPTX  {pptx_path}")
    print(f"  DOCX  {docx_path}")
    return {
        "md":   md_path,
        "pdf":  pdf_path,
        "pptx": pptx_path if pptx_path.exists() else None,
        "docx": docx_path if docx_path.exists() else None,
    }


# ── CLI ────────────────────────────────────────────────────────────────────────

def _build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(description="FanGetFameFast — Campaign Report Generator")
    p.add_argument("--case-id",     required=True, metavar="ID")
    p.add_argument("--hostname",    required=True, metavar="HOST")
    p.add_argument("--reports-dir", default=None,  metavar="DIR")
    p.add_argument("--output-dir",  default=None,  metavar="DIR")
    p.add_argument("--case-dir",    default=None,  metavar="DIR",  help="Case root dir (reports/<case_id>/); campaign MD lands here")
    p.add_argument("--docs-dir",    default=None,  metavar="DIR",  help="Documents dir (reports/<case_id>/documents/); PDF/PPTX/DOCX land here")
    return p


if __name__ == "__main__":
    args = _build_parser().parse_args()
    paths = generate(
        case_id     = args.case_id,
        hostname    = args.hostname,
        reports_dir = Path(args.reports_dir) if args.reports_dir else None,
        output_dir  = Path(args.output_dir)  if args.output_dir  else None,
        case_dir    = Path(args.case_dir)    if args.case_dir    else None,
        docs_dir    = Path(args.docs_dir)    if args.docs_dir    else None,
    )
    print("[combined] Report suite complete:")
    for fmt, p in paths.items():
        if p:
            print(f"  {fmt.upper():4s}  {p}")
