#!/usr/bin/env python3
# SPDX-License-Identifier: MIT OR Apache-2.0
# SPDX-FileCopyrightText: 2026 Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin · Joost Beekman
"""
board_deck.py — Campaign board-deck PPTX renderer.

Renders ``<case_id>_campaign_presentation.pptx`` using the board-deck design:
a cover slide with a severity ring and case-status cards, KPI strips,
numbered executive panels, eyebrow/headline/body card grids, a 5-step
incident timeline, a primary-cause banner with contributing-factor chips,
and prioritized recommendation rows.

The design is driven by a small bullet grammar (see
docs/campaign_report_template.md) layered on top of the hand-authored
campaign report markdown:

    | Severity   | High — Historical exposure |
    | Case status | Reconstructed |

    #### Board KPIs
    | Label | Value | Detail |
    |-------|-------|--------|
    | ... |

    ##### Executive Summary
    - **Headline** — Body text.

    ##### Key Findings / Impact Assessment / Residual Risks / Response Actions / Key Insights
    - **EYEBROW — Headline:** Body text.

    ##### Contributing Factors
    - **EYEBROW:** Headline

    ##### Incident Timeline
    - **HH:MM — Headline** — Body text.

    ##### Priority Actions
    - **[P0] Headline** — Body text.

Sections or subsections that don't follow the grammar fall back to the
generic bullet/table/text rendering from ``md_to_pptx``, so unstructured
campaign reports still render.

Python API:
    from lib.board_deck import convert
    pptx_path = convert(md_path, output_path, case_id="CASE-2026-001")
"""
from __future__ import annotations

import argparse
import re
import sys
from datetime import datetime, timezone
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import path_guard  # noqa: E402  write-path policy enforcement
from md_to_pptx import (  # noqa: E402
    _Block, _split_blocks, _parse_table_rows, _strip_md, _normalize_heading,
    _render_block, _TABLE_ROW_RE,
)
from generate_pptx_report import (  # noqa: E402
    _rgb, _set_bg, _rect, _text, _text_lines, _header_bar,
    _WHITE, _LIGHT,
)

# ── Palette (extracted from the hand-edited NIST-HACK-2026-V2 board deck) ──

_CARD_BG     = (0x14, 0x2A, 0x47)
_COVER_PANEL = (0x0F, 0x1F, 0x3D)
_TEXT_BODY   = (0xCB, 0xD5, 0xE1)
_TEXT_FAINT  = (0x94, 0xA3, 0xB8)

_RED    = (0xEF, 0x44, 0x44)
_AMBER  = (0xF5, 0x9E, 0x0B)
_GREEN  = (0x10, 0xB9, 0x81)
_CYAN   = (0x00, 0xBC, 0xD4)
_PURPLE = (0x8B, 0x5C, 0xF6)
_BLUE2  = (0x29, 0x79, 0xFF)

_RED_LIGHT    = (0xFC, 0xA5, 0xA5)
_AMBER_LIGHT  = (0xFD, 0xE6, 0x8A)
_AMBER_LIGHT2 = (0xFC, 0xD3, 0x4D)
_GREEN_LIGHT  = (0x6E, 0xE7, 0xB7)
_CYAN_LIGHT   = (0x67, 0xE8, 0xF9)
_PURPLE_LIGHT = (0xC4, 0xB5, 0xFD)
_BLUE_LIGHT   = (0x90, 0xCA, 0xF9)

_ACCENTS = [_RED, _AMBER, _GREEN, _CYAN, _PURPLE, _BLUE2]

_SEV_COLORS = {
    "critical": _RED, "high": _RED, "medium": _AMBER, "low": _GREEN, "info": _CYAN,
}

# ── Bullet grammar ───────────────────────────────────────────────────────

_RAW_BULLET_RE = re.compile(r"^[-*]\s+(.*)$")
_BOLD_RE = re.compile(r"^\*\*(.+?)\*\*\s*(.*)$")
_PRIORITY_RE = re.compile(r"^\[(P\d)\]\s*(.+)$")
_TIME_RE = re.compile(r"^(\d{1,2}:\d{2})\s*[—-]\s*(.+)$")
_EYEBROW_RE = re.compile(r"^([A-Z✓][A-Z0-9 /✓.]{1,28}?)\s*[—-]\s*(.+)$")


def _collect_raw_bullets(lines: list[str]) -> list[str]:
    """Collect ``- ...`` bullets, joining wrapped continuation lines."""
    items: list[str] = []
    i, n = 0, len(lines)
    while i < n:
        s = lines[i].strip()
        m = _RAW_BULLET_RE.match(s)
        if m:
            text = m.group(1)
            i += 1
            while i < n:
                cont = lines[i].strip()
                if not cont or cont.startswith(("#", "|", ">", "-", "*")) or re.match(r"^\d+[.)]\s", cont):
                    break
                text += " " + cont
                i += 1
            items.append(text)
        else:
            i += 1
    return items


def _parse_card_bullet(raw: str) -> dict | None:
    """Parse a ``- **...** — ...`` bullet into eyebrow/priority/time/headline/body."""
    m = _BOLD_RE.match(raw.strip())
    if not m:
        return None
    bold, rest = m.group(1).strip(), m.group(2).strip()
    had_colon = bold.endswith(":")
    bold = bold.rstrip(":").strip()
    rest = re.sub(r"^[—-]\s*", "", rest)

    eyebrow = priority = time = None
    headline = bold

    pm = _PRIORITY_RE.match(bold)
    tm = _TIME_RE.match(bold)
    em = _EYEBROW_RE.match(bold)
    if pm:
        priority, headline = pm.group(1), pm.group(2)
    elif tm:
        time, headline = tm.group(1), tm.group(2)
    elif em:
        eyebrow, headline = em.group(1), em.group(2)
    elif had_colon:
        eyebrow, headline = bold, rest
        rest = ""

    return {
        "eyebrow": eyebrow,
        "priority": priority,
        "time": time,
        "headline": _strip_md(headline.strip()),
        "body": _strip_md(rest.strip()),
    }


def _parse_card_bullets(lines: list[str]) -> list[dict]:
    cards = []
    for raw in _collect_raw_bullets(lines):
        card = _parse_card_bullet(raw)
        if card is not None:
            cards.append(card)
    return cards


_SUBHEAD_RE = re.compile(r"^(#{4,6})\s+(.*)$")


def _split_subblocks(lines: list[str]) -> dict[str, list[str]]:
    """Split a block's body lines into ####/##### sub-sections, keyed by
    normalized heading text. The "" key holds lines before the first
    sub-heading."""
    subs: dict[str, list[str]] = {"": []}
    cur = ""
    in_fence = False
    for line in lines:
        if line.strip().startswith("```"):
            in_fence = not in_fence
            subs[cur].append(line)
            continue
        m = _SUBHEAD_RE.match(line.strip()) if not in_fence else None
        if m:
            cur = _normalize_heading(_strip_md(m.group(2)))
            subs[cur] = []
        else:
            subs[cur].append(line)
    return subs


def _header_fields(blocks: list[_Block]) -> dict[str, str]:
    """Extract the ``| Field | Value |`` header table from the H1 block's body."""
    h1_block = next((b for b in blocks if b.level == 1), None)
    if h1_block is None:
        return {}
    table_lines = [l for l in h1_block.lines if _TABLE_ROW_RE.match(l.strip())]
    rows = _parse_table_rows(table_lines)
    fields: dict[str, str] = {}
    for row in rows:
        if len(row) >= 2:
            fields[_strip_md(row[0]).strip().lower()] = _strip_md(row[1]).strip()
    return fields


# ── Drawing helpers ──────────────────────────────────────────────────────

def _oval(slide, left, top, width, height, fill):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.fill.background()
    return shape


def _new_slide(prs, title: str, case_id: str, W, H, kicker: str | None = None,
               kicker_color=_BLUE_LIGHT, kicker_tick=_BLUE2):
    from pptx.util import Inches
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, (0x0a, 0x16, 0x28))
    _header_bar(slide, title[:90], case_id, W, H)
    if kicker:
        _rect(slide, Inches(0.35), Inches(0.95), Inches(0.11), Inches(0.04), fill=kicker_tick)
        _text(slide, kicker, Inches(0.50), Inches(0.87), Inches(8.75), Inches(0.35),
              size=10, bold=True, color=kicker_color)
    return slide


_THEMES = {
    "management_summary": (_BLUE_LIGHT, _BLUE2),
    "business_impact":    (_RED_LIGHT, _RED),
    "board_timeline":     (_CYAN_LIGHT, _CYAN),
    "root_cause":         (_AMBER_LIGHT, _AMBER),
    "response":           (_GREEN_LIGHT, _GREEN),
    "recommendations":    (_BLUE_LIGHT, _BLUE2),
    "lessons_learned":    (_PURPLE_LIGHT, _PURPLE),
}


# ── Cover slide ──────────────────────────────────────────────────────────

_TITLE_SPLIT_RE = re.compile(r"^(.+?)-(\d.*)$")


def _render_cover(prs, case_id: str, description: str, date_str: str,
                  severity: str, severity_detail: str, case_status: str, W, H) -> None:
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, (0x0a, 0x16, 0x28))

    # Right-side panel
    _rect(slide, Inches(8.53), 0, Inches(4.80), H, fill=_COVER_PANEL)
    _rect(slide, Inches(8.53), 0, Inches(0.02), H, fill=_BLUE2)
    # Left accent bar
    _rect(slide, 0, 0, Inches(0.12), H, fill=(0x15, 0x65, 0xc0))

    _text(slide, "DIGITAL FORENSICS  ·  INCIDENT RESPONSE",
          Inches(0.44), Inches(0.59), Inches(7.87), Inches(0.27),
          size=10, bold=True, color=_BLUE2)
    _rect(slide, Inches(0.44), Inches(1.03), Inches(0.66), Inches(0.04), fill=_BLUE2)

    m = _TITLE_SPLIT_RE.match(case_id)
    if m:
        line1, line2 = m.group(1), m.group(2)
    else:
        line1, line2 = case_id, ""
    _text(slide, line1, Inches(0.44), Inches(1.31), Inches(7.87), Inches(1.25),
          size=68, bold=True, color=_WHITE)
    if line2:
        _text(slide, line2, Inches(0.44), Inches(2.35), Inches(7.87), Inches(0.98),
              size=52, bold=True, color=_BLUE2)

    if description:
        _text(slide, description, Inches(0.44), Inches(3.39), Inches(7.87), Inches(0.47),
              size=22, color=_TEXT_BODY)

    # CONFIDENTIAL chip
    _rect(slide, Inches(0.44), Inches(4.27), Inches(4.16), Inches(0.55), fill=(0x3B, 0x0F, 0x12))
    _text(slide, "● CONFIDENTIAL — RESTRICTED", Inches(0.44), Inches(4.27),
          Inches(4.16), Inches(0.55), size=12, bold=True, color=_RED_LIGHT)

    # REPORT DATE / CASE STATUS cards
    _rect(slide, Inches(0.44), Inches(5.14), Inches(3.83), Inches(0.98), fill=_CARD_BG)
    _rect(slide, Inches(0.44), Inches(5.14), Inches(0.04), Inches(0.98), fill=_CYAN)
    _text(slide, "REPORT DATE", Inches(0.59), Inches(5.21), Inches(3.61), Inches(0.33),
          size=9, bold=True, color=_CYAN_LIGHT)
    _text(slide, date_str, Inches(0.59), Inches(5.52), Inches(3.61), Inches(0.55),
          size=20, bold=True, color=_WHITE)

    _rect(slide, Inches(4.48), Inches(5.14), Inches(3.83), Inches(0.98), fill=_CARD_BG)
    _rect(slide, Inches(4.48), Inches(5.14), Inches(0.04), Inches(0.98), fill=_AMBER)
    _text(slide, "CASE STATUS", Inches(4.64), Inches(5.21), Inches(3.61), Inches(0.33),
          size=9, bold=True, color=_AMBER_LIGHT2)
    _text(slide, case_status or "—", Inches(4.64), Inches(5.52), Inches(3.61), Inches(0.55),
          size=20, bold=True, color=_WHITE)

    # Severity ring
    sev_color = _SEV_COLORS.get((severity or "").strip().lower(), _CYAN)
    _oval(slide, Inches(9.62), Inches(0.77), Inches(3.06), Inches(3.06), fill=_CARD_BG)
    _oval(slide, Inches(9.95), Inches(1.09), Inches(2.41), Inches(2.41), fill=_COVER_PANEL)
    _text(slide, "SEVERITY", Inches(9.62), Inches(1.86), Inches(3.06), Inches(0.55),
          size=10, bold=True, color=_CYAN_LIGHT)
    _text(slide, (severity or "—").upper(), Inches(9.62), Inches(2.08), Inches(3.06), Inches(0.87),
          size=60, bold=True, color=_WHITE, align=PP_ALIGN.LEFT)
    if severity_detail:
        _text(slide, severity_detail, Inches(9.62), Inches(3.06), Inches(3.06), Inches(0.33),
              size=10, color=_TEXT_FAINT)

    # Prepared by
    _text(slide, "PREPARED BY", Inches(8.75), Inches(4.21), Inches(4.37), Inches(0.33),
          size=9, bold=True, color=_BLUE_LIGHT)
    _rect(slide, Inches(8.75), Inches(4.57), Inches(0.44), Inches(0.02), fill=_BLUE2)
    _text(slide, "Richard de Vries | Jeffrey Everling | Malin Janssen | Suzanne Maquelin | Joost Beekman",
          Inches(8.75), Inches(4.68), Inches(4.37), Inches(1.64), size=12, bold=True, color=_WHITE)

    # Footer
    _rect(slide, 0, H - Inches(0.55), W, Inches(0.55), fill=(0x08, 0x0f, 0x1e))
    _text(slide, "FanGetFameFast  ·  Forensics Agent Network",
          Inches(0.44), H - Inches(0.46), Inches(7.66), Inches(0.38), size=9, color=_TEXT_FAINT)
    _text(slide, f"Generated {date_str} UTC",
          Inches(9.30), H - Inches(0.46), Inches(3.83), Inches(0.38),
          size=9, color=_TEXT_FAINT, align=PP_ALIGN.RIGHT)


# ── Management Summary: KPI strip + executive panels ────────────────────

def _render_kpi_strip(prs, kpis: list[list[str]], title: str, case_id: str, W, H) -> None:
    from pptx.util import Inches
    theme_color, theme_tick = _THEMES["management_summary"]
    slide = _new_slide(prs, title, case_id, W, H, "EXECUTIVE SUMMARY", theme_color, theme_tick)

    card_w, card_h = Inches(3.09), Inches(1.31)
    xs = [Inches(0.35), Inches(3.53), Inches(6.71), Inches(9.89)]
    y = Inches(0.87)
    for i, row in enumerate(kpis[:4]):
        label = _strip_md(row[0]) if len(row) > 0 else ""
        value = _strip_md(row[1]) if len(row) > 1 else ""
        detail = _strip_md(row[2]) if len(row) > 2 else ""
        accent = _ACCENTS[i % len(_ACCENTS)]
        x = xs[i]
        _rect(slide, x, y, card_w, card_h, fill=_CARD_BG)
        _rect(slide, x, y, Inches(0.07), card_h, fill=accent)
        _text(slide, label.upper(), x, y + Inches(0.03), card_w, Inches(0.31),
              size=9, bold=True, color=accent)
        _text(slide, value, x, y + Inches(0.31), card_w, Inches(0.55), size=20, bold=True, color=_WHITE)
        if detail:
            _text(slide, detail, x, y + Inches(0.90), card_w, Inches(0.33), size=10, color=_TEXT_FAINT)
    return slide


def _render_exec_panels(slide, panels: list[dict], W, H) -> None:
    from pptx.util import Inches
    card_w, card_h = Inches(12.63), Inches(1.29)
    y = Inches(2.76)
    pitch = Inches(1.37)
    for i, panel in enumerate(panels[:3]):
        accent = _ACCENTS[i % len(_ACCENTS)]
        x = Inches(0.35)
        cy = y + i * pitch
        _rect(slide, x, cy, card_w, card_h, fill=_CARD_BG)
        _rect(slide, x, cy, Inches(0.07), card_h, fill=accent)
        _text(slide, f"{i + 1:02d}", x + Inches(0.11), cy + Inches(0.11), Inches(0.77), Inches(0.87),
              size=36, bold=True, color=accent)
        _text(slide, panel["headline"], x + Inches(0.96), cy + Inches(0.15), Inches(11.48), Inches(0.46),
              size=16, bold=True, color=_WHITE)
        if panel["body"]:
            _text(slide, panel["body"][:600], x + Inches(0.96), cy + Inches(0.61), Inches(11.48), Inches(0.66),
                  size=11, color=_TEXT_BODY)


# ── Generic eyebrow/headline/body card grids ─────────────────────────────

def _render_card_grid(prs, cards: list[dict], title: str, case_id: str, W, H,
                      theme_key: str, layout: str = "2x2") -> None:
    """layout: '2x2' (4 cards, 6.26x2.62), '2x2_tall' (4 cards, 6.26x2.73),
    '3+2' (5 cards, 4.15x2.62)."""
    from pptx.util import Inches
    theme_color, theme_tick = _THEMES[theme_key]
    kicker = title.upper()
    slide = _new_slide(prs, title, case_id, W, H, kicker, theme_color, theme_tick)

    if layout == "3+2":
        positions = [
            (Inches(0.35), Inches(1.20)), (Inches(4.59), Inches(1.20)), (Inches(8.83), Inches(1.20)),
            (Inches(0.35), Inches(3.92)), (Inches(4.59), Inches(3.92)),
        ]
        card_w, card_h = Inches(4.15), Inches(2.62)
        body_h = Inches(1.37)
    elif layout == "2x2_tall":
        positions = [
            (Inches(0.35), Inches(1.20)), (Inches(6.72), Inches(1.20)),
            (Inches(0.35), Inches(4.05)), (Inches(6.72), Inches(4.05)),
        ]
        card_w, card_h = Inches(6.26), Inches(2.73)
        body_h = Inches(1.48)
    else:  # 2x2
        positions = [
            (Inches(0.35), Inches(1.31)), (Inches(6.72), Inches(1.31)),
            (Inches(0.35), Inches(4.05)), (Inches(6.72), Inches(4.05)),
        ]
        card_w, card_h = Inches(6.26), Inches(2.62)
        body_h = Inches(1.37)

    for i, card in enumerate(cards[:len(positions)]):
        x, y = positions[i]
        accent = _ACCENTS[i % len(_ACCENTS)]
        _rect(slide, x, y, card_w, card_h, fill=_CARD_BG)
        _rect(slide, x, y, card_w, Inches(0.05), fill=accent)
        if card.get("eyebrow"):
            _text(slide, card["eyebrow"].upper(), x, y + Inches(0.13), card_w, Inches(0.31),
                  size=9, bold=True, color=accent)
        _text(slide, card["headline"], x, y + Inches(0.46), card_w, Inches(0.66), size=14, bold=True, color=_WHITE)
        if card.get("body"):
            _text(slide, card["body"][:400], x, y + card_h - body_h, card_w, body_h, size=11, color=_TEXT_BODY)


# ── Board Timeline: 5-step columns ───────────────────────────────────────

def _render_timeline(prs, steps: list[dict], title: str, case_id: str, W, H, date_label: str = "") -> None:
    from pptx.util import Inches
    theme_color, theme_tick = _THEMES["board_timeline"]
    kicker = "INCIDENT TIMELINE" + (f" — {date_label}" if date_label else "")
    slide = _new_slide(prs, title, case_id, W, H, kicker, theme_color, theme_tick)

    card_w, card_h = Inches(2.47), Inches(5.25)
    pitch = Inches(2.54)
    x0, y0 = Inches(0.35), Inches(1.29)
    for i, step in enumerate(steps[:5]):
        accent = _ACCENTS[i % len(_ACCENTS)]
        x = x0 + i * pitch
        _rect(slide, x, y0, card_w, card_h, fill=_CARD_BG)
        _rect(slide, x, y0, card_w, Inches(0.09), fill=accent)
        if step.get("time"):
            _text(slide, step["time"], x, y0 + Inches(0.20), card_w, Inches(0.66),
                  size=24, bold=True, color=accent)
        _text(slide, f"STEP {i + 1:02d}", x, y0 + Inches(0.96), card_w, Inches(0.35),
              size=9, bold=True, color=_TEXT_FAINT)
        _text(slide, step["headline"], x + Inches(0.09), y0 + Inches(1.40), Inches(2.30), Inches(0.77),
              size=13, bold=True, color=_WHITE)
        if step.get("body"):
            _text(slide, step["body"][:400], x + Inches(0.09), y0 + Inches(2.30), Inches(2.30), Inches(2.73),
                  size=10, color=_TEXT_BODY)


# ── Root Cause & Risk: primary cause banner + contributing-factor chips ──

def _render_root_cause(prs, primary: dict | None, factors: list[dict], title: str, case_id: str, W, H) -> None:
    from pptx.util import Inches
    theme_color, theme_tick = _THEMES["root_cause"]
    slide = _new_slide(prs, title, case_id, W, H, "ROOT CAUSE ANALYSIS", theme_color, theme_tick)

    if primary:
        x, y = Inches(0.35), Inches(1.42)
        w, h = Inches(12.63), Inches(2.62)
        _rect(slide, x, y, w, h, fill=_CARD_BG)
        _rect(slide, x, y, Inches(0.09), h, fill=_AMBER)
        _text(slide, (primary.get("eyebrow") or "PRIMARY CAUSE").upper(), x + Inches(0.22), y + Inches(0.13),
              Inches(4.37), Inches(0.35), size=11, bold=True, color=_AMBER)
        _text(slide, primary["headline"], x + Inches(0.22), y + Inches(0.53), Inches(12.19), Inches(1.86),
              size=24, bold=True, color=_WHITE)
        if primary.get("body"):
            _text(slide, primary["body"][:600], x + Inches(0.22), y + Inches(1.07), Inches(12.19), Inches(1.31),
                  size=14, color=_TEXT_BODY)

    chip_w, chip_h = Inches(4.14), Inches(1.86)
    xs = [Inches(0.35), Inches(4.60), Inches(8.84)]
    y2 = Inches(4.37)
    chip_accents = [_PURPLE, _RED, _AMBER]
    for i, factor in enumerate(factors[:3]):
        accent = chip_accents[i % len(chip_accents)]
        x = xs[i]
        _rect(slide, x, y2, chip_w, chip_h, fill=_CARD_BG)
        _rect(slide, x + Inches(0.09), y2 + Inches(0.22), Inches(0.39), Inches(0.39), fill=accent)
        if factor.get("eyebrow"):
            _text(slide, factor["eyebrow"].upper(), x + Inches(0.61), y2 + Inches(0.22), Inches(3.48), Inches(0.35),
                  size=10, bold=True, color=accent)
        _text(slide, factor["headline"], x + Inches(0.61), y2 + Inches(0.62), Inches(3.48), Inches(0.77),
              size=17, bold=True, color=_WHITE)


# ── Recommendations: prioritized rows ────────────────────────────────────

_PRIORITY_ACCENT = {"P0": _RED, "P1": _AMBER, "P2": _CYAN}


def _render_priority_rows(prs, rows: list[dict], title: str, case_id: str, W, H) -> None:
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN
    theme_color, theme_tick = _THEMES["recommendations"]
    slide = _new_slide(prs, title, case_id, W, H, "PRIORITY ACTIONS", theme_color, theme_tick)

    card_w, card_h = Inches(12.63), Inches(1.01)
    pitch = Inches(1.07)
    x0, y0 = Inches(0.35), Inches(1.29)
    for i, row in enumerate(rows[:5]):
        priority = row.get("priority") or ""
        accent = _PRIORITY_ACCENT.get(priority, _ACCENTS[i % len(_ACCENTS)])
        y = y0 + i * pitch
        _rect(slide, x0, y, card_w, card_h, fill=_CARD_BG)
        _rect(slide, x0, y, Inches(0.07), card_h, fill=accent)
        _text(slide, f"{i + 1:02d}", x0 + Inches(0.13), y + Inches(0.22), Inches(0.77), Inches(0.66),
              size=28, bold=True, color=accent)
        if priority:
            _rect(slide, x0 + Inches(1.09), y + Inches(0.33), Inches(0.42), Inches(0.35), fill=accent)
            _text(slide, priority, x0 + Inches(1.09), y + Inches(0.33), Inches(0.42), Inches(0.35),
                  size=10, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
        _text(slide, row["headline"], x0 + Inches(1.64), y + Inches(0.22), Inches(10.66), Inches(0.44),
              size=15, bold=True, color=_WHITE)
        if row.get("body"):
            _text(slide, row["body"][:400], x0 + Inches(1.64), y + Inches(0.59), Inches(10.66), Inches(0.39),
                  size=11, color=_TEXT_BODY)


# ── Section dispatch ─────────────────────────────────────────────────────

def _render_management_summary(prs, block: _Block, case_id: str, W, H) -> bool:
    subs = _split_subblocks(block.lines)
    kpi_lines = subs.get("board kpis")
    exec_lines = subs.get("executive summary")
    findings_lines = subs.get("key findings")

    if not (kpi_lines and exec_lines):
        return False

    rows = _parse_table_rows(kpi_lines)
    kpis = rows[1:] if rows and rows[0] and rows[0][0].lower() == "label" else rows
    panels = _parse_card_bullets(exec_lines)
    if not kpis or not panels:
        return False

    slide = _render_kpi_strip(prs, kpis, block.heading, case_id, W, H)
    _render_exec_panels(slide, panels, W, H)

    if findings_lines:
        cards = _parse_card_bullets(findings_lines)
        if cards:
            _render_card_grid(prs, cards, block.heading + " (cont.)", case_id, W, H,
                              "management_summary", layout="2x2_tall")
    return True


def _render_business_impact(prs, block: _Block, case_id: str, W, H) -> bool:
    subs = _split_subblocks(block.lines)
    impact_lines = subs.get("impact assessment")
    if not impact_lines:
        return False
    cards = _parse_card_bullets(impact_lines)
    if not cards:
        return False
    _render_card_grid(prs, cards, block.heading, case_id, W, H, "business_impact", layout="3+2")
    return True


def _render_board_timeline(prs, block: _Block, case_id: str, W, H, date_label: str = "") -> bool:
    subs = _split_subblocks(block.lines)
    timeline_lines = subs.get("incident timeline")
    if not timeline_lines:
        return False
    steps = _parse_card_bullets(timeline_lines)
    if not steps:
        return False
    _render_timeline(prs, steps, block.heading, case_id, W, H, date_label)
    return True


def _render_root_cause_section(prs, block: _Block, case_id: str, W, H) -> bool:
    subs = _split_subblocks(block.lines)
    rca_lines = subs.get("root cause analysis")
    factor_lines = subs.get("contributing factors")
    risk_lines = subs.get("residual risks")
    if not (rca_lines and factor_lines):
        return False
    primary_cards = _parse_card_bullets(rca_lines)
    factors = _parse_card_bullets(factor_lines)
    if not primary_cards or not factors:
        return False
    _render_root_cause(prs, primary_cards[0], factors, block.heading, case_id, W, H)

    if risk_lines:
        risks = _parse_card_bullets(risk_lines)
        if risks:
            _render_card_grid(prs, risks, block.heading + " (cont.)", case_id, W, H, "root_cause", layout="2x2")
    return True


def _render_response_section(prs, block: _Block, case_id: str, W, H) -> bool:
    subs = _split_subblocks(block.lines)
    response_lines = subs.get("response actions")
    if not response_lines:
        return False
    cards = _parse_card_bullets(response_lines)
    if not cards:
        return False
    _render_card_grid(prs, cards, block.heading, case_id, W, H, "response", layout="2x2")
    return True


def _render_recommendations_section(prs, block: _Block, case_id: str, W, H) -> bool:
    subs = _split_subblocks(block.lines)
    priority_lines = subs.get("priority actions")
    if not priority_lines:
        return False
    rows = _parse_card_bullets(priority_lines)
    if not rows:
        return False
    _render_priority_rows(prs, rows, block.heading, case_id, W, H)
    return True


def _render_lessons_learned_section(prs, block: _Block, case_id: str, W, H) -> bool:
    subs = _split_subblocks(block.lines)
    insight_lines = subs.get("key insights")
    if not insight_lines:
        return False
    cards = _parse_card_bullets(insight_lines)
    if not cards:
        return False
    _render_card_grid(prs, cards, block.heading, case_id, W, H, "lessons_learned", layout="2x2")
    return True


_BOARD_SECTIONS: dict[str, tuple] = {
    "management summary": _render_management_summary,
    "business impact": _render_business_impact,
    "board timeline": _render_board_timeline,
    "root cause & risk": _render_root_cause_section,
    "response & containment": _render_response_section,
    "recommendations": _render_recommendations_section,
    "lessons learned": _render_lessons_learned_section,
}

# Headings that get the structured board-deck treatment; if a heading
# matches but parsing fails, fall back to the generic renderer instead.
_BOARD_HEADING_KEYS = set(_BOARD_SECTIONS) | {"1. management summary"}


def _board_section_handler(heading: str):
    key = _normalize_heading(heading)
    if key.startswith("management summary") or key.endswith("management summary"):
        return _render_management_summary
    return _BOARD_SECTIONS.get(key)


# ── Public API ───────────────────────────────────────────────────────────

_DEFAULT_BOARD_SECTIONS = [
    "Management Summary",
    "Business Impact",
    "Board Timeline",
    "Root Cause & Risk",
    "Response & Containment",
    "Recommendations",
    "Lessons Learned",
]


def convert(md_path: Path, output_path: Path, case_id: str = "",
            title: str = "", date_str: str = "",
            board_sections: list[str] | None = None) -> Path:
    """Render *md_path* (a hand-authored campaign report) as the board-deck
    PPTX. Returns the output path.

    Only ``##``/``###`` blocks whose heading (after stripping any leading
    numbering, case-insensitive) matches an entry in *board_sections*
    (default: the seven standard board-briefing sections) are rendered —
    this keeps the campaign deck a concise board-level summary rather than
    the full technical report. Sections matching the board bullet grammar
    (Management Summary, Business Impact, Board Timeline, Root Cause & Risk,
    Response & Containment, Recommendations, Lessons Learned) are rendered
    with the board-deck card/KPI/timeline/priority layouts; any matching
    section that doesn't follow the grammar falls back to the generic
    bullet/table/text rendering."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise SystemExit(
            "[board_deck] 'python-pptx' package not found.\n"
            "Install: pip3 install python-pptx"
        )

    md_path = Path(md_path)
    if not md_path.exists():
        raise FileNotFoundError(f"Markdown file not found: {md_path}")

    out = path_guard.assert_writable(Path(output_path))
    path_guard.guard_output_dir(out.parent)

    if not date_str:
        date_str = datetime.now(timezone.utc).strftime("%Y-%m-%d")

    text = md_path.read_text(encoding="utf-8")
    blocks = _split_blocks(text)
    fields = _header_fields(blocks)

    h1_block = next((b for b in blocks if b.level == 1), None)
    h1_title = h1_block.heading if h1_block else (title or md_path.stem.replace("_", " ").title())
    cover_title = case_id or title or h1_title

    severity_raw = fields.get("severity", "")
    if "—" in severity_raw or "-" in severity_raw:
        parts = re.split(r"\s*[—-]\s*", severity_raw, maxsplit=1)
        severity, severity_detail = parts[0], (parts[1] if len(parts) > 1 else "")
    else:
        severity, severity_detail = severity_raw, ""
    case_status = fields.get("case status", "")

    # Derive a date label (e.g. "27 AUGUST 2004") for the timeline kicker
    # from the header table's "Incident date" field, if present.
    incident_date = fields.get("incident date", "")
    date_label = ""
    dm = re.search(r"(\d{4}-\d{2}-\d{2})", incident_date)
    if dm:
        try:
            date_label = datetime.strptime(dm.group(1), "%Y-%m-%d").strftime("%-d %B %Y").upper()
        except ValueError:
            date_label = ""

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W, H = prs.slide_width, prs.slide_height

    cover_date = datetime.now(timezone.utc).strftime("%-d %B %Y")
    _render_cover(prs, cover_title, "Campaign Forensics Report", cover_date,
                  severity, severity_detail, case_status, W, H)

    sections = board_sections if board_sections is not None else _DEFAULT_BOARD_SECTIONS
    allowlist = {_normalize_heading(s) for s in sections}

    for block in blocks:
        if block.level not in (2, 3):
            continue
        if _normalize_heading(block.heading) not in allowlist:
            continue
        handler = _board_section_handler(block.heading)
        rendered = False
        if handler is not None:
            try:
                rendered = handler(prs, block, case_id, W, H, date_label) if handler is _render_board_timeline \
                    else handler(prs, block, case_id, W, H)
            except Exception as exc:  # noqa: BLE001
                print(f"[board_deck] WARNING: structured render failed for "
                      f"'{block.heading}' ({exc}); falling back to generic layout")
                rendered = False
        if not rendered:
            _render_block(prs, block, case_id, W, H)

    prs.save(str(out))
    print(f"[board_deck] PPTX saved: {out}")
    return out


# ── CLI ───────────────────────────────────────────────────────────────────

def _build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(description="Render a Markdown campaign report as a board-deck PPTX")
    p.add_argument("markdown", help="Path to input .md file")
    p.add_argument("--output", "-o", required=True, metavar="PPTX", help="Output PPTX path")
    p.add_argument("--case-id", metavar="ID", default="", help="Case ID shown on cover/header")
    p.add_argument("--title", metavar="TITLE", default="", help="Cover title (default: case ID or H1)")
    p.add_argument("--date", metavar="YYYY-MM-DD", default="", help="Report date (default: today UTC)")
    return p


if __name__ == "__main__":
    args = _build_parser().parse_args()
    convert(
        md_path=Path(args.markdown),
        output_path=Path(args.output),
        case_id=args.case_id,
        title=args.title,
        date_str=args.date,
    )
