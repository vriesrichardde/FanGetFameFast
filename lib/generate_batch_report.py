#!/usr/bin/env python3
# SPDX-License-Identifier: MIT OR Apache-2.0
# SPDX-FileCopyrightText: 2026 Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin
"""
generate_batch_report.py — Overall batch report generator for multi-file investigations.

After batch_analyze.sh has processed all evidence files, this module reads every
per-case module report from ./reports/, aggregates findings, and produces a single
batch investigation report (Markdown + PDF + PPTX + DOCX).

Claude: enhance and elaborate when necessary throughout the batch report to surface
cross-host patterns, common IOCs, and shared TTPs that no single per-case report
would surface alone.

Usage (CLI):
    python3 lib/generate_batch_report.py \\
        --batch-id   BATCH-20260528-130000 \\
        --manifest   ./batch_work/BATCH-20260528-130000/manifest.json \\
        [--reports-dir ./reports] \\
        [--output-dir  ./reports] \\
        [--no-upload]

Python API:
    from lib.generate_batch_report import generate
    paths = generate(batch_id="BATCH-20260528-130000",
                     manifest_path=Path("batch_work/.../manifest.json"))
"""
from __future__ import annotations

import argparse
import json
import re
import subprocess
import sys
from datetime import datetime, timezone
try:
    from zoneinfo import ZoneInfo
    _CET = ZoneInfo("Europe/Amsterdam")
except ImportError:
    _CET = timezone.utc
from pathlib import Path
from typing import Any

PROJECT_ROOT = Path(__file__).parent.parent

_DARK_NAVY  = (0x0f, 0x17, 0x2a)
_MID_NAVY   = (0x1e, 0x3a, 0x5f)
_BLUE       = (0x1d, 0x4e, 0xd8)
_LIGHT_BLUE = (0x93, 0xc5, 0xfd)
_WHITE      = (0xff, 0xff, 0xff)
_LIGHT_BG   = (0xf8, 0xfa, 0xfc)
_ROW_ALT    = (0xf1, 0xf5, 0xf9)
_TEXT_DARK  = (0x1f, 0x29, 0x37)
_TEXT_MID   = (0x6b, 0x72, 0x80)
_AMBER      = (0xfb, 0xbf, 0x24)
_GREEN      = (0x22, 0xc5, 0x5e)
_RED        = (0xef, 0x44, 0x44)


# ── Manifest loading ───────────────────────────────────────────────────────────

def _load_manifest(manifest_path: Path) -> dict[str, Any]:
    with open(manifest_path) as f:
        return json.load(f)


# ── Per-case report discovery ──────────────────────────────────────────────────

def _discover_case_reports(reports_dir: Path, case_id: str) -> dict[str, str]:
    """Find all module reports for a single case ID."""
    stem = case_id.replace(" ", "_")
    found: dict[str, str] = {}
    patterns = {
        "fame_md":     [f"{stem}_fame_report.md"],
        "fast_md":     [f"{stem}_fast_report.md"],
        "fan_md":      [f"{stem}_incident_report.md", f"{stem}_fan_report.md"],
        "combined_md": [f"{stem}_combined_report.md"],
    }
    for key, filenames in patterns.items():
        for fn in filenames:
            candidate = reports_dir / fn
            if candidate.exists():
                found[key] = candidate.read_text(errors="replace")
                break
    return found


# ── Section extraction helpers ────────────────────────────────────────────────

def _extract_section(md: str, marker: str, max_chars: int = 2000) -> str:
    lines = md.splitlines()
    in_section = False
    section_level = 0
    collected: list[str] = []
    for line in lines:
        if marker.lower() in line.lower() and line.startswith("#"):
            in_section = True
            section_level = len(line) - len(line.lstrip("#"))
            continue
        if in_section:
            if line.startswith("#") and (len(line) - len(line.lstrip("#"))) <= section_level:
                break
            collected.append(line)
    result = "\n".join(collected).strip()
    return result[:max_chars] if len(result) > max_chars else result


def _extract_summary(md: str) -> str:
    return _extract_section(md, "Management Summary", max_chars=600)


def _extract_mitre(md: str) -> str:
    return _extract_section(md, "MITRE ATT", max_chars=2000)


def _extract_iocs(md: str) -> str:
    return _extract_section(md, "Indicators of Compromise", max_chars=2000)


def _extract_recommendations(md: str) -> str:
    return _extract_section(md, "Recommendations", max_chars=1500)


def _extract_top_finding(md: str) -> str:
    """Pull the first non-empty paragraph from any module report as a top finding."""
    for line in md.splitlines():
        stripped = line.strip()
        if stripped and not stripped.startswith("#") and not stripped.startswith("|") \
                and not stripped.startswith(">") and len(stripped) > 30:
            clean = re.sub(r"\*\*(.*?)\*\*", r"\1", stripped)
            clean = re.sub(r"\*(.*?)\*", r"\1", clean)
            return clean[:120]
    return "No findings available."


# ── Unified data aggregation ──────────────────────────────────────────────────

def _aggregate_mitre(
    all_cases: list[dict[str, Any]],
) -> list[tuple[str, str, str, str, int]]:
    """
    Return deduplicated list of (technique_id, name, tactic, modules, freq)
    across all cases, sorted by frequency descending.
    """
    freq: dict[str, dict] = {}
    for case in all_cases:
        for key in ("fame_md", "fast_md", "fan_md", "combined_md"):
            md = case.get("reports", {}).get(key, "")
            if not md:
                continue
            mitre = _extract_mitre(md)
            for line in mitre.splitlines():
                if re.match(r"\|\s*\[?T\d", line):
                    parts = [p.strip() for p in line.strip("|").split("|")]
                    if len(parts) >= 4:
                        tid = re.sub(r"\[.*?\]\(.*?\)", lambda m: re.search(r"\[(.*?)\]", m.group()).group(1), parts[0])
                        name = parts[1][:50]
                        tactic = parts[2][:30]
                        if tid not in freq:
                            freq[tid] = {"name": name, "tactic": tactic, "modules": set(), "count": 0}
                        freq[tid]["modules"].add(key.replace("_md", "").upper())
                        freq[tid]["count"] += 1
    return [
        (tid, v["name"], v["tactic"], "+".join(sorted(v["modules"])), v["count"])
        for tid, v in sorted(freq.items(), key=lambda x: -x[1]["count"])
    ]


def _aggregate_iocs(all_cases: list[dict[str, Any]]) -> list[tuple[str, str, set[str]]]:
    """
    Return deduplicated list of (ioc_value, ioc_type, {hostnames}).
    """
    iocs: dict[str, dict] = {}
    for case in all_cases:
        hostname = case.get("hostname", "unknown")
        for key in ("fame_md", "fast_md", "fan_md", "combined_md"):
            md = case.get("reports", {}).get(key, "")
            if not md:
                continue
            ioc_section = _extract_iocs(md)
            for line in ioc_section.splitlines():
                if "|" not in line or "---" in line or "IOC" in line:
                    continue
                parts = [p.strip() for p in line.strip("|").split("|")]
                if len(parts) >= 2:
                    ioc_val = parts[0].strip("`").strip()
                    ioc_type = parts[1] if len(parts) > 1 else "unknown"
                    if ioc_val and len(ioc_val) > 3 and not ioc_val.startswith("-"):
                        if ioc_val not in iocs:
                            iocs[ioc_val] = {"type": ioc_type, "hosts": set()}
                        iocs[ioc_val]["hosts"].add(hostname)
    return [(v, d["type"], d["hosts"]) for v, d in iocs.items()]


def _aggregate_recs(all_cases: list[dict[str, Any]]) -> list[str]:
    """Return deduplicated list of recommendations across all cases."""
    seen: set[str] = set()
    result: list[str] = []
    for case in all_cases:
        for key in ("fame_md", "fast_md", "fan_md", "combined_md"):
            md = case.get("reports", {}).get(key, "")
            if not md:
                continue
            rec_section = _extract_recommendations(md)
            for line in rec_section.splitlines():
                line = line.strip()
                if line and (line[0].isdigit() or line.startswith("-")):
                    cleaned = re.sub(r"^\d+\.\s*", "", re.sub(r"^-\s*", "", line))
                    cleaned = re.sub(r"\*\*(.*?)\*\*", r"\1", cleaned).strip()
                    if cleaned and cleaned not in seen:
                        seen.add(cleaned)
                        result.append(cleaned[:130])
    return result


# ── Markdown ───────────────────────────────────────────────────────────────────

def _build_markdown(
    manifest: dict[str, Any],
    all_cases: list[dict[str, Any]],
    batch_id: str,
    generated_utc: str,
) -> str:
    """
    Build the overall batch investigation report.

    Claude: enhance and elaborate when necessary — this is the single document
    that integrates findings from all investigated hosts.  Cross-host patterns,
    common IOCs, and shared TTPs should be called out explicitly.
    """
    total = len(all_cases)
    succeeded = sum(1 for c in all_cases if c.get("status") == "success")
    failed     = total - succeeded
    fame_count = sum(1 for c in all_cases if "fame_md" in c.get("reports", {}))
    fast_count = sum(1 for c in all_cases if "fast_md" in c.get("reports", {}))
    fan_count  = sum(1 for c in all_cases if "fan_md"  in c.get("reports", {}))

    lines: list[str] = []
    a = lines.append

    a("# FanGetFameFast — Batch Investigation Report")
    a("")
    a("| Field | Value |")
    a("|-------|-------|")
    a(f"| Batch ID | `{batch_id}` |")
    a(f"| Evidence directory | `{manifest.get('evidence_dir', 'N/A')}` |")
    a(f"| Hosts analysed | {total} ({succeeded} succeeded, {failed} failed) |")
    a(f"| FAME (memory) | {fame_count} host(s) |")
    a(f"| FAST (storage) | {fast_count} host(s) |")
    a(f"| FAN (network) | {fan_count} PCAP(s) |")
    a(f"| Generated (UTC) | {generated_utc} |")
    a(f"| Prepared by | Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin |")
    a("")
    a("> **Claude: enhance and elaborate when necessary** — this batch report aggregates")
    a("> forensic findings from all investigated hosts. Cross-host IOC matches, shared")
    a("> TTPs, and common threat actor indicators should be highlighted.")
    a("")

    # ── 1. Management executive summary ───────────────────────────────────────
    a("---")
    a("")
    a("## 1. Management executive summary")
    a("")
    a("> **Audience:** CISO, Legal, Law Enforcement, Internal Audit — no technical identifiers.")
    a("> Claude: enhance and elaborate when necessary.")
    a("")
    a(f"A batch investigation was conducted across {total} host(s). "
      f"Memory forensics (FAME) was run on {fame_count} host(s), "
      f"storage forensics (FAST) on {fast_count} host(s), and "
      f"network forensics (FAN) on {fan_count} capture(s). "
      f"Of {total} host(s) analysed, {succeeded} completed successfully and {failed} encountered errors.")
    a("")
    a("> Claude: synthesise the following per-host management summaries into a unified")
    a("> executive narrative. Group hosts by observed behaviour. Do not include technical identifiers.")
    a("")
    for case in all_cases:
        hostname = case.get("hostname", "unknown")
        for key, label in [("combined_md", "Combined"), ("fame_md", "Memory"),
                            ("fast_md", "Storage"), ("fan_md", "Network")]:
            md = case.get("reports", {}).get(key, "")
            if not md:
                continue
            summary = _extract_summary(md)
            if summary:
                a(f"**{hostname} ({label}):**")
                a("")
                a(summary[:400])
                a("")
            break  # use the best available report per case
    a("")

    # ── 2. Per-case findings table ─────────────────────────────────────────────
    a("---")
    a("")
    a("## 2. Per-case findings summary")
    a("")
    a("| Hostname | Case ID | Modules run | Status | Top finding |")
    a("|----------|---------|-------------|--------|-------------|")
    for case in all_cases:
        hostname = case.get("hostname", "unknown")
        case_id  = case.get("case_id", "unknown")
        status   = case.get("status",  "unknown")
        reports  = case.get("reports", {})
        modules_run = []
        if "fan_md"  in reports: modules_run.append("FAN")
        if "fame_md" in reports: modules_run.append("FAME")
        if "fast_md" in reports: modules_run.append("FAST")
        modules_str = ", ".join(modules_run) if modules_run else "none"
        status_icon = "✓" if status == "success" else "✗"

        top_md = next(
            (reports[k] for k in ("combined_md", "fame_md", "fast_md", "fan_md") if k in reports),
            "",
        )
        top = _extract_top_finding(top_md) if top_md else "No report available"
        a(f"| `{hostname}` | `{case_id}` | {modules_str} | {status_icon} {status} | {top} |")
    a("")

    # ── 3. Unified IOCs ────────────────────────────────────────────────────────
    a("---")
    a("")
    a("## 3. Unified indicators of compromise")
    a("")
    a("> Claude: enhance and elaborate when necessary. IOCs present on multiple hosts")
    a("> indicate lateral movement, shared infrastructure, or a common threat actor.")
    a("")
    iocs = _aggregate_iocs(all_cases)
    if iocs:
        a("| Indicator | Type | Hosts |")
        a("|-----------|------|-------|")
        for ioc_val, ioc_type, hosts in sorted(iocs, key=lambda x: -len(x[2])):
            hosts_str = ", ".join(sorted(hosts))
            a(f"| `{ioc_val}` | {ioc_type} | {hosts_str} |")
    else:
        a("*No IOCs extracted from module reports.*")
    a("")

    # ── 4. Unified MITRE ATT&CK ────────────────────────────────────────────────
    a("---")
    a("")
    a("## 4. Unified MITRE ATT&CK coverage")
    a("")
    a("> Claude: enhance and elaborate when necessary. Techniques observed on multiple")
    a("> hosts indicate systematic attacker behaviour or automated tooling.")
    a("")
    techniques = _aggregate_mitre(all_cases)
    if techniques:
        a("| Technique | Name | Tactic | Modules | Frequency |")
        a("|-----------|------|--------|---------|-----------|")
        for tid, name, tactic, modules_str, freq in techniques:
            a(f"| `{tid}` | {name} | {tactic} | {modules_str} | {freq}× |")
    else:
        a("*No MITRE ATT&CK techniques extracted from module reports.*")
    a("")

    # ── 5. Unified recommendations ────────────────────────────────────────────
    a("---")
    a("")
    a("## 5. Recommendations")
    a("")
    a("> Claude: enhance and elaborate when necessary — merge, de-duplicate, and re-prioritise.")
    a("")
    recs = _aggregate_recs(all_cases)
    if recs:
        for i, rec in enumerate(recs, 1):
            a(f"{i}. {rec}")
    else:
        a("*No recommendations extracted from module reports.*")
    a("")

    # ── 6. Full per-case reports ───────────────────────────────────────────────
    a("---")
    a("")
    a("## 6. Individual case reports")
    a("")
    a("> The following section embeds the full report for each host.")
    a("> Claude: enhance and elaborate when necessary.")
    a("")
    for case in all_cases:
        hostname = case.get("hostname", "unknown")
        case_id  = case.get("case_id",  "unknown")
        a(f"### {hostname}  (`{case_id}`)")
        a("")
        reports = case.get("reports", {})
        for key, label in [
            ("combined_md", "Combined FAN+FAME+FAST"),
            ("fame_md",  "FAME — Memory forensics"),
            ("fast_md",  "FAST — Storage forensics"),
            ("fan_md",   "FAN — Network forensics"),
        ]:
            md = reports.get(key, "")
            if not md:
                continue
            a(f"#### {label}")
            a("")
            a(md)
            a("")

    a("---")
    a("")
    a("*End of batch report. Evidence integrity preserved. All findings cited to their source module.*")
    a("")

    return "\n".join(lines)


# ── PPTX ───────────────────────────────────────────────────────────────────────

def _build_pptx(
    manifest: dict[str, Any],
    all_cases: list[dict[str, Any]],
    batch_id: str,
    generated_utc: str,
    output_path: Path,
) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("[batch] WARNING: python-pptx not installed — skipping PPTX.")
        return

    total      = len(all_cases)
    succeeded  = sum(1 for c in all_cases if c.get("status") == "success")
    fame_count = sum(1 for c in all_cases if "fame_md" in c.get("reports", {}))
    fast_count = sum(1 for c in all_cases if "fast_md" in c.get("reports", {}))
    fan_count  = sum(1 for c in all_cases if "fan_md"  in c.get("reports", {}))
    techniques = _aggregate_mitre(all_cases)
    iocs       = _aggregate_iocs(all_cases)
    recs       = _aggregate_recs(all_cases)
    evidence_dir = manifest.get("evidence_dir", "N/A")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def _rgb(t):
        return RGBColor(*t)

    def _rect(slide, l, t, w, h, fill):
        s = slide.shapes.add_shape(1, l, t, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = _rgb(fill)
        s.line.fill.background()
        return s

    def _txt(slide, text, l, t, w, h, sz, bold=False, color=_WHITE, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)[:500]
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)

    W = prs.slide_width
    H = prs.slide_height
    M = Inches(0.4)

    # Slide 1 — Cover
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, W, H, _DARK_NAVY)
    _rect(s, 0, 0, W, Inches(0.08), _BLUE)
    _rect(s, 0, H - Inches(0.08), W, Inches(0.08), _BLUE)
    _txt(s, "Fan Get Fame Fast", M, Inches(0.9), W - 2*M, Inches(1.2),
         52, bold=True, color=_LIGHT_BLUE, align=PP_ALIGN.CENTER)
    _txt(s, "Batch Investigation Report", M, Inches(2.0), W - 2*M, Inches(0.7),
         28, color=_WHITE, align=PP_ALIGN.CENTER)
    _txt(s, "FAN  ·  FAME  ·  FAST", M, Inches(2.65), W - 2*M, Inches(0.5),
         18, color=_LIGHT_BLUE, align=PP_ALIGN.CENTER)
    _rect(s, Inches(3), Inches(3.4), W - Inches(6), Inches(0.04), _BLUE)
    _txt(s, f"Batch: {batch_id}  |  {total} host(s)  |  {generated_utc[:10]}",
         M, Inches(3.7), W - 2*M, Inches(0.5), 14, color=_TEXT_MID, align=PP_ALIGN.CENTER)
    _txt(s, f"Evidence: {evidence_dir}",
         M, Inches(4.2), W - 2*M, Inches(0.4), 12, color=_TEXT_MID, align=PP_ALIGN.CENTER)
    _txt(s, "Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin",
         M, Inches(4.7), W - 2*M, Inches(0.4), 11, color=_TEXT_MID, align=PP_ALIGN.CENTER)
    _txt(s, "CONFIDENTIAL — FOR AUTHORISED PERSONNEL ONLY",
         M, H - Inches(0.7), W - 2*M, Inches(0.4), 11, color=_TEXT_MID, align=PP_ALIGN.CENTER)

    # Slide 2 — Batch overview (statistics)
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, W, Inches(1.1), _MID_NAVY)
    _txt(s, "Batch investigation overview", M, Inches(0.2), W, Inches(0.8),
         28, bold=True, color=_WHITE)
    _rect(s, 0, Inches(1.1), W, H - Inches(1.1), _LIGHT_BG)

    stats = [
        ("Hosts analysed",   str(total),                    _MID_NAVY),
        ("Succeeded",        str(succeeded),                 _GREEN),
        ("Failed",           str(total - succeeded),         _RED if total - succeeded > 0 else _TEXT_MID),
        ("FAME (memory)",    f"{fame_count} host(s)",        _BLUE),
        ("FAST (storage)",   f"{fast_count} host(s)",        _BLUE),
        ("FAN (network)",    f"{fan_count} capture(s)",      _BLUE),
        ("Unique IOCs",      str(len(iocs)),                 _AMBER),
        ("MITRE techniques", str(len(techniques)),           _AMBER),
    ]
    col_w = (W - 2*M - Inches(0.2)) // 4
    for i, (label, value, color) in enumerate(stats):
        col = i % 4
        row = i // 4
        cx = M + col * (col_w + Inches(0.07))
        cy = Inches(1.3) + row * Inches(2.6)
        _rect(s, cx, cy, col_w, Inches(2.4), _WHITE)
        _rect(s, cx, cy, col_w, Inches(0.08), color)
        _txt(s, value, cx + Inches(0.1), cy + Inches(0.3), col_w - Inches(0.2), Inches(1.2),
             42, bold=True, color=_rgb(color) and _TEXT_DARK, align=PP_ALIGN.CENTER)
        _txt(s, label, cx + Inches(0.1), cy + Inches(1.6), col_w - Inches(0.2), Inches(0.7),
             12, color=_TEXT_MID, align=PP_ALIGN.CENTER)
    _txt(s, "Claude: enhance and elaborate when necessary", M, H - Inches(0.4), W - 2*M, Inches(0.3),
         9, color=_TEXT_MID)

    # Slide 3 — Per-host summary table
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, W, Inches(1.1), _MID_NAVY)
    _txt(s, "Per-host findings summary", M, Inches(0.2), W, Inches(0.8),
         28, bold=True, color=_WHITE)

    headers = ["Hostname", "Case ID", "Modules", "Status", "Top finding"]
    col_ws  = [Inches(2.0), Inches(2.4), Inches(1.5), Inches(1.2), W - M - Inches(7.5)]
    row_h   = Inches(0.63)
    hx = M
    for h, cw in zip(headers, col_ws):
        _rect(s, hx, Inches(1.2), cw - Inches(0.05), row_h - Inches(0.04), _MID_NAVY)
        _txt(s, h, hx + Inches(0.08), Inches(1.25), cw, row_h, 12, bold=True, color=_WHITE)
        hx += cw
    for i, case in enumerate(all_cases[:8]):
        y = Inches(1.2) + (i + 1) * row_h
        bg = _LIGHT_BG if i % 2 == 0 else _ROW_ALT
        rx = M
        hostname = case.get("hostname", "?")[:25]
        case_id_short = case.get("case_id", "?")[-25:]
        reports = case.get("reports", {})
        modules = "+".join(k.replace("_md","").upper() for k in ("fan_md","fame_md","fast_md") if k in reports)
        status = "✓" if case.get("status") == "success" else "✗"
        top_md = next((reports[k] for k in ("combined_md","fame_md","fast_md","fan_md") if k in reports), "")
        top = _extract_top_finding(top_md)[:60] if top_md else "—"
        for val, cw in zip([hostname, case_id_short, modules or "—", status, top], col_ws):
            _rect(s, rx, y, cw - Inches(0.05), row_h - Inches(0.04), bg)
            _txt(s, val, rx + Inches(0.08), y + Inches(0.08), cw - Inches(0.13), row_h,
                 9, color=_TEXT_DARK)
            rx += cw
    if len(all_cases) > 8:
        _txt(s, f"… and {len(all_cases)-8} more — see full report",
             M, Inches(1.2) + 9 * row_h + Inches(0.05), W - 2*M, Inches(0.4),
             10, color=_TEXT_MID)
    _txt(s, "Claude: enhance and elaborate when necessary", M, H - Inches(0.4), W - 2*M, Inches(0.3),
         9, color=_TEXT_MID)

    # Slide 4 — Unified IOCs
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, W, Inches(1.1), _MID_NAVY)
    _txt(s, "Unified indicators of compromise", M, Inches(0.2), W, Inches(0.8),
         28, bold=True, color=_WHITE)
    if iocs:
        multi_host = [(v, t, h) for v, t, h in iocs if len(h) > 1]
        single_host = [(v, t, h) for v, t, h in iocs if len(h) == 1]
        shown = sorted(multi_host, key=lambda x: -len(x[2])) + single_host
        headers2 = ["Indicator", "Type", "Hosts (count)", "Hostnames"]
        col_ws2  = [Inches(3.5), Inches(1.5), Inches(1.2), W - M - Inches(6.6)]
        row_h2   = Inches(0.58)
        hx = M
        for h, cw in zip(headers2, col_ws2):
            _rect(s, hx, Inches(1.2), cw - Inches(0.05), row_h2 - Inches(0.04), _MID_NAVY)
            _txt(s, h, hx + Inches(0.08), Inches(1.25), cw, row_h2, 12, bold=True, color=_WHITE)
            hx += cw
        for i, (val, ioc_type, hosts) in enumerate(shown[:9]):
            y = Inches(1.2) + (i + 1) * row_h2
            bg = _LIGHT_BG if i % 2 == 0 else _ROW_ALT
            highlight = len(hosts) > 1
            rx = M
            for cell_val, cw in zip(
                [val[:45], ioc_type[:18], str(len(hosts)), ", ".join(sorted(hosts))[:60]],
                col_ws2,
            ):
                _rect(s, rx, y, cw - Inches(0.05), row_h2 - Inches(0.04),
                      _AMBER if highlight else bg)
                _txt(s, cell_val, rx + Inches(0.08), y + Inches(0.07), cw - Inches(0.13), row_h2,
                     9, color=_TEXT_DARK)
                rx += cw
    else:
        _txt(s, "No IOCs extracted from module reports.",
             M, Inches(2.0), W - 2*M, Inches(1.0), 16, color=_TEXT_MID)
    _txt(s, "Claude: enhance and elaborate when necessary", M, H - Inches(0.4), W - 2*M, Inches(0.3),
         9, color=_TEXT_MID)

    # Slide 5 — Unified MITRE ATT&CK
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, W, Inches(1.1), _MID_NAVY)
    _txt(s, "Unified MITRE ATT&CK coverage", M, Inches(0.2), W, Inches(0.8),
         28, bold=True, color=_WHITE)
    if techniques:
        headers3 = ["Technique", "Name", "Tactic", "Modules", "Frequency"]
        col_ws3  = [Inches(1.3), Inches(2.2), Inches(2.0), Inches(1.5), W - M - Inches(7.8)]
        row_h3   = Inches(0.63)
        hx = M
        for h, cw in zip(headers3, col_ws3):
            _rect(s, hx, Inches(1.2), cw - Inches(0.05), row_h3 - Inches(0.04), _MID_NAVY)
            _txt(s, h, hx + Inches(0.08), Inches(1.25), cw, row_h3, 12, bold=True, color=_WHITE)
            hx += cw
        for i, (tid, name, tactic, mods, freq) in enumerate(techniques[:8]):
            y = Inches(1.2) + (i + 1) * row_h3
            bg = _LIGHT_BG if i % 2 == 0 else _ROW_ALT
            rx = M
            for val, cw in zip([tid, name[:30], tactic[:25], mods, f"{freq}×"], col_ws3):
                _rect(s, rx, y, cw - Inches(0.05), row_h3 - Inches(0.04), bg)
                _txt(s, val, rx + Inches(0.08), y + Inches(0.08), cw - Inches(0.13), row_h3,
                     10, color=_TEXT_DARK)
                rx += cw
    else:
        _txt(s, "No MITRE ATT&CK techniques mapped.", M, Inches(2.0), W - 2*M, Inches(1.0),
             16, color=_TEXT_MID)
    _txt(s, "Claude: enhance and elaborate when necessary", M, H - Inches(0.4), W - 2*M, Inches(0.3),
         9, color=_TEXT_MID)

    # Slide 6 — Recommendations
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, W, Inches(1.1), _MID_NAVY)
    _txt(s, "Recommendations", M, Inches(0.2), W, Inches(0.8), 28, bold=True, color=_WHITE)
    if recs:
        row_h4 = Inches(0.62)
        for i, rec in enumerate(recs[:8]):
            y = Inches(1.2) + i * row_h4
            _rect(s, M, y, Inches(0.45), row_h4 - Inches(0.08), _BLUE)
            _txt(s, str(i+1), M + Inches(0.08), y + Inches(0.08), Inches(0.3), row_h4,
                 14, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
            _txt(s, rec, M + Inches(0.55), y + Inches(0.08), W - M - Inches(0.95), row_h4,
                 12, color=_TEXT_DARK)
    else:
        _txt(s, "No recommendations extracted from module reports.",
             M, Inches(2.0), W - 2*M, Inches(1.0), 16, color=_TEXT_MID)
    _txt(s, "Claude: enhance and elaborate when necessary", M, H - Inches(0.4), W - 2*M, Inches(0.3),
         9, color=_TEXT_MID)

    # Slide 7 — Case index
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 0, W, Inches(1.1), _MID_NAVY)
    _txt(s, f"Investigated cases ({total})", M, Inches(0.2), W, Inches(0.8),
         28, bold=True, color=_WHITE)
    _rect(s, 0, Inches(1.1), W, H - Inches(1.1), _LIGHT_BG)
    col_count = 2
    items_per_col = (len(all_cases) + col_count - 1) // col_count
    col_w7 = (W - 2*M - Inches(0.4)) / col_count
    for i, case in enumerate(all_cases):
        col = i // items_per_col
        row = i % items_per_col
        cx = M + col * (col_w7 + Inches(0.4))
        cy = Inches(1.3) + row * Inches(0.46)
        hostname = case.get("hostname", "?")
        reports  = case.get("reports", {})
        mods = "+".join(k.replace("_md","").upper() for k in ("fan_md","fame_md","fast_md") if k in reports)
        status_dot = "●" if case.get("status") == "success" else "○"
        _txt(s, f"{status_dot}  {hostname}  [{mods or '—'}]",
             cx, cy, col_w7, Inches(0.4), 11, color=_TEXT_DARK)

    prs.save(str(output_path))
    print(f"[batch] PPTX saved: {output_path}")


# ── DOCX ───────────────────────────────────────────────────────────────────────

def _build_docx(
    manifest: dict[str, Any],
    all_cases: list[dict[str, Any]],
    batch_id: str,
    generated_utc: str,
    output_path: Path,
) -> None:
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        print("[batch] WARNING: python-docx not installed — skipping DOCX.")
        return

    total      = len(all_cases)
    succeeded  = sum(1 for c in all_cases if c.get("status") == "success")
    fame_count = sum(1 for c in all_cases if "fame_md" in c.get("reports", {}))
    fast_count = sum(1 for c in all_cases if "fast_md" in c.get("reports", {}))
    fan_count  = sum(1 for c in all_cases if "fan_md"  in c.get("reports", {}))
    iocs       = _aggregate_iocs(all_cases)
    techniques = _aggregate_mitre(all_cases)
    recs       = _aggregate_recs(all_cases)

    doc = Document()
    for section in doc.sections:
        section.top_margin    = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin   = Inches(1.2)
        section.right_margin  = Inches(1.2)

    def _xml_safe(text: str) -> str:
        return re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", str(text))

    def _h(text, level):
        p = doc.add_heading(_xml_safe(text), level=level)
        p.runs[0].font.color.rgb = RGBColor(0x0f, 0x17, 0x2a)

    def _p(text, bold=False, italic=False):
        p = doc.add_paragraph()
        run = p.add_run(_xml_safe(text))
        run.bold   = bold
        run.italic = italic

    def _note(text):
        p = doc.add_paragraph()
        run = p.add_run(_xml_safe(text))
        run.italic = True
        run.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)

    # Cover
    doc.add_paragraph()
    t = doc.add_heading("Fan Get Fame Fast — Batch Investigation Report", 0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = sub.add_run("FAN · FAME · FAST  |  Multi-host Batch Report")
    r.font.size = Pt(14)
    r.font.color.rgb = RGBColor(0x1d, 0x4e, 0xd8)
    doc.add_paragraph()

    tbl = doc.add_table(rows=7, cols=2)
    tbl.style = "Table Grid"
    meta_rows = [
        ("Batch ID",        batch_id),
        ("Evidence dir",    manifest.get("evidence_dir", "N/A")),
        ("Hosts analysed",  f"{total} ({succeeded} succeeded, {total-succeeded} failed)"),
        ("FAME (memory)",   f"{fame_count} host(s)"),
        ("FAST (storage)",  f"{fast_count} host(s)"),
        ("FAN (network)",   f"{fan_count} capture(s)"),
        ("Generated UTC",   generated_utc),
    ]
    for i, (label, value) in enumerate(meta_rows):
        tbl.rows[i].cells[0].text = label
        tbl.rows[i].cells[1].text = value

    doc.add_paragraph()
    conf = doc.add_paragraph("CONFIDENTIAL — FOR AUTHORISED PERSONNEL ONLY")
    conf.alignment = WD_ALIGN_PARAGRAPH.CENTER
    conf.runs[0].font.bold = True
    conf.runs[0].font.color.rgb = RGBColor(0xef, 0x44, 0x44)
    doc.add_page_break()

    # 1. Management summary
    _h("1. Management executive summary", 1)
    _note("Audience: CISO, Legal, Law Enforcement, Internal Audit. "
          "Claude: enhance and elaborate when necessary.")
    _p(f"A batch investigation was conducted across {total} host(s). "
       f"Memory forensics (FAME) was run on {fame_count} host(s), "
       f"storage forensics (FAST) on {fast_count} host(s), and "
       f"network forensics (FAN) on {fan_count} capture(s). "
       f"Of {total} host(s) analysed, {succeeded} completed successfully.")
    doc.add_paragraph()
    for case in all_cases:
        hostname = case.get("hostname", "unknown")
        for key, label in [("combined_md","Combined"), ("fame_md","Memory"),
                            ("fast_md","Storage"), ("fan_md","Network")]:
            md = case.get("reports", {}).get(key, "")
            if not md:
                continue
            summary = _extract_summary(md)
            if summary:
                _h(f"{hostname} ({label})", 2)
                _p(summary[:400].strip())
                doc.add_paragraph()
            break
    doc.add_page_break()

    # 2. Per-case table
    _h("2. Per-case findings summary", 1)
    tbl2 = doc.add_table(rows=1 + len(all_cases), cols=4)
    tbl2.style = "Table Grid"
    for j, hdr in enumerate(["Hostname", "Modules", "Status", "Top finding"]):
        tbl2.rows[0].cells[j].text = hdr
    for i, case in enumerate(all_cases):
        hostname = case.get("hostname", "unknown")
        reports  = case.get("reports", {})
        mods = "+".join(k.replace("_md","").upper() for k in ("fan_md","fame_md","fast_md") if k in reports)
        status = "success" if case.get("status") == "success" else "failed"
        top_md = next((reports[k] for k in ("combined_md","fame_md","fast_md","fan_md") if k in reports), "")
        top = _extract_top_finding(top_md) if top_md else "N/A"
        for j, val in enumerate([hostname, mods or "—", status, top]):
            tbl2.rows[i+1].cells[j].text = val
    doc.add_page_break()

    # 3. Unified IOCs
    _h("3. Unified indicators of compromise", 1)
    _note("Claude: enhance and elaborate when necessary. "
          "IOCs on multiple hosts indicate lateral movement or shared infrastructure.")
    if iocs:
        tbl3 = doc.add_table(rows=1 + len(iocs), cols=3)
        tbl3.style = "Table Grid"
        for j, hdr in enumerate(["Indicator", "Type", "Hosts"]):
            tbl3.rows[0].cells[j].text = hdr
        for i, (val, ioc_type, hosts) in enumerate(sorted(iocs, key=lambda x: -len(x[2]))):
            tbl3.rows[i+1].cells[0].text = val
            tbl3.rows[i+1].cells[1].text = ioc_type
            tbl3.rows[i+1].cells[2].text = ", ".join(sorted(hosts))
    else:
        _p("No IOCs extracted from module reports.")
    doc.add_page_break()

    # 4. Unified MITRE
    _h("4. Unified MITRE ATT&CK coverage", 1)
    _note("Claude: enhance and elaborate when necessary.")
    if techniques:
        tbl4 = doc.add_table(rows=1 + len(techniques), cols=5)
        tbl4.style = "Table Grid"
        for j, hdr in enumerate(["Technique", "Name", "Tactic", "Modules", "Frequency"]):
            tbl4.rows[0].cells[j].text = hdr
        for i, (tid, name, tactic, mods, freq) in enumerate(techniques):
            for j, val in enumerate([tid, name, tactic, mods, str(freq)]):
                tbl4.rows[i+1].cells[j].text = val
    else:
        _p("No MITRE ATT&CK techniques extracted.")
    doc.add_page_break()

    # 5. Recommendations
    _h("5. Recommendations", 1)
    _note("Claude: enhance and elaborate when necessary.")
    if recs:
        for i, rec in enumerate(recs, 1):
            p = doc.add_paragraph(style="List Number")
            p.add_run(rec)
    else:
        _p("No recommendations extracted.")
    doc.add_page_break()

    # 6. Individual case reports
    _h("6. Individual case reports", 1)
    for case in all_cases:
        hostname = case.get("hostname", "unknown")
        case_id  = case.get("case_id",  "unknown")
        _h(f"{hostname}  ({case_id})", 2)
        reports = case.get("reports", {})
        for key, label in [
            ("combined_md", "Combined FAN+FAME+FAST"),
            ("fame_md",  "FAME — Memory forensics"),
            ("fast_md",  "FAST — Storage forensics"),
            ("fan_md",   "FAN — Network forensics"),
        ]:
            md = reports.get(key, "")
            if not md:
                continue
            _h(label, 3)
            _note("Claude: enhance and elaborate when necessary.")
            for line in md.splitlines():
                stripped = line.strip()
                if not stripped:
                    doc.add_paragraph()
                elif stripped.startswith("# "):
                    doc.add_heading(stripped[2:], 4)
                elif stripped.startswith("## "):
                    doc.add_heading(stripped[3:], 4)
                elif stripped.startswith("```"):
                    pass
                elif stripped.startswith("|") and "---" not in stripped:
                    _p(stripped)
                elif not stripped.startswith("|"):
                    clean = re.sub(r"\*\*(.*?)\*\*", r"\1", stripped)
                    clean = re.sub(r"\*(.*?)\*",     r"\1", clean)
                    clean = re.sub(r"`(.*?)`",        r"\1", clean)
                    if clean:
                        _p(clean)
            break  # only embed the best available report per case
        doc.add_page_break()

    doc.save(str(output_path))
    print(f"[batch] DOCX saved: {output_path}")


# ── Public API ─────────────────────────────────────────────────────────────────

def generate(
    batch_id: str,
    manifest_path: Path,
    reports_dir: Path | None = None,
    output_dir: Path | None = None,
    no_upload: bool = False,
) -> dict[str, Path | None]:
    reports_dir = reports_dir or (PROJECT_ROOT / "reports")
    output_dir  = output_dir  or reports_dir
    output_dir.mkdir(parents=True, exist_ok=True)

    manifest = _load_manifest(manifest_path)

    # Attach discovered reports to each case entry
    all_cases: list[dict[str, Any]] = []
    for entry in manifest.get("cases", []):
        case_id = entry.get("case_id", "")
        entry["reports"] = _discover_case_reports(reports_dir, case_id)
        all_cases.append(entry)

    if not all_cases:
        print("[batch] WARNING: No cases found in manifest — batch report will be empty.")

    generated_utc = datetime.now(_CET).strftime("%Y-%m-%d %H:%M CET")
    stem = batch_id.replace(" ", "_")

    # Markdown
    md_text  = _build_markdown(manifest, all_cases, batch_id, generated_utc)
    md_path  = output_dir / f"{stem}_batch_report.md"
    md_path.write_text(md_text)
    print(f"[batch] Markdown saved: {md_path}")

    # PDF
    pdf_path: Path | None = None
    try:
        sys.path.insert(0, str(PROJECT_ROOT / "lib"))
        from md_to_pdf import convert as md2pdf
        pdf_path = output_dir / f"{stem}_batch_report.pdf"
        md2pdf(md_path, pdf_path)
        print(f"[batch] PDF saved: {pdf_path}")
    except Exception as exc:
        print(f"[batch] WARNING: PDF generation failed: {exc}")

    # PPTX
    pptx_path = output_dir / f"{stem}_batch_presentation.pptx"
    _build_pptx(manifest, all_cases, batch_id, generated_utc, pptx_path)

    # DOCX
    docx_path = output_dir / f"{stem}_batch_report.docx"
    _build_docx(manifest, all_cases, batch_id, generated_utc, docx_path)

    # Upload
    if not no_upload:
        try:
            cmd = [
                sys.executable,
                str(PROJECT_ROOT / "lib" / "investigations_upload.py"),
                "--case-id", batch_id,
                "--md",      str(md_path),
            ]
            if pdf_path and pdf_path.exists():
                cmd += ["--pdf", str(pdf_path)]
            if pptx_path.exists():
                cmd += ["--pptx", str(pptx_path)]
            if docx_path.exists():
                cmd += ["--docx", str(docx_path)]
            subprocess.run(cmd, check=True)
            print(f"[batch] Uploaded to vault under case ID: {batch_id}")
        except Exception as exc:
            print(f"[batch] WARNING: Upload failed: {exc}")

    return {
        "md":   md_path,
        "pdf":  pdf_path,
        "pptx": pptx_path if pptx_path.exists() else None,
        "docx": docx_path if docx_path.exists() else None,
    }


# ── CLI ────────────────────────────────────────────────────────────────────────

def _build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(description="FanGetFameFast — Batch Report Generator")
    p.add_argument("--batch-id",    required=True,  metavar="ID")
    p.add_argument("--manifest",    required=True,  metavar="FILE")
    p.add_argument("--reports-dir", default=None,   metavar="DIR")
    p.add_argument("--output-dir",  default=None,   metavar="DIR")
    p.add_argument("--no-upload",   action="store_true")
    return p


if __name__ == "__main__":
    args = _build_parser().parse_args()
    paths = generate(
        batch_id      = args.batch_id,
        manifest_path = Path(args.manifest),
        reports_dir   = Path(args.reports_dir) if args.reports_dir else None,
        output_dir    = Path(args.output_dir)  if args.output_dir  else None,
        no_upload     = args.no_upload,
    )
    print("[batch] Report suite complete:")
    for fmt, p in paths.items():
        if p:
            print(f"  {fmt.upper():4s}  {p}")
