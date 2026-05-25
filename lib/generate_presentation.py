#!/usr/bin/env python3
# SPDX-License-Identifier: MIT OR Apache-2.0
# SPDX-FileCopyrightText: 2026 Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin
"""
generate_presentation.py — PCAP Incident PowerPoint Presentation Generator

Produces a styled .pptx executive briefing from the analysis outputs written
by the 22 FAN modules.  Reuses load_all_data() and _overall_severity() from
generate_pcap_report so there is no duplication of data-loading logic.

Usage (standalone):
    python3 lib/generate_presentation.py --stem <stem> [options]

Via shell wrapper:
    scripts/generate_pcap_presentation.sh --stem <stem> --case-id <id> ...
"""

from __future__ import annotations

import argparse
import json
import sys
from datetime import datetime, timezone
from pathlib import Path
from textwrap import wrap

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

# ── project path ─────────────────────────────────────────────────────────────
PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT / "lib"))

from generate_pcap_report import (  # noqa: E402
    _overall_severity,
    _sev_rank,
    load_all_data,
)

# ── colour palette ────────────────────────────────────────────────────────────
C_NAVY    = RGBColor(0x0D, 0x1B, 0x2A)   # slide header bar background
C_BLUE    = RGBColor(0x1B, 0x4F, 0x72)   # accent / table header
C_SILVER  = RGBColor(0xF0, 0xF3, 0xF4)   # alternating row background
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK   = RGBColor(0x00, 0x00, 0x00)
C_RED     = RGBColor(0xC0, 0x39, 0x2B)   # CRITICAL
C_ORANGE  = RGBColor(0xCA, 0x6F, 0x1E)   # HIGH
C_YELLOW  = RGBColor(0xB7, 0x95, 0x0B)   # MEDIUM
C_GREEN   = RGBColor(0x1E, 0x8B, 0x4C)   # LOW / clean
C_GREY    = RGBColor(0x71, 0x7D, 0x7E)   # INFO / none

SEV_COLOUR = {
    "critical": C_RED,
    "high":     C_ORANGE,
    "medium":   C_YELLOW,
    "low":      C_GREEN,
    "info":     C_GREY,
    "none":     C_GREEN,
}

# Slide canvas: widescreen 16:9
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

HEADER_H  = Cm(2.2)
FOOTER_H  = Cm(0.8)
MARGIN_L  = Cm(1.2)
MARGIN_R  = Cm(1.2)
CONTENT_Y = HEADER_H + Cm(0.3)
CONTENT_H = SLIDE_H - HEADER_H - FOOTER_H - Cm(0.3)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R


# ── low-level helpers ─────────────────────────────────────────────────────────

def _rgb(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(r, g, b)


def _add_textbox(slide, left, top, width, height,
                 text: str = "",
                 font_size: int = 12,
                 bold: bool = False,
                 italic: bool = False,
                 colour: RGBColor = C_BLACK,
                 bg: RGBColor | None = None,
                 align=PP_ALIGN.LEFT,
                 wrap_text: bool = True) -> object:
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap_text
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = colour
    if bg is not None:
        from pptx.oxml.ns import qn
        from lxml import etree
        sp = txb._element
        spPr = sp.find(qn("p:spPr"))
        if spPr is None:
            spPr = etree.SubElement(sp, qn("p:spPr"))
        solidFill = etree.SubElement(spPr, qn("a:solidFill"))
        srgbClr   = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", f"{bg.rgb:06X}")
    return txb


def _header(slide, title: str, subtitle: str = "") -> None:
    """Navy header bar across the full slide width."""
    from pptx.oxml.ns import qn
    from lxml import etree

    # background rectangle
    rect = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, SLIDE_W, HEADER_H,
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = C_NAVY
    rect.line.fill.background()

    # title text
    _add_textbox(slide,
                 MARGIN_L, Cm(0.3),
                 SLIDE_W - MARGIN_L * 2, Cm(1.2),
                 text=title,
                 font_size=22, bold=True, colour=C_WHITE,
                 align=PP_ALIGN.LEFT)

    if subtitle:
        _add_textbox(slide,
                     MARGIN_L, Cm(1.3),
                     SLIDE_W - MARGIN_L * 2, Cm(0.8),
                     text=subtitle,
                     font_size=12, colour=RGBColor(0xAE, 0xD6, 0xF1),
                     align=PP_ALIGN.LEFT)


def _footer(slide, case_id: str, page_num: int) -> None:
    y = SLIDE_H - FOOTER_H
    _add_textbox(slide, MARGIN_L, y, Cm(12), FOOTER_H,
                 text=f"Case: {case_id}  ·  CONFIDENTIAL",
                 font_size=8, colour=C_GREY)
    _add_textbox(slide, SLIDE_W - Cm(4), y, Cm(3.8), FOOTER_H,
                 text=f"Slide {page_num}  ·  FAN",
                 font_size=8, colour=C_GREY, align=PP_ALIGN.RIGHT)


def _severity_badge(slide, sev: str, left, top, width=Cm(2.8), height=Cm(0.55)) -> None:
    colour = SEV_COLOUR.get(sev.lower(), C_GREY)
    rect = slide.shapes.add_shape(1, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = colour
    rect.line.fill.background()
    tf = rect.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = sev.upper()
    run.font.size  = Pt(8)
    run.font.bold  = True
    run.font.color.rgb = C_WHITE


def _bullet_slide(prs, title: str, subtitle: str, items: list[str],
                  case_id: str, page_num: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _header(slide, title, subtitle)
    _footer(slide, case_id, page_num)

    tf_left = MARGIN_L
    tf_top  = CONTENT_Y
    tf_w    = CONTENT_W
    tf_h    = CONTENT_H

    txb = slide.shapes.add_textbox(tf_left, tf_top, tf_w, tf_h)
    tf  = txb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p   = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text           = f"• {item}"
        run.font.size      = Pt(13)
        run.font.color.rgb = C_BLACK
        p.space_after      = Pt(4)


def _table_slide(prs, title: str, subtitle: str,
                 headers: list[str], rows: list[list[str]],
                 case_id: str, page_num: int,
                 col_widths: list[float] | None = None) -> None:
    """Generic table slide with navy header row and alternating row colours."""
    if not rows:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, title, subtitle)
    _footer(slide, case_id, page_num)

    ncols  = len(headers)
    nrows  = len(rows) + 1  # +1 for header

    # Cap rows that would overflow the slide
    max_rows = 18
    rows = rows[:max_rows]
    nrows = len(rows) + 1

    if col_widths is None:
        col_widths = [1.0 / ncols] * ncols

    row_h   = Cm(0.6)
    total_w = CONTENT_W
    col_cms = [int(total_w * w) for w in col_widths]

    table = slide.shapes.add_table(
        nrows, ncols,
        MARGIN_L, CONTENT_Y,
        total_w, int(Cm(row_h.cm * nrows)),
    ).table

    # column widths
    for ci, cw in enumerate(col_cms):
        table.columns[ci].width = cw

    # header row
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_BLUE
        tf_h = cell.text_frame
        tf_h.paragraphs[0].clear()
        tf_h.paragraphs[0].alignment = PP_ALIGN.CENTER
        run_h = tf_h.paragraphs[0].add_run()
        run_h.text = hdr
        run_h.font.size  = Pt(9)
        run_h.font.bold  = True
        run_h.font.color.rgb = C_WHITE

    # data rows
    for ri, row in enumerate(rows):
        bg = C_SILVER if ri % 2 == 0 else C_WHITE
        for ci, val in enumerate(row[:ncols]):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf3 = cell.text_frame
            tf3.paragraphs[0].clear()
            run3 = tf3.paragraphs[0].add_run()
            run3.text = str(val)[:120]
            run3.font.size = Pt(8)
            run3.font.color.rgb = C_BLACK


# ── slide builders ────────────────────────────────────────────────────────────

def _slide_title(prs, stem: str, case_id: str, overall_sev: str,
                 stats: dict, generated: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # full-slide navy background
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_NAVY
    bg.line.fill.background()

    # accent bar
    bar = slide.shapes.add_shape(1, 0, Cm(7.5), SLIDE_W, Cm(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = SEV_COLOUR.get(overall_sev.lower(), C_ORANGE)
    bar.line.fill.background()

    _add_textbox(slide, MARGIN_L, Cm(2.5), SLIDE_W - MARGIN_L * 2, Cm(2.2),
                 text="PCAP Incident Briefing",
                 font_size=16, colour=RGBColor(0xAE, 0xD6, 0xF1),
                 align=PP_ALIGN.CENTER)

    _add_textbox(slide, MARGIN_L, Cm(4.2), SLIDE_W - MARGIN_L * 2, Cm(2.8),
                 text=stem,
                 font_size=36, bold=True, colour=C_WHITE,
                 align=PP_ALIGN.CENTER)

    _add_textbox(slide, MARGIN_L, Cm(7.8), SLIDE_W - MARGIN_L * 2, Cm(1.0),
                 text=f"Overall Severity: {overall_sev.upper()}",
                 font_size=18, bold=True,
                 colour=SEV_COLOUR.get(overall_sev.lower(), C_ORANGE),
                 align=PP_ALIGN.CENTER)

    stat_line = (
        f"{stats.get('flows', 0)} flows  ·  "
        f"{stats.get('ips', 0)} IPs  ·  "
        f"{stats.get('fqdns', 0)} FQDNs  ·  "
        f"{stats.get('packets', 0)} packets  ·  "
        f"{stats.get('mb', '?')} MB  ·  "
        f"{stats.get('duration', '?')} duration"
    )
    _add_textbox(slide, MARGIN_L, Cm(9.2), SLIDE_W - MARGIN_L * 2, Cm(0.7),
                 text=stat_line,
                 font_size=11, colour=RGBColor(0xAE, 0xD6, 0xF1),
                 align=PP_ALIGN.CENTER)

    _add_textbox(slide, MARGIN_L, Cm(10.5), SLIDE_W - MARGIN_L * 2, Cm(0.6),
                 text=f"Case ID: {case_id}",
                 font_size=11, colour=C_SILVER,
                 align=PP_ALIGN.CENTER)

    _add_textbox(slide, MARGIN_L, Cm(11.3), SLIDE_W - MARGIN_L * 2, Cm(0.6),
                 text=f"Generated: {generated}  ·  FAN — Forensics Agent Network",
                 font_size=9, colour=C_GREY,
                 align=PP_ALIGN.CENTER)

    _add_textbox(slide, MARGIN_L, SLIDE_H - Cm(0.8), SLIDE_W - MARGIN_L * 2, Cm(0.6),
                 text="CONFIDENTIAL — For authorised recipients only",
                 font_size=8, colour=C_GREY, italic=True,
                 align=PP_ALIGN.CENTER)


def _slide_executive_summary(prs, data: dict, case_id: str, overall_sev: str,
                              page: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "Executive Summary", "Management overview — no technical identifiers")
    _footer(slide, case_id, page)

    sev_colour = SEV_COLOUR.get(overall_sev.lower(), C_ORANGE)

    # Severity box
    rect = slide.shapes.add_shape(1, MARGIN_L, CONTENT_Y, Cm(8), Cm(2.0))
    rect.fill.solid()
    rect.fill.fore_color.rgb = sev_colour
    rect.line.fill.background()
    _add_textbox(slide, MARGIN_L, CONTENT_Y, Cm(8), Cm(2.0),
                 text=f"Overall Severity\n{overall_sev.upper()}",
                 font_size=20, bold=True, colour=C_WHITE,
                 align=PP_ALIGN.CENTER)

    # Stat boxes
    stats_data = [
        ("Flows",    str(len(data.get("netflow", [])))),
        ("IPs",      str(len(data.get("unique_ips", [])))),
        ("FQDNs",    str(len(data.get("unique_fqdns", [])))),
    ]
    box_w = Cm(5.5)
    for i, (label, value) in enumerate(stats_data):
        bx = SLIDE_W - MARGIN_R - box_w * (3 - i) - Cm(0.3) * (2 - i)
        rect2 = slide.shapes.add_shape(1, bx, CONTENT_Y, box_w, Cm(2.0))
        rect2.fill.solid()
        rect2.fill.fore_color.rgb = C_BLUE
        rect2.line.fill.background()
        _add_textbox(slide, bx, CONTENT_Y, box_w, Cm(2.0),
                     text=f"{value}\n{label}",
                     font_size=20, bold=True, colour=C_WHITE,
                     align=PP_ALIGN.CENTER)

    # Summary bullets
    bullets = _build_executive_bullets(data, overall_sev)
    y = CONTENT_Y + Cm(2.4)
    txb = slide.shapes.add_textbox(MARGIN_L, y, CONTENT_W, CONTENT_H - Cm(2.4))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p   = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text           = f"• {b}"
        run.font.size      = Pt(13)
        run.font.color.rgb = C_BLACK
        p.space_after      = Pt(5)


def _build_executive_bullets(data: dict, overall_sev: str) -> list[str]:
    bullets: list[str] = []

    # HTTP threats
    http = data.get("http_results", {})
    sus_ua = http.get("suspicious_ua", {})
    if sus_ua.get("findings"):
        bullets.append(
            f"Suspicious network behaviour detected: unusual HTTP client signature "
            f"observed from {len(sus_ua['findings'])} connection(s)."
        )

    # File hashes
    fh = data.get("fh_data", {})
    malicious = [r for r in data.get("fh_records", []) if r.get("verdict", "").lower() == "malicious"]
    suspicious = [r for r in data.get("fh_records", []) if r.get("verdict", "").lower() == "suspicious"]
    if malicious:
        bullets.append(
            f"{len(malicious)} file(s) extracted from the capture were identified as malicious "
            f"by threat intelligence services."
        )
    elif suspicious:
        bullets.append(
            f"{len(suspicious)} file(s) extracted from the capture are flagged as suspicious "
            f"and require further review."
        )

    # Suricata
    alerts = data.get("suricata_alerts", [])
    if alerts:
        bullets.append(
            f"Intrusion detection system raised {len(alerts)} alert(s) against the captured traffic."
        )

    # YARA
    yara_matches = [r for r in data.get("yara_records", []) if r.get("rule")]
    if yara_matches:
        bullets.append(
            f"Threat hunting rules matched {len(yara_matches)} pattern(s) in the captured data."
        )

    # Catch-all for other threats
    threat_modules = [
        ("dns_results",  "Domain name system abuse"),
        ("icmp_results", "Network-level covert channel activity"),
        ("arp_results",  "Address resolution protocol manipulation"),
        ("tcp_results",  "Transmission control protocol anomalies"),
    ]
    for key, label in threat_modules:
        res = data.get(key, {})
        findings_found = any(
            v.get("findings") for v in res.values() if isinstance(v, dict)
        )
        if findings_found:
            bullets.append(f"{label} detected in the captured traffic.")

    if not bullets:
        bullets.append("No high-severity threats were automatically detected in this capture.")
        bullets.append(
            "Manual review of network flows and extracted files is recommended to rule "
            "out low-and-slow or encrypted threats."
        )

    bullets.append(
        "All network capture data and analysis artefacts have been preserved for "
        "chain-of-custody purposes."
    )
    return bullets


def _slide_key_findings(prs, data: dict, case_id: str, page: int) -> None:
    """One slide listing every category with findings, coloured by severity."""
    findings: list[tuple[str, str, str]] = []  # (severity, category, description)

    # HTTP
    for cat_key, cat_label in [
        ("suspicious_ua",   "Suspicious User-Agent"),
        ("suspicious_uri",  "Suspicious URI Pattern"),
        ("large_upload",    "Large HTTP Upload"),
        ("http_beaconing",  "HTTP Beaconing"),
        ("scanning_codes",  "HTTP Scanning / Error Flood"),
    ]:
        cat = data.get("http_results", {}).get(cat_key, {})
        if cat.get("findings"):
            findings.append((
                cat.get("severity", "high"),
                cat_label,
                f"{len(cat['findings'])} finding(s) detected",
            ))

    # Suricata
    for a in data.get("suricata_alerts", [])[:5]:
        findings.append((
            a.get("severity", "high"),
            f"Suricata: {a.get('signature', 'IDS Alert')}",
            a.get("category", ""),
        ))

    # YARA
    for r in data.get("yara_records", [])[:3]:
        findings.append((
            r.get("severity", "high"),
            f"YARA: {r.get('rule', 'Rule match')}",
            r.get("description", ""),
        ))

    # File hashes
    for r in data.get("fh_records", []):
        if r.get("verdict", "").lower() in ("malicious", "suspicious"):
            findings.append((
                "critical" if r["verdict"].lower() == "malicious" else "high",
                f"File: {r.get('filename', 'extracted file')}",
                f"Verdict: {r.get('verdict', '')} — {r.get('sha256', '')[:16]}…",
            ))

    # Generic protocol modules
    for key, label in [
        ("dns_results",     "DNS"),
        ("icmp_results",    "ICMP"),
        ("arp_results",     "ARP"),
        ("tcp_results",     "TCP"),
        ("udp_results",     "UDP"),
        ("tls_results",     "TLS"),
    ]:
        res = data.get(key, {})
        for cat_key, cat_val in res.items():
            if not isinstance(cat_val, dict):
                continue
            if cat_val.get("findings"):
                findings.append((
                    cat_val.get("severity", "medium"),
                    f"{label}: {cat_key.replace('_', ' ').title()}",
                    f"{len(cat_val['findings'])} finding(s)",
                ))

    if not findings:
        findings.append(("info", "No automatic threat findings", "Manual review recommended"))

    # Sort by severity
    findings.sort(key=lambda x: _sev_rank(x[0]))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header(slide, "Key Findings", f"{len(findings)} finding category(s) detected")
    _footer(slide, case_id, page)

    y = CONTENT_Y
    row_h = Cm(0.85)
    for sev, cat, desc in findings[:16]:
        colour = SEV_COLOUR.get(sev.lower(), C_GREY)

        # severity pill
        pill = slide.shapes.add_shape(1, MARGIN_L, y + Cm(0.12), Cm(2.5), Cm(0.55))
        pill.fill.solid()
        pill.fill.fore_color.rgb = colour
        pill.line.fill.background()
        tf2 = pill.text_frame
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
        run2 = tf2.paragraphs[0].add_run()
        run2.text = sev.upper()
        run2.font.size = Pt(7)
        run2.font.bold = True
        run2.font.color.rgb = C_WHITE

        # category + description
        _add_textbox(slide,
                     MARGIN_L + Cm(2.7), y,
                     CONTENT_W - Cm(2.7), row_h,
                     text=f"{cat}  —  {desc}",
                     font_size=11, colour=C_BLACK)
        y += row_h


def _slide_network_overview(prs, data: dict, case_id: str, page: int) -> None:
    netflow = data.get("netflow", [])
    ips     = data.get("unique_ips", [])
    fqdns   = data.get("unique_fqdns", [])

    # Protocol distribution
    proto_count: dict[str, int] = {}
    for row in netflow:
        p = row.get("protocol", row.get("Protocol", "Unknown"))
        proto_count[p] = proto_count.get(p, 0) + 1

    proto_lines = [f"{p}: {c} flows" for p, c in
                   sorted(proto_count.items(), key=lambda x: -x[1])]

    items = (
        [f"Total flows: {len(netflow)}",
         f"Unique IPs: {len(ips)}",
         f"Unique FQDNs: {len(fqdns)}",
         ""]
        + ["Protocol distribution:"]
        + [f"    {l}" for l in proto_lines]
        + [""]
        + ["Unique FQDNs:"]
        + [f"    {f}" for f in fqdns]
        + [""]
        + ["Unique IPs:"]
        + [f"    {ip}" for ip in ips]
    )
    _bullet_slide(prs, "Network Overview", "Traffic composition and unique endpoints",
                  items, case_id, page)


def _slide_top_flows(prs, data: dict, case_id: str, page: int) -> None:
    netflow = data.get("netflow", [])
    if not netflow:
        return

    # Detect columns
    sample = netflow[0] if netflow else {}
    has_bytes  = "bytes" in sample or "length" in sample
    has_proto  = "protocol" in sample or "Protocol" in sample
    has_dst    = "dst_ip" in sample or "ip.dst" in sample

    rows = []
    for r in netflow[:20]:
        src   = r.get("src_ip", r.get("ip.src", ""))
        dst   = r.get("dst_ip", r.get("ip.dst", ""))
        proto = r.get("protocol", r.get("Protocol", ""))
        dport = r.get("dst_port", r.get("tcp.dstport", r.get("udp.dstport", "")))
        nbytes = r.get("bytes", r.get("length", r.get("frame.len", "")))
        rows.append([src, dst, str(dport), proto, str(nbytes)])

    _table_slide(prs,
                 "Top Network Flows", "Source → Destination",
                 ["Source IP", "Destination IP", "Dst Port", "Protocol", "Bytes"],
                 rows, case_id, page,
                 col_widths=[0.22, 0.22, 0.12, 0.12, 0.12])


def _slide_http_findings(prs, data: dict, case_id: str, page: int) -> None:
    http = data.get("http_results", {})
    rows = []
    for cat_key in http:
        cat = http[cat_key]
        if not isinstance(cat, dict):
            continue
        for f in cat.get("findings", []):
            if isinstance(f, dict):
                rows.append([
                    cat.get("severity", ""),
                    cat_key.replace("_", " ").title(),
                    f.get("src_ip", ""),
                    f.get("host", f.get("dst_ip", "")),
                    f.get("uri", f.get("indicator", ""))[:60],
                ])
            else:
                rows.append([cat.get("severity", ""), cat_key, "", "", str(f)[:60]])
    if not rows:
        return
    _table_slide(prs, "HTTP Threat Findings", "Unusual HTTP patterns detected",
                 ["Sev", "Category", "Source", "Host", "Detail"],
                 rows, case_id, page,
                 col_widths=[0.08, 0.18, 0.16, 0.20, 0.38])


def _slide_tls_sessions(prs, data: dict, case_id: str, page: int) -> None:
    sessions = data.get("tls_sessions", [])
    if not sessions:
        return
    rows = []
    for s in sessions[:15]:
        rows.append([
            s.get("timestamp", s.get("frame.time", ""))[:19],
            s.get("src_ip",    s.get("ip.src", "")),
            s.get("dst_ip",    s.get("ip.dst", "")),
            s.get("server_name", s.get("sni", "")),
            s.get("tls_version", ""),
        ])
    _table_slide(prs, "TLS Sessions", "Encrypted session inventory",
                 ["Timestamp", "Client IP", "Server IP", "SNI", "TLS Version"],
                 rows, case_id, page,
                 col_widths=[0.20, 0.16, 0.16, 0.30, 0.18])


def _slide_dns_records(prs, data: dict, case_id: str, page: int) -> None:
    dns = data.get("dns_flows", [])
    if not dns:
        return
    rows = []
    for r in dns[:18]:
        rows.append([
            r.get("timestamp", r.get("frame.time", ""))[:19],
            r.get("query",     r.get("dns.qry.name", "")),
            r.get("response",  r.get("dns.a", "")),
            r.get("src_ip",    r.get("ip.src", "")),
        ])
    _table_slide(prs, "DNS Records", "All DNS queries and responses",
                 ["Timestamp", "Query", "Response", "Requester IP"],
                 rows, case_id, page,
                 col_widths=[0.22, 0.34, 0.24, 0.20])


def _slide_file_hashes(prs, data: dict, case_id: str, page: int) -> None:
    records = data.get("fh_records", [])
    if not records:
        return
    rows = []
    for r in records[:15]:
        rows.append([
            r.get("filename", r.get("name", ""))[:30],
            r.get("protocol", ""),
            r.get("size", r.get("file_size", "")),
            r.get("md5",  "")[:12] + "…" if r.get("md5") else "",
            r.get("verdict", r.get("disposition", "—")),
        ])
    _table_slide(prs, "File Hash Analysis", "Files extracted from PCAP and checked against OSINT",
                 ["Filename", "Protocol", "Size (B)", "MD5 (partial)", "Verdict"],
                 rows, case_id, page,
                 col_widths=[0.30, 0.10, 0.10, 0.20, 0.10])


def _slide_suricata(prs, data: dict, case_id: str, page: int) -> None:
    alerts = data.get("suricata_alerts", [])
    if not alerts:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _header(slide, "Suricata IDS", "Signature-based intrusion detection results")
        _footer(slide, case_id, page)
        _add_textbox(slide, MARGIN_L, CONTENT_Y, CONTENT_W, Cm(2),
                     text="No Suricata alerts triggered for this capture.",
                     font_size=14, colour=C_GREY, italic=True)
        return
    rows = []
    for a in alerts[:15]:
        rows.append([
            a.get("severity",  ""),
            a.get("signature", a.get("alert", ""))[:50],
            a.get("src_ip",    ""),
            a.get("dst_ip",    ""),
            a.get("timestamp", "")[:19],
        ])
    _table_slide(prs, "Suricata IDS Alerts", f"{len(alerts)} alert(s) triggered",
                 ["Sev", "Signature", "Src IP", "Dst IP", "Timestamp"],
                 rows, case_id, page,
                 col_widths=[0.08, 0.44, 0.16, 0.16, 0.16])


def _slide_yara(prs, data: dict, case_id: str, page: int) -> None:
    matches = data.get("yara_records", [])
    if not matches:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _header(slide, "YARA Threat Hunting", "Rule-based pattern matching results")
        _footer(slide, case_id, page)
        _add_textbox(slide, MARGIN_L, CONTENT_Y, CONTENT_W, Cm(2),
                     text="No YARA rule matches in this capture.",
                     font_size=14, colour=C_GREY, italic=True)
        return
    rows = [[r.get("rule",""), r.get("target",""), r.get("tags",""), r.get("description","")[:50]]
            for r in matches[:15]]
    _table_slide(prs, "YARA Rule Matches", f"{len(matches)} match(es) found",
                 ["Rule", "Target", "Tags", "Description"],
                 rows, case_id, page,
                 col_widths=[0.25, 0.20, 0.15, 0.40])


def _slide_mitre(prs, data: dict, case_id: str, page: int) -> None:
    """Collect all MITRE ATT&CK technique IDs referenced in any module output."""
    techniques: dict[str, tuple[str, str, str]] = {}  # id -> (name, tactic, severity)

    def _add(tid: str, name: str, tactic: str, sev: str) -> None:
        if tid not in techniques or _sev_rank(sev) < _sev_rank(techniques[tid][2]):
            techniques[tid] = (name, tactic, sev)

    # HTTP
    for cat_key, cat in data.get("http_results", {}).items():
        if not isinstance(cat, dict):
            continue
        for tid in cat.get("mitre_ids", cat.get("mitre", [])):
            _add(tid, cat_key.replace("_", " ").title(), "Command and Control", cat.get("severity", "medium"))

    # Generic protocol modules
    for mod_key in ["dns_results","icmp_results","ntp_results","arp_results",
                    "tcp_results","udp_results","dhcp_results","mdns_results",
                    "quic_results","snmp_results","nbns_results","llmnr_results",
                    "stun_results","ssdp_results","netbios_results"]:
        for cat_key, cat in data.get(mod_key, {}).items():
            if not isinstance(cat, dict):
                continue
            for tid in cat.get("mitre_ids", cat.get("mitre", [])):
                _add(tid, cat_key.replace("_", " ").title(), mod_key.replace("_results","").upper(), cat.get("severity","medium"))

    # Suricata
    for a in data.get("suricata_alerts", []):
        for tid in a.get("mitre_ids", []):
            _add(tid, a.get("signature","Suricata alert"), "Detect", a.get("severity","medium"))

    if not techniques:
        # fallback: show a "none detected" slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _header(slide, "MITRE ATT&CK Coverage", "Techniques observed in this capture")
        _footer(slide, case_id, page)
        _add_textbox(slide, MARGIN_L, CONTENT_Y, CONTENT_W, Cm(2),
                     text="No MITRE ATT&CK techniques were automatically mapped in this capture.\n"
                          "Manual analysis may identify additional coverage.",
                     font_size=13, colour=C_GREY, italic=True)
        return

    rows = [
        [tid, name, tactic, sev.upper()]
        for tid, (name, tactic, sev) in
        sorted(techniques.items(), key=lambda x: _sev_rank(x[1][2]))
    ]
    _table_slide(prs, "MITRE ATT&CK Coverage",
                 f"{len(rows)} technique(s) observed",
                 ["Technique ID", "Name", "Tactic", "Severity"],
                 rows, case_id, page,
                 col_widths=[0.14, 0.38, 0.30, 0.18])


def _slide_iocs(prs, data: dict, case_id: str, page: int) -> None:
    """Aggregate IOCs from all modules."""
    iocs: list[tuple[str, str, str, str]] = []  # (sev, type, value, source)

    # CTI enrichment
    for row in data.get("fan_correlation", []):
        val = row.get("indicator", row.get("value", ""))
        if val:
            iocs.append((row.get("severity","medium"), row.get("type",""), val, "CTI"))

    for row in data.get("fan_ip", []):
        ip = row.get("ip", "")
        if ip:
            iocs.append((row.get("severity","medium"), "ip", ip, "CTI"))

    # HTTP findings
    for cat in data.get("http_results", {}).values():
        if not isinstance(cat, dict):
            continue
        for f in cat.get("findings", []):
            if not isinstance(f, dict):
                continue
            for field in ("src_ip", "dst_ip", "host", "domain"):
                val = f.get(field, "")
                if val:
                    iocs.append((cat.get("severity","medium"), field.replace("_"," "), val, "HTTP"))

    # Deduplicate by value
    seen: set[str] = set()
    unique_iocs: list[tuple[str,str,str,str]] = []
    for sev, typ, val, src in iocs:
        if val not in seen:
            seen.add(val)
            unique_iocs.append((sev, typ, val, src))

    unique_iocs.sort(key=lambda x: _sev_rank(x[0]))

    if not unique_iocs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _header(slide, "Indicators of Compromise", "Extracted from all analysis modules")
        _footer(slide, case_id, page)
        _add_textbox(slide, MARGIN_L, CONTENT_Y, CONTENT_W, Cm(2),
                     text="No indicators of compromise were automatically extracted.",
                     font_size=13, colour=C_GREY, italic=True)
        return

    rows = [[sev.upper(), typ, val[:60], src] for sev, typ, val, src in unique_iocs]
    _table_slide(prs, "Indicators of Compromise",
                 f"{len(rows)} unique indicator(s)",
                 ["Severity", "Type", "Value", "Source"],
                 rows, case_id, page,
                 col_widths=[0.12, 0.12, 0.56, 0.10])


def _slide_recommendations(prs, data: dict, case_id: str,
                            overall_sev: str, page: int) -> None:
    recs: list[str] = []

    http = data.get("http_results", {})
    if any(c.get("findings") for c in http.values() if isinstance(c, dict)):
        recs.append("Investigate source host(s) for suspicious HTTP activity and review proxy logs.")

    alerts = data.get("suricata_alerts", [])
    if alerts:
        recs.append(f"Investigate {len(alerts)} Suricata alert(s) — triage against known-good baselines.")
    else:
        recs.append("Update Suricata rules (run update_suricata_rules.sh) to improve detection coverage.")

    yara_matches = data.get("yara_records", [])
    if yara_matches:
        recs.append(f"Review {len(yara_matches)} YARA match(es) — confirm whether targets are malicious.")

    malicious_files = [r for r in data.get("fh_records",[]) if r.get("verdict","").lower()=="malicious"]
    if malicious_files:
        recs.append(f"Isolate host(s) involved — {len(malicious_files)} malicious file(s) extracted from the capture.")

    suspicious_files = [r for r in data.get("fh_records",[]) if r.get("verdict","").lower()=="suspicious"]
    if suspicious_files:
        recs.append(f"Submit {len(suspicious_files)} suspicious file(s) to sandbox for dynamic analysis.")

    if overall_sev.lower() in ("critical", "high"):
        recs.append("Preserve all capture artefacts for chain-of-custody and forensic follow-up.")
        recs.append("Consider blocking identified malicious IPs and domains at the perimeter firewall.")

    recs.append("Run /obsidian-record after the investigation to persist findings to the vault.")

    _bullet_slide(prs, "Recommendations", "Prioritised response actions",
                  recs, case_id, page)


# ── main builder ──────────────────────────────────────────────────────────────

def build_presentation(stem: str, case_id: str, output_path: Path,
                       base_dir: Path | None = None) -> Path:
    """
    Build a .pptx investigation briefing for the given PCAP stem.
    Returns the path to the written file.
    """
    import generate_pcap_report as rpt
    if base_dir:
        rpt.ANALYSIS_DIR = base_dir

    data        = load_all_data(stem)
    overall_sev = _overall_severity(data)
    generated   = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")

    # Capture stats
    netflow  = data.get("netflow", [])
    total_bytes = sum(int(r.get("bytes", r.get("length", r.get("frame.len", 0)) or 0))
                      for r in netflow)
    stats = {
        "flows":    len(netflow),
        "ips":      len(data.get("unique_ips", [])),
        "fqdns":    len(data.get("unique_fqdns", [])),
        "packets":  len(netflow),
        "mb":       f"{total_bytes / 1_048_576:.1f} MB" if total_bytes else "?",
        "duration": "see report",
    }

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    # Ensure we have a blank layout
    if len(prs.slide_layouts) < 7:
        for _ in range(7 - len(prs.slide_layouts)):
            prs.slide_layouts[0]

    page = 1
    _slide_title(prs, stem, case_id, overall_sev, stats, generated)
    page += 1; _slide_executive_summary(prs, data, case_id, overall_sev, page)
    page += 1; _slide_key_findings(prs, data, case_id, page)
    page += 1; _slide_network_overview(prs, data, case_id, page)
    page += 1; _slide_top_flows(prs, data, case_id, page)
    if data.get("has_dns"):
        page += 1; _slide_dns_records(prs, data, case_id, page)
    if data.get("has_http"):
        page += 1; _slide_http_findings(prs, data, case_id, page)
    if data.get("has_tls"):
        page += 1; _slide_tls_sessions(prs, data, case_id, page)
    if data.get("has_fh"):
        page += 1; _slide_file_hashes(prs, data, case_id, page)
    page += 1; _slide_suricata(prs, data, case_id, page)
    page += 1; _slide_yara(prs, data, case_id, page)
    page += 1; _slide_mitre(prs, data, case_id, page)
    page += 1; _slide_iocs(prs, data, case_id, page)
    page += 1; _slide_recommendations(prs, data, case_id, overall_sev, page)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


# ── CLI ───────────────────────────────────────────────────────────────────────

def main() -> None:
    ap = argparse.ArgumentParser(
        description="PCAP Incident PowerPoint Presentation Generator"
    )
    ap.add_argument("--stem",         required=True, help="PCAP file stem")
    ap.add_argument("--case-id",      default="",    help="Case ID")
    ap.add_argument("--output-dir",   default=None,  help="Output directory")
    ap.add_argument("--base-dir",     default=None,  help="Analysis base directory")
    ap.add_argument("--report-version", default="1", help="Version stamp")
    args = ap.parse_args()

    import generate_pcap_report as rpt

    base_dir = Path(args.base_dir) if args.base_dir else rpt.ANALYSIS_DIR
    out_dir  = Path(args.output_dir) if args.output_dir else (PROJECT_ROOT / "reports")
    fname    = f"{args.stem}_incident_briefing_v{args.report_version}.pptx"
    out_path = out_dir / fname

    print(f"[presentation] Stem     : {args.stem}")
    print(f"[presentation] Case ID  : {args.case_id}")
    print(f"[presentation] Output   : {out_path}")

    result = build_presentation(
        stem=args.stem,
        case_id=args.case_id or args.stem,
        output_path=out_path,
        base_dir=base_dir,
    )
    print(f"[presentation] Done     : {result}")


if __name__ == "__main__":
    main()
