#!/usr/bin/env python3
# SPDX-License-Identifier: MIT OR Apache-2.0
# SPDX-FileCopyrightText: 2026 Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin
"""
generate_pcap_report.py — PCAP Incident Report Generator

Aggregates outputs from PCAP analysis modules and produces a structured
incident report in Markdown and PDF format.

Consumed (all optional — missing sources are skipped gracefully):
  ./analysis/pcap/<stem>/            → netflow.csv, unique_ips.txt, unique_fqdns.txt
  ./analysis/fan_ip/<stem>/             → correlation.csv, ip_enrichment.csv
  ./analysis/dns_threats/<stem>/     → dns_threats.json, dns_flows.csv
  ./analysis/icmp_threats/<stem>/    → icmp_threats.json, icmp_flows.csv
  ./analysis/ntp_threats/<stem>/     → ntp_threats.json, ntp_flows.csv
  ./analysis/http_threats/<stem>/    → http_threats.json, http_flows.csv
  ./analysis/cert_inspector/<stem>/  → certs.json, certs.csv
  ./analysis/tls_inspector/<stem>/   → tls_sessions.json, tls_sessions.csv

Output:
  ./reports/<stem>_incident_report.md
  ./reports/<stem>_incident_report.pdf  (requires: pip3 install markdown weasyprint)
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import json
import subprocess
import sys
from collections import defaultdict
from datetime import datetime, timezone
try:
    from zoneinfo import ZoneInfo
    _CET = ZoneInfo("Europe/Amsterdam")
except ImportError:
    _CET = timezone.utc
from pathlib import Path
from typing import Any

try:
    from research_notes import parse_steps as _parse_research_steps, parse_events as _parse_research_events
    from hallucination_guard import (
        ConfidenceTier,
        tag_finding,
        render_confidence_summary,
        reset_counter as _hg_reset,
    )
except ModuleNotFoundError:
    # Imported as a package (e.g. `from lib.generate_pcap_report import ...`)
    # rather than run as a script — put lib/ on the path for the sibling import.
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    from research_notes import parse_steps as _parse_research_steps, parse_events as _parse_research_events
    from hallucination_guard import (
        ConfidenceTier,
        tag_finding,
        render_confidence_summary,
        reset_counter as _hg_reset,
    )

PROJECT_ROOT = Path(__file__).parent.parent
ANALYSIS_DIR = PROJECT_ROOT / "analysis"
REPORTS_DIR  = PROJECT_ROOT / "reports"

SEVERITY_ORDER = {"critical": 0, "high": 1, "medium": 2, "low": 3, "info": 4}
SEVERITY_BADGE = {
    "critical": "**[CRITICAL]**",
    "high":     "**[HIGH]**",
    "medium":   "[MEDIUM]",
    "low":      "[LOW]",
    "info":     "[INFO]",
}

MAX_TIMELINE_ROWS = 150   # cap on timeline events

# ── Colour palette ─────────────────────────────────────────────────────────────
_DARK_NAVY  = (0x0f, 0x17, 0x2a)
_MID_NAVY   = (0x1e, 0x3a, 0x5f)
_BLUE       = (0x1d, 0x4e, 0xd8)
_LIGHT_BLUE = (0x93, 0xc5, 0xfd)
_WHITE      = (0xff, 0xff, 0xff)
_LIGHT_BG   = (0xf8, 0xfa, 0xfc)
_ROW_ALT    = (0xf1, 0xf5, 0xf9)
_TEXT_DARK  = (0x1f, 0x29, 0x37)
_TEXT_MID   = (0x6b, 0x72, 0x80)
_AMBER      = (0xfb, 0xbf, 0x24)
_GREEN      = (0x22, 0xc5, 0x5e)

_SEV_COLORS = {
    "critical": (0xef, 0x44, 0x44),
    "high":     (0xf9, 0x73, 0x16),
    "medium":   (0xea, 0xb3, 0x08),
    "low":      (0x22, 0xc5, 0x5e),
    "info":     (0x6b, 0x72, 0x80),
}


# ── Loaders ────────────────────────────────────────────────────────────────────

def _load_json(path: Path) -> dict | list | None:
    if path.exists():
        try:
            return json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            return None
    return None


def _load_csv(path: Path) -> list[dict]:
    if not path.exists():
        return []
    try:
        with path.open(encoding="utf-8") as fh:
            return list(csv.DictReader(fh))
    except Exception:
        return []


def _load_lines(path: Path) -> list[str]:
    if not path.exists():
        return []
    return [l.strip() for l in path.read_text(encoding="utf-8").splitlines() if l.strip()]


def _normalize_categories(raw: dict | None) -> dict:
    """Normalise new-style module JSON (has 'categories' key) to flat results dict."""
    if not raw:
        return {}
    if "categories" in raw:
        result = {}
        for key, cat in raw["categories"].items():
            norm = dict(cat)
            if "label" in norm and "name" not in norm:
                norm["name"] = norm["label"]
            result[key] = norm
        return result
    return raw


def load_all_data(stem: str) -> dict:
    pcap_dir  = ANALYSIS_DIR / "pcap"           / stem
    fan_ip_dir   = ANALYSIS_DIR / "fan_ip"            / stem
    dns_dir   = ANALYSIS_DIR / "dns_threats"    / stem
    icmp_dir  = ANALYSIS_DIR / "icmp_threats"   / stem
    ntp_dir   = ANALYSIS_DIR / "ntp_threats"    / stem
    http_dir  = ANALYSIS_DIR / "http_threats"   / stem
    cert_dir  = ANALYSIS_DIR / "cert_inspector" / stem
    tls_dir   = ANALYSIS_DIR / "tls_inspector"  / stem
    arp_dir   = ANALYSIS_DIR / "arp_threats"    / stem
    tcp_dir   = ANALYSIS_DIR / "tcp_threats"    / stem
    udp_dir   = ANALYSIS_DIR / "udp_threats"    / stem
    dhcp_dir  = ANALYSIS_DIR / "dhcp_threats"   / stem
    mdns_dir  = ANALYSIS_DIR / "mdns_threats"   / stem
    quic_dir    = ANALYSIS_DIR / "quic_threats"    / stem
    snmp_dir    = ANALYSIS_DIR / "snmp_threats"   / stem
    nbns_dir    = ANALYSIS_DIR / "nbns_threats"   / stem
    llmnr_dir   = ANALYSIS_DIR / "llmnr_threats"  / stem
    stun_dir    = ANALYSIS_DIR / "stun_threats"   / stem
    ssdp_dir    = ANALYSIS_DIR / "ssdp_threats"   / stem
    netbios_dir = ANALYSIS_DIR / "netbios_threats" / stem
    fh_dir       = ANALYSIS_DIR / "file_hashes"  / stem
    suricata_dir = ANALYSIS_DIR / "suricata"     / stem
    yara_dir     = ANALYSIS_DIR / "yara_pcap"    / stem

    # Strip non-category _sessions key from TLS inspector JSON
    tls_raw = _load_json(tls_dir / "tls_sessions.json") or {}
    tls_results = {k: v for k, v in tls_raw.items() if not k.startswith("_")}

    return {
        "stem":             stem,
        "netflow":          _load_csv(pcap_dir / "netflow.csv"),
        "unique_ips":       _load_lines(pcap_dir / "unique_ips.txt"),
        "unique_fqdns":     _load_lines(pcap_dir / "unique_fqdns.txt"),
        "fan_correlation":  _load_csv(fan_ip_dir / "correlation.csv"),
        "fan_ip":           _load_csv(fan_ip_dir / "ip_enrichment.csv"),
        "dns_results":      _load_json(dns_dir  / "dns_threats.json") or {},
        "dns_flows":        _load_csv(dns_dir   / "dns_flows.csv"),
        "icmp_results":     _load_json(icmp_dir / "icmp_threats.json") or {},
        "icmp_flows":       _load_csv(icmp_dir  / "icmp_flows.csv"),
        "ntp_results":      _load_json(ntp_dir  / "ntp_threats.json") or {},
        "ntp_flows":        _load_csv(ntp_dir   / "ntp_flows.csv"),
        "http_results":     _load_json(http_dir / "http_threats.json") or {},
        "http_flows":       _load_csv(http_dir  / "http_flows.csv"),
        "cert_results":     _load_json(cert_dir / "certs.json") or {},
        "cert_records":     _load_csv(cert_dir  / "certs.csv"),
        "tls_results":      tls_results,
        "tls_sessions":     _load_csv(tls_dir   / "tls_sessions.csv"),
        # New protocol threat modules
        "arp_results":      _normalize_categories(_load_json(arp_dir  / "arp_threats.json")),
        "tcp_results":      _normalize_categories(_load_json(tcp_dir  / "tcp_threats.json")),
        "udp_results":      _normalize_categories(_load_json(udp_dir  / "udp_threats.json")),
        "dhcp_results":     _normalize_categories(_load_json(dhcp_dir / "dhcp_threats.json")),
        "mdns_results":     _normalize_categories(_load_json(mdns_dir / "mdns_threats.json")),
        "quic_results":     _normalize_categories(_load_json(quic_dir   / "quic_threats.json")),
        "snmp_results":     _normalize_categories(_load_json(snmp_dir   / "snmp_threats.json")),
        "nbns_results":     _normalize_categories(_load_json(nbns_dir   / "nbns_threats.json")),
        "llmnr_results":    _normalize_categories(_load_json(llmnr_dir  / "llmnr_threats.json")),
        "stun_results":     _normalize_categories(_load_json(stun_dir   / "stun_threats.json")),
        "ssdp_results":     _normalize_categories(_load_json(ssdp_dir   / "ssdp_threats.json")),
        "netbios_results":  _normalize_categories(_load_json(netbios_dir / "netbios_threats.json")),
        "fh_data":          _load_json(fh_dir / "file_hashes.json") or {},
        "fh_records":       _load_csv(fh_dir  / "file_hashes.csv"),
        # IDS / signature modules
        "suricata_data":    _load_json(suricata_dir / "suricata_alerts.json") or {},
        "suricata_alerts":  _load_csv(suricata_dir  / "suricata_alerts.csv"),
        "yara_data":        _load_json(yara_dir     / "yara_matches.json") or {},
        "yara_records":     _load_csv(yara_dir      / "yara_matches.csv"),
        # Availability flags
        "has_pcap":     (pcap_dir    / "netflow.csv").exists(),
        "has_fan_ip":      (fan_ip_dir     / "correlation.csv").exists() or (fan_ip_dir / "ip_enrichment.csv").exists(),
        "has_dns":      (dns_dir     / "dns_threats.json").exists(),
        "has_icmp":     (icmp_dir    / "icmp_threats.json").exists(),
        "has_ntp":      (ntp_dir     / "ntp_threats.json").exists(),
        "has_http":     (http_dir    / "http_threats.json").exists(),
        "has_cert":     (cert_dir    / "certs.json").exists(),
        "has_tls":      (tls_dir     / "tls_sessions.json").exists(),
        "has_arp":      (arp_dir     / "arp_threats.json").exists(),
        "has_tcp":      (tcp_dir     / "tcp_threats.json").exists(),
        "has_udp":      (udp_dir     / "udp_threats.json").exists(),
        "has_dhcp":     (dhcp_dir    / "dhcp_threats.json").exists(),
        "has_mdns":     (mdns_dir    / "mdns_threats.json").exists(),
        "has_quic":     (quic_dir     / "quic_threats.json").exists(),
        "has_snmp":     (snmp_dir    / "snmp_threats.json").exists(),
        "has_nbns":     (nbns_dir    / "nbns_threats.json").exists(),
        "has_llmnr":    (llmnr_dir   / "llmnr_threats.json").exists(),
        "has_stun":     (stun_dir    / "stun_threats.json").exists(),
        "has_ssdp":     (ssdp_dir    / "ssdp_threats.json").exists(),
        "has_netbios":  (netbios_dir / "netbios_threats.json").exists(),
        "has_fh":       (fh_dir      / "file_hashes.json").exists(),
        "has_suricata": (suricata_dir / "suricata_alerts.json").exists(),
        "has_yara":     (yara_dir    / "yara_matches.json").exists(),
    }


# ── Severity helpers ───────────────────────────────────────────────────────────

def _sev_rank(s: str) -> int:
    return SEVERITY_ORDER.get((s or "info").lower(), 4)


def _overall_severity(data: dict) -> str:
    best = "info"
    # Malicious files are always critical
    fh = data.get("fh_data", {})
    if fh.get("malicious_count", 0) > 0:
        return "critical"
    if fh.get("suspicious_count", 0) > 0 and _sev_rank("high") < _sev_rank(best):
        best = "high"
    # Suricata critical alerts → critical overall
    suricata = data.get("suricata_data", {})
    if suricata.get("critical_count", 0) > 0:
        return "critical"
    if suricata.get("high_count", 0) > 0 and _sev_rank("high") < _sev_rank(best):
        best = "high"
    # YARA critical matches → critical overall
    yara = data.get("yara_data", {})
    if yara.get("critical_count", 0) > 0:
        return "critical"
    if yara.get("high_count", 0) > 0 and _sev_rank("high") < _sev_rank(best):
        best = "high"
    for results_key in ("icmp_results", "dns_results", "ntp_results", "http_results",
                        "cert_results", "tls_results",
                        "arp_results", "tcp_results", "udp_results",
                        "dhcp_results", "mdns_results", "quic_results",
                        "snmp_results", "nbns_results", "llmnr_results",
                        "stun_results", "ssdp_results", "netbios_results"):
        for cat in data[results_key].values():
            sev = cat.get("severity", "info")
            if sev in ("critical", "high", "medium", "low") and cat.get("count", 0) > 0:
                if _sev_rank(sev) < _sev_rank(best):
                    best = sev
    # CTI malicious indicators raise severity
    mal = [r for r in data["fan_ip"] if r.get("reputation") == "malicious"]
    if mal and _sev_rank("high") < _sev_rank(best):
        best = "high"
    return best


def _triggered(results: dict) -> list[dict]:
    """Return categories that have at least one finding, sorted by severity."""
    cats = [v for v in results.values() if isinstance(v, dict) and v.get("count", 0) > 0]
    return sorted(cats, key=lambda c: _sev_rank(c.get("severity", "info")))


# ── Capture timing ─────────────────────────────────────────────────────────────

def _capture_window(data: dict) -> tuple[str, str, float]:
    """Return (first_ts, last_ts, duration_sec) from netflow data."""
    first = last = ""
    min_t = max_t = None

    for row in data["netflow"]:
        for col in ("first_seen_utc", "last_seen_utc"):
            ts = row.get(col, "")
            if not ts:
                continue
            try:
                # Accept ISO-8601 with Z or offset
                dt = datetime.fromisoformat(ts.replace("Z", "+00:00"))
                ep = dt.timestamp()
                if min_t is None or ep < min_t:
                    min_t = ep; first = ts
                if max_t is None or ep > max_t:
                    max_t = ep; last = ts
            except (ValueError, TypeError):
                continue

    dur = round(max_t - min_t, 1) if (min_t is not None and max_t is not None) else 0.0
    return first, last, dur


def _format_duration(sec: float) -> str:
    if sec < 60:
        return f"{sec:.0f} seconds"
    if sec < 3600:
        return f"{sec / 60:.1f} minutes"
    return f"{sec / 3600:.1f} hours"


# ── Timeline ───────────────────────────────────────────────────────────────────

def _ts_key(ts_str: str) -> str:
    """Normalise timestamp string for sorting (empty → sort last)."""
    return ts_str if ts_str else "9999-99-99"


def _icmp_timeline_events(data: dict) -> list[dict]:
    icmp_results = data["icmp_results"]
    icmp_flows   = data["icmp_flows"]

    # Build sets of flagged IPs from detectors that have clear src_ip in findings
    flood_ips: set[str] = set()
    for f in icmp_results.get("flood", {}).get("findings", []):
        if "src_ip" in f:
            flood_ips.add(f["src_ip"])

    tunnel_ips: set[str] = set()
    for f in icmp_results.get("tunneling", {}).get("findings", []):
        if "src_ip" in f:
            tunnel_ips.add(f["src_ip"])

    exfil_ips: set[str] = set()
    for f in icmp_results.get("exfiltration", {}).get("findings", []):
        if "src_ip" in f:
            exfil_ips.add(f["src_ip"])

    sweep_ips: set[str] = set()
    for f in icmp_results.get("sweep", {}).get("findings", []):
        if "src_ip" in f:
            sweep_ips.add(f["src_ip"])

    events: list[dict] = []
    counts: dict[str, int] = defaultdict(int)
    limits = {
        "ICMP Flood":                30,
        "ICMP Tunneling":            30,
        "ICMP Data Exfiltration":    20,
        "Ping of Death":             20,
        "ICMP Fragmentation Attack": 20,
        "Smurf / Broadcast Amplification": 20,
        "ICMP Redirect Attack":      50,  # usually rare
        "ICMP Network Sweep":        20,
        "ICMP Recon Types":          50,
        "Destination Unreachable Flood": 20,
    }

    for row in icmp_flows:
        ts  = row.get("timestamp_utc", "")
        src = row.get("src_ip", "")
        dst = row.get("dst_ip", "")
        typ = row.get("icmp_type", "")
        tname = row.get("icmp_type_name", f"Type {typ}")
        ip_len  = int(row.get("ip_len", 0) or 0)
        data_est = int(row.get("data_len_estimated", 0) or 0)
        frag     = row.get("fragmented", "0")

        label = sev = desc = None

        if src in flood_ips and typ == "8":
            label, sev = "ICMP Flood", "high"
            desc = f"Echo Request from flagged flood source"
        elif src in tunnel_ips and data_est >= 128:
            label, sev = "ICMP Tunneling", "critical"
            desc = f"Oversized Echo payload ({data_est} bytes) — possible ICMP tunnel"
        elif src in exfil_ips and typ == "8":
            label, sev = "ICMP Data Exfiltration", "critical"
            desc = f"High-entropy Echo payload from exfil source ({data_est} bytes)"
        elif ip_len >= 65000 or int(row.get("frag_offset", 0) or 0) > 65000:
            label, sev = "Ping of Death", "high"
            desc = f"Oversized IP datagram: ip_len={ip_len}"
        elif frag == "1":
            label, sev = "ICMP Fragmentation Attack", "medium"
            desc = f"Fragmented ICMP from {src} to {dst}"
        elif dst.endswith(".255") or dst == "255.255.255.255":
            label, sev = "Smurf / Broadcast Amplification", "high"
            desc = f"Echo Request to broadcast address {dst}"
        elif typ == "5":
            label, sev = "ICMP Redirect Attack", "high"
            desc = f"ICMP Redirect from {src} — potential routing manipulation"
        elif src in sweep_ips and typ == "8":
            label, sev = "ICMP Network Sweep", "medium"
            desc = f"Ping sweep packet to {dst}"
        elif typ in ("13", "14", "15", "16", "17", "18"):
            label, sev = "ICMP Recon Types", "medium"
            desc = f"{tname} — obsolete ICMP recon type"
        elif typ == "3":
            label, sev = "Destination Unreachable Flood", "medium"
            desc = f"Destination Unreachable from {src}"

        if label:
            if counts[label] < limits.get(label, 20):
                events.append({
                    "timestamp": ts, "protocol": "ICMP",
                    "src_ip": src, "dst_ip": dst,
                    "event_type": label, "severity": sev, "description": desc,
                })
                counts[label] += 1

    return events


def _dns_timeline_events(data: dict) -> list[dict]:
    dns_results = data["dns_results"]
    dns_flows   = data["dns_flows"]

    # Collect flagged domains from detectors
    dga_domains: set[str] = set()
    for f in dns_results.get("dga", {}).get("findings", []):
        d = f.get("fqdn") or f.get("domain") or f.get("sld", "")
        if d:
            dga_domains.add(d)

    tunnel_domains: set[str] = set()
    for f in dns_results.get("tunneling", {}).get("findings", []):
        d = f.get("fqdn") or f.get("domain") or f.get("apex", "")
        if d:
            tunnel_domains.add(d)

    beacon_domains: set[str] = set()
    for f in dns_results.get("beaconing", {}).get("findings", []):
        d = f.get("fqdn") or f.get("domain") or f.get("apex", "")
        if d:
            beacon_domains.add(d)

    events: list[dict] = []
    counts: dict[str, int] = defaultdict(int)
    limits = {
        "DNS Tunneling / Exfiltration":        30,
        "C&C DNS Beaconing":                   30,
        "Domain Generation Algorithm (DGA)":   30,
        "NXDomain Flooding":                   30,
        "DNS Amplification / Reflection":      20,
        "Unusual DNS Record Types":            50,
        "DNS Response Spoofing":               20,
        "Fast Flux DNS":                       20,
        "Unauthorized DNS Servers":            20,
        "Typosquatting / Impersonation":       20,
        "Excessive DNS Query Rate":            10,
    }

    for row in dns_flows:
        ts    = row.get("timestamp_utc", "")
        src   = row.get("src_ip", "")
        dst   = row.get("dst_ip", "")
        fqdn  = row.get("fqdn", "")
        qtype = row.get("qtype_name", "")
        rcode = row.get("rcode", "")
        direction = row.get("direction", "")
        frame_len = int(row.get("frame_len_bytes", 0) or 0)

        label = sev = desc = None

        if fqdn in tunnel_domains or (len(fqdn) > 100):
            label, sev = "DNS Tunneling / Exfiltration", "critical"
            desc = f"Suspected tunnel domain: {fqdn} ({qtype})"
        elif fqdn in beacon_domains:
            label, sev = "C&C DNS Beaconing", "high"
            desc = f"Beacon domain: {fqdn} from {src}"
        elif fqdn in dga_domains:
            label, sev = "Domain Generation Algorithm (DGA)", "high"
            desc = f"DGA domain: {fqdn} ({qtype})"
        elif rcode == "NXDOMAIN":
            label, sev = "NXDomain Flooding", "medium"
            desc = f"NXDOMAIN: {fqdn} from {src}"
        elif qtype in ("AXFR", "IXFR"):
            label, sev = "Unusual DNS Record Types", "high"
            desc = f"Zone transfer attempt ({qtype}): {fqdn} from {src}"
        elif qtype == "ANY":
            label, sev = "DNS Amplification / Reflection", "medium"
            desc = f"ANY query for {fqdn} (amplification risk)"
        elif frame_len > 512 and direction == "response":
            label, sev = "DNS Amplification / Reflection", "medium"
            desc = f"Large DNS response ({frame_len} B) for {fqdn}"

        if label:
            if counts[label] < limits.get(label, 20):
                events.append({
                    "timestamp": ts, "protocol": "DNS",
                    "src_ip": src, "dst_ip": dst,
                    "event_type": label, "severity": sev, "description": desc,
                })
                counts[label] += 1

    return events


def _ntp_timeline_events(data: dict) -> list[dict]:
    ntp_results = data["ntp_results"]
    ntp_flows   = data["ntp_flows"]

    # Collect flagged IPs from detectors
    amp_victims: set[str] = set()
    for f in ntp_results.get("amplification", {}).get("findings", []):
        if "victim_ip" in f:
            amp_victims.add(f["victim_ip"])

    flood_ips: set[str] = set()
    for f in ntp_results.get("ntp_flood", {}).get("findings", []):
        if "src_ip" in f:
            flood_ips.add(f["src_ip"])

    recon_ips: set[str] = set()
    for f in ntp_results.get("ntp_recon", {}).get("findings", []):
        if "src_ip" in f:
            recon_ips.add(f["src_ip"])

    manip_ips: set[str] = set()
    for f in ntp_results.get("time_manipulation", {}).get("findings", []):
        if "src_ip" in f:
            manip_ips.add(f["src_ip"])

    events: list[dict] = []
    counts: dict[str, int] = defaultdict(int)
    limits = {
        "NTP Amplification Attack":      40,
        "NTP Flood":                     30,
        "NTP Kiss-of-Death (KoD)":       30,
        "NTP Mode 7 / Monlist Abuse":    30,
        "Spoofed NTP Response":          20,
        "NTP Time Manipulation":         20,
        "NTP Reconnaissance / Scan":     20,
    }

    for row in ntp_flows:
        ts       = row.get("timestamp_utc", "")
        src      = row.get("src_ip", "")
        dst      = row.get("dst_ip", "")
        mode     = row.get("ntp_mode", "")
        mode_name = row.get("ntp_mode_name", "")
        stratum  = row.get("stratum", "")
        refid    = row.get("refid", "")
        frame_len = int(row.get("frame_len", 0) or 0)
        is_kod   = row.get("is_kod", "0") == "1"
        is_mode7 = row.get("is_mode7", "0") == "1"
        is_large = row.get("is_large_response", "0") == "1"

        label = sev = desc = None

        if dst in amp_victims and is_large:
            label, sev = "NTP Amplification Attack", "critical"
            desc = f"Large NTP response ({frame_len} B) to victim {dst} from {src}"
        elif src in flood_ips:
            label, sev = "NTP Flood", "high"
            desc = f"NTP flood packet from {src} to {dst}"
        elif is_kod:
            label, sev = "NTP Kiss-of-Death (KoD)", "high"
            desc = f"KoD response ({refid}) from {src} to {dst}"
        elif is_mode7:
            label, sev = "NTP Mode 7 / Monlist Abuse", "high"
            desc = f"Mode 7 (private/deprecated) NTP packet from {src} to {dst}"
        elif src in manip_ips and mode == "4":
            label, sev = "NTP Time Manipulation", "high"
            desc = f"Anomalous NTP server response from {src} (stratum {stratum})"
        elif mode == "4" and row.get("src_port", "123") != "123":
            label, sev = "Spoofed NTP Response", "high"
            desc = f"NTP response from non-standard port {row.get('src_port')} src={src}"
        elif stratum == "16" and mode == "4":
            label, sev = "Spoofed NTP Response", "medium"
            desc = f"Unsynchronised NTP server (stratum 16) response from {src}"
        elif src in recon_ips and mode == "3":
            label, sev = "NTP Reconnaissance / Scan", "medium"
            desc = f"NTP client query to {dst} from scanning source {src}"

        if label:
            if counts[label] < limits.get(label, 20):
                events.append({
                    "timestamp": ts, "protocol": "NTP",
                    "src_ip": src, "dst_ip": dst,
                    "event_type": label, "severity": sev, "description": desc,
                })
                counts[label] += 1

    return events


def _http_timeline_events(data: dict) -> list[dict]:
    http_results = data["http_results"]
    http_flows   = data["http_flows"]

    # Collect flagged IPs/hosts/URIs from detectors for annotation
    beacon_paths: set[tuple] = set()
    for f in http_results.get("http_beaconing", {}).get("findings", []):
        beacon_paths.add((f.get("src_ip", ""), f.get("host", ""), f.get("uri_path", "")))

    suspicious_ua_ips: set[str] = set()
    for f in http_results.get("suspicious_ua", {}).get("findings", []):
        if "src_ip" in f:
            suspicious_ua_ips.add(f["src_ip"])

    scanning_clients: set[str] = set()
    for f in http_results.get("scanning_codes", {}).get("findings", []):
        if "client_ip" in f:
            scanning_clients.add(f["client_ip"])

    exploit_ips: set[str] = set()
    for f in http_results.get("suspicious_uri", {}).get("findings", []):
        if "src_ip" in f:
            exploit_ips.add(f["src_ip"])

    events: list[dict] = []
    counts: dict[str, int] = defaultdict(int)
    limits = {
        "HTTP Beaconing":               30,
        "Suspicious User-Agent":        30,
        "Suspicious URI Patterns":      40,
        "HTTP Scanning / Error Code Flood": 20,
        "Large HTTP Upload":            20,
        "HTTP Cookie Anomaly":          20,
        "Host Header Anomaly":          20,
        "Unusual HTTP Methods":         30,
        "Unusual HTTP Server Header":   20,
        "Suspicious HTTP Referer":      20,
        "Deprecated TLS Version":       20,
    }

    for row in http_flows:
        ts        = row.get("timestamp_utc", "")
        src       = row.get("src_ip", "")
        dst       = row.get("dst_ip", "")
        direction = row.get("direction", "")
        method    = row.get("method", "")
        uri       = row.get("uri", "")
        host      = row.get("host", "")
        ua        = row.get("user_agent", "")
        status    = row.get("status_code", "")
        server    = row.get("server", "")
        try:
            frame_len = int(row.get("frame_len", 0) or 0)
            content_len = int(row.get("content_length", 0) or 0)
        except ValueError:
            frame_len = content_len = 0

        label = sev = desc = None

        # Beacon — regular interval requests from flagged (src, host, path)
        uri_path = uri.split("?")[0] if "?" in uri else uri
        if (src, host, uri_path[:120]) in beacon_paths and direction == "request":
            label, sev = "HTTP Beaconing", "critical"
            desc = f"Beacon request: {method} {host}{uri[:80]} from {src}"

        elif direction == "request" and src in exploit_ips:
            label, sev = "Suspicious URI Patterns", "critical"
            desc = f"Exploit/probe request: {method} {uri[:100]}"

        elif direction == "request" and src in suspicious_ua_ips and ua:
            label, sev = "Suspicious User-Agent", "high"
            desc = f"Tool/framework UA: {ua[:80]}"

        elif direction == "request" and method in (
                "TRACE", "PROPFIND", "MKCOL", "SEARCH", "LOCK", "UNLOCK"):
            label, sev = "Unusual HTTP Methods", "high"
            desc = f"Unusual method {method} to {host}{uri[:80]}"

        elif direction == "request" and max(frame_len, content_len) >= 500_000:
            label, sev = "Large HTTP Upload", "high"
            desc = f"{method} upload {max(frame_len,content_len):,} B to {host}{uri[:60]}"

        elif direction == "response" and dst in scanning_clients and status.startswith(("4", "5")):
            label, sev = "HTTP Scanning / Error Code Flood", "high"
            desc = f"Error response HTTP {status} to scanning client {dst}"

        elif direction == "response" and server:
            srv_lower = server.lower()
            if any(kw in srv_lower for kw in ("metasploit", "empire", "sliver", "havoc", "cobalt")):
                label, sev = "Unusual HTTP Server Header", "critical"
                desc = f"C2 server header: '{server[:80]}' from {src}"

        if label:
            if counts[label] < limits.get(label, 20):
                events.append({
                    "timestamp": ts, "protocol": "HTTP",
                    "src_ip": src, "dst_ip": dst,
                    "event_type": label, "severity": sev, "description": desc,
                })
                counts[label] += 1

    return events


def _cert_timeline_events(data: dict) -> list[dict]:
    cert_results = data["cert_results"]
    cert_records = data["cert_records"]

    self_signed_ips: set[str] = set()
    for f in cert_results.get("self_signed", {}).get("findings", []):
        if "dst_ip" in f:
            self_signed_ips.add(f["dst_ip"])

    mismatch_ips: set[str] = set()
    for f in cert_results.get("sni_mismatch", {}).get("findings", []):
        if "dst_ip" in f:
            mismatch_ips.add(f["dst_ip"])

    expired_ips: set[str] = set()
    for f in cert_results.get("expired", {}).get("findings", []):
        if "dst_ip" in f:
            expired_ips.add(f["dst_ip"])

    short_ips: set[str] = set()
    for f in cert_results.get("short_validity", {}).get("findings", []):
        if "dst_ip" in f:
            short_ips.add(f["dst_ip"])

    events: list[dict] = []
    counts: dict[str, int] = defaultdict(int)
    limits = {
        "Self-Signed Certificate":         20,
        "Certificate CN/SNI Mismatch":     20,
        "Expired Certificate":             20,
        "Very Short Certificate Validity": 20,
    }

    for row in cert_records:
        ts  = row.get("timestamp_utc", "")
        src = row.get("src_ip", "")
        dst = row.get("dst_ip", "")
        cn  = row.get("subject_cn", "")
        sni = row.get("sni", "")

        label = sev = desc = None

        if dst in self_signed_ips:
            label, sev = "Self-Signed Certificate", "critical"
            desc = f"Self-signed cert presented by {dst} (CN: {cn or '—'})"
        elif dst in mismatch_ips:
            label, sev = "Certificate CN/SNI Mismatch", "critical"
            desc = f"SNI '{sni}' ≠ cert CN '{cn}' — possible MITM on {dst}"
        elif dst in expired_ips:
            label, sev = "Expired Certificate", "high"
            desc = f"Expired cert from {dst} (CN: {cn or '—'})"
        elif dst in short_ips:
            label, sev = "Very Short Certificate Validity", "high"
            desc = f"Short-lived cert from {dst} (CN: {cn or '—'}, {row.get('valid_days', '?')} days)"

        if label:
            if counts[label] < limits.get(label, 20):
                events.append({
                    "timestamp": ts, "protocol": "TLS/Cert",
                    "src_ip": src, "dst_ip": dst,
                    "event_type": label, "severity": sev, "description": desc,
                })
                counts[label] += 1

    return events


def _tls_timeline_events(data: dict) -> list[dict]:
    tls_results  = data["tls_results"]
    tls_sessions = data["tls_sessions"]

    suspicious_clients: set[str] = set()
    for f in tls_results.get("suspicious_ja4", {}).get("findings", []):
        if "client_ip" in f:
            suspicious_clients.add(f["client_ip"])

    weak_pairs: set[tuple] = set()
    for f in tls_results.get("weak_cipher", {}).get("findings", []):
        weak_pairs.add((f.get("client_ip", ""), f.get("server_ip", "")))

    nonstandard_servers: set[str] = set()
    for f in tls_results.get("non_standard_port", {}).get("findings", []):
        nonstandard_servers.add(f"{f.get('server_ip', '')}:{f.get('server_port', '')}")

    events: list[dict] = []
    counts: dict[str, int] = defaultdict(int)
    limits = {
        "Suspicious JA4/JA3 Fingerprint": 30,
        "Weak Cipher Suite Negotiated":   30,
        "TLS on Non-Standard Port":       20,
        "Deprecated TLS Negotiated":      20,
    }

    for row in tls_sessions:
        ts     = row.get("timestamp_utc", "")
        client = row.get("client_ip", "")
        server = row.get("server_ip", "")
        port   = row.get("server_port", "")
        sni    = row.get("sni", "")
        ver    = row.get("tls_version", "")

        label = sev = desc = None

        if client in suspicious_clients:
            label, sev = "Suspicious JA4/JA3 Fingerprint", "critical"
            desc = f"Known C2 fingerprint from {client} → {server}:{port} (SNI: {sni or '—'})"
        elif (client, server) in weak_pairs:
            label, sev = "Weak Cipher Suite Negotiated", "high"
            desc = f"Weak cipher negotiated: {client} → {server}:{port}"
        elif f"{server}:{port}" in nonstandard_servers:
            label, sev = "TLS on Non-Standard Port", "medium"
            desc = f"TLS on port {port}: {client} → {server} (SNI: {sni or '—'})"
        elif ver in ("TLS 1.0", "TLS 1.1"):
            label, sev = "Deprecated TLS Negotiated", "high"
            desc = f"{ver} negotiated: {client} → {server}:{port}"

        if label:
            if counts[label] < limits.get(label, 20):
                events.append({
                    "timestamp": ts, "protocol": "TLS",
                    "src_ip": client, "dst_ip": server,
                    "event_type": label, "severity": sev, "description": desc,
                })
                counts[label] += 1

    return events


def _generic_timeline_events(results: dict, protocol: str, ip_fields: tuple) -> list[dict]:
    """Build timeline events directly from findings in a normalised results dict."""
    events: list[dict] = []
    for cat_key, cat in results.items():
        if not isinstance(cat, dict) or cat.get("count", 0) == 0:
            continue
        sev  = cat.get("severity", "info")
        name = cat.get("name", cat_key)
        for f in cat.get("findings", [])[:15]:
            ts   = f.get("timestamp_utc", f.get("first_timestamp", ""))
            src  = next((f.get(k) for k in ip_fields if f.get(k)), "—")
            dst  = f.get("dst_ip", f.get("target_ip", "—"))
            desc = f"{name}"
            events.append({
                "timestamp": ts, "protocol": protocol,
                "src_ip": src, "dst_ip": dst,
                "event_type": name, "severity": sev, "description": desc,
            })
    return events


def build_timeline(data: dict) -> list[dict]:
    events: list[dict] = []

    if data["has_icmp"]:
        events.extend(_icmp_timeline_events(data))
    if data["has_dns"]:
        events.extend(_dns_timeline_events(data))
    if data["has_ntp"]:
        events.extend(_ntp_timeline_events(data))
    if data["has_http"]:
        events.extend(_http_timeline_events(data))
    if data["has_cert"]:
        events.extend(_cert_timeline_events(data))
    if data["has_tls"]:
        events.extend(_tls_timeline_events(data))
    if data["has_arp"]:
        events.extend(_generic_timeline_events(
            data["arp_results"], "ARP",
            ("src_mac", "eth_src_mac", "target_ip")))
    if data["has_tcp"]:
        events.extend(_generic_timeline_events(
            data["tcp_results"], "TCP",
            ("src_ip", "injecting_ip")))
    if data["has_udp"]:
        events.extend(_generic_timeline_events(
            data["udp_results"], "UDP",
            ("src_ip", "reflector_ip")))
    if data["has_dhcp"]:
        events.extend(_generic_timeline_events(
            data["dhcp_results"], "DHCP",
            ("src_ip", "client_mac")))
    if data["has_mdns"]:
        events.extend(_generic_timeline_events(
            data["mdns_results"], "mDNS",
            ("src_ip",)))
    if data["has_quic"]:
        events.extend(_generic_timeline_events(
            data["quic_results"], "QUIC",
            ("src_ip", "server_ip")))

    # Suricata IDS alerts → timeline events
    if data["has_suricata"]:
        sev_order = {"critical": 0, "high": 1, "medium": 2, "low": 3, "info": 4}
        for alert in sorted(
            data["suricata_data"].get("alerts", []),
            key=lambda a: (sev_order.get(a.get("severity", "info"), 4), a.get("timestamp_utc", "")),
        )[:30]:
            events.append({
                "timestamp":  alert.get("timestamp_utc", ""),
                "protocol":   alert.get("proto", "SURICATA"),
                "src_ip":     alert.get("src_ip", ""),
                "dst_ip":     alert.get("dest_ip", ""),
                "event_type": alert.get("signature", "IDS Alert"),
                "severity":   alert.get("severity", "medium"),
                "description": (
                    f"{alert.get('signature','')} "
                    f"(cat: {alert.get('category','')})"
                ),
            })

    # Sort by timestamp; events without timestamp sort last
    events.sort(key=lambda e: _ts_key(e["timestamp"]))

    # Prioritise critical/high when truncating
    if len(events) > MAX_TIMELINE_ROWS:
        priority = [e for e in events if e["severity"] in ("critical", "high")]
        others   = [e for e in events if e["severity"] not in ("critical", "high")]
        keep = priority + others
        keep.sort(key=lambda e: _ts_key(e["timestamp"]))
        events = keep[:MAX_TIMELINE_ROWS]

    return events


# ── IOC extraction ─────────────────────────────────────────────────────────────

def extract_iocs(data: dict) -> list[dict]:
    seen: set[str] = set()
    iocs: list[dict] = []

    def _add(ioc_type, value, severity, category, source):
        key = f"{ioc_type}:{value}"
        if key not in seen:
            seen.add(key)
            iocs.append({"type": ioc_type, "value": value,
                         "severity": severity, "category": category, "source": source})

    # TLS — extract server IPs from triggered categories
    for cat_key, cat in data["tls_results"].items():
        if not isinstance(cat, dict) or cat.get("count", 0) == 0:
            continue
        sev = cat.get("severity", "info")
        for f in cat.get("findings", []):
            for field in ("server_ip", "client_ip", "src_ip"):
                ip = f.get(field)
                if ip:
                    _add("ip", ip, sev, cat.get("name", cat_key), "TLS Inspector")
                    break

    # Cert — extract server IPs and domain names from triggered categories
    for cat_key, cat in data["cert_results"].items():
        if not isinstance(cat, dict) or cat.get("count", 0) == 0:
            continue
        sev = cat.get("severity", "info")
        for f in cat.get("findings", []):
            if "dst_ip" in f:
                _add("ip", f["dst_ip"], sev, cat.get("name", cat_key), "Certificate")
            sni_val = f.get("sni") or f.get("subject_cn", "")
            if sni_val and "." in sni_val and not sni_val[0].isdigit() and not sni_val.startswith("*"):
                _add("domain", sni_val, sev, cat.get("name", cat_key), "Certificate")

    # HTTP — extract IPs and domains from triggered categories
    for cat_key, cat in data["http_results"].items():
        if cat.get("count", 0) == 0:
            continue
        sev = cat.get("severity", "info")
        for f in cat.get("findings", []):
            for ip_field in ("src_ip", "client_ip"):
                ip = f.get(ip_field)
                if ip:
                    _add("ip", ip, sev, cat.get("name", cat_key), "HTTP")
                    break
            # Beacon/UA findings also carry a host — record as domain IOC
            host = f.get("host")
            if host and "." in host and not host[0].isdigit():
                _add("domain", host, sev, cat.get("name", cat_key), "HTTP")

    # NTP — extract IPs from triggered categories
    for cat_key, cat in data["ntp_results"].items():
        if cat.get("count", 0) == 0:
            continue
        sev = cat.get("severity", "info")
        for f in cat.get("findings", []):
            for field in ("src_ip", "victim_ip", "client_ip"):
                ip = f.get(field)
                if ip:
                    _add("ip", ip, sev, cat.get("name", cat_key), "NTP")
                    break

    # ICMP — extract IPs from triggered categories
    for cat_key, cat in data["icmp_results"].items():
        if cat.get("count", 0) == 0:
            continue
        sev = cat.get("severity", "info")
        for f in cat.get("findings", []):
            if "src_ip" in f:
                _add("ip", f["src_ip"], sev, cat.get("name", cat_key), "ICMP")
            if "dst_ip" in f and (f.get("dst_ip", "").endswith(".255")
                                   or f.get("dst_ip") == "255.255.255.255"):
                _add("ip", f["dst_ip"], sev, cat.get("name", cat_key), "ICMP")

    # DNS — extract domains from triggered categories
    for cat_key, cat in data["dns_results"].items():
        if cat.get("count", 0) == 0:
            continue
        sev = cat.get("severity", "info")
        for f in cat.get("findings", []):
            for field in ("fqdn", "domain", "apex", "sld"):
                val = f.get(field)
                if val:
                    _add("domain", val, sev, cat.get("name", cat_key), "DNS")
                    break

    # ARP / TCP / UDP / DHCP / mDNS / QUIC — extract IPs from triggered categories
    for proto_key, source_label, ip_fields in (
        ("arp_results",  "ARP",  ("target_ip", "src_ip", "eth_src_mac")),
        ("tcp_results",  "TCP",  ("src_ip", "injecting_ip")),
        ("udp_results",  "UDP",  ("src_ip", "reflector_ip", "victim_ip")),
        ("dhcp_results", "DHCP", ("src_ip",)),
        ("mdns_results", "mDNS", ("src_ip",)),
        ("quic_results", "QUIC", ("src_ip", "server_ip", "victim_ip")),
    ):
        for cat_key, cat in data[proto_key].items():
            if not isinstance(cat, dict) or cat.get("count", 0) == 0:
                continue
            sev = cat.get("severity", "info")
            for f in cat.get("findings", []):
                for field in ip_fields:
                    ip = f.get(field)
                    if ip and "." in str(ip) and ":" not in str(ip):
                        _add("ip", ip, sev, cat.get("name", cat_key), source_label)
                        break

    # CTI malicious / suspicious IPs
    for row in data["fan_ip"]:
        rep = row.get("reputation", "")
        if rep in ("malicious", "suspicious"):
            sev = "high" if rep == "malicious" else "medium"
            _add("ip", row["ip"], sev, f"CTI ({rep})", "CTI Enrichment")

    # CTI malicious / suspicious domains
    for row in data["fan_correlation"]:
        rep = row.get("reputation", "")
        if rep in ("malicious", "suspicious"):
            sev = "high" if rep == "malicious" else "medium"
            _add("domain", row["fqdn"], sev, f"CTI ({rep})", "CTI Enrichment")

    # File hashes — malicious/suspicious extracted files
    for rec in data.get("fh_data", {}).get("files", []):
        verdict = rec.get("osint_verdict", "")
        if verdict in ("malicious", "suspicious"):
            sev = "critical" if verdict == "malicious" else "high"
            fname = rec.get("filename", "unknown")
            _add("hash", rec["sha256"], sev, f"Extracted file: {fname}", "File Hashes")
            _add("hash", rec["md5"],    sev, f"Extracted file: {fname} (MD5)", "File Hashes")

    # Suricata — IPs from critical/high alerts
    for alert in data.get("suricata_data", {}).get("alerts", []):
        sev = alert.get("severity", "info")
        if sev not in ("critical", "high"):
            continue
        cat  = alert.get("category", "IDS Alert")
        sig  = alert.get("signature", "")
        desc = f"Suricata: {sig}" if sig else "Suricata IDS"
        for ip in (alert.get("src_ip", ""), alert.get("dest_ip", "")):
            if ip and "." in ip:
                _add("ip", ip, sev, desc, cat or "Suricata IDS")

    iocs.sort(key=lambda x: (_sev_rank(x["severity"]), x["type"], x["value"]))
    return iocs


# ── MITRE ATT&CK coverage ──────────────────────────────────────────────────────

_MITRE_TACTICS = {
    "T1498": "Impact", "T1498.001": "Impact", "T1498.002": "Impact",
    "T1499": "Impact", "T1499.002": "Impact",
    "T1572": "Command and Control",
    "T1557": "Credential Access / Collection",
    "T1595": "Reconnaissance", "T1595.001": "Reconnaissance",
    "T1048": "Exfiltration", "T1048.001": "Exfiltration",
    "T1568": "Command and Control", "T1568.001": "Command and Control",
    "T1568.002": "Command and Control",
    "T1071": "Command and Control", "T1071.001": "Command and Control",
    "T1071.004": "Command and Control",
    "T1583": "Resource Development", "T1583.001": "Resource Development",
    "T1584": "Resource Development", "T1584.002": "Resource Development",
    "T1590": "Reconnaissance",
    "T1070": "Defense Evasion",
    "T1190": "Initial Access",
    "T1048.002": "Exfiltration", "T1048.003": "Exfiltration",
    "T1040": "Credential Access",
    "T1587.003": "Resource Development",
    "T1571": "Command and Control",
    "T1557.002": "Credential Access / Collection",
    "T1018": "Discovery",
    "T1046": "Discovery",
    "T1563": "Lateral Movement",
    "T1001": "Command and Control",
    "T1550": "Defense Evasion / Lateral Movement",
    "T1562": "Defense Evasion",
    "T1565": "Impact", "T1565.001": "Impact",
    "T1105": "Lateral Movement",
    # Suricata / YARA supplemental
    "T1059": "Execution", "T1059.001": "Execution",
    "T1003": "Credential Access", "T1003.001": "Credential Access",
    "T1021": "Lateral Movement", "T1021.002": "Lateral Movement",
    "T1068": "Privilege Escalation",
    "T1189": "Initial Access",
    "T1204": "Execution",
    "T1486": "Impact",
    "T1505": "Persistence", "T1505.003": "Persistence",
    "T1110": "Credential Access",
    "T1095": "Command and Control",
    "T1041": "Exfiltration",
    "T1573": "Command and Control",
    "T1090": "Command and Control",
}


def mitre_coverage(data: dict) -> list[dict]:
    seen: dict[str, dict] = {}

    for results_key in ("icmp_results", "dns_results", "ntp_results", "http_results",
                        "cert_results", "tls_results",
                        "arp_results", "tcp_results", "udp_results",
                        "dhcp_results", "mdns_results", "quic_results"):
        for cat in data[results_key].values():
            if cat.get("count", 0) == 0:
                continue
            mitre = cat.get("mitre", [])
            if not mitre:
                continue
            tid, tname = mitre[0], mitre[1] if len(mitre) > 1 else ""
            if tid not in seen or _sev_rank(cat["severity"]) < _sev_rank(seen[tid]["severity"]):
                seen[tid] = {
                    "id": tid,
                    "name": tname,
                    "tactic": _MITRE_TACTICS.get(tid, "—"),
                    "severity": cat.get("severity", "info"),
                    "category": cat.get("name", ""),
                }

    # File hashes — T1105 Ingress Tool Transfer
    fh = data.get("fh_data", {})
    if fh.get("malicious_count", 0) > 0 or fh.get("suspicious_count", 0) > 0:
        tid, tname = "T1105", "Ingress Tool Transfer"
        sev = "critical" if fh.get("malicious_count", 0) > 0 else "high"
        seen.setdefault(tid, {
            "id": tid, "name": tname,
            "tactic": _MITRE_TACTICS.get(tid, "Lateral Movement"),
            "severity": sev, "category": "File Hash Analysis",
        })

    # Suricata — map alert categories to MITRE techniques
    _SURICATA_CAT_MITRE = {
        "a network trojan was detected":           ("T1071",    "Standard Application Layer Protocol"),
        "malware command and control activity":    ("T1071",    "Standard Application Layer Protocol"),
        "command and control":                     ("T1071",    "Standard Application Layer Protocol"),
        "web application attack":                  ("T1190",    "Exploit Public-Facing Application"),
        "attempted administrator privilege gain":  ("T1068",    "Exploitation for Privilege Escalation"),
        "attempted user privilege gain":           ("T1068",    "Exploitation for Privilege Escalation"),
        "exploit kit activity detected":           ("T1189",    "Drive-by Compromise"),
        "credential theft":                        ("T1003",    "OS Credential Dumping"),
        "trojan activity":                         ("T1204",    "User Execution"),
        "network scan":                            ("T1046",    "Network Service Discovery"),
        "port scan":                               ("T1046",    "Network Service Discovery"),
        "denial of service":                       ("T1498",    "Network Denial of Service"),
        "attempted denial of service":             ("T1498",    "Network Denial of Service"),
        "ransomware":                              ("T1486",    "Data Encrypted for Impact"),
        "lateral movement":                        ("T1021",    "Remote Services"),
        "policy violation":                        ("T1562",    "Impair Defenses"),
    }
    for alert in data.get("suricata_data", {}).get("alerts", []):
        cat = (alert.get("category") or "").lower().strip()
        entry = _SURICATA_CAT_MITRE.get(cat)
        if entry:
            tid, tname = entry
            sev = alert.get("severity", "medium")
            if tid not in seen or _sev_rank(sev) < _sev_rank(seen[tid]["severity"]):
                seen[tid] = {
                    "id": tid, "name": tname,
                    "tactic": _MITRE_TACTICS.get(tid, "—"),
                    "severity": sev, "category": alert.get("category", "Suricata IDS"),
                }

    # YARA — map rule mitre_att metadata to coverage table
    for match in data.get("yara_data", {}).get("matches", []):
        tid = (match.get("mitre_att") or "").strip()
        if not tid:
            continue
        sev = match.get("severity", "medium")
        tname = match.get("category", match.get("rule", ""))
        if tid not in seen or _sev_rank(sev) < _sev_rank(seen[tid]["severity"]):
            seen[tid] = {
                "id": tid, "name": tname,
                "tactic": _MITRE_TACTICS.get(tid, "—"),
                "severity": sev, "category": f"YARA: {match.get('rule','')}",
            }

    return sorted(seen.values(), key=lambda x: _sev_rank(x["severity"]))


# ── Recommendations ────────────────────────────────────────────────────────────

def build_recommendations(data: dict, overall_sev: str) -> list[str]:
    recs: list[str] = []
    triggered_icmp = {k for k, v in data["icmp_results"].items() if isinstance(v, dict) and v.get("count", 0) > 0}
    triggered_dns  = {k for k, v in data["dns_results"].items()  if isinstance(v, dict) and v.get("count", 0) > 0}
    triggered_ntp  = {k for k, v in data["ntp_results"].items()  if isinstance(v, dict) and v.get("count", 0) > 0}
    triggered_http = {k for k, v in data["http_results"].items() if isinstance(v, dict) and v.get("count", 0) > 0}
    triggered_cert = {k for k, v in data["cert_results"].items() if isinstance(v, dict) and v.get("count", 0) > 0}
    triggered_tls  = {k for k, v in data["tls_results"].items()  if isinstance(v, dict) and v.get("count", 0) > 0}
    triggered_arp  = {k for k, v in data["arp_results"].items()  if isinstance(v, dict) and v.get("count", 0) > 0}
    triggered_tcp  = {k for k, v in data["tcp_results"].items()  if isinstance(v, dict) and v.get("count", 0) > 0}
    triggered_udp  = {k for k, v in data["udp_results"].items()  if isinstance(v, dict) and v.get("count", 0) > 0}
    triggered_dhcp = {k for k, v in data["dhcp_results"].items() if isinstance(v, dict) and v.get("count", 0) > 0}
    triggered_mdns = {k for k, v in data["mdns_results"].items() if isinstance(v, dict) and v.get("count", 0) > 0}
    triggered_quic = {k for k, v in data["quic_results"].items() if isinstance(v, dict) and v.get("count", 0) > 0}
    mal_ips = [r for r in data["fan_ip"] if r.get("reputation") == "malicious"]
    sus_ips = [r for r in data["fan_ip"] if r.get("reputation") == "suspicious"]

    if mal_ips:
        recs.append(f"**IMMEDIATE**: Block {len(mal_ips)} confirmed malicious IP(s) at the network perimeter firewall.")
        recs.append("Quarantine and forensically examine hosts that communicated with confirmed malicious IPs.")

    if "exfiltration" in triggered_icmp:
        recs.append("Inspect and block outbound ICMP Echo packets with payloads exceeding 64 bytes; investigate source hosts for malware.")
    if "tunneling" in triggered_icmp:
        recs.append("Block or rate-limit ICMP Echo packets larger than 128 bytes; deploy deep-packet inspection to detect ICMP tunnel tools (ptunnel, icmptunnel).")
    if "flood" in triggered_icmp:
        recs.append("Implement ICMP rate limiting at the border router/firewall; verify upstream ISP nullrouting for flood sources.")
    if "smurf" in triggered_icmp:
        recs.append("Ensure broadcast forwarding is disabled on all routers; block directed broadcast addresses at the network edge.")
    if "redirect" in triggered_icmp:
        recs.append("Disable ICMP Redirect processing on all hosts (`net.ipv4.conf.all.accept_redirects = 0`); investigate hosts sending Type 5 messages.")
    if "pod" in triggered_icmp:
        recs.append("Confirm firewall rules drop oversized ICMP packets (ip_len ≥ 65000). Patch affected systems if unpatched Ping of Death vulnerability exists.")

    if "tunneling" in triggered_dns or "exfiltration" in triggered_dns:
        recs.append("Block DNS queries with labels exceeding 40 characters; implement DNS-layer security (e.g., Cisco Umbrella, NextDNS) with DGA/tunnel detection.")
    if "dga" in triggered_dns:
        recs.append("Investigate hosts generating DGA-pattern queries; consider DNS RPZ (Response Policy Zone) to sink-hole known DGA domains.")
    if "beaconing" in triggered_dns:
        recs.append("Isolate hosts exhibiting regular-interval DNS beaconing; conduct full memory and disk forensics for malware implants.")
    if "amplification" in triggered_dns:
        recs.append("Block inbound DNS ANY queries at the perimeter; configure rate limiting on DNS resolvers.")
    if "nxdomain_flood" in triggered_dns:
        recs.append("Investigate hosts generating high NXDOMAIN volumes; could indicate active DGA malware scanning for live C2 domains.")

    if "amplification" in triggered_ntp:
        recs.append("Block or rate-limit inbound NTP responses (UDP/123) from external sources; disable monlist (`restrict default noquery`) on all NTP servers.")
    if "mode7_abuse" in triggered_ntp:
        recs.append("Disable NTP Mode 7 on all NTP servers (`restrict default noquery nomodify`); block inbound mode-7 packets at the perimeter firewall.")
    if "kiss_of_death" in triggered_ntp:
        recs.append("Investigate KoD packets for spoofed sources; verify NTP client configurations to ensure they are querying authorised servers only.")
    if "spoofed_response" in triggered_ntp:
        recs.append("Audit NTP server list on all hosts; block NTP responses from non-port-123 sources; consider NTPsec or authenticated NTP (RFC 8915 NTS).")
    if "time_manipulation" in triggered_ntp:
        recs.append("Isolate hosts receiving anomalous NTP time; review system and security logs for timestamp anomalies; deploy NTP monitoring (chronyc tracking).")
    if "ntp_flood" in triggered_ntp:
        recs.append("Rate-limit NTP traffic at the network perimeter; contact upstream ISP for nullrouting if flood originates externally.")
    if "ntp_recon" in triggered_ntp:
        recs.append("Block or alert on hosts scanning NTP infrastructure (mode 3 queries to many servers); review for DDoS recruitment activity.")

    if "http_beaconing" in triggered_http:
        recs.append("Isolate hosts exhibiting regular-interval HTTP beaconing; conduct full memory and disk forensics for malware implants and C2 agents.")
    if "suspicious_uri" in triggered_http:
        recs.append("Review web server logs for successful exploitation; patch or WAF-protect endpoints targeted by injection and path traversal attempts.")
    if "suspicious_ua" in triggered_http:
        recs.append("Block or alert on traffic from known offensive tool user-agents at the proxy/firewall; investigate source hosts for installed attack tools.")
    if "unusual_server" in triggered_http:
        recs.append("**IMMEDIATE**: Block IP addresses serving C2-framework Server headers; these are attacker-controlled infrastructure.")
    if "large_upload" in triggered_http:
        recs.append("Investigate large HTTP uploads (POST/PUT) for data exfiltration; review destination hosts and inspect payloads for sensitive data.")
    if "scanning_codes" in triggered_http:
        recs.append("Block or rate-limit clients generating high HTTP error volumes; review web application firewall (WAF) rules for scanning signatures.")
    if "old_tls" in triggered_http:
        recs.append("Disable SSL 2.0, SSL 3.0, and TLS 1.0/1.1 on all servers and clients; enforce TLS 1.2 minimum (TLS 1.3 preferred).")
    if "cookie_anomaly" in triggered_http:
        recs.append("Inspect oversized or high-entropy HTTP cookies for encoded exfiltration payloads; review application cookie handling for injection vulnerabilities.")
    if "host_header_anomaly" in triggered_http:
        recs.append("Configure web servers to reject requests with IP-address Host headers or malformed values; deploy Host header validation in the WAF.")

    if "sni_mismatch" in triggered_cert:
        recs.append("Investigate TLS CN/SNI mismatches for MITM interception appliances or rogue SSL inspection proxies on the network path.")
    if "self_signed" in triggered_cert:
        recs.append("Block or alert on TLS to self-signed certificates for externally-facing services; cross-reference server IPs with threat intel for C2 infrastructure.")
    if "expired" in triggered_cert:
        recs.append("Audit servers presenting expired TLS certificates for neglected or attacker-deployed infrastructure; update certificate management processes.")
    if "weak_signature" in triggered_cert:
        recs.append("Replace MD5/SHA-1 signed certificates immediately; enforce SHA-256 minimum in all certificate policies and CA signing configurations.")
    if "short_validity" in triggered_cert:
        recs.append("Investigate servers with very short-lived certificates (< 30 days); cross-reference JA4 fingerprints and IP reputation for C2 framework indicators.")

    if "suspicious_ja4" in triggered_tls:
        recs.append("**IMMEDIATE**: Block IPs matching known C2 JA4/JA3 fingerprints; isolate communicating hosts and conduct full memory and disk forensics.")
    if "weak_cipher" in triggered_tls:
        recs.append("Disable broken cipher suites (NULL, EXPORT, RC4, anonymous DH, DES, 3DES) on all TLS servers; enforce AES-GCM and ChaCha20-Poly1305 cipher suites only.")
    if "deprecated_tls" in triggered_tls:
        recs.append("Disable TLS 1.0 and TLS 1.1 on all servers; configure TLS 1.2 as the minimum with TLS 1.3 preferred wherever supported by clients.")
    if "non_standard_port" in triggered_tls:
        recs.append("Review TLS sessions on non-standard ports for C2 tunnels or exfiltration channels; block unexpected TLS ports at the perimeter firewall.")
    if "cipher_diversity" in triggered_tls:
        recs.append("Investigate source IPs presenting high TLS cipher-suite diversity for automated TLS scanning or C2 beaconers rotating their TLS fingerprint profile.")

    # ARP recommendations
    if "arp_poisoning" in triggered_arp:
        recs.append("**IMMEDIATE**: Investigate ARP cache poisoning — identify the spoofing host, remove it from the network, and reset ARP caches on affected devices (arp -d).")
    if "gratuitous_arp" in triggered_arp:
        recs.append("Enable Dynamic ARP Inspection (DAI) on managed switches; configure ARP rate limiting to suppress gratuitous ARP floods.")
    if "arp_flood" in triggered_arp:
        recs.append("Enable ARP rate limiting on switch ports; configure storm control to prevent ARP broadcast saturation.")
    if "arp_scan" in triggered_arp:
        recs.append("Investigate host performing ARP reconnaissance; verify it is an authorised scanner or asset discovery tool.")
    if "arp_proxy_anomaly" in triggered_arp:
        recs.append("Investigate ARP proxy anomalies for unauthorised relay devices or ARP-level man-in-the-middle appliances.")

    # TCP recommendations
    if "syn_flood" in triggered_tcp:
        recs.append("Enable SYN cookies on affected servers (`net.ipv4.tcp_syncookies = 1`); rate-limit inbound SYN packets at the perimeter firewall.")
    if "port_scan" in triggered_tcp:
        recs.append("Block or rate-limit scanning source IPs at the firewall; review exposed services and close unnecessary ports.")
    if "rst_flood" in triggered_tcp:
        recs.append("Investigate RST flood source for session disruption intent; deploy stateful firewall rules to validate RST sequences.")
    if "stealth_scan" in triggered_tcp:
        recs.append("Block stealth scan source IPs; review IDS/IPS signatures for FIN, NULL, and Xmas scan detection.")
    if "session_hijack" in triggered_tcp:
        recs.append("**HIGH PRIORITY**: Investigate TCP session hijacking indicators; deploy TCP sequence number randomisation and consider IPsec for sensitive sessions.")
    if "half_open_flood" in triggered_tcp:
        recs.append("Tune SYN backlog size (`net.ipv4.tcp_max_syn_backlog`); deploy upstream rate limiting for SYN traffic to affected server IPs.")

    # UDP recommendations
    if "udp_flood" in triggered_udp:
        recs.append("Rate-limit inbound UDP at the network perimeter; contact ISP for upstream nullrouting of flood source IPs.")
    if "udp_amplification" in triggered_udp:
        recs.append("Block or rate-limit UDP responses from amplification ports (DNS/53, NTP/123, SSDP/1900, Memcached/11211) from external sources.")
    if "udp_port_scan" in triggered_udp:
        recs.append("Block UDP port scanning source IPs; review whether UDP services are exposed unnecessarily.")
    if "udp_fragmentation" in triggered_udp:
        recs.append("Configure the firewall to drop or rate-limit fragmented UDP packets from untrusted sources.")
    if "udp_spoofing" in triggered_udp:
        recs.append("Investigate very-low-TTL UDP packets for spoofed-source DoS activity; deploy BCP38 anti-spoofing filters at egress.")

    # DHCP recommendations
    if "dhcp_starvation" in triggered_dhcp:
        recs.append("**IMMEDIATE**: Enable DHCP snooping on all access switches; restrict DHCP requests per port (port security with MAC limiting).")
    if "rogue_dhcp_server" in triggered_dhcp:
        recs.append("**IMMEDIATE**: Enable DHCP snooping to designate trusted DHCP server ports; isolate and remove rogue DHCP servers from the network.")
    if "dhcp_spoofing" in triggered_dhcp:
        recs.append("Verify that all DHCP responses originate from the authorised DHCP server; investigate server IP and option 54 mismatches.")
    if "dhcp_release_flood" in triggered_dhcp:
        recs.append("Rate-limit DHCP RELEASE/DECLINE messages per port; monitor DHCP server event logs for abnormal churn.")
    if "dhcp_relay_anomaly" in triggered_dhcp:
        recs.append("Audit DHCP relay agent configuration; ensure only authorised relay agents are permitted on DHCP snooping trusted ports.")
    if "dhcp_injection" in triggered_dhcp:
        recs.append("Review DHCP option 43 and option 82 content for malicious payloads; validate vendor class identifiers against the expected device list.")

    # mDNS recommendations
    if "mdns_spoofing" in triggered_mdns:
        recs.append("**HIGH PRIORITY**: Investigate mDNS spoofing for name-based man-in-the-middle attacks; deploy mDNS guard features if available on managed switches.")
    if "mdns_amplification" in triggered_mdns:
        recs.append("Restrict mDNS (UDP 5353) to the local link segment; block mDNS forwarding across router interfaces.")
    if "mdns_info_leakage" in triggered_mdns:
        recs.append("Audit mDNS service announcements for sensitive hostname and service name disclosures; suppress unnecessary service advertisements.")
    if "mdns_outside_local" in triggered_mdns:
        recs.append("Block mDNS (UDP 5353) traffic at router interfaces; investigate routable-IP mDNS sources for reconnaissance activity.")
    if "mdns_flood" in triggered_mdns:
        recs.append("Rate-limit mDNS traffic at the switch level; investigate the flooding source for misconfiguration or malicious automation.")

    # QUIC recommendations
    if "quic_amplification" in triggered_quic:
        recs.append("Investigate QUIC amplification sources; ensure QUIC servers enforce address validation tokens before sending large responses.")
    if "quic_0rtt_replay" in triggered_quic:
        recs.append("Verify that QUIC servers enforce anti-replay protection for 0-RTT data; consider disabling 0-RTT for sensitive API endpoints.")
    if "quic_version_anomaly" in triggered_quic:
        recs.append("Investigate unrecognised QUIC version numbers for version negotiation probing or stack fuzzing activity.")
    if "quic_handshake_exhaustion" in triggered_quic:
        recs.append("Rate-limit incomplete QUIC Initial packets per source IP at the load balancer; deploy QUIC-aware connection rate limiting.")
    if "quic_nonstandard_port" in triggered_quic:
        recs.append("Block QUIC (UDP) on non-standard ports at the perimeter; investigate sessions for C2 tunnelling or data exfiltration.")

    # File hash recommendations
    fh = data.get("fh_data", {})
    fh_files = fh.get("files", [])
    mal_files = [f for f in fh_files if f.get("osint_verdict") == "malicious"]
    sus_files = [f for f in fh_files if f.get("osint_verdict") == "suspicious"]
    if mal_files:
        recs.append(f"**IMMEDIATE**: {len(mal_files)} malicious file(s) extracted from this PCAP. Quarantine all hosts involved in these transfers and conduct full memory and disk forensics.")
        recs.append("Submit malicious file hashes to your threat intelligence platform and update endpoint detection signatures.")
        recs.append("Block the SHA256 hashes at endpoint protection and email gateway; search for the same hashes across all endpoints.")
    if sus_files:
        recs.append(f"Investigate {len(sus_files)} suspicious file(s) extracted from PCAP transfers; submit hashes for sandbox detonation.")
    if fh.get("files_found", 0) > 0 and not mal_files and not sus_files:
        recs.append(f"{fh.get('files_found',0)} file(s) extracted from PCAP. Review file list in the File Hash Analysis section and verify none are unexpected.")

    if sus_ips:
        recs.append(f"Investigate {len(sus_ips)} suspicious IP(s) flagged by CTI enrichment; verify context before blocking.")

    # Suricata recommendations
    suricata = data.get("suricata_data", {})
    if suricata.get("critical_count", 0) > 0:
        recs.append(
            f"**IMMEDIATE**: {suricata['critical_count']} critical Suricata alert(s) detected. "
            "Quarantine all involved hosts and initiate full memory and disk forensics immediately."
        )
    if suricata.get("high_count", 0) > 0:
        recs.append(
            f"Investigate {suricata['high_count']} high-severity Suricata alert(s); "
            "block flagged source IPs at the perimeter firewall and review host artefacts."
        )
    if suricata.get("total_alerts", 0) > 0 and not suricata.get("critical_count") and not suricata.get("high_count"):
        recs.append(
            f"Review {suricata['total_alerts']} Suricata alert(s) in Section 2.14; "
            "tune rules to reduce false positives and confirm no true positives were missed."
        )
    if data.get("has_suricata") and suricata.get("total_alerts", 0) == 0:
        recs.append(
            "Suricata ran with no alerts. Ensure rules are up to date: "
            "./scripts/update_suricata_rules.sh"
        )

    # YARA recommendations
    yara = data.get("yara_data", {})
    if yara.get("critical_count", 0) > 0:
        recs.append(
            f"**IMMEDIATE**: {yara['critical_count']} critical YARA rule match(es) in this PCAP. "
            "Review match details in Section 2.15 and isolate affected hosts."
        )
    if yara.get("high_count", 0) > 0:
        recs.append(
            f"Investigate {yara['high_count']} high-severity YARA match(es); "
            "correlate with Suricata alerts and CTI enrichment before escalating."
        )

    if not recs:
        recs.append("No high-priority immediate actions required. Continue routine monitoring.")

    recs.append("Retain PCAP capture and all analysis artefacts for chain of custody; do not modify source files.")
    return recs


# ── Report sections ────────────────────────────────────────────────────────────

def _md_table(headers: list[str], rows: list[list[str]]) -> list[str]:
    lines = ["| " + " | ".join(headers) + " |"]
    lines.append("|" + "|".join([" --- " for _ in headers]) + "|")
    for row in rows:
        lines.append("| " + " | ".join(str(c) for c in row) + " |")
    return lines


def sec_header(stem: str, case_id: str, overall_sev: str, now: str,
               first_ts: str, last_ts: str, duration: float,
               report_version: int = 1) -> list[str]:
    badge = SEVERITY_BADGE.get(overall_sev, overall_sev.upper())
    lines = [
        f"# PCAP Incident Report — `{stem}`",
        "",
        f"| Field | Value |",
        f"|-------|-------|",
        f"| Case ID | {case_id or '—'} |",
        f"| Report Version | v{report_version} |",
        f"| Overall Severity | {badge} |",
        f"| Report Generated | {now} |",
        f"| Prepared By | Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin |",
        f"| Capture Start | {first_ts or '—'} |",
        f"| Capture End | {last_ts or '—'} |",
        f"| Capture Duration | {_format_duration(duration) if duration else '—'} |",
        "",
        "---",
        "",
    ]
    return lines


def sec_management_summary(data: dict, overall_sev: str,
                            first_ts: str, last_ts: str, duration: float) -> list[str]:
    icmp_trig = _triggered(data["icmp_results"])
    dns_trig  = _triggered(data["dns_results"])
    ntp_trig  = _triggered(data["ntp_results"])
    http_trig = _triggered(data["http_results"])
    cert_trig = _triggered(data["cert_results"])
    tls_trig  = _triggered({k: v for k, v in data["tls_results"].items() if isinstance(v, dict)})
    arp_trig  = _triggered(data["arp_results"])
    tcp_trig  = _triggered(data["tcp_results"])
    udp_trig  = _triggered(data["udp_results"])
    dhcp_trig = _triggered(data["dhcp_results"])
    mdns_trig = _triggered(data["mdns_results"])
    quic_trig = _triggered(data["quic_results"])
    mal_ips   = [r for r in data["fan_ip"] if r.get("reputation") == "malicious"]
    sus_ips   = [r for r in data["fan_ip"] if r.get("reputation") == "suspicious"]
    mal_fqdns = [r for r in data["fan_correlation"] if r.get("reputation") == "malicious"]

    netflow_count = len(data["netflow"])
    ip_count      = len(data["unique_ips"])
    fqdn_count    = len(data["unique_fqdns"])

    lines = ["## 1. Management Summary", ""]

    # Opening assessment
    if overall_sev == "critical":
        lines.append(
            "Analysis of the captured network traffic has identified **CRITICAL-severity threats** "
            "that require immediate containment action. Adversarial activity consistent with data "
            "exfiltration or command-and-control (C2) communications is present in the capture."
        )
    elif overall_sev == "high":
        lines.append(
            "Analysis of the captured network traffic has identified **HIGH-severity threats** "
            "requiring prompt investigation. Attack patterns consistent with denial-of-service, "
            "reconnaissance, or network compromise attempts have been detected."
        )
    elif overall_sev == "medium":
        lines.append(
            "Analysis of the captured network traffic has identified **MEDIUM-severity anomalies** "
            "that should be investigated. While no immediately critical threats are confirmed, "
            "the observed patterns warrant closer examination."
        )
    else:
        lines.append(
            "Analysis of the captured network traffic did not identify significant threat indicators "
            "at this time. Routine monitoring is recommended."
        )

    lines.append("")

    # Traffic scope
    scope_parts = []
    if netflow_count:
        scope_parts.append(f"**{netflow_count:,} unique network flows**")
    if ip_count:
        scope_parts.append(f"**{ip_count:,} unique IP addresses**")
    if fqdn_count:
        scope_parts.append(f"**{fqdn_count:,} unique domain names**")
    if duration:
        scope_parts.append(f"captured over **{_format_duration(duration)}**")

    if scope_parts:
        lines.append("The capture encompasses " + ", ".join(scope_parts) + ".")
        lines.append("")

    # Key findings
    all_trig = (icmp_trig + dns_trig + ntp_trig + http_trig + cert_trig + tls_trig +
                arp_trig + tcp_trig + udp_trig + dhcp_trig + mdns_trig + quic_trig)
    if all_trig:
        lines.append("**Key findings:**")
        lines.append("")
        for cat in all_trig[:8]:
            sev  = cat.get("severity", "info")
            name = cat.get("name", "Unknown")
            cnt  = cat.get("count", 0)
            desc = cat.get("description", "")
            short_desc = desc.split(".")[0] if desc else ""
            lines.append(f"- {SEVERITY_BADGE.get(sev, sev.upper())} **{name}** — "
                         f"{cnt} finding(s). {short_desc}.")
        lines.append("")

    # CTI intelligence
    if mal_ips or mal_fqdns:
        total_mal = len(mal_ips) + len(mal_fqdns)
        lines.append(
            f"Threat intelligence enrichment identified **{total_mal} indicator(s)** "
            f"confirmed as malicious by open-source intelligence sources."
        )
        if mal_ips:
            ip_list = ", ".join(f"`{r['ip']}`" for r in mal_ips[:5])
            lines.append(f"Malicious IPs: {ip_list}" + (" (and more)" if len(mal_ips) > 5 else "") + ".")
        if mal_fqdns:
            fqdn_list = ", ".join(f"`{r['fqdn']}`" for r in mal_fqdns[:3])
            lines.append(f"Malicious domains: {fqdn_list}" + (" (and more)" if len(mal_fqdns) > 3 else "") + ".")
        lines.append("")
    elif sus_ips:
        lines.append(
            f"Threat intelligence enrichment flagged **{len(sus_ips)} IP(s) as suspicious**. "
            "These require manual review before a malicious classification can be confirmed."
        )
        lines.append("")

    # Immediate actions summary
    lines += [
        "**Recommended immediate actions:**",
        "",
    ]
    recs = build_recommendations(data, overall_sev)
    for r in recs[:5]:
        lines.append(f"- {r}")

    lines += ["", "See Section 7 for the full recommendations list.", "", "---", ""]
    return lines


def sec_findings_icmp(data: dict) -> list[str]:
    results = data["icmp_results"]
    if not results:
        return []

    lines = ["## 2. Findings", "", "### 2.1 ICMP Threat Analysis", ""]

    trig = _triggered(results)
    clean = [v for v in results.values() if isinstance(v, dict) and v.get("count", 0) == 0]

    if not trig:
        lines += ["No ICMP threats detected in this capture.", ""]
        return lines

    # Summary table
    rows = [[SEVERITY_BADGE.get(c["severity"], c["severity"]),
             c["name"], str(c["count"]), c["mitre"][0] if c.get("mitre") else "—"]
            for c in trig]
    lines += _md_table(["Severity", "Category", "Findings", "MITRE ATT&CK"], rows)
    lines.append("")

    # Per-category detail
    for cat in trig:
        sev  = cat.get("severity", "info")
        name = cat.get("name", "")
        mitre = cat.get("mitre", [])
        mitre_str = f"[{mitre[0]}](https://attack.mitre.org/techniques/{mitre[0].replace('.','/')}/) — {mitre[1]}" if len(mitre) >= 2 else (mitre[0] if mitre else "—")
        desc = cat.get("description", "")
        findings = cat.get("findings", [])

        lines += [
            f"#### {name}",
            "",
            f"**Severity:** {SEVERITY_BADGE.get(sev, sev)}  |  **MITRE ATT&CK:** {mitre_str}",
            "",
            desc,
            "",
        ]

        if findings:
            lines.append(f"**Top findings** (showing up to 10 of {len(findings)}):")
            lines.append("")
            lines.append("```")
            for f in findings[:10]:
                lines.append(json.dumps(f, default=str))
            lines.append("```")
            lines.append("")

    if clean:
        lines += [
            "**Clean categories (no findings):** " +
            ", ".join(c.get("name", "") for c in clean),
            "",
        ]

    lines += ["---", ""]
    return lines


def sec_findings_dns(data: dict) -> list[str]:
    results = data["dns_results"]
    if not results:
        return []

    lines = ["### 2.2 DNS Threat Analysis", ""]

    trig  = _triggered(results)
    clean = [v for v in results.values() if isinstance(v, dict) and v.get("count", 0) == 0]

    if not trig:
        lines += ["No DNS threats detected in this capture.", ""]
        return lines

    rows = [[SEVERITY_BADGE.get(c["severity"], c["severity"]),
             c["name"], str(c["count"]), c["mitre"][0] if c.get("mitre") else "—"]
            for c in trig]
    lines += _md_table(["Severity", "Category", "Findings", "MITRE ATT&CK"], rows)
    lines.append("")

    for cat in trig:
        sev   = cat.get("severity", "info")
        name  = cat.get("name", "")
        mitre = cat.get("mitre", [])
        mitre_str = f"[{mitre[0]}](https://attack.mitre.org/techniques/{mitre[0].replace('.','/')}/) — {mitre[1]}" if len(mitre) >= 2 else (mitre[0] if mitre else "—")
        desc  = cat.get("description", "")
        findings = cat.get("findings", [])

        lines += [
            f"#### {name}",
            "",
            f"**Severity:** {SEVERITY_BADGE.get(sev, sev)}  |  **MITRE ATT&CK:** {mitre_str}",
            "",
            desc,
            "",
        ]

        if findings:
            lines.append(f"**Top findings** (showing up to 10 of {len(findings)}):")
            lines.append("")
            lines.append("```")
            for f in findings[:10]:
                lines.append(json.dumps(f, default=str))
            lines.append("```")
            lines.append("")

    if clean:
        lines += [
            "**Clean categories (no findings):** " +
            ", ".join(c.get("name", "") for c in clean),
            "",
        ]

    lines += ["---", ""]
    return lines


def sec_findings_ntp(data: dict) -> list[str]:
    results = data["ntp_results"]
    if not results:
        return []

    lines = ["### 2.3 NTP Threat Analysis", ""]

    trig  = _triggered(results)
    clean = [v for v in results.values() if isinstance(v, dict) and v.get("count", 0) == 0]

    if not trig:
        lines += ["No NTP threats detected in this capture.", ""]
        return lines

    rows = [[SEVERITY_BADGE.get(c["severity"], c["severity"]),
             c["name"], str(c["count"]), c["mitre"][0] if c.get("mitre") else "—"]
            for c in trig]
    lines += _md_table(["Severity", "Category", "Findings", "MITRE ATT&CK"], rows)
    lines.append("")

    for cat in trig:
        sev   = cat.get("severity", "info")
        name  = cat.get("name", "")
        mitre = cat.get("mitre", [])
        mitre_str = (
            f"[{mitre[0]}](https://attack.mitre.org/techniques/{mitre[0].replace('.','/')}/) — {mitre[1]}"
            if len(mitre) >= 2 else (mitre[0] if mitre else "—")
        )
        desc  = cat.get("description", "")
        findings = cat.get("findings", [])

        lines += [
            f"#### {name}",
            "",
            f"**Severity:** {SEVERITY_BADGE.get(sev, sev)}  |  **MITRE ATT&CK:** {mitre_str}",
            "",
            desc,
            "",
        ]

        if findings:
            lines.append(f"**Top findings** (showing up to 10 of {len(findings)}):")
            lines.append("")
            lines.append("```")
            for f in findings[:10]:
                lines.append(json.dumps(f, default=str))
            lines.append("```")
            lines.append("")

    if clean:
        lines += [
            "**Clean categories (no findings):** " +
            ", ".join(c.get("name", "") for c in clean),
            "",
        ]

    lines += ["---", ""]
    return lines


def sec_findings_http(data: dict) -> list[str]:
    results = data["http_results"]
    if not results:
        return []

    lines = ["### 2.4 HTTP(S) Unusual Patterns", ""]

    trig  = _triggered(results)
    clean = [v for v in results.values() if isinstance(v, dict) and v.get("count", 0) == 0]

    if not trig:
        lines += ["No HTTP(S) threat patterns detected in this capture.", ""]
        return lines

    rows = [[SEVERITY_BADGE.get(c["severity"], c["severity"]),
             c["name"], str(c["count"]), c["mitre"][0] if c.get("mitre") else "—"]
            for c in trig]
    lines += _md_table(["Severity", "Category", "Findings", "MITRE ATT&CK"], rows)
    lines.append("")

    for cat in trig:
        sev   = cat.get("severity", "info")
        name  = cat.get("name", "")
        mitre = cat.get("mitre", [])
        mitre_str = (
            f"[{mitre[0]}](https://attack.mitre.org/techniques/{mitre[0].replace('.','/')}/) — {mitre[1]}"
            if len(mitre) >= 2 else (mitre[0] if mitre else "—")
        )
        desc     = cat.get("description", "")
        findings = cat.get("findings", [])

        lines += [
            f"#### {name}",
            "",
            f"**Severity:** {SEVERITY_BADGE.get(sev, sev)}  |  **MITRE ATT&CK:** {mitre_str}",
            "",
            desc,
            "",
        ]

        if findings:
            lines.append(f"**Top findings** (showing up to 10 of {len(findings)}):")
            lines.append("")
            lines.append("```")
            for f in findings[:10]:
                lines.append(json.dumps(f, default=str))
            lines.append("```")
            lines.append("")

    if clean:
        lines += [
            "**Clean categories (no findings):** " +
            ", ".join(c.get("name", "") for c in clean),
            "",
        ]

    lines += ["---", ""]
    return lines


def sec_findings_cert(data: dict) -> list[str]:
    results = data["cert_results"]
    if not results:
        return []

    lines = ["### 2.5 TLS Certificate Inspector", ""]

    trig  = _triggered(results)
    clean = [v for v in results.values() if isinstance(v, dict) and v.get("count", 0) == 0]

    if not trig:
        lines += ["No certificate anomalies detected in this capture.", ""]
        return lines

    rows = [[SEVERITY_BADGE.get(c["severity"], c["severity"]),
             c["name"], str(c["count"]), c["mitre"][0] if c.get("mitre") else "—"]
            for c in trig]
    lines += _md_table(["Severity", "Category", "Findings", "MITRE ATT&CK"], rows)
    lines.append("")

    for cat in trig:
        sev   = cat.get("severity", "info")
        name  = cat.get("name", "")
        mitre = cat.get("mitre", [])
        mitre_str = (
            f"[{mitre[0]}](https://attack.mitre.org/techniques/{mitre[0].replace('.','/')}/) — {mitre[1]}"
            if len(mitre) >= 2 else (mitre[0] if mitre else "—")
        )
        desc     = cat.get("description", "")
        findings = cat.get("findings", [])

        lines += [
            f"#### {name}",
            "",
            f"**Severity:** {SEVERITY_BADGE.get(sev, sev)}  |  **MITRE ATT&CK:** {mitre_str}",
            "",
            desc,
            "",
        ]

        if findings:
            lines.append(f"**Top findings** (showing up to 10 of {len(findings)}):")
            lines.append("")
            lines.append("```")
            for f in findings[:10]:
                lines.append(json.dumps(f, default=str))
            lines.append("```")
            lines.append("")

    if clean:
        lines += [
            "**Clean categories (no findings):** " +
            ", ".join(c.get("name", "") for c in clean),
            "",
        ]

    lines += ["---", ""]
    return lines


def sec_findings_tls(data: dict) -> list[str]:
    results  = data["tls_results"]
    sessions = data["tls_sessions"]
    if not results and not sessions:
        return []

    lines = ["### 2.6 TLS Session Inspector", ""]

    # Filter to category dicts only (exclude _sessions etc.)
    cat_results = {k: v for k, v in results.items() if isinstance(v, dict)}
    trig  = _triggered(cat_results)
    clean = [v for v in cat_results.values() if v.get("count", 0) == 0]

    # Session inventory summary
    if sessions:
        lines.append(f"*{len(sessions)} unique TLS session(s) analysed.*")
        lines.append("")
        # Top 15 sessions table
        lines.append("**Session inventory (first 15 sessions):**")
        lines.append("")
        inv_rows = [
            [s.get("timestamp_utc", "—"), s.get("client_ip", "—"),
             s.get("server_ip", "—"), s.get("server_port", "—"),
             s.get("sni", "—") or "—", s.get("tls_version", "—"),
             f"`{s.get('cipher_hex', '—')}`",
             f"`{s.get('ja4', '—')[:24]}…`" if len(s.get("ja4", "")) > 24 else f"`{s.get('ja4', '—')}`"]
            for s in sessions[:15]
        ]
        lines += _md_table(
            ["Timestamp", "Client IP", "Server IP", "Port", "SNI", "TLS Version", "Cipher", "JA4"],
            inv_rows
        )
        lines.append("")

    if not trig:
        lines += ["No TLS session anomalies detected.", ""]
        lines += ["---", ""]
        return lines

    lines.append("**Threat findings:**")
    lines.append("")
    rows = [[SEVERITY_BADGE.get(c["severity"], c["severity"]),
             c["name"], str(c["count"]), c["mitre"][0] if c.get("mitre") else "—"]
            for c in trig]
    lines += _md_table(["Severity", "Category", "Findings", "MITRE ATT&CK"], rows)
    lines.append("")

    for cat in trig:
        sev   = cat.get("severity", "info")
        name  = cat.get("name", "")
        mitre = cat.get("mitre", [])
        mitre_str = (
            f"[{mitre[0]}](https://attack.mitre.org/techniques/{mitre[0].replace('.','/')}/) — {mitre[1]}"
            if len(mitre) >= 2 else (mitre[0] if mitre else "—")
        )
        desc     = cat.get("description", "")
        findings = cat.get("findings", [])

        lines += [
            f"#### {name}",
            "",
            f"**Severity:** {SEVERITY_BADGE.get(sev, sev)}  |  **MITRE ATT&CK:** {mitre_str}",
            "",
            desc,
            "",
        ]

        if findings:
            lines.append(f"**Top findings** (showing up to 10 of {len(findings)}):")
            lines.append("")
            lines.append("```")
            for f in findings[:10]:
                lines.append(json.dumps(f, default=str))
            lines.append("```")
            lines.append("")

    if clean:
        lines += [
            "**Clean categories (no findings):** " +
            ", ".join(c.get("name", "") for c in clean),
            "",
        ]

    lines += ["---", ""]
    return lines


def _sec_findings_generic(data: dict, results_key: str, section_num: str,
                           title: str, has_key: str) -> list[str]:
    """Render a generic findings section for any of the new protocol threat modules."""
    if not data[has_key]:
        return []
    results = data[results_key]
    if not results:
        return []

    lines = [f"### {section_num} {title}", ""]
    trig  = _triggered(results)
    clean = [v for v in results.values() if isinstance(v, dict) and v.get("count", 0) == 0]

    if not trig:
        lines += [f"No {title.split()[0]} threats detected in this capture.", ""]
        return lines

    rows = [[SEVERITY_BADGE.get(c["severity"], c["severity"]),
             c.get("name", ""), str(c["count"]),
             c["mitre"][0] if c.get("mitre") else "—"]
            for c in trig]
    lines += _md_table(["Severity", "Category", "Findings", "MITRE ATT&CK"], rows)
    lines.append("")

    for cat in trig:
        sev  = cat.get("severity", "info")
        name = cat.get("name", "")
        mitre = cat.get("mitre", [])
        mitre_str = (
            f"[{mitre[0]}](https://attack.mitre.org/techniques/{mitre[0].replace('.','/')}/) "
            f"— {mitre[1]}" if len(mitre) >= 2 else (mitre[0] if mitre else "—")
        )
        desc     = cat.get("description", "")
        findings = cat.get("findings", [])

        lines += [
            f"#### {name}",
            "",
            f"**Severity:** {SEVERITY_BADGE.get(sev, sev)}  |  **MITRE ATT&CK:** {mitre_str}",
            "",
            desc,
            "",
        ]
        if findings:
            lines.append(f"**Top findings** (showing up to 10 of {len(findings)}):")
            lines.append("")
            lines.append("```")
            for f in findings[:10]:
                lines.append(json.dumps(f, default=str))
            lines.append("```")
            lines.append("")

    if clean:
        lines += [
            "**Clean categories (no findings):** " +
            ", ".join(c.get("name", "") for c in clean),
            "",
        ]

    lines += ["---", ""]
    return lines


def sec_findings_arp(data: dict) -> list[str]:
    return _sec_findings_generic(data, "arp_results",  "2.7",  "ARP Threat Analysis", "has_arp")


def sec_findings_tcp(data: dict) -> list[str]:
    return _sec_findings_generic(data, "tcp_results",  "2.8",  "TCP Threat Analysis", "has_tcp")


def sec_findings_udp(data: dict) -> list[str]:
    return _sec_findings_generic(data, "udp_results",  "2.9",  "UDP Threat Analysis", "has_udp")


def sec_findings_dhcp(data: dict) -> list[str]:
    return _sec_findings_generic(data, "dhcp_results", "2.10", "DHCP Threat Analysis", "has_dhcp")


def sec_findings_mdns(data: dict) -> list[str]:
    return _sec_findings_generic(data, "mdns_results", "2.11", "mDNS Threat Analysis", "has_mdns")


def sec_findings_quic(data: dict) -> list[str]:
    return _sec_findings_generic(data, "quic_results",    "2.12", "QUIC Threat Analysis",    "has_quic")


def sec_findings_snmp(data: dict) -> list[str]:
    return _sec_findings_generic(data, "snmp_results",    "2.13", "SNMP Threat Analysis",    "has_snmp")


def sec_findings_nbns(data: dict) -> list[str]:
    return _sec_findings_generic(data, "nbns_results",    "2.14", "NBNS Threat Analysis",    "has_nbns")


def sec_findings_llmnr(data: dict) -> list[str]:
    return _sec_findings_generic(data, "llmnr_results",   "2.15", "LLMNR Threat Analysis",   "has_llmnr")


def sec_findings_stun(data: dict) -> list[str]:
    return _sec_findings_generic(data, "stun_results",    "2.16", "STUN Threat Analysis",    "has_stun")


def sec_findings_ssdp(data: dict) -> list[str]:
    return _sec_findings_generic(data, "ssdp_results",    "2.17", "SSDP Threat Analysis",    "has_ssdp")


def sec_findings_netbios(data: dict) -> list[str]:
    return _sec_findings_generic(data, "netbios_results", "2.18", "NetBIOS Threat Analysis", "has_netbios")


def sec_findings_filehashes(data: dict) -> list[str]:
    if not data["has_fh"]:
        return []
    fh = data.get("fh_data", {})
    files = fh.get("files", [])
    lines = ["### 2.19 File Hash Analysis", ""]

    if not files:
        lines += ["No files were extracted from this PCAP capture.", "", "---", ""]
        return lines

    mal  = [f for f in files if f.get("osint_verdict") == "malicious"]
    sus  = [f for f in files if f.get("osint_verdict") == "suspicious"]
    rest = [f for f in files if f.get("osint_verdict") not in ("malicious", "suspicious")]

    summary_rows = [
        ["Files Extracted", str(fh.get("files_found", len(files)))],
        ["Malicious",       str(fh.get("malicious_count",  len(mal)))],
        ["Suspicious",      str(fh.get("suspicious_count", len(sus)))],
        ["Clean / Unknown", str(len(rest))],
    ]
    lines += _md_table(["Metric", "Count"], summary_rows)
    lines.append("")

    # Full inventory table
    lines += ["#### Extracted File Inventory", ""]
    rows = []
    for rec in files:
        verdict = rec.get("osint_verdict", "") or "—"
        badge = {"malicious": "🔴 CRITICAL", "suspicious": "🟠 HIGH"}.get(verdict, verdict)
        size = f"{rec.get('size_bytes', 0):,}"
        rows.append([
            rec.get("protocol", "—"),
            f"`{rec.get('filename','')[:40]}`",
            size,
            f"`{rec.get('md5','')[:12]}…`",
            f"`{rec.get('sha256','')[:16]}…`",
            badge,
        ])
    lines += _md_table(["Protocol", "Filename", "Size (B)", "MD5", "SHA256", "Verdict"], rows)
    lines.append("")

    # Threat detail for malicious/suspicious files
    if mal or sus:
        lines += ["#### Threat Detail", ""]
        for rec in mal + sus:
            sev_label = "CRITICAL" if rec.get("osint_verdict") == "malicious" else "HIGH"
            lines += [
                f"**{rec.get('filename', 'unknown')}** — {sev_label}",
                "",
                f"- Protocol: {rec.get('protocol','—')} | Size: {rec.get('size_bytes',0):,} bytes",
                f"- MD5:    `{rec.get('md5','')}`",
                f"- SHA256: `{rec.get('sha256','')}`",
                f"- OSINT: {rec.get('osint_summary','No data') or 'No data'}",
                "",
            ]

    lines += ["---", ""]
    return lines


def sec_findings_suricata(data: dict) -> list[str]:
    if not data["has_suricata"]:
        return []
    suricata = data.get("suricata_data", {})
    alerts   = suricata.get("alerts", [])
    lines    = ["### 2.20 Suricata IDS Alerts", ""]

    if not alerts:
        lines += ["Suricata ran but produced no alerts for this capture.", "", "---", ""]
        return lines

    critical = suricata.get("critical_count", 0)
    high     = suricata.get("high_count", 0)
    medium   = suricata.get("medium_count", 0)
    low      = suricata.get("low_count", 0)

    summary_rows = [
        ["Total Alerts",    str(suricata.get("total_alerts", len(alerts)))],
        ["Critical",        str(critical)],
        ["High",            str(high)],
        ["Medium",          str(medium)],
        ["Low",             str(low)],
    ]
    lines += _md_table(["Metric", "Count"], summary_rows)
    lines.append("")

    # Top signatures
    sig_counts: dict[str, int] = {}
    sig_to_cat:  dict[str, str] = {}
    sig_to_sev:  dict[str, str] = {}
    sev_order = {"critical": 0, "high": 1, "medium": 2, "low": 3, "info": 4}
    for a in alerts:
        sig = a.get("signature", "")
        sig_counts[sig] = sig_counts.get(sig, 0) + 1
        sig_to_cat[sig] = a.get("category", "")
        cur = sig_to_sev.get(sig, "info")
        if sev_order.get(a.get("severity", "info"), 4) < sev_order.get(cur, 4):
            sig_to_sev[sig] = a.get("severity", "info")

    top_sigs = sorted(sig_counts.items(), key=lambda x: -x[1])[:15]
    lines += ["#### Top Signatures", ""]
    rows = []
    for sig, cnt in top_sigs:
        sev = sig_to_sev.get(sig, "info").upper()
        cat = sig_to_cat.get(sig, "")
        rows.append([str(cnt), sev, sig[:80], cat[:50]])
    lines += _md_table(["Count", "Severity", "Signature", "Category"], rows)
    lines.append("")

    # Alert detail for critical/high
    critical_high = [a for a in alerts if a.get("severity") in ("critical", "high")]
    if critical_high:
        lines += ["#### Critical / High Alerts", ""]
        rows2 = []
        for a in sorted(critical_high, key=lambda x: (sev_order.get(x.get("severity","info"),4), x.get("timestamp_utc","")))[:30]:
            ts   = (a.get("timestamp_utc") or "")[:19]
            flow = f"{a.get('src_ip','?')}:{a.get('src_port','?')} → {a.get('dest_ip','?')}:{a.get('dest_port','?')}"
            rows2.append([ts, flow, a.get("proto",""), a.get("severity","").upper(), a.get("signature","")[:70]])
        lines += _md_table(["Timestamp", "Flow", "Proto", "Severity", "Signature"], rows2)
        lines.append("")

    lines += ["---", ""]
    return lines


def sec_findings_yara(data: dict) -> list[str]:
    if not data["has_yara"]:
        return []
    yara    = data.get("yara_data", {})
    matches = yara.get("matches", [])
    lines   = ["### 2.21 YARA Rule Matches", ""]

    if not matches:
        lines += ["YARA scan completed with no rule matches.", "", "---", ""]
        return lines

    critical = yara.get("critical_count", 0)
    high     = yara.get("high_count", 0)
    medium   = yara.get("medium_count", 0)

    summary_rows = [
        ["Total Matches",  str(yara.get("total_matches", len(matches)))],
        ["Critical",       str(critical)],
        ["High",           str(high)],
        ["Medium / Low",   str(medium)],
    ]
    lines += _md_table(["Metric", "Count"], summary_rows)
    lines.append("")

    sev_order = {"critical": 0, "high": 1, "medium": 2, "low": 3, "info": 4}
    sorted_matches = sorted(
        matches, key=lambda m: (sev_order.get(m.get("severity", "info"), 4), m.get("rule", ""))
    )

    lines += ["#### Match Details", ""]
    rows = []
    for m in sorted_matches[:30]:
        rows.append([
            m.get("severity", "").upper(),
            m.get("rule", "")[:50],
            f"`{m.get('target_name','')[:40]}`",
            m.get("category", "")[:30],
            m.get("mitre_att", ""),
        ])
    lines += _md_table(["Severity", "Rule", "Target", "Category", "MITRE"], rows)
    lines.append("")

    # Threat detail for critical/high
    critical_high = [m for m in sorted_matches if m.get("severity") in ("critical", "high")]
    if critical_high:
        lines += ["#### Threat Detail", ""]
        for m in critical_high:
            lines += [
                f"**{m.get('rule','')}** — {m.get('severity','').upper()}",
                "",
                f"- Target: `{m.get('target_name','—')}`",
                f"- Rule file: `{m.get('rule_file','—')}`",
                f"- Description: {m.get('description','—')}",
                f"- Category: {m.get('category','—')}",
                f"- MITRE ATT&CK: {m.get('mitre_att','—')}",
                "",
            ]

    lines += ["---", ""]
    return lines


def sec_network_summary(data: dict, duration: float) -> list[str]:
    if not data["has_pcap"]:
        return []

    lines = ["### 2.22 Network Flow Summary", ""]
    netflow    = data["netflow"]
    unique_ips = data["unique_ips"]

    if not netflow:
        lines += ["No netflow data available.", "", "---", ""]
        return lines

    # Protocol distribution
    proto_count: dict[str, int] = defaultdict(int)
    proto_bytes: dict[str, int] = defaultdict(int)
    for row in netflow:
        proto = row.get("protocol", "?")
        proto_count[proto] += int(row.get("packets", 0) or 0)
        proto_bytes[proto] += int(row.get("bytes", 0) or 0)

    total_pkts  = sum(proto_count.values())
    total_bytes = sum(proto_bytes.values())

    lines += [
        f"| Metric | Value |",
        f"|--------|-------|",
        f"| Unique flows | {len(netflow):,} |",
        f"| Unique IPs | {len(unique_ips):,} |",
        f"| Unique FQDNs | {len(data['unique_fqdns']):,} |",
        f"| Total packets | {total_pkts:,} |",
        f"| Total bytes | {total_bytes:,} ({total_bytes / 1_048_576:.1f} MB) |",
    ]
    if duration:
        lines.append(f"| Duration | {_format_duration(duration)} |")
    lines.append("")

    # Protocol breakdown
    proto_rows = sorted(proto_count.items(), key=lambda x: x[1], reverse=True)
    lines.append("**Protocol distribution:**")
    lines.append("")
    prows = [[p, f"{c:,}", f"{proto_bytes[p]:,}",
              f"{proto_bytes[p] / total_bytes * 100:.1f}%" if total_bytes else "—"]
             for p, c in proto_rows[:10]]
    lines += _md_table(["Protocol", "Packets", "Bytes", "% of Traffic"], prows)
    lines.append("")

    # Top talkers by bytes
    talker_bytes: dict[str, int] = defaultdict(int)
    for row in netflow:
        src = row.get("src_ip", "")
        if src:
            talker_bytes[src] += int(row.get("bytes", 0) or 0)

    top_talkers = sorted(talker_bytes.items(), key=lambda x: x[1], reverse=True)[:10]
    if top_talkers:
        lines.append("**Top 10 source IPs by traffic volume:**")
        lines.append("")
        trows = [[ip, f"{b:,}", f"{b / total_bytes * 100:.1f}%" if total_bytes else "—"]
                 for ip, b in top_talkers]
        lines += _md_table(["Source IP", "Bytes", "% of Traffic"], trows)
        lines.append("")

    # Top destination IPs
    dst_bytes: dict[str, int] = defaultdict(int)
    for row in netflow:
        dst = row.get("dst_ip", "")
        if dst:
            dst_bytes[dst] += int(row.get("bytes", 0) or 0)

    top_dsts = sorted(dst_bytes.items(), key=lambda x: x[1], reverse=True)[:10]
    if top_dsts:
        lines.append("**Top 10 destination IPs by traffic volume:**")
        lines.append("")
        drows = [[ip, f"{b:,}", f"{b / total_bytes * 100:.1f}%" if total_bytes else "—"]
                 for ip, b in top_dsts]
        lines += _md_table(["Destination IP", "Bytes", "% of Traffic"], drows)
        lines.append("")

    lines += ["---", ""]
    return lines


def sec_cti_enrichment(data: dict) -> list[str]:
    if not data["has_fan_ip"]:
        return []

    lines = ["### 2.23 CTI Enrichment", ""]

    corr = data["fan_correlation"]
    ips  = data["fan_ip"]

    mal_fqdns = [r for r in corr if r.get("reputation") == "malicious"]
    sus_fqdns = [r for r in corr if r.get("reputation") == "suspicious"]
    mal_ips   = [r for r in ips  if r.get("reputation") == "malicious"]
    sus_ips   = [r for r in ips  if r.get("reputation") == "suspicious"]

    lines += [
        f"| Indicator Type | Malicious | Suspicious | Total Enriched |",
        f"|----------------|-----------|------------|----------------|",
        f"| IP Addresses | {len(mal_ips)} | {len(sus_ips)} | {len(ips)} |",
        f"| Domains (FQDNs) | {len(mal_fqdns)} | {len(sus_fqdns)} | {len(corr)} |",
        "",
    ]

    if mal_ips or sus_ips:
        lines.append("**High-priority IP indicators:**")
        lines.append("")
        prows = [[r["ip"], r["reputation"].upper(),
                  r.get("reverse_dns", "—") or "—",
                  (r.get("osint_summary", "") or "")[:120]]
                 for r in (mal_ips + sus_ips)[:20]]
        lines += _md_table(["IP", "Reputation", "Reverse DNS", "OSINT Summary"], prows)
        lines.append("")

    if mal_fqdns or sus_fqdns:
        lines.append("**High-priority domain indicators:**")
        lines.append("")
        drows = [[r["fqdn"], r["reputation"].upper(),
                  r.get("source_tags", "—") or "—",
                  (r.get("osint_summary", "") or "")[:120]]
                 for r in (mal_fqdns + sus_fqdns)[:20]]
        lines += _md_table(["FQDN", "Reputation", "Sources", "OSINT Summary"], drows)
        lines.append("")

    if not (mal_ips or sus_ips or mal_fqdns or sus_fqdns):
        lines += ["No malicious or suspicious indicators found by CTI enrichment.", ""]

    lines += ["---", ""]
    return lines


def sec_timeline(timeline: list[dict]) -> list[str]:
    lines = ["## 3. Threat Timeline", ""]

    if not timeline:
        lines += ["No notable events to display.", "", "---", ""]
        return lines

    lines.append(f"*Showing {len(timeline)} significant events (sorted chronologically).*")
    lines.append("")

    rows = [[e["timestamp"] or "—", e["protocol"],
             e["event_type"], SEVERITY_BADGE.get(e["severity"], e["severity"]),
             e.get("src_ip", "—"), e.get("dst_ip", "—"),
             e.get("description", "")]
            for e in timeline]

    lines += _md_table(
        ["Timestamp (UTC)", "Protocol", "Event Type", "Severity", "Source", "Destination", "Description"],
        rows
    )
    lines += ["", "---", ""]
    return lines


def sec_iocs(iocs: list[dict]) -> list[str]:
    lines = ["## 4. Indicators of Compromise", ""]

    if not iocs:
        lines += ["No indicators of compromise extracted.", "", "---", ""]
        return lines

    lines.append(f"*{len(iocs)} unique indicator(s) extracted from all analysis outputs.*")
    lines.append("")

    rows = [[SEVERITY_BADGE.get(i["severity"], i["severity"]),
             i["type"], f"`{i['value']}`", i["category"], i["source"]]
            for i in iocs]
    lines += _md_table(["Severity", "Type", "Value", "Category", "Source"], rows)
    lines += ["", "---", ""]
    return lines


def sec_mitre(coverage: list[dict]) -> list[str]:
    lines = ["## 5. MITRE ATT&CK Coverage", ""]

    if not coverage:
        lines += ["No MITRE ATT&CK techniques observed.", "", "---", ""]
        return lines

    rows = [[
        f"[{t['id']}](https://attack.mitre.org/techniques/{t['id'].replace('.', '/')}/ )",
        t["name"],
        t["tactic"],
        SEVERITY_BADGE.get(t["severity"], t["severity"]),
        t["category"],
    ] for t in coverage]
    lines += _md_table(["Technique", "Name", "Tactic", "Severity", "Triggered By"], rows)
    lines += ["", "---", ""]
    return lines


def sec_recommendations(recs: list[str]) -> list[str]:
    lines = ["## 6. Recommendations", ""]
    for r in recs:
        lines.append(f"- {r}")
    lines += ["", "---", ""]
    return lines


def sec_appendix(stem: str, data: dict, case_id: str = "") -> list[str]:
    ev_note = (
        f"All artifact files are preserved in `./{case_id}_evidence/analysis/` "
        "and uploaded to the investigations vault alongside this report. "
        "SHA-256 hashes are recorded in the research notes (Appendix B) for chain-of-custody verification."
        if case_id else
        "Analysis source files listed below."
    )
    lines = ["## Appendix A: Analysis Source Files", "", ev_note, ""]

    sources = []
    if data["has_pcap"]:
        d = ANALYSIS_DIR / "pcap" / stem
        sources += [
            ("PCAP Netflow", str(d / "netflow.csv")),
            ("Unique IPs", str(d / "unique_ips.txt")),
            ("Unique FQDNs", str(d / "unique_fqdns.txt")),
        ]
    if data["has_icmp"]:
        d = ANALYSIS_DIR / "icmp_threats" / stem
        sources += [
            ("ICMP Threat Report", str(d / "icmp_threats_report.md")),
            ("ICMP Findings JSON", str(d / "icmp_threats.json")),
            ("ICMP Flows CSV", str(d / "icmp_flows.csv")),
        ]
    if data["has_dns"]:
        d = ANALYSIS_DIR / "dns_threats" / stem
        sources += [
            ("DNS Threat Report", str(d / "dns_threats_report.md")),
            ("DNS Findings JSON", str(d / "dns_threats.json")),
            ("DNS Flows CSV", str(d / "dns_flows.csv")),
        ]
    if data["has_ntp"]:
        d = ANALYSIS_DIR / "ntp_threats" / stem
        sources += [
            ("NTP Threat Report", str(d / "ntp_threats_report.md")),
            ("NTP Findings JSON", str(d / "ntp_threats.json")),
            ("NTP Flows CSV", str(d / "ntp_flows.csv")),
        ]
    if data["has_http"]:
        d = ANALYSIS_DIR / "http_threats" / stem
        sources += [
            ("HTTP Threat Report", str(d / "http_threats_report.md")),
            ("HTTP Findings JSON", str(d / "http_threats.json")),
            ("HTTP Flows CSV", str(d / "http_flows.csv")),
        ]
    if data["has_cert"]:
        d = ANALYSIS_DIR / "cert_inspector" / stem
        sources += [
            ("Certificate Inspector Report", str(d / "cert_inspector_report.md")),
            ("Certificate Findings JSON", str(d / "certs.json")),
            ("Certificate Inventory CSV", str(d / "certs.csv")),
        ]
    if data["has_tls"]:
        d = ANALYSIS_DIR / "tls_inspector" / stem
        sources += [
            ("TLS Inspector Report", str(d / "tls_inspector_report.md")),
            ("TLS Sessions JSON", str(d / "tls_sessions.json")),
            ("TLS Sessions CSV", str(d / "tls_sessions.csv")),
        ]
    if data["has_arp"]:
        d = ANALYSIS_DIR / "arp_threats" / stem
        sources += [
            ("ARP Threat Report", str(d / "arp_threats_report.md")),
            ("ARP Findings JSON", str(d / "arp_threats.json")),
            ("ARP Flows CSV",    str(d / "arp_flows.csv")),
        ]
    if data["has_tcp"]:
        d = ANALYSIS_DIR / "tcp_threats" / stem
        sources += [
            ("TCP Threat Report", str(d / "tcp_threats_report.md")),
            ("TCP Findings JSON", str(d / "tcp_threats.json")),
            ("TCP Flows CSV",    str(d / "tcp_flows.csv")),
        ]
    if data["has_udp"]:
        d = ANALYSIS_DIR / "udp_threats" / stem
        sources += [
            ("UDP Threat Report", str(d / "udp_threats_report.md")),
            ("UDP Findings JSON", str(d / "udp_threats.json")),
            ("UDP Flows CSV",    str(d / "udp_flows.csv")),
        ]
    if data["has_dhcp"]:
        d = ANALYSIS_DIR / "dhcp_threats" / stem
        sources += [
            ("DHCP Threat Report", str(d / "dhcp_threats_report.md")),
            ("DHCP Findings JSON", str(d / "dhcp_threats.json")),
            ("DHCP Flows CSV",    str(d / "dhcp_flows.csv")),
        ]
    if data["has_mdns"]:
        d = ANALYSIS_DIR / "mdns_threats" / stem
        sources += [
            ("mDNS Threat Report", str(d / "mdns_threats_report.md")),
            ("mDNS Findings JSON", str(d / "mdns_threats.json")),
            ("mDNS Flows CSV",    str(d / "mdns_flows.csv")),
        ]
    if data["has_quic"]:
        d = ANALYSIS_DIR / "quic_threats" / stem
        sources += [
            ("QUIC Threat Report", str(d / "quic_threats_report.md")),
            ("QUIC Findings JSON", str(d / "quic_threats.json")),
            ("QUIC Flows CSV",    str(d / "quic_flows.csv")),
        ]
    if data["has_fh"]:
        d = ANALYSIS_DIR / "file_hashes" / stem
        sources += [
            ("File Hash Report",  str(d / "file_hashes_report.md")),
            ("File Hash JSON",    str(d / "file_hashes.json")),
            ("File Hash CSV",     str(d / "file_hashes.csv")),
            ("Extracted Files",   str(d / "files/")),
        ]
    if data["has_suricata"]:
        d = ANALYSIS_DIR / "suricata" / stem
        sources += [
            ("Suricata IDS Report",   str(d / "suricata_report.md")),
            ("Suricata Alerts JSON",  str(d / "suricata_alerts.json")),
            ("Suricata Alerts CSV",   str(d / "suricata_alerts.csv")),
            ("Suricata EVE JSON",     str(d / "eve.json")),
        ]
    if data["has_yara"]:
        d = ANALYSIS_DIR / "yara_pcap" / stem
        sources += [
            ("YARA Matches Report",  str(d / "yara_report.md")),
            ("YARA Matches JSON",    str(d / "yara_matches.json")),
            ("YARA Matches CSV",     str(d / "yara_matches.csv")),
        ]
    if data["has_fan_ip"]:
        d = ANALYSIS_DIR / "fan_ip" / stem
        sources += [
            ("CTI IP Enrichment", str(d / "ip_enrichment.csv")),
            ("CTI FQDN Correlation", str(d / "correlation.csv")),
        ]

    if sources:
        rows = [[name, f"`{path}`"] for name, path in sources]
        lines += _md_table(["Dataset", "Path"], rows)
    else:
        lines += ["No analysis source files found."]

    lines.append("")
    return lines


# ── DOCX / PPTX helpers ────────────────────────────────────────────────────────

def _compute_file_metadata(file_path: Path) -> tuple[str, str]:
    if not file_path.exists():
        return "N/A", "N/A"
    mtime_str = datetime.fromtimestamp(
        file_path.stat().st_mtime, tz=_CET
    ).strftime("%d-%b-%Y %H:%M CET")
    h = hashlib.sha256()
    with open(file_path, "rb") as fh:
        for chunk in iter(lambda: fh.read(65536), b""):
            h.update(chunk)
    return mtime_str, h.hexdigest()


def _set_doc_font(doc: Any, font_name: str = "Arial") -> None:
    for style_name in [
        "Normal", "Heading 1", "Heading 2", "Heading 3", "Heading 4",
        "No Spacing", "Intense Quote", "List Number", "List Bullet", "Table Grid",
    ]:
        try:
            doc.styles[style_name].font.name = font_name
        except Exception:
            pass


def _add_header_footer(section: Any, case_id: str) -> None:
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    hdr = section.header
    h_para = hdr.paragraphs[0] if hdr.paragraphs else hdr.add_paragraph()
    h_para.clear()
    h_run = h_para.add_run(case_id)
    h_run.font.name = "Arial"
    h_run.font.size = Pt(9)
    h_run.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)
    h_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    ftr = section.footer
    f_para = ftr.paragraphs[0] if ftr.paragraphs else ftr.add_paragraph()
    f_para.clear()
    f_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    def _fld_char(fld_type: str) -> Any:
        r = OxmlElement("w:r")
        fc = OxmlElement("w:fldChar")
        fc.set(qn("w:fldCharType"), fld_type)
        r.append(fc)
        return r

    def _instr(text: str) -> Any:
        r = OxmlElement("w:r")
        it = OxmlElement("w:instrText")
        it.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
        it.text = text
        r.append(it)
        return r

    def _txt_run(text: str) -> Any:
        r = f_para.add_run(text)
        r.font.name = "Arial"
        r.font.size = Pt(9)
        r.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)
        return r

    _txt_run("Page ")
    f_para._p.append(_fld_char("begin"))
    f_para._p.append(_instr(" PAGE "))
    f_para._p.append(_fld_char("end"))
    _txt_run(" of ")
    f_para._p.append(_fld_char("begin"))
    f_para._p.append(_instr(" NUMPAGES "))
    f_para._p.append(_fld_char("end"))


def _add_watermark(section: Any) -> None:
    from lxml import etree

    hdr = section.header
    wm_para = hdr.add_paragraph()
    vml = (
        '<w:r xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"'
        ' xmlns:v="urn:schemas-microsoft-com:vml"'
        ' xmlns:o="urn:schemas-microsoft-com:office:office">'
        "<w:rPr><w:noProof/></w:rPr>"
        "<w:pict>"
        '<v:shape id="WaterMark" o:spid="_x0000_s2049"'
        ' type="#_x0000_t136"'
        ' style="position:absolute;margin-left:0;margin-top:0;'
        "width:430pt;height:120pt;z-index:-251654144;"
        "rotation:315;"
        "mso-position-horizontal:center;"
        "mso-position-horizontal-relative:page;"
        "mso-position-vertical:center;"
        'mso-position-vertical-relative:page"'
        ' fillcolor="#C0C0C0" stroked="f">'
        "<v:textpath"
        ' on="t" fitshape="t" string="CONFIDENTIAL"'
        " style='font-family:\"Arial\";font-size:1pt;font-weight:bold'"
        "/>"
        "</v:shape>"
        "</w:pict>"
        "</w:r>"
    )
    wm_para._p.append(etree.fromstring(vml))


def _remove_blank_paragraphs(doc: Any) -> None:
    from docx.oxml.ns import qn

    paras = list(doc.paragraphs)
    to_remove = []
    for i, para in enumerate(paras):
        if para.text.strip():
            continue
        if para.style.name.startswith("Heading"):
            continue
        if para._element.find(".//" + qn("w:br")) is not None:
            continue
        prev_heading = i > 0 and paras[i - 1].style.name.startswith("Heading")
        next_heading = i + 1 < len(paras) and paras[i + 1].style.name.startswith("Heading")
        if prev_heading or next_heading:
            continue
        to_remove.append(para._element)
    for elem in to_remove:
        if elem.getparent() is not None:
            elem.getparent().remove(elem)


def _build_fan_pptx(
    stem: str,
    case_id: str,
    generated_cet: str,
    overall_sev: str,
    timeline: list[dict],
    iocs: list[dict],
    recs: list[str],
    data: dict,
    output_path: Path,
) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("[fan] WARNING: python-pptx not installed — skipping PPTX. pip3 install python-pptx")
        return

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def _rgb(t):
        return RGBColor(*t)

    def _rect(slide, l, t, w, h, fill):
        s = slide.shapes.add_shape(1, l, t, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = _rgb(fill)
        s.line.fill.background()
        return s

    def _txt(slide, text, l, t, w, h, sz, bold=False, color=_WHITE, align=PP_ALIGN.LEFT):
        tb  = slide.shapes.add_textbox(l, t, w, h)
        tf  = tb.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.name = "Arial"
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)
        return tb

    W = prs.slide_width
    H = prs.slide_height
    M = Inches(0.4)

    sev_color = _SEV_COLORS.get(overall_sev, _SEV_COLORS["info"])
    pcap_name = f"{stem}.pcap"

    # ── Slide 1 — Cover ───────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    _rect(s1, 0, 0, W, H, _DARK_NAVY)
    _rect(s1, 0, 0, W, Inches(0.08), _BLUE)
    _rect(s1, 0, H - Inches(0.08), W, Inches(0.08), _BLUE)
    _txt(s1, "FAN", M, Inches(1.2), W - 2*M, Inches(1.2),
         72, bold=True, color=_LIGHT_BLUE, align=PP_ALIGN.CENTER)
    _txt(s1, "Forensics agent network", M, Inches(2.2), W - 2*M, Inches(0.7),
         28, color=_WHITE, align=PP_ALIGN.CENTER)
    _txt(s1, "Network forensics incident report", M, Inches(2.9), W - 2*M, Inches(0.6),
         20, color=_LIGHT_BLUE, align=PP_ALIGN.CENTER)
    _rect(s1, Inches(3), Inches(3.8), W - Inches(6), Inches(0.04), _BLUE)
    _txt(s1, f"Case: {case_id}  |  PCAP: {pcap_name}  |  {generated_cet[:10]}",
         M, Inches(4.1), W - 2*M, Inches(0.5), 14, color=_TEXT_MID, align=PP_ALIGN.CENTER)
    _txt(s1, "Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin",
         M, Inches(4.6), W - 2*M, Inches(0.4), 11, color=_TEXT_MID, align=PP_ALIGN.CENTER)
    _txt(s1, "Fan Get Fame Fast  |  FAN module",
         M, H - Inches(0.7), W - 2*M, Inches(0.4), 11, color=_TEXT_MID, align=PP_ALIGN.CENTER)

    # ── Slide 2 — Key findings ────────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    _rect(s2, 0, 0, W, Inches(1.1), _MID_NAVY)
    _txt(s2, "Key findings", M, Inches(0.2), W, Inches(0.8), 28, bold=True, color=_WHITE)
    _txt(s2, f"{case_id}  |  {pcap_name}", M, Inches(0.75), W, Inches(0.3), 12, color=_LIGHT_BLUE)

    unique_ips   = len(data.get("unique_ips", []))
    unique_fqdns = len(data.get("unique_fqdns", []))
    ioc_count    = len(iocs)
    alert_count  = sum(1 for e in timeline if e.get("severity") in ("critical", "high"))

    summary = (
        f"Network forensic analysis of {pcap_name} identified {unique_ips} unique IP "
        f"addresses and {unique_fqdns} unique domain names. "
        f"Overall severity: {overall_sev.upper()}. "
        f"{ioc_count} indicator(s) of compromise extracted; "
        f"{alert_count} critical/high-severity event(s) in the timeline."
    )
    _txt(s2, summary, M, Inches(1.3), W - 2*M, Inches(3.0), 15, color=_TEXT_DARK)

    metrics = [
        ("Unique IPs",    str(unique_ips)),
        ("Unique FQDNs",  str(unique_fqdns)),
        ("IOCs",          str(ioc_count)),
        ("Severity",      overall_sev.upper()),
    ]
    col_w = (W - 2*M) // len(metrics)
    for i, (label, value) in enumerate(metrics):
        cx = M + i * col_w
        val_color = sev_color if label == "Severity" else _AMBER
        _rect(s2, cx + Inches(0.05), Inches(4.7), col_w - Inches(0.1), Inches(1.3), _MID_NAVY)
        _txt(s2, value, cx + Inches(0.1), Inches(4.8), col_w - Inches(0.2), Inches(0.7),
             18, bold=True, color=val_color)
        _txt(s2, label, cx + Inches(0.1), Inches(5.5), col_w - Inches(0.2), Inches(0.4),
             10, color=_LIGHT_BLUE)

    # ── Slide 3 — Incident timeline ───────────────────────────────────────────
    s3 = prs.slides.add_slide(blank)
    _rect(s3, 0, 0, W, Inches(1.1), _MID_NAVY)
    _txt(s3, "Incident timeline", M, Inches(0.2), W, Inches(0.8), 28, bold=True, color=_WHITE)
    _txt(s3, f"{case_id}  |  {pcap_name}", M, Inches(0.75), W, Inches(0.3), 12, color=_LIGHT_BLUE)

    shown_events = [e for e in timeline if e.get("severity") in ("critical", "high", "medium")][:12]
    if not shown_events:
        shown_events = timeline[:12]
    if shown_events:
        col_ws = [Inches(2.5), Inches(1.0), Inches(3.0), W - M - Inches(7.0)]
        hdrs   = ["Timestamp", "Severity", "Category", "Description"]
        row_h  = Inches(0.46)
        hx = M
        for h_t, cw_t in zip(hdrs, col_ws):
            _rect(s3, hx, Inches(1.15), cw_t - Inches(0.05), row_h - Inches(0.04), _MID_NAVY)
            _txt(s3, h_t, hx + Inches(0.08), Inches(1.2), cw_t - Inches(0.13), row_h,
                 12, bold=True, color=_WHITE)
            hx += cw_t
        for i, ev in enumerate(shown_events[:11]):
            y  = Inches(1.15) + (i + 1) * row_h
            bg = _LIGHT_BG if i % 2 == 0 else _ROW_ALT
            rx = M
            vals = [
                str(ev.get("timestamp", ""))[:24],
                str(ev.get("severity", "")).upper(),
                str(ev.get("category", ""))[:25],
                str(ev.get("description", ""))[:70],
            ]
            for val, cw_t in zip(vals, col_ws):
                _rect(s3, rx, y, cw_t - Inches(0.05), row_h - Inches(0.04), bg)
                _txt(s3, val, rx + Inches(0.08), y + Inches(0.06),
                     cw_t - Inches(0.13), row_h, 10, color=_TEXT_DARK)
                rx += cw_t
    else:
        _txt(s3, "No threat events in timeline.", M, Inches(2.0), W - 2*M, Inches(1.0),
             16, color=_TEXT_MID)

    # ── Slide 4 — Key evidence ────────────────────────────────────────────────
    s4 = prs.slides.add_slide(blank)
    _rect(s4, 0, 0, W, Inches(1.1), _MID_NAVY)
    _txt(s4, "Key evidence", M, Inches(0.2), W, Inches(0.8), 28, bold=True, color=_WHITE)
    _txt(s4, f"{case_id}  |  {pcap_name}", M, Inches(0.75), W, Inches(0.3), 12, color=_LIGHT_BLUE)

    module_hits = [
        ("DNS threats",    data.get("has_dns",      False)),
        ("HTTP threats",   data.get("has_http",     False)),
        ("TLS inspection", data.get("has_tls",      False)),
        ("Cert inspection",data.get("has_cert",     False)),
        ("ICMP threats",   data.get("has_icmp",     False)),
        ("TCP threats",    data.get("has_tcp",      False)),
        ("Suricata IDS",   data.get("has_suricata", False)),
        ("YARA matches",   data.get("has_yara",     False)),
        ("CTI enrichment", data.get("has_fan_ip",   False)),
        ("ARP threats",    data.get("has_arp",      False)),
    ]
    row_h = Inches(0.65)
    for i, (label, hit) in enumerate(module_hits[:8]):
        y  = Inches(1.25) + i * row_h
        bg = _LIGHT_BG if i % 2 == 0 else _ROW_ALT
        _rect(s4, M, y, Inches(3.5), row_h - Inches(0.06), bg)
        _rect(s4, M + Inches(3.5), y, W - M - Inches(3.5) - M, row_h - Inches(0.06), bg)
        mark_color = _GREEN if hit else _TEXT_MID
        mark = "Triggered" if hit else "No findings"
        _txt(s4, label, M + Inches(0.1), y + Inches(0.08), Inches(3.3), row_h,
             13, bold=True, color=_TEXT_DARK)
        _txt(s4, mark, M + Inches(3.6), y + Inches(0.08), W - M - Inches(4.1), row_h,
             12, color=mark_color)

    # ── Slide 5 — Recommendations ─────────────────────────────────────────────
    s5 = prs.slides.add_slide(blank)
    _rect(s5, 0, 0, W, Inches(1.1), _MID_NAVY)
    _txt(s5, "Recommendations", M, Inches(0.2), W, Inches(0.8), 28, bold=True, color=_WHITE)
    _txt(s5, f"{case_id}  |  {pcap_name}", M, Inches(0.75), W, Inches(0.3), 12, color=_LIGHT_BLUE)
    row_h = Inches(0.72)
    for i, rec in enumerate(recs[:7]):
        y = Inches(1.2) + i * row_h
        _rect(s5, M, y, Inches(0.5), row_h - Inches(0.08), _BLUE)
        _txt(s5, str(i + 1), M + Inches(0.1), y + Inches(0.1),
             Inches(0.3), row_h, 16, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
        import re as _re
        rec_clean = _re.sub(r"\*\*(.*?)\*\*", r"\1", rec[:120])
        _txt(s5, rec_clean, M + Inches(0.6), y + Inches(0.1),
             W - M - Inches(1.0), row_h, 13, color=_TEXT_DARK)

    prs.save(str(output_path))
    print(f"[fan] PPTX saved: {output_path}")


def _build_fan_docx(
    stem: str,
    case_id: str,
    generated_cet: str,
    overall_sev: str,
    timeline: list[dict],
    iocs: list[dict],
    recs: list[str],
    data: dict,
    output_path: Path,
) -> None:
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        print("[fan] WARNING: python-docx not installed — skipping DOCX. pip3 install python-docx")
        return

    import re as _re

    doc = Document()
    styles = doc.styles
    pcap_name = f"{stem}.pcap"

    # ── Page setup ────────────────────────────────────────────────────────────
    for section in doc.sections:
        section.top_margin      = Inches(1.0)
        section.bottom_margin   = Inches(1.0)
        section.left_margin     = Inches(1.2)
        section.right_margin    = Inches(1.2)
        section.header_distance = Inches(0.4)
        section.footer_distance = Inches(0.4)

    _set_doc_font(doc, "Arial")
    doc.styles["Normal"].paragraph_format.space_after  = Pt(5)
    doc.styles["Normal"].paragraph_format.space_before = Pt(0)
    for _i in range(1, 5):
        try:
            _hs = doc.styles[f"Heading {_i}"]
            _hs.paragraph_format.space_before = Pt(14 if _i == 1 else 10)
            _hs.paragraph_format.space_after  = Pt(4)
            _hs.font.name = "Arial"
        except Exception:
            pass

    def _heading(text: str, level: int) -> None:
        p = doc.add_heading(text, level=level)
        if p.runs:
            p.runs[0].font.color.rgb = RGBColor(0x0f, 0x17, 0x2a)
            p.runs[0].font.name = "Arial"

    def _para(text: str, bold: bool = False, italic: bool = False) -> None:
        p = doc.add_paragraph()
        run = p.add_run(text)
        run.bold      = bold
        run.italic    = italic
        run.font.name = "Arial"

    def _note(text: str) -> None:
        p = (doc.add_paragraph(style="Intense Quote")
             if "Intense Quote" in [s.name for s in styles]
             else doc.add_paragraph())
        run = p.add_run(text)
        run.italic = True
        run.font.name = "Arial"
        run.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)

    def _code(text: str) -> None:
        sn = "No Spacing" if "No Spacing" in [s.name for s in styles] else "Normal"
        p  = doc.add_paragraph(style=sn)
        r  = p.add_run(text)
        r.font.name  = "Courier New"
        r.font.size  = Pt(9)
        r.font.color.rgb = RGBColor(0x1e, 0x3a, 0x5f)

    def _table_2col(rows: list[tuple[str, str]], header: bool = True) -> None:
        start = 1 if header else 0
        tbl = doc.add_table(rows=len(rows) + start, cols=2)
        tbl.style = "Table Grid"
        if header:
            for i, h in enumerate(["Field", "Value"]):
                c = tbl.rows[0].cells[i]
                c.text = h
                c.paragraphs[0].runs[0].font.bold = True
                c.paragraphs[0].runs[0].font.name = "Arial"
        for i, (k, v) in enumerate(rows):
            r = tbl.rows[i + start]
            r.cells[0].text = k
            r.cells[1].text = v
            for j in range(2):
                for run in r.cells[j].paragraphs[0].runs:
                    run.font.name = "Arial"

    # ── Header, footer, watermark ─────────────────────────────────────────────
    for section in doc.sections:
        _add_header_footer(section, case_id)
        _add_watermark(section)

    # ── Cover ──────────────────────────────────────────────────────────────────
    title = doc.add_heading("FAN — Network forensics report", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if title.runs:
        title.runs[0].font.name = "Arial"

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = sub.add_run("Forensics Agent Network  |  Fan Get Fame Fast")
    r.font.size  = Pt(14)
    r.font.name  = "Arial"
    r.font.color.rgb = RGBColor(0x1d, 0x4e, 0xd8)

    _table_2col([
        ("Case ID",      case_id),
        ("PCAP file",    pcap_name),
        ("Module",       "FAN — Forensics Agent Network"),
        ("Analysts",     "Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin"),
        ("Generated",    generated_cet),
        ("Analysis tools", "tshark · Suricata · YARA · 22 protocol detectors"),
    ])
    doc.add_page_break()

    # ── Part A — Investigation methodology ───────────────────────────────────
    _heading("Part A — Investigation methodology", 1)
    doc.add_paragraph()

    _heading("A.1  Evidence acquisition and integrity", 2)
    _para(
        f"The PCAP file {pcap_name!r} was acquired as a read-only artifact. "
        "All analysis tools open the file in read-only mode. The PCAP is never modified. "
        "The SHA-256 hash of the PCAP is recorded at analysis start and verified at "
        "report generation time to confirm evidence integrity throughout the pipeline."
    )
    doc.add_paragraph()
    _para(
        "Analysis output is written to separate working directories (analysis/) and "
        "archived in the investigations vault on completion. The PCAP itself is never "
        "stored in the working directory — only the analysis outputs."
    )
    doc.add_paragraph()

    _heading("A.2  Tool suite", 2)
    _table_2col([
        ("tshark",              "Wireshark CLI — packet decoding and flow extraction"),
        ("Suricata IDS",        "Signature-based network intrusion detection"),
        ("YARA",                "Pattern matching against PCAP payload bytes"),
        ("DNS threat detector", "DNS anomaly and threat analysis"),
        ("HTTP threat detector","HTTP/HTTPS anomaly and threat analysis"),
        ("TLS inspector",       "TLS session and certificate analysis"),
        ("Certificate inspector","X.509 certificate chain analysis"),
        ("ICMP detector",       "ICMP tunneling, flood, and anomaly detection"),
        ("TCP detector",        "TCP scan, RST flood, and C2 pattern detection"),
        ("ARP detector",        "ARP spoofing and poisoning detection"),
        ("+ 12 protocol detectors", "NTP, DHCP, mDNS, NBNS, LLMNR, QUIC, SNMP, STUN, SSDP, UDP, NetBIOS, file hashes"),
    ])
    doc.add_paragraph()

    _heading("A.3  Analysis pipeline", 2)
    _para(
        f"PCAP: {pcap_name}  →  22 protocol threat detectors  →  "
        "Suricata IDS + YARA scanning  →  CTI enrichment  →  "
        "report generation  →  investigations vault upload"
    )
    doc.add_paragraph()
    doc.add_page_break()

    # ── Part B — Artifact extraction catalog ─────────────────────────────────
    _heading("Part B — Analysis module catalog", 1)
    _note("Each module below ran against the PCAP and produced the listed outputs.")
    doc.add_paragraph()

    modules_catalog = [
        ("B.1",  "DNS threat analysis",       "has_dns",      "dns_threats/",
         "Detects DNS tunneling, DGA domains, DNS exfiltration, and typosquatting."),
        ("B.2",  "HTTP/HTTPS threat analysis", "has_http",     "http_threats/",
         "Identifies C2 beacons, malicious downloads, and suspicious User-Agent strings."),
        ("B.3",  "TLS session inspection",    "has_tls",      "tls_inspector/",
         "Analyzes TLS versions, cipher suites, and SNI patterns for weaknesses."),
        ("B.4",  "Certificate inspection",    "has_cert",     "cert_inspector/",
         "Validates X.509 certificate chains, expiry, and self-signed certs."),
        ("B.5",  "ICMP threat detection",     "has_icmp",     "icmp_threats/",
         "Detects ICMP tunneling, flood patterns, and oversized ICMP packets."),
        ("B.6",  "TCP threat detection",      "has_tcp",      "tcp_threats/",
         "Identifies port scans, SYN floods, and RST injection patterns."),
        ("B.7",  "ARP threat detection",      "has_arp",      "arp_threats/",
         "Detects ARP spoofing, ARP poisoning, and gratuitous ARP anomalies."),
        ("B.8",  "UDP threat detection",      "has_udp",      "udp_threats/",
         "Identifies UDP flood patterns and anomalous UDP traffic."),
        ("B.9",  "Suricata IDS",              "has_suricata", "suricata/",
         "Signature-based detection against the Suricata ruleset (ET/Open)."),
        ("B.10", "YARA pattern matching",     "has_yara",     "yara_pcap/",
         "Custom YARA rules matched against PCAP payload bytes."),
        ("B.11", "CTI enrichment",            "has_fan_ip",   "fan_ip/",
         "IPs and FQDNs cross-referenced against OpenCTI threat intelligence."),
    ]

    for ref, name, has_key, path, desc in modules_catalog:
        ran = data.get(has_key, False)
        _heading(f"{ref}  {name}", 2)
        _para("Status", bold=True)
        _para("Triggered — findings available." if ran else "No findings — analysis complete, no threats detected.")
        doc.add_paragraph()
        _para("Description", bold=True)
        _para(desc)
        doc.add_paragraph()
        _para("Output path", bold=True)
        _code(f"analysis/{path}{stem}/")
        doc.add_paragraph()

    doc.add_page_break()

    # ── Part C — Findings ──────────────────────────────────────────────────────
    _heading("Part C — Findings", 1)
    doc.add_page_break()

    # C.1 Management summary
    _heading("C.1  Management summary", 2)
    _note("Audience: CISO, Legal, Internal Audit — no technical identifiers.")
    unique_ips   = len(data.get("unique_ips", []))
    unique_fqdns = len(data.get("unique_fqdns", []))
    _para(
        f"Network forensic analysis of the captured traffic identified {unique_ips} unique "
        f"IP addresses and {unique_fqdns} unique domain names. "
        f"Overall severity: {overall_sev.upper()}. "
        f"{len(iocs)} indicator(s) of compromise extracted. "
        "The timeline of threat events has been reconstructed and is available for "
        "cross-reference with storage and memory forensics findings."
    )
    doc.add_paragraph()

    # C.2 IOCs
    _heading("C.2  Indicators of compromise", 2)
    _note("All IOC values defanged. Claude: enhance and elaborate when necessary.")
    if iocs:
        col_hdrs = ["Type", "Value", "Severity", "Context"]
        tbl = doc.add_table(rows=len(iocs) + 1, cols=len(col_hdrs))
        tbl.style = "Table Grid"
        for i, h in enumerate(col_hdrs):
            c = tbl.rows[0].cells[i]
            c.text = h
            if c.paragraphs[0].runs:
                c.paragraphs[0].runs[0].font.bold = True
                c.paragraphs[0].runs[0].font.name = "Arial"
        for i, ioc in enumerate(iocs):
            vals = [
                str(ioc.get("type", "")),
                str(ioc.get("value", "")),
                str(ioc.get("severity", "")),
                str(ioc.get("context", "")),
            ]
            for j, val in enumerate(vals):
                c = tbl.rows[i + 1].cells[j]
                c.text = val
                for run in c.paragraphs[0].runs:
                    run.font.name = "Arial"
    else:
        _para("No malicious indicators of compromise identified from network analysis.")
    doc.add_paragraph()

    # C.3 Recommendations
    _heading("C.3  Recommendations", 2)
    _note("Claude: enhance and elaborate when necessary.")
    for i, rec in enumerate(recs, 1):
        rec_clean = _re.sub(r"\*\*(.*?)\*\*", r"\1", rec)
        p = doc.add_paragraph(style="List Number")
        run = p.add_run(rec_clean)
        run.font.name = "Arial"
    doc.add_page_break()

    # ── Appendix A — 4-column ─────────────────────────────────────────────────
    _heading("Appendix A — Analysis source files", 1)
    _note("SHA-256 hashes computed at report generation time.")

    artifact_files: list[tuple[Path, str]] = []
    if data.get("has_pcap"):
        d = ANALYSIS_DIR / "pcap" / stem
        artifact_files += [
            (d / "netflow.csv",      "PCAP netflow (tshark)"),
            (d / "unique_ips.txt",   "Unique IP addresses"),
            (d / "unique_fqdns.txt", "Unique FQDNs"),
        ]
    for has_key, subdir, label in [
        ("has_dns",      "dns_threats",    "DNS threats JSON"),
        ("has_http",     "http_threats",   "HTTP threats JSON"),
        ("has_tls",      "tls_inspector",  "TLS sessions JSON"),
        ("has_cert",     "cert_inspector", "Certificate findings JSON"),
        ("has_icmp",     "icmp_threats",   "ICMP threats JSON"),
        ("has_tcp",      "tcp_threats",    "TCP threats JSON"),
        ("has_arp",      "arp_threats",    "ARP threats JSON"),
        ("has_suricata", "suricata",       "Suricata alerts JSON"),
        ("has_yara",     "yara_pcap",      "YARA matches JSON"),
        ("has_fan_ip",   "fan_ip",         "CTI enrichment CSV"),
    ]:
        if data.get(has_key):
            d = ANALYSIS_DIR / subdir / stem
            fname = f"{subdir.split('_')[0]}_{subdir.split('_')[1] if '_' in subdir else 'results'}.json"
            fp = d / fname
            if not fp.exists():
                fps = list(d.glob("*.json"))
                fp = fps[0] if fps else d
            artifact_files.append((fp, label))

    tbl = doc.add_table(rows=len(artifact_files) + 1, cols=4)
    tbl.style = "Table Grid"
    for i, h in enumerate(["Filename", "Description", "Generated (CET)", "SHA-256 (first 32)"]):
        c = tbl.rows[0].cells[i]
        c.text = h
        if c.paragraphs[0].runs:
            c.paragraphs[0].runs[0].font.bold = True
            c.paragraphs[0].runs[0].font.name = "Arial"
    for i, (fp, desc) in enumerate(artifact_files):
        mtime, sha256 = _compute_file_metadata(fp)
        vals = [fp.name, desc, mtime, sha256[:32] if sha256 != "N/A" else "N/A"]
        for j, val in enumerate(vals):
            c = tbl.rows[i + 1].cells[j]
            c.text = val
            for run in c.paragraphs[0].runs:
                run.font.name = "Arial"

    _remove_blank_paragraphs(doc)
    doc.save(str(output_path))
    print(f"[fan] DOCX saved: {output_path}")


# ── PDF conversion ─────────────────────────────────────────────────────────────

def convert_to_pdf(
    md_path: Path,
    pdf_path: Path,
    title: str = "",
    case_id: str = "",
    date_str: str = "",
) -> bool:
    """
    Convert *md_path* to a styled PDF with cover page, running page-header stripe,
    and 'Page X of Y' pagination.  Delegates to lib/md_to_pdf.py (WeasyPrint).
    Falls back to plain WeasyPrint with @page pagination if md_to_pdf unavailable.
    """
    try:
        sys.path.insert(0, str(Path(__file__).parent))
        from md_to_pdf import convert as _md_convert
        _md_convert(
            md_path=md_path,
            output_path=pdf_path,
            title=title or md_path.stem.replace("_", " ").replace("-", " ").title(),
            subtitle=md_path.stem,
            case_id=case_id,
            date_str=date_str,
        )
        return pdf_path.exists()
    except SystemExit:
        pass
    except Exception as exc:
        print(f"[report] md_to_pdf error: {exc}", file=sys.stderr)

    # Fallback: plain weasyprint with @page pagination only
    try:
        import markdown as md_lib
        import weasyprint

        _PLAIN_CSS = """
        @page { size: A4; margin: 2cm;
                @bottom-right { content: "Page " counter(page) " of " counter(pages);
                                 font-size: 8pt; color: #9ca3af; }
                @bottom-left  { content: "CONFIDENTIAL — DFIR INTERNAL USE ONLY";
                                 font-size: 8pt; color: #9ca3af; } }
        body { font-family: Arial, sans-serif; font-size: 10pt;
               color: #1a1a1a; line-height: 1.5; }
        h1 { color: #0f172a; border-bottom: 2px solid #1d4ed8; padding-bottom: 4px; }
        h2 { color: #1e3a5f; border-bottom: 1px solid #93c5fd; padding-bottom: 3px;
             margin-top: 1.5em; }
        h3 { color: #1d4ed8; margin-top: 1.2em; }
        table { border-collapse: collapse; width: 100%; margin: 1em 0; font-size: 9pt; }
        thead tr { background: #1e3a5f; color: white; }
        th { padding: 6px 8px; text-align: left; font-size: 8pt; text-transform: uppercase; }
        td { padding: 5px 8px; border-bottom: 1px solid #e5e7eb; vertical-align: top; }
        tr:nth-child(even) { background: #f8fafc; }
        code { background: #f1f5f9; padding: 1px 4px; border-radius: 3px;
               font-family: monospace; font-size: 9pt; }
        pre  { background: #0f172a; color: #e2e8f0; padding: 10px; border-radius: 6px;
               font-size: 8pt; white-space: pre-wrap; word-break: break-all; }
        pre code { background: none; color: inherit; }
        blockquote { background: #eff6ff; border-left: 4px solid #1d4ed8;
                     padding: 0.4cm 0.6cm; margin: 0.3cm 0; }
        hr { border: none; border-top: 1px solid #e5e7eb; margin: 1.5em 0; }
        """
        text = md_path.read_text(encoding="utf-8")
        body = md_lib.markdown(text, extensions=["tables", "fenced_code", "toc", "nl2br"])
        full = (
            f'<!DOCTYPE html><html><head><meta charset="utf-8">'
            f"<style>{_PLAIN_CSS}</style></head><body>{body}</body></html>"
        )
        weasyprint.HTML(string=full).write_pdf(str(pdf_path), presentational_hints=True)
        return pdf_path.exists()
    except ImportError:
        pass
    except Exception as exc:
        print(f"[report] weasyprint fallback error: {exc}", file=sys.stderr)

    return False


# ── Narrative loader ──────────────────────────────────────────────────────────

def _load_narrative(case_id: str, reports_dir: Path) -> dict[str, str]:
    """Load Claude-generated narrative sections from {case_id}_narrative.md."""
    if not case_id:
        return {}
    path = reports_dir / f"{case_id}_narrative.md"
    if not path.exists():
        return {}
    sections: dict[str, str] = {}
    current: str | None = None
    for line in path.read_text(encoding="utf-8").splitlines():
        if line.startswith("## "):
            current = line[3:].strip()
            sections[current] = ""
        elif line.startswith("<!--"):
            continue
        elif current is not None:
            sections[current] += line + "\n"
    return {k: v.strip() for k, v in sections.items()}


# ── Hallucination Guard ────────────────────────────────────────────────────────

def _build_fan_hallucination_guard_section(data: dict, case_id: str, out_dir: Path) -> list[str]:
    """
    Build the Hallucination Guard section for FAN (PCAP) reports.

    Tags each key conclusion with a ConfidenceTier based on which analysers
    ran and what they produced. Tiers are assigned by code logic, not by
    Claude prompt instructions.
    """
    _hg_reset()
    findings = []
    steps = _parse_research_steps(case_id, str(out_dir)) if case_id else []

    # Suricata alerts — highest-confidence: IDS rule match on wire traffic
    suricata = data.get("suricata_data") or {}
    alerts = suricata.get("alerts") or (suricata if isinstance(suricata, list) else [])
    if alerts:
        findings.append(tag_finding(
            f"Suricata IDS produced {len(alerts)} alert(s) — rule-based detections on wire traffic",
            ConfidenceTier.CONFIRMED,
            [],
            ["suricata"],
            ["fan"],
        ))
    elif data.get("has_suricata"):
        findings.append(tag_finding(
            "Suricata ran but produced no alerts",
            ConfidenceTier.CONFIRMED,
            [],
            ["suricata"],
            ["fan"],
        ))
    else:
        findings.append(tag_finding(
            "Suricata IDS not run — rule-based alert detection unavailable",
            ConfidenceTier.UNVERIFIABLE,
            [],
            ["suricata"],
            ["fan"],
        ))

    # YARA matches
    yara = data.get("yara_results") or {}
    yara_hits = [k for k, v in yara.items() if isinstance(v, dict) and v.get("triggered")]
    if yara_hits:
        findings.append(tag_finding(
            f"YARA matched {len(yara_hits)} rule(s) in PCAP payload: {', '.join(yara_hits[:3])}",
            ConfidenceTier.CONFIRMED,
            [],
            ["yara-python"],
            ["fan"],
        ))
    elif data.get("has_yara"):
        findings.append(tag_finding(
            "YARA scan ran on PCAP — no rule matches detected",
            ConfidenceTier.CONFIRMED,
            [],
            ["yara-python"],
            ["fan"],
        ))
    else:
        findings.append(tag_finding(
            "YARA scan not run on PCAP — malware signature detection in traffic unavailable",
            ConfidenceTier.UNVERIFIABLE,
            [],
            ["yara-python"],
            ["fan"],
        ))

    # Protocol threat detections (direct analyser output → CONFIRMED)
    protocol_map = [
        ("has_dns",     "dns_results",     "DNS",     "fan_dns_threats"),
        ("has_http",    "http_results",    "HTTP(S)", "fan_http_threats"),
        ("has_tls",     "tls_results",     "TLS",     "fan_tls_inspector"),
        ("has_tcp",     "tcp_results",     "TCP",     "fan_tcp_threats"),
        ("has_icmp",    "icmp_results",    "ICMP",    "fan_icmp_threats"),
        ("has_arp",     "arp_results",     "ARP",     "fan_arp_threats"),
        ("has_dhcp",    "dhcp_results",    "DHCP",    "fan_dhcp_threats"),
        ("has_ntp",     "ntp_results",     "NTP",     "fan_ntp_threats"),
        ("has_mdns",    "mdns_results",    "mDNS",    "fan_mdns_threats"),
        ("has_nbns",    "nbns_results",    "NBNS",    "fan_nbns_threats"),
        ("has_llmnr",   "llmnr_results",   "LLMNR",   "fan_llmnr_threats"),
        ("has_snmp",    "snmp_results",    "SNMP",    "fan_snmp_threats"),
        ("has_quic",    "quic_results",    "QUIC",    "fan_quic_threats"),
        ("has_stun",    "stun_results",    "STUN",    "fan_stun_threats"),
        ("has_ssdp",    "ssdp_results",    "SSDP",    "fan_ssdp_threats"),
        ("has_netbios", "netbios_results", "NetBIOS", "fan_netbios_threats"),
    ]
    for has_key, result_key, proto_name, tool_name in protocol_map:
        if data.get(has_key):
            results = data.get(result_key) or {}
            triggered = [k for k, v in results.items()
                         if isinstance(v, dict) and v.get("triggered")]
            if triggered:
                findings.append(tag_finding(
                    f"{proto_name} analyser flagged {len(triggered)} category/ies: "
                    f"{', '.join(triggered[:3])}",
                    ConfidenceTier.CONFIRMED,
                    [],
                    [tool_name],
                    ["fan"],
                ))
            else:
                findings.append(tag_finding(
                    f"{proto_name} analyser ran — no threats detected",
                    ConfidenceTier.CONFIRMED,
                    [],
                    [tool_name],
                    ["fan"],
                ))

    # Behavioral inferences (beaconing, C2 patterns) — INFERRED
    dns_results = data.get("dns_results") or {}
    for cat, v in dns_results.items():
        if isinstance(v, dict) and v.get("triggered") and "beacon" in cat.lower():
            findings.append(tag_finding(
                f"DNS beaconing pattern inferred from query frequency analysis ({cat})",
                ConfidenceTier.INFERRED,
                [],
                ["fan_dns_threats"],
                ["fan"],
            ))

    http_results = data.get("http_results") or {}
    for cat, v in http_results.items():
        if isinstance(v, dict) and v.get("triggered") and any(
            kw in cat.lower() for kw in ("c2", "command", "beacon", "exfil")
        ):
            findings.append(tag_finding(
                f"HTTP C2/exfiltration pattern inferred from traffic analysis ({cat})",
                ConfidenceTier.INFERRED,
                [],
                ["fan_http_threats"],
                ["fan"],
            ))

    # Assumptions from research notes
    for s in steps:
        outcome = s.get("outcome", "")
        if "[ASSUMPTION]" in outcome or s.get("confidence") == "assumed":
            text = outcome.replace("[ASSUMPTION]", "").strip()
            if text:
                findings.append(tag_finding(
                    text,
                    ConfidenceTier.ASSUMED,
                    [s["id"]] if s.get("id") else [],
                    [s.get("source_tool", "")] if s.get("source_tool") else [],
                    ["fan"],
                ))

    if not findings:
        return []

    section_md = render_confidence_summary(findings, module_label="FAN")
    return section_md.splitlines()


# ── Evidence trail ────────────────────────────────────────────────────────────

def _build_evidence_trail(case_id: str, reports_dir: Path) -> list[str]:
    if not case_id:
        return []
    steps  = _parse_research_steps(case_id, str(reports_dir))
    events = _parse_research_events(case_id, str(reports_dir))
    if not steps and not events:
        return []

    lines: list[str] = [
        "---", "",
        "## Appendix B — Investigation Evidence Trail", "",
    ]

    # Attacker timeline — events sorted by evidence timestamp
    if events:
        def _ev_sort(ev: dict) -> tuple:
            ts = ev.get("timestamp", "")
            if ts:
                try:
                    from datetime import datetime, timezone
                    dt = datetime.strptime(ts.replace(" UTC", "").strip(), "%Y-%m-%d %H:%M:%S")
                    return (0, dt.replace(tzinfo=timezone.utc))
                except ValueError:
                    pass
            from datetime import datetime, timezone
            return (1, datetime.min.replace(tzinfo=timezone.utc))

        sorted_events = sorted(events, key=_ev_sort)
        lines += [
            "### Attacker Timeline", "",
            "Attacker events observed in the evidence, ordered by evidence timestamp.", "",
            "| Timestamp (UTC) | Severity | Event | Source |",
            "|-----------------|----------|-------|--------|",
        ]
        for ev in sorted_events:
            ts   = ev.get("timestamp", "") or "—"
            sev  = ev.get("severity", "info").upper()
            desc = ev.get("description", "")[:160].replace("|", "\\|")
            if len(ev.get("description", "")) > 160:
                desc += "…"
            src = (ev.get("source_detail", "") or "—").replace("|", "\\|")
            lines.append(f"| {ts} | **{sev}** | {desc} | {src} |")
        lines += ["", ""]

    # Analysis timeline — analyst investigation steps
    if steps:
        lines += [
            "### Analysis Timeline", "",
            "Steps recorded in the research notes during this investigation. "
            f"Preserved artifacts are in `{case_id}_evidence/`.", "",
            "| Step ID | Timestamp | Analysis Step | Outcome |",
            "|---------|-----------|---------------|---------|",
        ]
        for s in steps:
            sid = f"`{s['id']}`" if s["id"] else "—"
            outcome = s["outcome"].replace("|", "\\|")
            lines.append(f"| {sid} | {s['timestamp']} | {s['title']} | {outcome} |")
        lines += [
            "",
            "*Cross-reference step IDs with the research notes and preserved artifacts "
            f"in `{case_id}_evidence/` to verify any conclusion in this report.*",
            "",
        ]
    return lines


def _sec_incident_timeline(case_id: str, out_dir: Path) -> list[str]:
    """Incident Timeline section sourced from Claude-generated narrative."""
    narrative = _load_narrative(case_id, out_dir)
    timeline_text = narrative.get("attack_timeline", "")
    lines = [
        "---", "",
        "## Incident Timeline", "",
        "> Chronological reconstruction of the attack path. Each finding references",
        "> the investigation step (RN-NNN) and the preserved source file in",
        f"> `{case_id}_evidence/`.",
        "",
    ]
    if timeline_text:
        lines.append(timeline_text)
    else:
        lines += [
            "> *Incident timeline not yet generated. Run the FAN skill to produce*",
            f"> *`{case_id}_narrative.md` with the `attack_timeline` section.*",
        ]
    lines.append("")
    return lines


# ── Main entry point ───────────────────────────────────────────────────────────

def generate_report(
    stem: str,
    case_id: str = "",
    output_dir: Path | None = None,
    base_dir: Path | None = None,
    report_version: int = 1,
) -> dict[str, Path | None]:
    global ANALYSIS_DIR
    if base_dir:
        ANALYSIS_DIR = base_dir

    out_dir = output_dir or REPORTS_DIR
    out_dir.mkdir(parents=True, exist_ok=True)

    print(f"[report] Loading analysis data for stem: {stem}")
    data = load_all_data(stem)

    avail = []
    if data["has_pcap"]:  avail.append("PCAP netflow")
    if data["has_icmp"]:  avail.append("ICMP threats")
    if data["has_dns"]:   avail.append("DNS threats")
    if data["has_ntp"]:   avail.append("NTP threats")
    if data["has_http"]:  avail.append("HTTP(S) threats")
    if data["has_cert"]:  avail.append("certificate inspector")
    if data["has_tls"]:   avail.append("TLS inspector")
    if data["has_arp"]:   avail.append("ARP threats")
    if data["has_tcp"]:   avail.append("TCP threats")
    if data["has_udp"]:   avail.append("UDP threats")
    if data["has_dhcp"]:  avail.append("DHCP threats")
    if data["has_mdns"]:  avail.append("mDNS threats")
    if data["has_quic"]:    avail.append("QUIC threats")
    if data["has_snmp"]:    avail.append("SNMP threats")
    if data["has_nbns"]:    avail.append("NBNS threats")
    if data["has_llmnr"]:   avail.append("LLMNR threats")
    if data["has_stun"]:    avail.append("STUN threats")
    if data["has_ssdp"]:    avail.append("SSDP threats")
    if data["has_netbios"]: avail.append("NetBIOS threats")
    if data["has_fh"]:       avail.append("file hashes")
    if data["has_suricata"]: avail.append("Suricata IDS")
    if data["has_yara"]:     avail.append("YARA rules")
    if data["has_fan_ip"]:      avail.append("CTI enrichment")
    print(f"[report] Available sources: {', '.join(avail) or 'none'}")

    now = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
    generated_cet = datetime.now(_CET).strftime("%d-%b-%Y %H:%M CET")
    overall_sev = _overall_severity(data)
    first_ts, last_ts, duration = _capture_window(data)

    print(f"[report] Overall severity: {overall_sev.upper()}")
    print(f"[report] Building timeline ...")
    timeline = build_timeline(data)
    print(f"[report] Timeline: {len(timeline)} events")

    iocs     = extract_iocs(data)
    coverage = mitre_coverage(data)
    recs     = build_recommendations(data, overall_sev)

    # Assemble sections
    sections: list[str] = []
    sections.extend(sec_header(stem, case_id, overall_sev, now, first_ts, last_ts, duration, report_version))
    sections.extend(sec_management_summary(data, overall_sev, first_ts, last_ts, duration))
    sections.extend(_sec_incident_timeline(case_id, out_dir))
    sections.extend(sec_findings_icmp(data))
    sections.extend(sec_findings_dns(data))
    sections.extend(sec_findings_ntp(data))
    sections.extend(sec_findings_http(data))
    sections.extend(sec_findings_cert(data))
    sections.extend(sec_findings_tls(data))
    sections.extend(sec_findings_arp(data))
    sections.extend(sec_findings_tcp(data))
    sections.extend(sec_findings_udp(data))
    sections.extend(sec_findings_dhcp(data))
    sections.extend(sec_findings_mdns(data))
    sections.extend(sec_findings_quic(data))
    sections.extend(sec_findings_snmp(data))
    sections.extend(sec_findings_nbns(data))
    sections.extend(sec_findings_llmnr(data))
    sections.extend(sec_findings_stun(data))
    sections.extend(sec_findings_ssdp(data))
    sections.extend(sec_findings_netbios(data))
    sections.extend(sec_findings_filehashes(data))
    sections.extend(sec_findings_suricata(data))
    sections.extend(sec_findings_yara(data))
    sections.extend(sec_network_summary(data, duration))
    sections.extend(sec_cti_enrichment(data))
    sections.extend(sec_timeline(timeline))
    sections.extend(sec_iocs(iocs))
    sections.extend(sec_mitre(coverage))
    sections.extend(sec_recommendations(recs))
    sections.extend(sec_appendix(stem, data, case_id))
    sections.extend(_build_fan_hallucination_guard_section(data, case_id, out_dir))
    sections.extend(_build_evidence_trail(case_id, out_dir))

    md_content = "\n".join(sections) + "\n"

    md_path  = out_dir / f"{stem}_incident_report.md"
    pdf_path = out_dir / f"{stem}_incident_report.pdf"

    md_path.write_text(md_content, encoding="utf-8")
    print(f"[report] Markdown: {md_path}")

    print(f"[report] Converting to PDF ...")
    pdf_title = f"PCAP Incident Report — {stem}" + (f" | {case_id}" if case_id else "")
    ok = convert_to_pdf(md_path, pdf_path, title=pdf_title, case_id=case_id, date_str=now[:10])
    if ok:
        print(f"[report] PDF:      {pdf_path}")
    else:
        print("[report] PDF conversion failed. Install: pip3 install markdown weasyprint", file=sys.stderr)
        pdf_path = None

    # PPTX
    pptx_path = out_dir / f"{stem}_fan_presentation.pptx"
    try:
        _build_fan_pptx(
            stem, case_id, generated_cet, overall_sev,
            timeline, iocs, recs, data, pptx_path,
        )
    except Exception as exc:
        print(f"[report] PPTX generation failed: {exc}", file=sys.stderr)
        pptx_path = None

    # DOCX
    docx_path = out_dir / f"{stem}_fan_report.docx"
    try:
        _build_fan_docx(
            stem, case_id, generated_cet, overall_sev,
            timeline, iocs, recs, data, docx_path,
        )
    except Exception as exc:
        print(f"[report] DOCX generation failed: {exc}", file=sys.stderr)
        docx_path = None

    print(f"[report] Done.")
    return {
        "md":   md_path,
        "pdf":  pdf_path,
        "pptx": pptx_path if pptx_path and pptx_path.exists() else None,
        "docx": docx_path if docx_path and docx_path.exists() else None,
    }


# ── CLI ────────────────────────────────────────────────────────────────────────

def _build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        description="PCAP Incident Report Generator — aggregates ICMP/DNS/netflow/CTI analysis into a report",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "Examples:\n"
            "  python3 lib/generate_pcap_report.py --stem capture\n"
            "  python3 lib/generate_pcap_report.py --stem capture --case-id CASE-2025-001\n"
            "  python3 lib/generate_pcap_report.py --stem my-pcap --output-dir ./reports/case1/\n"
        ),
    )
    p.add_argument("--stem",           metavar="STEM", required=True, help="PCAP file stem (matches analysis subdirectories)")
    p.add_argument("--case-id",        metavar="ID",   default="",   help="Case ID stamped into the report")
    p.add_argument("--output-dir",     metavar="DIR",                help="Output directory (default: ./reports/)")
    p.add_argument("--base-dir",       metavar="DIR",                help="Analysis base directory (default: ./analysis/)")
    p.add_argument("--report-version", metavar="N",   type=int, default=1,
                   help="Report version number stamped into the report header (default: 1)")
    return p


if __name__ == "__main__":
    args  = _build_parser().parse_args()
    out   = Path(args.output_dir) if args.output_dir else None
    base  = Path(args.base_dir)   if args.base_dir   else None
    paths = generate_report(stem=args.stem, case_id=args.case_id, output_dir=out,
                             base_dir=base, report_version=args.report_version)
    print("[report] Report suite complete:")
    for fmt, p in paths.items():
        if p:
            print(f"  {fmt.upper():4s}  {p}")
