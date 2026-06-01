#!/usr/bin/env python3
# SPDX-License-Identifier: MIT OR Apache-2.0
# SPDX-FileCopyrightText: 2026 Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin
"""
generate_campaign_report.py — Cross-case campaign report for FanGetFameFast.

Reads all {case_id}_narrative.md and {case_id}_research_notes.md files in
reports_dir and produces a single unified investigation report covering all
analysed hosts:

  CAMPAIGN_{campaign_id}_report.md
  CAMPAIGN_{campaign_id}_report.pdf
  CAMPAIGN_{campaign_id}_board_deck.pptx

Usage:
    python3 lib/generate_campaign_report.py \
        --campaign-id SHIELDBASE-2026 \
        --title "Operation ShieldBase" \
        [--reports-dir ./reports] \
        [--output-dir ./reports]
"""
from __future__ import annotations

import argparse
import re
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

PROJECT_ROOT = Path(__file__).parent.parent
REPORTS_DIR  = PROJECT_ROOT / "reports"

sys.path.insert(0, str(PROJECT_ROOT / "lib"))

# Severity keywords (ranked, highest first)
_SEV_KEYWORDS = {
    "critical": [
        "rootkit", "mnemosyne", "mimikatz", "meterpreter", "dkom",
        "ransomware", "c2 framework", "command-and-control", "credential theft",
        "code injection", "malfind", "sliver", "cobalt strike", "empire",
        "CRITICAL",
    ],
    "high": [
        "suspicious process", "anomalous", "malware", "injection", "yara match",
        "hidden process", "rdp from non-domain", "staging", "exfil",
        "HIGH",
    ],
    "medium": [
        "deviation", "baseline diff", "legacy", "outdated", "expired cert",
        "MEDIUM",
    ],
    "low": ["no threats", "clean", "legitimate", "LOW"],
}

_SEV_RANK = {"critical": 0, "high": 1, "medium": 2, "low": 3, "info": 4}

# ── Helpers ───────────────────────────────────────────────────────────────────

def _read(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace") if path.exists() else ""


def _detect_severity(text: str) -> str:
    t = text.lower()
    for sev, keywords in _SEV_KEYWORDS.items():
        if any(kw.lower() in t for kw in keywords):
            return sev
    return "medium"


def _extract_section(markdown: str, heading: str) -> str:
    pattern = rf"##\s+{re.escape(heading)}(.*?)(?=\n##\s|\Z)"
    m = re.search(pattern, markdown, re.DOTALL | re.IGNORECASE)
    return m.group(1).strip() if m else ""


def _detect_module(case_id: str) -> str:
    p = case_id.upper()
    if p.startswith("FAME"):  return "FAME"
    if p.startswith("FAST"):  return "FAST"
    if p.startswith("FAN"):   return "FAN"
    return "UNKNOWN"


def _extract_hostname(notes_text: str) -> str:
    m = re.search(r"\*\*Hostname:\*\*\s*([^\s|`]+)", notes_text)
    return m.group(1).strip("*`|") if m else "—"


def _first_paragraph(text: str, max_chars: int = 300) -> str:
    """Return the first meaningful paragraph (skip >-quote lines)."""
    lines = []
    for ln in text.splitlines():
        stripped = ln.strip()
        if stripped.startswith(">") or stripped.startswith("---") or not stripped:
            if lines:
                break
            continue
        lines.append(stripped)
    result = " ".join(lines)
    result = re.sub(r"\*\*(.*?)\*\*", r"\1", result)
    result = re.sub(r"`(.*?)`", r"\1", result)
    return result[:max_chars].rstrip(" ,;")


# ── Top 3 Immediate Actions derivation ───────────────────────────────────────

_WHO_HINTS: list[tuple[list[str], str]] = [
    (["isolat", "quarantin", "contain", "disconnect", "network"], "SOC / IR Team"),
    (["creden", "password", "account", "reset", "privilege"],     "IT Admin"),
    (["notify", "escalat", "report", "brief", "ciso", "legal"],   "CISO / Management"),
    (["patch", "update", "remediati", "fix", "mitigat"],          "IT / Systems Admin"),
    (["forensic", "collect", "preserv", "image", "dump"],         "IR Team"),
]

_WHEN_MAP = {"critical": "**Immediately**", "high": "Within 24 hours", "medium": "Within 7 days"}


def _infer_who(text: str) -> str:
    t = text.lower()
    for keywords, role in _WHO_HINTS:
        if any(kw in t for kw in keywords):
            return role
    return "SOC / IT Admin"


def _build_top3_actions(cases: list) -> str:
    """Derive up to 3 immediate actions from the highest-severity cases."""
    seen: set[str] = set()
    actions: list[dict] = []

    for c in sorted(cases, key=lambda x: (_SEV_RANK.get(x.severity, 4), x.case_id)):
        recs = [
            ln.strip().lstrip("•-").strip()
            for ln in c.pptx_recommendations.splitlines()
            if ln.strip().lstrip("•-").strip() and len(ln.strip()) > 10
        ]
        if not recs:
            # Fall back: synthesise a containment action from the key finding
            recs = [f"Investigate and contain findings on {c.hostname}"]

        for rec in recs:
            key = rec.lower()[:60]
            if key in seen:
                continue
            seen.add(key)

            # Split rec into "what" (first sentence) and "how" (remainder)
            parts = re.split(r"[.;]", rec, maxsplit=1)
            what = parts[0].strip()
            how  = parts[1].strip(" .") if len(parts) > 1 and parts[1].strip() else what
            why  = f"{c.key_finding[:100]}{'…' if len(c.key_finding) > 100 else ''} ({c.severity.upper()})"
            who  = _infer_who(rec)
            when = _WHEN_MAP.get(c.severity, "Within 7 days")

            actions.append({"what": what, "why": why, "who": who, "when": when, "how": how})
            if len(actions) == 3:
                break
        if len(actions) == 3:
            break

    if not actions:
        return ""

    lines: list[str] = []
    lines.append("## 1b. Top 3 Immediate Actions")
    lines.append("")
    lines.append("> **Audience:** CISO, IR Lead, IT Admin — act on these now.")
    lines.append("")
    lines.append("| # | What | Why | Who | When | How |")
    lines.append("|---|------|-----|-----|------|-----|")
    for i, a in enumerate(actions, 1):
        what = a["what"].replace("|", "\\|")
        why  = a["why"].replace("|", "\\|")
        who  = a["who"]
        when = a["when"]
        how  = a["how"].replace("|", "\\|")
        lines.append(f"| {i} | {what} | {why} | {who} | {when} | {how} |")
    lines.append("")
    return "\n".join(lines)


# ── Case data loading ─────────────────────────────────────────────────────────

class CaseInfo:
    def __init__(self, case_id: str, reports_dir: Path):
        self.case_id  = case_id
        self.module   = _detect_module(case_id)

        narrative_path = reports_dir / f"{case_id}_narrative.md"
        notes_path     = reports_dir / f"{case_id}_research_notes.md"

        self.narrative_text = _read(narrative_path)
        self.notes_text     = _read(notes_path)

        self.hostname = _extract_hostname(self.notes_text)

        inv_summary_raw = _extract_section(self.notes_text, "Investigation Summary")
        self.inv_summary = " ".join(
            ln.lstrip(">").strip()
            for ln in inv_summary_raw.splitlines()
            if ln.strip() and not ln.strip() == "---"
        )

        self.severity  = _detect_severity(self.inv_summary or self.narrative_text)
        self.key_finding = _first_paragraph(self.inv_summary, 250) or "No summary available."

        # Collect research note steps and attacker events for timelines
        from research_notes import parse_steps, parse_events
        self.steps  = parse_steps(case_id, str(reports_dir))
        self.events = parse_events(case_id, str(reports_dir))

        # First meaningful timestamp from steps
        self.first_ts: datetime | None = None
        for s in self.steps:
            ts_str = s.get("timestamp", "")
            try:
                dt = datetime.strptime(ts_str.replace(" UTC", ""), "%Y-%m-%d %H:%M:%S")
                self.first_ts = dt.replace(tzinfo=timezone.utc)
                break
            except ValueError:
                continue

        # Narrative sections for PPTX
        self.attack_timeline   = _extract_section(self.narrative_text, "attack_timeline")
        self.pptx_risk         = _extract_section(self.narrative_text, "pptx_risk")
        self.pptx_recommendations = _extract_section(self.narrative_text, "pptx_recommendations")


def _load_all_cases(reports_dir: Path) -> list[CaseInfo]:
    notes_files = sorted(reports_dir.glob("*_research_notes.md"))
    cases = []
    for nf in notes_files:
        case_id = nf.name.replace("_research_notes.md", "")
        cases.append(CaseInfo(case_id, reports_dir))
    # Sort: severity (critical first), then module, then case_id
    cases.sort(key=lambda c: (_SEV_RANK.get(c.severity, 4), c.module, c.case_id))
    return cases


# ── Markdown generation ───────────────────────────────────────────────────────

_SEV_EMOJI = {"critical": "🔴", "high": "🟠", "medium": "🟡", "low": "🟢", "info": "⚪"}

def _build_campaign_markdown(
    cases: list[CaseInfo],
    campaign_id: str,
    title: str,
    generated_utc: str,
    reports_dir: Path | None = None,
) -> str:
    lines: list[str] = []
    a = lines.append

    # Header
    a(f"# {title}")
    a("")
    a("| Field | Value |")
    a("|-------|-------|")
    a(f"| Campaign ID | `{campaign_id}` |")
    a(f"| Hosts investigated | {len(cases)} |")
    a(f"| Generated (UTC) | {generated_utc} |")
    a("")

    # Management summary
    critical_cases = [c for c in cases if c.severity == "critical"]
    high_cases     = [c for c in cases if c.severity == "high"]
    clean_cases    = [c for c in cases if c.severity == "low"]

    a("---")
    a("")
    a("## 1. Management Summary")
    a("")
    a("> **Audience:** CISO, Legal, Internal Audit — no technical identifiers.")
    a("")
    a(f"A forensic investigation was conducted across **{len(cases)} systems** "
      f"within the organisation. The investigation spanned memory forensics, "
      f"disk forensics, and network forensics.")
    a("")
    a(f"**{len(critical_cases)} of {len(cases)} systems** returned critical-severity findings "
      f"indicating active compromise. "
      f"{'No' if not high_cases else str(len(high_cases))} additional systems returned "
      f"high-severity anomalies requiring further review. "
      f"{'No systems were assessed as clean.' if not clean_cases else ''}")
    a("")
    a("The findings indicate a coordinated, multi-stage intrusion affecting "
      "both endpoint workstations and critical infrastructure (domain controller, "
      "file server, remote desktop servers). Immediate containment and remediation "
      "action is required.")
    a("")

    # Top 3 Immediate Actions
    top3 = _build_top3_actions(cases)
    if top3:
        a("---")
        a("")
        a(top3)

    # Host inventory table
    a("---")
    a("")
    a("## 2. Host Investigation Summary")
    a("")
    a("| Severity | Host | Module | Key Finding |")
    a("|----------|------|--------|-------------|")
    for c in cases:
        sev_label = c.severity.upper()
        finding   = c.key_finding[:120].replace("|", "\\|") + ("…" if len(c.key_finding) > 120 else "")
        a(f"| **{sev_label}** | `{c.hostname}` | {c.module} | {finding} |")
    a("")

    # Batch synthesis — embed if present
    rpts = reports_dir or REPORTS_DIR
    synthesis_path = rpts / f"{campaign_id}_batch_synthesis.md"
    if synthesis_path.exists():
        synthesis_text = synthesis_path.read_text(encoding="utf-8").strip()
        if synthesis_text:
            a("---")
            a("")
            a("## 2b. Batch Synthesis")
            a("")
            a("> *Cross-case analyst synthesis — patterns, outliers, and open leads spanning the full host set.*")
            a("")
            a(synthesis_text)
            a("")

    # Campaign attack timeline — two tables
    _EXAM_TYPE = {
        "FAME": "Memory image analysis",
        "FAST": "Disk image analysis",
        "FAN":  "Network capture analysis",
    }

    a("---")
    a("")
    a("## 3. Campaign Attack Timeline")
    a("")

    # Table 1 — Investigation Steps (what was done)
    a("### Investigation Steps")
    a("")
    a("| Timestamp | Host | Module | Case ID | Examination Type |")
    a("|-----------|------|--------|---------|-----------------|")
    timed = [(c.first_ts, c) for c in cases if c.first_ts]
    timed.sort(key=lambda x: x[0])
    untimed = [c for c in cases if not c.first_ts]
    for ts, c in timed:
        ts_str = ts.strftime("%Y-%m-%d %H:%M UTC")
        exam   = _EXAM_TYPE.get(c.module, "Forensic analysis")
        a(f"| {ts_str} | `{c.hostname}` | {c.module} | `{c.case_id}` | {exam} |")
    for c in untimed:
        exam = _EXAM_TYPE.get(c.module, "Forensic analysis")
        a(f"| — | `{c.hostname}` | {c.module} | `{c.case_id}` | {exam} |")
    a("")

    # Table 2 — Attacker Timeline (events sorted by evidence timestamp)
    a("### Attacker Timeline")
    a("")
    a("| Timestamp (UTC) | Severity | Host | Module | Event |")
    a("|-----------------|----------|------|--------|-------|")

    all_events: list[tuple] = []
    for c in cases:
        for ev in c.events:
            all_events.append((c, ev))

    def _ev_sort(item: tuple) -> tuple:
        _, ev = item
        ts = ev.get("timestamp", "")
        if ts:
            try:
                dt = datetime.strptime(ts.replace(" UTC", "").strip(), "%Y-%m-%d %H:%M:%S")
                return (0, dt.replace(tzinfo=timezone.utc))
            except ValueError:
                pass
        return (1, datetime.min.replace(tzinfo=timezone.utc))

    all_events.sort(key=_ev_sort)

    for c, ev in all_events:
        ts   = ev.get("timestamp", "") or "—"
        sev  = ev.get("severity", "info").upper()
        desc = ev.get("description", "")[:160].replace("|", "\\|")
        if len(ev.get("description", "")) > 160:
            desc += "…"
        a(f"| {ts} | **{sev}** | `{c.hostname}` | {c.module} | {desc} |")
    a("")

    # Combined risks
    a("---")
    a("")
    a("## 4. Risk Assessment")
    a("")
    all_risks: list[str] = []
    seen_risks: set[str] = set()
    for c in cases:
        for ln in c.pptx_risk.splitlines():
            ln = ln.strip().lstrip("•").strip()
            if ln and ln.lower() not in seen_risks:
                seen_risks.add(ln.lower())
                report_link = f"{c.case_id}_{c.module.lower()}_report.md"
                all_risks.append(f"- ({c.hostname} · [{c.case_id}](./{report_link})) {ln}")
    for r in all_risks[:20]:
        a(r)
    a("")

    # Combined recommendations
    a("---")
    a("")
    a("## 5. Recommendations")
    a("")
    seen_recs: set[str] = set()
    all_recs: list[tuple[str, str, str]] = []  # (text, case_id, module)
    for c in sorted(cases, key=lambda x: _SEV_RANK.get(x.severity, 4)):
        for ln in c.pptx_recommendations.splitlines():
            ln = ln.strip().lstrip("•").strip()
            if ln and ln.lower() not in seen_recs and len(ln) > 10:
                seen_recs.add(ln.lower())
                all_recs.append((ln, c.case_id, c.module.lower()))
    for i, (rec, cid, mod) in enumerate(all_recs[:15], 1):
        report_link = f"{cid}_{mod}_report.md"
        a(f"{i}. {rec} — [{cid}](./{report_link})")
    a("")

    # MITRE coverage (extracted from inv_summary)
    mitre_hits: dict[str, tuple[str, list[str]]] = {}  # tid → (name, [case_ids])
    for c in cases:
        report_link = f"{c.case_id}_{c.module.lower()}_report.md"
        for m in re.finditer(r"(T\d{4}(?:\.\d{3})?)\s*\(([^)]+)\)", c.inv_summary):
            tid, tname = m.group(1), m.group(2).strip()
            if tid not in mitre_hits:
                mitre_hits[tid] = (tname, [])
            link = f"[{c.case_id}](./{report_link})"
            if link not in mitre_hits[tid][1]:
                mitre_hits[tid][1].append(link)

    if mitre_hits:
        a("---")
        a("")
        a("## 6. Campaign MITRE ATT&CK Coverage")
        a("")
        a("| Technique ID | Name | Source |")
        a("|-------------|------|--------|")
        for tid in sorted(mitre_hits):
            tname, sources = mitre_hits[tid]
            a(f"| `{tid}` | {tname} | {', '.join(sources)} |")
        a("")

    # Appendix — case index
    a("---")
    a("")
    a("## Appendix — Case Index")
    a("")
    a("| Case ID | Host | Module | Severity | Report |")
    a("|---------|------|--------|----------|--------|")
    for c in sorted(cases, key=lambda x: x.case_id):
        report_name = f"{c.case_id}_{c.module.lower()}_report.md"
        a(f"| `{c.case_id}` | `{c.hostname}` | {c.module} | {c.severity.upper()} |"
          f" [{report_name}](./{report_name}) |")
    a("")

    return "\n".join(lines)


# ── Timeline PNG for campaign ─────────────────────────────────────────────────

def _build_campaign_timeline_png(cases: list[CaseInfo], output_path: Path) -> Path | None:
    """Generate a campaign-level timeline PNG: one dot per case, coloured by severity."""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import matplotlib.patches as mpatches
    except ImportError:
        print("[campaign] WARNING: matplotlib not available — timeline PNG skipped")
        return None

    _SEV_COLOR = {
        "critical": "#FF5252",
        "high":     "#FF9800",
        "medium":   "#FFEB3B",
        "low":      "#43A047",
        "info":     "#B0BEC5",
    }
    _BG = "#0A1628"

    timed = [(c.first_ts, c) for c in cases if c.first_ts]
    timed.sort(key=lambda x: x[0])

    if not timed:
        return None

    dpi   = 150
    fig_w = 1600 / dpi
    fig_h = 520  / dpi

    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    fig.patch.set_facecolor(_BG)
    ax.set_facecolor(_BG)
    ax.axis("off")

    times = [t for t, _ in timed]
    t_min = times[0]
    t_max = times[-1]
    span_sec = max((t_max - t_min).total_seconds(), 1)
    margin_l, margin_r = 0.05, 0.95

    def _x(dt: datetime) -> float:
        return margin_l + ((dt - t_min).total_seconds() / span_sec) * (margin_r - margin_l)

    spine_y = 0.5
    ax.plot([margin_l - 0.02, margin_r + 0.02], [spine_y, spine_y],
            color="#1565C0", linewidth=2, transform=ax.transAxes, clip_on=False)

    n = len(timed)
    for i, (dt, c) in enumerate(timed):
        xpos  = _x(dt)
        above = (i % 2 == 0)
        color = _SEV_COLOR.get(c.severity, _SEV_COLOR["info"])

        # Tick
        tick_top    = spine_y + 0.2 if above else spine_y
        tick_bottom = spine_y if above else spine_y - 0.2
        ax.plot([xpos, xpos], [tick_bottom, tick_top],
                color=color, linewidth=1.5,
                transform=ax.transAxes, clip_on=False)

        # Dot
        ax.plot(xpos, spine_y, "o", color=color, markersize=7,
                transform=ax.transAxes, clip_on=False, zorder=5)

        # Label: hostname + module
        label_y = spine_y + 0.24 if above else spine_y - 0.32
        label   = f"{c.hostname}\n{c.module}"
        ax.text(xpos, label_y, label,
                ha="center", va="bottom" if above else "top",
                color="#FFFFFF", fontsize=6.5, linespacing=1.3,
                transform=ax.transAxes, clip_on=False)

    # Timestamp endpoints
    def _fmt(dt: datetime) -> str:
        return dt.strftime("%Y-%m-%d\n%H:%M UTC")
    for dt, ypos in [(t_min, _x(t_min)), (t_max, _x(t_max))]:
        ax.text(ypos, spine_y - 0.08, _fmt(dt),
                ha="center", va="top", color="#B0BEC5", fontsize=6,
                linespacing=1.3, transform=ax.transAxes, clip_on=False)

    # Legend
    patches = [mpatches.Patch(color=col, label=sev.upper())
               for sev, col in _SEV_COLOR.items() if sev != "info"]
    legend = ax.legend(handles=patches, loc="lower right",
                       framealpha=0.15, facecolor="#1A2A3A",
                       edgecolor="#2979FF", fontsize=7,
                       labelcolor="#FFFFFF", handlelength=1.2)
    legend.get_frame().set_linewidth(0.5)

    ax.set_title("Campaign Investigation Timeline", color="#FFFFFF",
                 fontsize=10, fontweight="bold", loc="left", pad=6)

    plt.tight_layout(pad=0.3)
    fig.savefig(str(output_path), dpi=dpi, bbox_inches="tight",
                facecolor=_BG, edgecolor="none")
    plt.close(fig)
    return output_path


# ── PPTX board deck ───────────────────────────────────────────────────────────

def _build_campaign_pptx(
    cases: list[CaseInfo],
    campaign_id: str,
    title: str,
    timeline_png: Path | None,
    output_path: Path,
) -> Path:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise SystemExit("[pptx] python-pptx not installed.")

    from generate_pptx_report import (
        _set_bg, _rect, _text, _text_lines, _header_bar, _sev_badge,
        _slide_cover,
        _DARK_NAVY, _BLUE, _ELECTRIC, _WHITE, _LIGHT, _ALERT, _MID_NAVY,
        _LIGHT_BLUE, _AMBER,
    )

    def _rgb(r, g, b): return RGBColor(r, g, b)

    critical_cases = [c for c in cases if c.severity == "critical"]
    high_cases     = [c for c in cases if c.severity == "high"]

    date_str = datetime.now(timezone.utc).strftime("%Y-%m-%d")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W, H = prs.slide_width, prs.slide_height

    # ── Slide 1 — Cover ───────────────────────────────────────────────────────
    _slide_cover(
        prs,
        title=title,
        case_id=campaign_id,
        description=f"{len(cases)} hosts investigated — {len(critical_cases)} critical",
        date_str=date_str,
        W=W, H=H,
    )

    # ── Slide 2 — Executive Summary ───────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, _DARK_NAVY)
    _header_bar(slide, "Executive Summary", campaign_id, W, H)

    CONTENT_TOP = Inches(0.82)
    CONTENT_L   = Inches(0.4)
    CONTENT_W   = Inches(12.5)

    overall_sev = "critical" if critical_cases else ("high" if high_cases else "medium")
    _sev_badge(slide, overall_sev, CONTENT_L, CONTENT_TOP, width=Inches(1.3), height=Inches(0.32))

    sev_line = f"{len(critical_cases)} of {len(cases)} hosts returned CRITICAL findings requiring immediate action."
    _text(slide, sev_line,
          CONTENT_L + Inches(1.45), CONTENT_TOP, CONTENT_W - Inches(1.5), Inches(0.35),
          size=13, bold=True, color=_WHITE)

    # Stat boxes
    BOX_TOP = CONTENT_TOP + Inches(0.45)
    BOX_H   = Inches(0.75)
    box_items = [
        (str(len(cases)),          "Hosts investigated"),
        (str(len(critical_cases)), "Critical severity"),
        (str(len(high_cases)),     "High severity"),
        (str(len(cases) - len(critical_cases) - len(high_cases)), "Lower / clean"),
    ]
    box_w  = Inches(2.9)
    box_gap = Inches(0.2)
    bx = CONTENT_L
    for val, lbl in box_items:
        _rect(slide, bx, BOX_TOP, box_w, BOX_H, fill=_MID_NAVY, line=_BLUE, line_width_pt=0.5)
        _text(slide, val, bx + Inches(0.12), BOX_TOP + Inches(0.04),
              box_w - Inches(0.24), Inches(0.38),
              size=22, bold=True, color=_rgb(*_ELECTRIC))
        _text(slide, lbl, bx + Inches(0.12), BOX_TOP + Inches(0.44),
              box_w - Inches(0.24), Inches(0.26),
              size=10, color=_LIGHT)
        bx += box_w + box_gap

    # Key finding bullets (top 5 critical)
    KF_TOP = BOX_TOP + BOX_H + Inches(0.2)
    _text(slide, "KEY FINDINGS", CONTENT_L, KF_TOP, CONTENT_W, Inches(0.28),
          size=9, bold=True, color=_ELECTRIC)
    _rect(slide, CONTENT_L, KF_TOP + Inches(0.28), CONTENT_W, Inches(0.02), fill=_ELECTRIC)

    bullets = []
    for c in cases[:6]:
        short = c.key_finding[:100] + ("…" if len(c.key_finding) > 100 else "")
        col   = _ALERT if c.severity == "critical" else _WHITE
        bullets.append((f"[{c.severity.upper():<8}] {c.hostname}: {short}", 10, False, col))

    _text_lines(slide, bullets,
                CONTENT_L, KF_TOP + Inches(0.35), CONTENT_W, Inches(3.8))

    # ── Slide 3 — Campaign Timeline ───────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, _DARK_NAVY)
    _header_bar(slide, "Campaign Investigation Timeline", campaign_id, W, H)

    CONTENT_TOP = Inches(0.78)
    IMG_H       = Inches(3.2)

    if timeline_png and Path(timeline_png).exists():
        slide.shapes.add_picture(
            str(timeline_png),
            Inches(0.3), CONTENT_TOP, W - Inches(0.6), IMG_H
        )

    # Host inventory table (below timeline)
    tbl_top = CONTENT_TOP + IMG_H + Inches(0.15)
    tbl_h   = H - tbl_top - Inches(0.2)
    tbl_rows = min(len(cases) + 1, 10)

    tbl = slide.shapes.add_table(
        tbl_rows, 3, Inches(0.3), tbl_top, W - Inches(0.6), tbl_h
    ).table
    tbl.columns[0].width = Inches(3.5)
    tbl.columns[1].width = Inches(1.5)
    tbl.columns[2].width = W - Inches(0.6) - Inches(3.5) - Inches(1.5)

    def _tcell(table_obj, row, col, text, bold=False, size=9, fg=_WHITE, bg=None, align=PP_ALIGN.LEFT):
        cell = table_obj.cell(row, col)
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*bg)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*fg)

    _tcell(tbl, 0, 0, "Host", bold=True, size=8, fg=_WHITE, bg=_BLUE)
    _tcell(tbl, 0, 1, "Module", bold=True, size=8, fg=_WHITE, bg=_BLUE, align=PP_ALIGN.CENTER)
    _tcell(tbl, 0, 2, "Key Finding", bold=True, size=8, fg=_WHITE, bg=_BLUE)

    for i, c in enumerate(cases[:tbl_rows - 1]):
        row = i + 1
        sev_colors = {"critical": _ALERT, "high": (0xFF, 0x98, 0x00),
                      "medium": (0xFF, 0xEB, 0x3B), "low": (0x43, 0xA0, 0x47)}
        sev_col = sev_colors.get(c.severity, _LIGHT)
        bg = _MID_NAVY if i % 2 == 0 else _DARK_NAVY
        _tcell(tbl, row, 0, f"{c.hostname} [{c.severity.upper()}]",
               bold=False, size=8, fg=sev_col, bg=bg)
        _tcell(tbl, row, 1, c.module, bold=False, size=8, fg=_LIGHT, bg=bg, align=PP_ALIGN.CENTER)
        finding = c.key_finding[:80] + ("…" if len(c.key_finding) > 80 else "")
        _tcell(tbl, row, 2, finding, bold=False, size=7.5, fg=_WHITE, bg=bg)

    # ── Slide 4 — Risk & Impact ───────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, _DARK_NAVY)
    _header_bar(slide, "Risk & Impact", campaign_id, W, H)

    seen: set[str] = set()
    risk_bullets: list[str] = []
    for c in cases:
        for ln in c.pptx_risk.splitlines():
            ln = ln.strip().lstrip("•").strip()
            if ln and ln.lower() not in seen and len(ln) > 10:
                seen.add(ln.lower())
                risk_bullets.append(ln)
    risk_bullets = risk_bullets[:10]

    CONTENT_TOP = Inches(0.85)
    CONTENT_L   = Inches(0.4)
    CONTENT_W   = Inches(12.5)
    _text(slide, "CAMPAIGN RISK SUMMARY", CONTENT_L, CONTENT_TOP,
          CONTENT_W, Inches(0.28), size=10, bold=True, color=_ELECTRIC)
    _rect(slide, CONTENT_L, CONTENT_TOP + Inches(0.28), CONTENT_W, Inches(0.02), fill=_ELECTRIC)

    item_h = (H - CONTENT_TOP - Inches(0.45)) / max(len(risk_bullets), 1)
    item_h = min(item_h, Inches(0.65))
    for idx, risk in enumerate(risk_bullets):
        row_top = CONTENT_TOP + Inches(0.35) + idx * item_h
        _rect(slide, CONTENT_L, row_top + Inches(0.13),
              Inches(0.1), Inches(0.1), fill=_ALERT)
        _text(slide, risk,
              CONTENT_L + Inches(0.2), row_top,
              CONTENT_W - Inches(0.25), item_h - Inches(0.04),
              size=11, color=_WHITE)

    # ── Slide 5 — Compromised Host Inventory ──────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, _DARK_NAVY)
    _header_bar(slide, "Compromised Host Inventory", campaign_id, W, H)

    CONTENT_TOP = Inches(0.78)
    tbl_h = H - CONTENT_TOP - Inches(0.2)
    num_rows = min(len(cases) + 1, 20)

    tbl2 = slide.shapes.add_table(
        num_rows, 4, Inches(0.3), CONTENT_TOP, W - Inches(0.6), tbl_h
    ).table
    tbl2.columns[0].width = Inches(3.8)
    tbl2.columns[1].width = Inches(1.3)
    tbl2.columns[2].width = Inches(1.5)
    tbl2.columns[3].width = W - Inches(0.6) - Inches(3.8) - Inches(1.3) - Inches(1.5)

    for col, hdr in enumerate(["Case ID", "Host", "Severity", "Key Finding"]):
        _tcell(tbl2, 0, col, hdr, bold=True, size=8, fg=_WHITE, bg=_BLUE)

    for i, c in enumerate(cases[:num_rows - 1]):
        row = i + 1
        sev_colors = {"critical": _ALERT, "high": (0xFF, 0x98, 0x00),
                      "medium": (0xFF, 0xEB, 0x3B), "low": (0x43, 0xA0, 0x47)}
        sev_col = sev_colors.get(c.severity, _LIGHT)
        bg = _MID_NAVY if i % 2 == 0 else _DARK_NAVY
        _tcell(tbl2, row, 0, c.case_id, bold=False, size=7.5, fg=_LIGHT, bg=bg)
        _tcell(tbl2, row, 1, c.hostname, bold=False, size=8, fg=_WHITE, bg=bg)
        _tcell(tbl2, row, 2, c.severity.upper(), bold=True, size=8, fg=sev_col, bg=bg, align=PP_ALIGN.CENTER)
        finding = c.key_finding[:80] + ("…" if len(c.key_finding) > 80 else "")
        _tcell(tbl2, row, 3, finding, bold=False, size=7, fg=_WHITE, bg=bg)

    # ── Slide 6 — Recommendations ─────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, _DARK_NAVY)
    _header_bar(slide, "Recommendations & Next Steps", campaign_id, W, H)

    seen_recs: set[str] = set()
    all_recs: list[str] = []
    for c in cases:
        for ln in c.pptx_recommendations.splitlines():
            ln = ln.strip().lstrip("•").strip()
            if ln and ln.lower() not in seen_recs and len(ln) > 10:
                seen_recs.add(ln.lower())
                all_recs.append(ln)

    if not all_recs:
        all_recs = ["Review full technical reports and implement findings — CISO / IT"]

    CONTENT_TOP = Inches(0.85)
    CONTENT_L   = Inches(0.4)
    CONTENT_W   = Inches(12.5)
    item_h = (H - CONTENT_TOP - Inches(0.3)) / max(len(all_recs[:8]), 1)
    item_h = min(item_h, Inches(0.72))

    for idx, rec in enumerate(all_recs[:8]):
        row_top = CONTENT_TOP + idx * item_h
        is_immediate = rec.upper().startswith("IMMEDIATE")
        num_color = _ALERT if is_immediate else _BLUE
        _rect(slide, CONTENT_L, row_top + Inches(0.05),
              Inches(0.36), Inches(0.36), fill=num_color)
        _text(slide, str(idx + 1),
              CONTENT_L + Inches(0.04), row_top + Inches(0.06),
              Inches(0.28), Inches(0.28),
              size=12, bold=True, color=_WHITE,
              align=PP_ALIGN.CENTER)
        text_color = _ALERT if is_immediate else _WHITE
        _text(slide, rec, CONTENT_L + Inches(0.48), row_top,
              CONTENT_W - Inches(0.55), item_h - Inches(0.06),
              size=11, bold=is_immediate, color=text_color)
        if idx < min(len(all_recs), 8) - 1:
            _rect(slide, CONTENT_L, row_top + item_h - Inches(0.02),
                  CONTENT_W, Inches(0.01), fill=_MID_NAVY)

    # ── Slide 7 — MITRE Coverage ──────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, _DARK_NAVY)
    _header_bar(slide, "MITRE ATT&CK Coverage — Campaign", campaign_id, W, H)

    mitre_hits: dict[str, str] = {}
    for c in cases:
        for m in re.finditer(r"(T\d{4}(?:\.\d{3})?)\s*\(([^)]+)\)", c.inv_summary):
            tid, tname = m.group(1), m.group(2)
            if tid not in mitre_hits:
                mitre_hits[tid] = tname.strip()

    CONTENT_TOP = Inches(0.82)
    CONTENT_L   = Inches(0.35)

    if not mitre_hits:
        _text(slide, "MITRE techniques extracted from individual case reports.",
              CONTENT_L, CONTENT_TOP + Inches(0.3), Inches(12), Inches(0.5),
              size=12, color=_LIGHT)
    else:
        items  = sorted(mitre_hits.items())
        COLS   = 2
        COL_W  = Inches(6.1)
        ROW_H  = Inches(0.34)
        ipc    = (len(items) + COLS - 1) // COLS
        for idx, (tid, tname) in enumerate(items[:24]):
            col = idx // ipc
            row = idx % ipc
            lx  = CONTENT_L + col * (COL_W + Inches(0.2))
            ly  = CONTENT_TOP + Inches(0.1) + row * ROW_H
            _rect(slide, lx, ly + Inches(0.1), Inches(0.12), Inches(0.12), fill=_ELECTRIC)
            _text(slide, f"{tid} — {tname}",
                  lx + Inches(0.2), ly, COL_W - Inches(0.25), ROW_H - Inches(0.04),
                  size=9.5, color=_WHITE)

        _text(slide, f"{len(mitre_hits)} unique techniques observed across {len(cases)} investigated hosts.",
              CONTENT_L, H - Inches(0.5), Inches(12), Inches(0.35),
              size=9, italic=True, color=_LIGHT)

    prs.save(str(output_path))
    print(f"[campaign] PPTX written: {output_path}")
    return output_path


# ── Main API ──────────────────────────────────────────────────────────────────

def generate(
    campaign_id: str,
    title: str = "",
    reports_dir: Path | None = None,
    output_dir: Path | None = None,
) -> dict[str, Path]:
    rpts_dir = reports_dir or REPORTS_DIR
    out_dir  = output_dir  or rpts_dir

    print(f"[campaign] Loading cases from {rpts_dir} ...")
    cases = _load_all_cases(rpts_dir)
    print(f"[campaign] {len(cases)} cases loaded.")

    if not title:
        title = (
            f"Investigation — {campaign_id}" if len(cases) == 1
            else f"Campaign Investigation — {campaign_id}"
        )

    generated_utc = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M:%S UTC")

    # Markdown
    md_content = _build_campaign_markdown(cases, campaign_id, title, generated_utc, rpts_dir)
    md_path = out_dir / f"CAMPAIGN_{campaign_id}_report.md"
    md_path.write_text(md_content, encoding="utf-8")
    print(f"[campaign] Markdown written: {md_path}")

    # Campaign timeline PNG
    timeline_path = out_dir / f"CAMPAIGN_{campaign_id}_timeline.png"
    timeline_png  = _build_campaign_timeline_png(cases, timeline_path)

    # PDF from markdown
    pdf_path = out_dir / f"CAMPAIGN_{campaign_id}_report.pdf"
    try:
        from md_to_pdf import convert
        convert(
            md_path=md_path,
            output_path=pdf_path,
            title=title,
            case_id=campaign_id,
        )
        print(f"[campaign] PDF written: {pdf_path}")
    except Exception as exc:
        print(f"[campaign] WARNING: PDF generation failed: {exc}")
        pdf_path = None

    # PPTX
    pptx_path = out_dir / f"CAMPAIGN_{campaign_id}_board_deck.pptx"
    _build_campaign_pptx(cases, campaign_id, title, timeline_png, pptx_path)

    return {
        "md":       md_path,
        "pdf":      pdf_path,
        "pptx":     pptx_path,
        "timeline": timeline_png,
    }


# ── CLI ───────────────────────────────────────────────────────────────────────

def _build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        description="Generate a campaign-level report across all investigated cases"
    )
    p.add_argument("--campaign-id",  required=True, metavar="ID",
                   help="Campaign identifier, e.g. SHIELDBASE-2026")
    p.add_argument("--title",        default="",    metavar="TEXT",
                   help="Report title (default: Campaign Investigation — <id>)")
    p.add_argument("--reports-dir",  default=str(REPORTS_DIR), metavar="DIR",
                   help="Directory containing case narrative + research notes files")
    p.add_argument("--output-dir",   default="",    metavar="DIR",
                   help="Output directory (default: same as reports-dir)")
    return p


if __name__ == "__main__":
    args = _build_parser().parse_args()
    out_dir = Path(args.output_dir) if args.output_dir else None
    generate(
        campaign_id  = args.campaign_id,
        title        = args.title,
        reports_dir  = Path(args.reports_dir),
        output_dir   = out_dir,
    )
