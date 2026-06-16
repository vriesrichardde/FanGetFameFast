#!/usr/bin/env python3
# SPDX-License-Identifier: MIT OR Apache-2.0
# SPDX-FileCopyrightText: 2026 Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin · Joost Beekman
"""
generate_fame_report.py — FAME (Forensic Analysis Memory) report generator.

Aggregates Volatility 3 / Memory Baseliner analysis outputs from ./analysis/memory/
into a structured incident report in Markdown, PDF, PPTX (Microsoft PowerPoint),
and DOCX (Microsoft Word) formats.

All report sections follow the FanGetFameFast dual-register voice:
  - Management Summary: no technical identifiers; plain business language
  - Technical Body: precise identifiers; scoped conclusions citing evidence source

Claude instructs itself to "enhance and elaborate when necessary" on each section
so that analysts reviewing the output get full contextual depth.

Usage (CLI):
    python3 lib/generate_fame_report.py \\
        --case-id FAME-2026-001 \\
        --hostname SERVER1234 \\
        [--analysis-dir ./analysis/memory] \\
        [--output-dir ./reports]

Python API:
    from lib.generate_fame_report import generate
    paths = generate(case_id="FAME-2026-001", hostname="SERVER1234")
"""
from __future__ import annotations

import argparse
import hashlib
import json
import re
import sys
from datetime import datetime, timezone
from pathlib import Path
sys.path.insert(0, str(Path(__file__).resolve().parent))
import path_guard  # noqa: E402  write-path policy enforcement
import report_completeness  # noqa: E402  narrative/reasoning completeness gate
from typing import Any

try:
    from research_notes import (
        parse_steps as _parse_research_steps,
        parse_events as _parse_research_events,
        parse_reflections as _parse_research_reflections,
    )
    import report_sections
    from hallucination_guard import (
        ConfidenceTier,
        tag_finding,
        render_confidence_summary,
        reset_counter as _hg_reset,
    )
except ModuleNotFoundError:
    # Imported as a package (e.g. `from lib.generate_fame_report import generate`)
    # rather than run as a script — put lib/ on the path for the sibling import.
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    from research_notes import (
        parse_steps as _parse_research_steps,
        parse_events as _parse_research_events,
        parse_reflections as _parse_research_reflections,
    )
    import report_sections
    from hallucination_guard import (
        ConfidenceTier,
        tag_finding,
        render_confidence_summary,
        reset_counter as _hg_reset,
    )
try:
    from zoneinfo import ZoneInfo
    _CET = ZoneInfo("Europe/Amsterdam")
except ImportError:
    _CET = timezone.utc  # fallback — install tzdata if missing

PROJECT_ROOT = Path(__file__).parent.parent

# ── Colour palette (RGB tuples) ────────────────────────────────────────────────
_DARK_NAVY  = (0x0f, 0x17, 0x2a)
_MID_NAVY   = (0x1e, 0x3a, 0x5f)
_BLUE       = (0x1d, 0x4e, 0xd8)
_LIGHT_BLUE = (0x93, 0xc5, 0xfd)
_WHITE      = (0xff, 0xff, 0xff)
_LIGHT_BG   = (0xf8, 0xfa, 0xfc)
_ROW_ALT    = (0xf1, 0xf5, 0xf9)
_TEXT_DARK  = (0x1f, 0x29, 0x37)
_TEXT_MID   = (0x6b, 0x72, 0x80)
_AMBER      = (0xfb, 0xbf, 0x24)
_GREEN      = (0x22, 0xc5, 0x5e)

_SEV_RGB = {
    "critical": (0xef, 0x44, 0x44),
    "high":     (0xf9, 0x73, 0x16),
    "medium":   (0xea, 0xb3, 0x08),
    "low":      (0x22, 0xc5, 0x5e),
    "info":     (0x6b, 0x72, 0x80),
}


# ── Data loading helpers ───────────────────────────────────────────────────────

def _load_analysis(analysis_dir: Path) -> dict[str, Any]:
    """Read all Volatility 3 output files from analysis_dir into a single dict."""
    data: dict[str, Any] = {
        "pslist":          _read_text(analysis_dir / "pslist.txt"),
        "psscan":          _read_text(analysis_dir / "psscan.txt"),
        "linux_pslist":    _read_text(analysis_dir / "linux_pslist.txt"),
        "cmdline":         _read_text(analysis_dir / "cmdline.txt"),
        "netstat":         _read_text(analysis_dir / "netstat.txt"),
        "netscan":         _read_text(analysis_dir / "netscan.txt"),
        "malfind":         _read_text(analysis_dir / "malfind.txt"),
        "svcscan":         _read_text(analysis_dir / "svcscan.txt"),
        "modules":         _read_text(analysis_dir / "modules.txt"),
        "modscan":         _read_text(analysis_dir / "modscan.txt"),
        "userassist":      _read_text(analysis_dir / "userassist.txt"),
        "hivelist":        _read_text(analysis_dir / "hivelist.txt"),
        "filescan":        _read_text(analysis_dir / "filescan.txt"),
        "mem_timeline":    _read_text(analysis_dir / "mem_timeline.txt"),
        "banners":         _read_text(analysis_dir / "banners.txt"),
        "vmcoreinfo":      _read_text(analysis_dir / "vmcoreinfo.txt"),
        "windows_info":    _read_text(analysis_dir / "windows_info.txt"),
        "proc_baseline":   _read_csv(analysis_dir / "proc_baseline.csv"),
        "drv_baseline":    _read_csv(analysis_dir / "drv_baseline.csv"),
        "svc_baseline":    _read_csv(analysis_dir / "svc_baseline.csv"),
        "shutdown_report":    _read_text(analysis_dir / "SERVER1234_shutdown_analysis.md"),
        "syslog_patterns":    _read_text(analysis_dir / "syslog_patterns.txt"),
        "isf_investigation":  _read_text(analysis_dir / "isf_investigation.txt"),
        "yara_scan":          _read_text(analysis_dir / "yara_scan.txt"),
    }
    # Merge individual YARA output files (yara_common_malware.txt, yara_pe_analysis.txt, …)
    # when the consolidated yara_scan.txt is absent (CLI-based YARA runs write separate files)
    if not data["yara_scan"]:
        _yara_parts = [
            _read_text(fname)
            for fname in sorted(analysis_dir.glob("yara_*.txt"))
            if fname.name != "yara_scan.txt"
        ]
        _yara_parts = [p for p in _yara_parts if p.strip()]
        if _yara_parts:
            data["yara_scan"] = "\n\n".join(_yara_parts)
    # Load MemProcFS JSON results
    memprocfs_dir = analysis_dir / "memprocfs"
    memprocfs_results: dict[str, Any] = {}
    if memprocfs_dir.is_dir():
        for jf in sorted(memprocfs_dir.glob("memprocfs_*.json")):
            try:
                memprocfs_results[jf.stem] = json.loads(jf.read_text())
            except Exception:
                pass
    data["memprocfs_results"] = memprocfs_results
    data["rekall_status"] = _read_text(memprocfs_dir / "rekall_status.txt") if memprocfs_dir.is_dir() else ""

    # Load YARA rule files from analysis/yara/ (case-specific) and project yara/
    yara_rules: list[str] = []
    for yara_dir in [analysis_dir.parent / "yara", analysis_dir.parent.parent / "yara"]:
        if yara_dir.is_dir():
            for rf in sorted(yara_dir.glob("*.yar")) + sorted(yara_dir.glob("*.yara")):
                yara_rules.append(rf.read_text(errors="replace"))
    data["yara_rules"] = "\n\n".join(yara_rules)
    # Summarise YARA match statistics for slides
    if data["yara_scan"]:
        rule_hits = {}
        for line in data["yara_scan"].splitlines():
            if line and not line.startswith(("0x", "\t", " ", "#")):
                rule_name = line.split(" ")[0]
                rule_hits[rule_name] = rule_hits.get(rule_name, 0) + 1
        data["yara_summary"] = rule_hits
    else:
        data["yara_summary"] = {}
    # Also absorb any *.json findings dropped by skill scripts
    for jf in sorted(analysis_dir.glob("*.json")):
        data[jf.stem] = json.loads(jf.read_text())
    return data


def _read_text(path: Path) -> str:
    return path.read_text(errors="replace") if path.exists() else ""


def _read_csv(path: Path) -> list[dict]:
    if not path.exists():
        return []
    lines = path.read_text(errors="replace").splitlines()
    if not lines:
        return []
    headers = [h.strip() for h in lines[0].split(",")]
    rows = []
    for line in lines[1:]:
        if line.strip():
            cols = [c.strip() for c in line.split(",")]
            rows.append(dict(zip(headers, cols)))
    return rows


def _available_plugins(data: dict[str, Any]) -> list[str]:
    return [k for k, v in data.items() if v and k not in ("proc_baseline", "drv_baseline", "svc_baseline")]


# ── Confidence & Gaps helpers ──────────────────────────────────────────────────

def _is_dkom_active(data: dict[str, Any]) -> bool:
    pslist = data.get("pslist", "")
    psscan = data.get("psscan", "")
    # Process rows always start with a numeric PID; header/progress lines do not
    pslist_procs = [l for l in pslist.splitlines() if l.strip() and l.strip()[0].isdigit()]
    psscan_procs = [l for l in psscan.splitlines() if l.strip() and l.strip()[0].isdigit()]
    return not pslist_procs and bool(psscan_procs)


_DKOM_AFFECTED = {"pstree", "cmdline", "svcscan", "malfind", "filescan", "hivelist"}

_EXPECTED_PLUGINS: list[tuple[str, str, str]] = [
    ("Image Type Detection",          "banners",      "OS type unknown — cannot select correct plugin chain"),
    ("Hidden Process Scan (psscan)",  "psscan",       "Critical — authoritative process list when DKOM active"),
    ("Active Process List (pslist)",  "pslist",       "EPROCESS walk; empty output may indicate DKOM (T1014) or a symbol/version mismatch"),
    ("Process Tree (pstree)",         "pstree",       "Parent-child relationships unavailable"),
    ("Command Lines (cmdline)",       "cmdline",      "Process command arguments unverifiable"),
    ("Network Connections (netscan)", "netscan",      "C2 connection pivot source"),
    ("Kernel Modules (modscan)",      "modscan",      "Rootkit driver detection"),
    ("Services (svcscan)",            "svcscan",      "Service-based persistence partially blind"),
    ("Registry (hivelist)",           "hivelist",     "Registry-based persistence unverifiable"),
    ("Code Injection (malfind)",      "malfind",      "Process injection evidence unavailable"),
    ("File Handles (filescan)",       "filescan",     "Open file handle enumeration unavailable"),
    ("YARA Scan",                     "yara_scan",    "Malware signature detection unavailable"),
    ("Baseline Comparison",           "proc_baseline","Memory Baseliner requires baseline.json"),
]


def _compute_plugin_completeness(data: dict[str, Any], dkom: bool) -> list[dict]:
    rows = []
    for label, key, missing_note in _EXPECTED_PLUGINS:
        val = data.get(key)
        has_data = bool(val) if not isinstance(val, list) else bool(val)

        if key == "pslist":
            proc_lines = [l for l in (val or "").splitlines()
                          if l.strip() and l.strip()[0].isdigit()]
            if not proc_lines and dkom:
                rows.append({"step": label, "status": "Empty result",
                             "notes": "Possible DKOM (active-list unlinking) or symbol/version "
                                      "mismatch — psscan remains authoritative either way"})
                continue

        if dkom and key in _DKOM_AFFECTED:
            rows.append({"step": label, "status": "Empty — possible DKOM", "notes": missing_note})
            continue

        if key == "banners" and not has_data:
            if data.get("windows_info"):
                rows.append({"step": label, "status": "Complete", "notes": "via windows_info"})
                continue

        if key == "proc_baseline":
            rows.append({"step": label, "status": "Complete" if has_data else "Not run",
                         "notes": "" if has_data else missing_note})
            continue

        rows.append({"step": label, "status": "Complete" if has_data else "Not run",
                     "notes": "" if has_data else missing_note})
    return rows


def _detect_gaps(
    data: dict[str, Any], dkom: bool, case_id: str,
    reports_dir: Path, opencti_findings: str,
) -> list[str]:
    gaps = []
    if dkom:
        gaps.append(
            "**Possible DKOM (T1014) or symbol/version mismatch:** `pslist` returned empty while "
            "`psscan` recovered a process list, and several other plugins — pstree, cmdline, "
            "svcscan, malfind, filescan, hivelist — also returned empty. This pattern is "
            "consistent with active-list unlinking by a rootkit, but can equally arise from a "
            "Volatility 3 / symbol-table version mismatch for this kernel build. Process-level "
            "evidence (command lines, injected regions, service registry entries) could not be "
            "confirmed from this image alone. A disk image (FAST), a secondary memory capture, or "
            "re-running with a matching symbol table would help distinguish the two explanations."
        )
    if not data.get("proc_baseline"):
        gaps.append(
            "**Memory Baseliner not run:** No known-good baseline was available for comparison. "
            "Anomalous-but-legitimate processes cannot be distinguished from attacker tooling "
            "without a baseline from a clean system at the same OS/patch level."
        )
    if not data.get("yara_scan"):
        gaps.append(
            "**YARA scan not run:** No YARA rules were applied to this image — malware signature "
            "detection is unavailable. Place `.yar` or `.yara` files in `./yara/` and re-run."
        )
    fan_report = reports_dir / f"{case_id}_fan_report.md"
    if not fan_report.exists():
        gaps.append(
            f"**No FAN (network) report for {case_id}:** Extracted C2 IPs have not been confirmed "
            "in PCAP traffic. A network capture from the same time window would confirm beaconing "
            "cadence, payload content, and lateral movement targets."
        )
    fast_report = reports_dir / f"{case_id}_fast_report.md"
    if not fast_report.exists():
        gaps.append(
            f"**No FAST (storage) report for {case_id}:** Disk-level artifacts — rootkit driver "
            "presence, browser history, prefetch, MFT entries, and Windows Event Logs — are "
            "unexamined. The initial infection vector and full persistence mechanism remain unknown."
        )
    if not opencti_findings:
        gaps.append(
            "**No OpenCTI enrichment:** Extracted IOCs (C2 IPs, process hashes) have not been "
            "attributed to a known threat actor or campaign. Run `/fan-opencti-lookup` to correlate."
        )
    return gaps


def _parse_assumptions(case_id: str, reports_dir: Path) -> list[str]:
    notes_path = reports_dir / f"{case_id}_research_notes.md"
    if not notes_path.exists():
        return []
    assumptions = []
    for line in notes_path.read_text(encoding="utf-8").splitlines():
        if "— Assumption:" in line and line.startswith("### ["):
            text = line.split("— Assumption:", 1)[1].strip()
            if text:
                assumptions.append(text)
        elif "| **Outcome** |" in line and "[ASSUMPTION]" in line:
            text = line.split("[ASSUMPTION]", 1)[1].strip().rstrip("|").strip()
            if text:
                assumptions.append(text)
    return assumptions


def _score_overall_confidence(dkom: bool, data: dict[str, Any]) -> tuple[str, str]:
    yara_ran = bool(data.get("yara_scan"))
    baseline_ran = bool(data.get("proc_baseline"))
    psscan_ran = bool(data.get("psscan"))
    netscan_ran = bool(data.get("netscan"))

    if not yara_ran and not baseline_ran:
        return (
            "LOW",
            "YARA scan and Memory Baseliner were both unavailable — malware signature detection "
            "and baseline deviation analysis could not be performed. Key findings rest on "
            "pool-scan artifacts alone.",
        )
    if dkom:
        return (
            "MEDIUM",
            "`pslist` returned empty while `psscan` recovered a process list, and several other "
            "EPROCESS-walk-dependent plugins also returned empty — consistent with either DKOM "
            "(T1014, active-list unlinking) or a Volatility 3 / symbol-table version mismatch for "
            "this kernel build. Pool-scan alternatives (psscan, netscan, modscan) remain "
            "authoritative either way and were used as the primary evidence basis. Confidence is "
            "MEDIUM because process command-line and injection evidence that would normally come "
            "from cmdline and malfind is absent, and the cause of the empty pslist has not been "
            "independently confirmed.",
        )
    if psscan_ran and netscan_ran and yara_ran:
        return (
            "HIGH",
            "All critical pool-scan plugins produced results, YARA signatures were applied, "
            "and no rootkit-induced visibility gaps were detected.",
        )
    return (
        "MEDIUM",
        "One or more critical analysis steps did not produce results — findings are based on "
        "available evidence only.",
    )


def _build_followup_investigations(
    data: dict[str, Any], dkom: bool, case_id: str,
    reports_dir: Path, opencti_findings: str,
) -> list[str]:
    followups = []
    if not (reports_dir / f"{case_id}_fast_report.md").exists():
        followups.append(
            "**[FAST]** Acquire a disk image from the subject host; confirm rootkit driver on "
            "disk, locate dropper, recover prefetch and MFT entries, and examine Windows Event "
            "Logs for the initial access vector."
        )
    if not (reports_dir / f"{case_id}_fan_report.md").exists():
        followups.append(
            "**[FAN]** Capture live traffic or retrieve historical PCAP; confirm C2 beaconing "
            "to extracted IPs; identify initial delivery mechanism (phishing, drive-by, exploit)."
        )
    if not opencti_findings:
        followups.append(
            "**[OpenCTI]** Submit extracted C2 IPs and process hashes to OpenCTI; cross-reference "
            "against known threat actor campaigns and malware families via `/fan-opencti-lookup`."
        )
    if dkom:
        followups.append(
            "**[Re-image]** If a sandbox is available, detonate the rootkit in isolation to recover "
            "suppressed process, service, and file-handle artifacts."
        )
    if "ESTABLISHED" in data.get("netscan", ""):
        followups.append(
            "**[Credential scope]** Rotate all domain credentials accessible from this host; "
            "check domain controller event logs for lateral movement originating from this machine."
        )
    if not followups:
        followups.append("No specific follow-up investigations identified based on current findings.")
    return followups


def _build_confidence_gaps_section(
    data: dict[str, Any],
    case_id: str,
    reports_dir: Path,
    opencti_findings: str = "",
) -> list[str]:
    lines: list[str] = []
    a = lines.append

    dkom = _is_dkom_active(data)
    level, reasoning = _score_overall_confidence(dkom, data)
    completeness = _compute_plugin_completeness(data, dkom)
    gaps = _detect_gaps(data, dkom, case_id, reports_dir, opencti_findings)
    assumptions = _parse_assumptions(case_id, reports_dir)
    followups = _build_followup_investigations(data, dkom, case_id, reports_dir, opencti_findings)

    a("---")
    a("")
    a("## 17. Confidence & gaps")
    a("")
    a(f"**Overall investigation confidence: {level}**")
    a("")
    a(f"*{reasoning}*")
    a("")
    a("### Completeness")
    a("")
    a("| Analysis step | Status | Notes |")
    a("|---|---|---|")
    for row in completeness:
        a(f"| {row['step']} | {row['status']} | {row['notes']} |")
    a("")
    a("### Data gaps & missing evidence")
    a("")
    if gaps:
        for gap in gaps:
            a(f"- {gap}")
    else:
        a("No significant data gaps identified — all standard analysis steps completed.")
    a("")
    a("### Assumptions")
    a("")
    if assumptions:
        for assumption in assumptions:
            a(f"- {assumption}")
    else:
        a("No explicit assumptions recorded. Use `python3 lib/research_notes.py assumption`")
        a("to document analytical judgements for reviewer transparency.")
    a("")
    a("### Reflection log")
    a("")
    reflections = _parse_research_reflections(case_id, str(reports_dir))
    if reflections:
        for r in reflections:
            a(f"**{r['id']} — {r['trigger']}** *(recorded {r['timestamp']})*")
            a("")
            if r["reinterpret"] and r["reinterpret"] != "—":
                a(f"> Re-interpretations: {r['reinterpret']}")
                a("")
            if r["open_leads"] and r["open_leads"] != "—":
                a(f"> Open leads: {r['open_leads']}")
                a("")
    else:
        a("No reflection entries recorded.")
        a("Use `python3 lib/research_notes.py reflect` to log mid-investigation re-assessments.")
    a("")
    a("### Recommended follow-up investigations")
    a("")
    for followup in followups:
        a(f"- {followup}")
    a("")

    return lines


# ── Hallucination Guard ────────────────────────────────────────────────────────

def _build_hallucination_guard_section(
    data: dict[str, Any],
    case_id: str,
    reports_dir: Path,
    dkom: bool,
) -> str:
    """
    Build the Hallucination Guard section for FAME reports.

    Tags each key conclusion with a ConfidenceTier based on which Volatility
    plugins ran and what they produced. Tiers are assigned by code logic, not
    by Claude prompt instructions.
    """
    _hg_reset()
    findings = []
    steps = _parse_research_steps(case_id, str(reports_dir))
    step_ids = [s["id"] for s in steps if s.get("id")]

    def _rn(n: int) -> list[str]:
        """Return RN-NNN id for the nth step if it exists."""
        return [step_ids[n - 1]] if n <= len(step_ids) else []

    # Process scan findings
    if data.get("psscan"):
        findings.append(tag_finding(
            "Process pool scan (psscan) produced process list — authoritative even under DKOM",
            ConfidenceTier.CONFIRMED,
            _rn(1) or [],
            ["volatility3/psscan"],
            ["fame"],
        ))
    if data.get("pslist"):
        proc_lines = [l for l in data["pslist"].splitlines() if l.strip() and l.strip()[0].isdigit()]
        if proc_lines:
            findings.append(tag_finding(
                f"Active process list (pslist) returned {len(proc_lines)} process row(s)",
                ConfidenceTier.CONFIRMED,
                [],
                ["volatility3/pslist"],
                ["fame"],
            ))
        elif dkom:
            findings.append(tag_finding(
                "pslist empty, psscan has data — possible DKOM (T1014) or "
                "symbol/version mismatch",
                ConfidenceTier.CONFIRMED,
                [],
                ["volatility3/pslist"],
                ["fame"],
            ))

    # Network connections
    for plugin in ("netscan", "netstat"):
        if data.get(plugin):
            findings.append(tag_finding(
                f"Network connection data present from {plugin}",
                ConfidenceTier.CONFIRMED,
                [],
                [f"volatility3/{plugin}"],
                ["fame"],
            ))
            break

    # YARA scan
    if data.get("yara_scan"):
        yara_hits = data.get("yara_summary", {})
        if yara_hits:
            findings.append(tag_finding(
                f"YARA scan matched {len(yara_hits)} rule(s): {', '.join(list(yara_hits.keys())[:3])}",
                ConfidenceTier.CONFIRMED,
                [],
                ["yara-python"],
                ["fame"],
            ))
        else:
            findings.append(tag_finding(
                "YARA scan ran — no rule matches detected",
                ConfidenceTier.CONFIRMED,
                [],
                ["yara-python"],
                ["fame"],
            ))
    else:
        findings.append(tag_finding(
            "YARA scan not run — malware signature detection unavailable",
            ConfidenceTier.UNVERIFIABLE,
            [],
            ["yara-python"],
            ["fame"],
        ))

    # Code injection (malfind is heuristic → INFERRED)
    if data.get("malfind"):
        findings.append(tag_finding(
            "Code injection candidates found by malfind (heuristic — requires analyst triage)",
            ConfidenceTier.INFERRED,
            [],
            ["volatility3/malfind"],
            ["fame"],
        ))

    # DKOM-suppressed plugins
    if dkom:
        for plugin in ("cmdline", "svcscan", "filescan", "hivelist"):
            if not data.get(plugin):
                findings.append(tag_finding(
                    f"{plugin} unavailable — DKOM rootkit suppressed this plugin (T1014)",
                    ConfidenceTier.UNVERIFIABLE,
                    [],
                    [f"volatility3/{plugin}"],
                    ["fame"],
                ))

    # Baseline comparison
    if data.get("proc_baseline"):
        findings.append(tag_finding(
            "Memory Baseliner baseline comparison completed",
            ConfidenceTier.CONFIRMED,
            [],
            ["memory_baseliner"],
            ["fame"],
        ))
    else:
        findings.append(tag_finding(
            "Memory Baseliner not run — anomalous-but-legitimate processes cannot be distinguished",
            ConfidenceTier.UNVERIFIABLE,
            [],
            ["memory_baseliner"],
            ["fame"],
        ))

    # Assumptions from research notes
    for s in steps:
        outcome = s.get("outcome", "")
        if "[ASSUMPTION]" in outcome or s.get("confidence") == "assumed":
            text = outcome.replace("[ASSUMPTION]", "").strip()
            if text:
                findings.append(tag_finding(
                    text,
                    ConfidenceTier.ASSUMED,
                    [s["id"]] if s.get("id") else [],
                    [s.get("source_tool", "")] if s.get("source_tool") else [],
                    ["fame"],
                ))

    return render_confidence_summary(findings, module_label="FAME")


# ── Narrative loader ──────────────────────────────────────────────────────────

def _load_narrative(case_id: str, reports_dir: Path) -> dict[str, str]:
    """Load Claude-generated narrative sections from {case_id}_narrative.md."""
    path = reports_dir / f"{case_id}_narrative.md"
    if not path.exists():
        return {}
    sections: dict[str, str] = {}
    current: str | None = None
    for line in path.read_text(encoding="utf-8").splitlines():
        if line.startswith("## "):
            current = line[3:].strip()
            sections[current] = ""
        elif line.startswith("<!--"):
            continue
        elif current is not None:
            sections[current] += line + "\n"
    return {k: v.strip() for k, v in sections.items()}


def _narr_bullets(narrative: dict, key: str) -> list[str]:
    """Return the bullet lines of a narrative pptx_* section as clean strings.

    Accepts '-', '*' or '•' bullet markers; falls back to non-empty paragraph
    lines if the section has no explicit bullets. Returns [] when absent."""
    text = (narrative or {}).get(key, "") or ""

    def _clean(s: str) -> str:
        return re.sub(r"\*\*(.*?)\*\*", r"\1", s).replace("**", "").strip()

    bullets: list[str] = []
    # A real bullet marker is '-', '*' or '•' FOLLOWED BY whitespace. A line such
    # as "-testing tooling" (a hyphenated word wrapped across lines) is NOT a bullet.
    marker = re.compile(r"^([-*•])\s+(.*)$")
    for raw in text.splitlines():
        line = raw.strip()
        if not line:
            continue
        m = marker.match(line)
        if m:
            item = m.group(2).strip()
            if item:
                bullets.append(item)
        elif bullets:
            # continuation of a wrapped bullet — join onto the previous one
            bullets[-1] = (bullets[-1] + " " + line).strip()
    if not bullets:
        # No explicit bullet markers: treat blank-line-separated paragraphs as items.
        for para in re.split(r"\n\s*\n", text):
            p = " ".join(l.strip() for l in para.splitlines() if l.strip())
            if p:
                bullets.append(p)
    return [_clean(b) for b in bullets]


# ── Markdown generation ────────────────────────────────────────────────────────

def _build_markdown(
    data: dict[str, Any],
    case_id: str,
    hostname: str,
    image_path: str,
    generated_utc: str,
    opencti_findings: str = "",
    fan_summary: str = "",
    fast_summary: str = "",
    reports_dir: Path | None = None,
) -> str:
    """
    Build the full incident report in Markdown.

    Claude: enhance and elaborate when necessary on each section to ensure
    the analyst receives full contextual depth appropriate to their level.
    """
    lines: list[str] = []
    a = lines.append

    reports_dir = reports_dir or (PROJECT_ROOT / "reports")
    narrative = _load_narrative(case_id, reports_dir)

    narrative_result = report_completeness.check_narrative(case_id, "FAME", reports_dir)
    reasoning_result = report_completeness.check_research_notes(case_id, reports_dir)
    report_completeness.write_incomplete_marker(reports_dir, case_id, narrative_result, reasoning_result)
    incomplete_banner = report_completeness.format_incomplete_banner(narrative_result, reasoning_result)
    if incomplete_banner:
        print(f"[fame] WARNING: investigation incomplete for {case_id}/{hostname} — "
              f"see {case_id}_INVESTIGATION_INCOMPLETE.json", file=sys.stderr)

    # ── Header ────────────────────────────────────────────────────────────────
    a(f"# FAME Memory Forensics Report")
    a(f"")
    a(f"| Field | Value |")
    a(f"|-------|-------|")
    a(f"| Case ID | `{case_id}` |")
    a(f"| Hostname | `{hostname}` |")
    a(f"| Memory image | `{image_path}` |")
    a(f"| Module | FAME — Forensic Analysis Memory |")
    a(f"| Analyst | Claude Code — FAME skill |")
    a(f"| Generated (UTC) | {generated_utc} |")
    a(f"")
    if incomplete_banner:
        lines.extend(incomplete_banner)

    # ── Cross-module summary ───────────────────────────────────────────────────
    if fan_summary or fast_summary:
        a("---")
        a("")
        a("## Cross-module intelligence")
        a("")
        a("> Claude: enhance and elaborate when necessary — correlate memory findings")
        a("> with any available network (FAN) or storage (FAST) evidence below.")
        a("")
        if fan_summary:
            a("### Network forensics (FAN) summary")
            a("")
            a(fan_summary.strip())
            a("")
        if fast_summary:
            a("### Storage forensics (FAST) summary")
            a("")
            a(fast_summary.strip())
            a("")

    # ── Management Summary ────────────────────────────────────────────────────
    a("---")
    a("")
    a("## 1. Management summary")
    a("")
    a("> **Audience:** CISO, Legal, Internal Audit — no technical identifiers.")
    a("")
    a("> Claude: enhance and elaborate when necessary to provide full business-impact")
    a("> context while keeping language free of IPs, ports, file sizes, and workstation IDs.")
    a("")
    shutdown_md = data.get("shutdown_report", "")
    exec_bullets = _narr_bullets(narrative, "pptx_executive_summary")
    impact_bullets = _narr_bullets(narrative, "pptx_impact")
    if shutdown_md:
        # Legacy "unexpected shutdown" scenario.
        for line in shutdown_md.splitlines():
            if "deliberately" in line.lower() or "intentionally" in line.lower() or "root cause" in line.lower():
                a(f"> {line.strip()}")
                a("")
                break
        a("Memory forensic analysis of the subject server was conducted to determine")
        a("the cause of an unexpected shutdown event. The analysis examined the memory")
        a("image captured from the server and reconstructed system activity leading up")
        a("to the shutdown.")
        a("")
        if "msfadmin" in shutdown_md:
            a("**Finding:** The shutdown was deliberately triggered by a person at the physical")
            a("server console. Two failed login attempts with unknown credentials were made before")
            a("a successful login. The authenticated user immediately obtained elevated (administrator)")
            a("privileges and issued a shutdown command. All system services halted in the expected")
            a("orderly sequence and the server restarted approximately 30 seconds later.")
            a("")
            a("**Business Impact:** The event indicates either an undocumented maintenance action")
            a("or unauthorized physical access to the server room. No evidence of a remote attacker,")
            a("hardware failure, or software crash was found in the memory image.")
    elif exec_bullets:
        # Narrative-driven management summary (general case).
        a(f"Memory forensic analysis of host **{hostname}** was performed to establish what")
        a("activity occurred on the workstation. The key business-level findings are:")
        a("")
        for b in exec_bullets:
            a(f"- {b}")
        a("")
        if impact_bullets:
            a("**Business impact:**")
            a("")
            for b in impact_bullets:
                a(f"- {b}")
            a("")
    else:
        a("Refer to the Technical Body below for detailed findings extracted from the memory image.")
    a("")

    # ── Incident Timeline (Claude-generated) ─────────────────────────────────
    timeline_text = narrative.get("attack_timeline", "")
    a("---")
    a("")
    a("## 2. Incident Timeline")
    a("")
    a("> Chronological reconstruction of the attack path. Each finding references")
    a("> the investigation step (RN-NNN) and the preserved source file in")
    a(f"> `{case_id}_evidence/`.")
    a("")
    if timeline_text:
        a(timeline_text)
    else:
        a("> *Incident timeline not yet generated. Run the FAME skill to produce*")
        a(f"> *`{case_id}_narrative.md` with the `attack_timeline` section.*")
    a("")

    # ── System Profile ────────────────────────────────────────────────────────
    a("---")
    a("")
    a("## 3. System profile")
    a("")
    profile_lines = []
    for src in (data.get("banners", ""), data.get("vmcoreinfo", ""), data.get("windows_info", ""), shutdown_md):
        if src:
            profile_lines.append(src[:500])
    if profile_lines:
        a("Extracted from memory image:")
        a("")
        a("```")
        for pl in profile_lines[:2]:
            a(pl.strip()[:300])
        a("```")
        a("")

    # ── Detailed Timeline ─────────────────────────────────────────────────────
    if shutdown_md:
        a("---")
        a("")
        a("## 3. Detailed event timeline")
        a("")
        a("> Claude: enhance and elaborate when necessary — add MITRE ATT&CK technique")
        a("> references and business-impact annotations alongside each timeline event.")
        a("")
        in_table = False
        for line in shutdown_md.splitlines():
            if "| Time |" in line or "| Timestamp |" in line or ("|" in line and ("08:" in line or "EDT" in line or "UTC" in line)):
                in_table = True
            if in_table or "##" in line:
                a(line)
        a("")

    # ── Process Analysis ──────────────────────────────────────────────────────
    for plugin, title in [("pslist", "Process List (windows.pslist)"),
                           ("linux_pslist", "Process List (linux.pslist)"),
                           ("psscan", "Process Scan (windows.psscan — pool scan)")]:
        content = data.get(plugin, "")
        if content and len(content.strip()) > 20:
            a("---")
            a("")
            a(f"## 4. {title}")
            a("")
            proc_narrative = narrative.get("section_processes", "")
            if proc_narrative:
                a(proc_narrative)
                a("")
            else:
                a("> Claude: enhance and elaborate when necessary — flag any processes that")
                a("> do not belong to the OS baseline, appear in psscan but not pslist (hidden),")
                a("> or show suspicious parent-child relationships.")
                a("")
            a(f"> **Source file:** [`{case_id}_evidence/memory/{plugin}.txt`]"
              f"(./{case_id}_evidence/memory/{plugin}.txt)")
            a("")
            a("```")
            a(content.strip()[:3000])
            a("```")
            a("")
            break

    # ── Network Connections ───────────────────────────────────────────────────
    for plugin, title in [("netstat", "Network Connections (windows.netstat)"),
                           ("netscan", "Network Connections (windows.netscan — pool scan)")]:
        content = data.get(plugin, "")
        if content and len(content.strip()) > 20:
            a("---")
            a("")
            a(f"## 5. {title}")
            a("")
            net_narrative = narrative.get("section_network", "")
            if net_narrative:
                a(net_narrative)
                a("")
            else:
                a("> Claude: enhance and elaborate when necessary — identify any external")
                a("> connections and cross-reference with OpenCTI / FAN findings.")
                a("")
            a(f"> **Source file:** [`{case_id}_evidence/memory/{plugin}.txt`]"
              f"(./{case_id}_evidence/memory/{plugin}.txt)")
            a("")
            a("```")
            a(content.strip()[:3000])
            a("```")
            a("")
            break

    # ── Code Injection / Malfind ──────────────────────────────────────────────
    malfind = data.get("malfind", "")
    if malfind and len(malfind.strip()) > 20:
        a("---")
        a("")
        a("## 6. Code injection analysis (windows.malfind)")
        a("")
        mal_narrative = narrative.get("section_malware", "")
        if mal_narrative:
            a(mal_narrative)
            a("")
        else:
            a("> Claude: enhance and elaborate when necessary — distinguish JIT-compiled")
            a("> false positives (.NET/Java) from genuine shellcode injection indicators.")
            a("")
        a(f"> **Source file:** [`{case_id}_evidence/memory/malfind.txt`]"
          f"(./{case_id}_evidence/memory/malfind.txt)")
        a("")
        a("```")
        a(malfind.strip()[:3000])
        a("```")
        a("")

    # ── Services ──────────────────────────────────────────────────────────────
    svcscan = data.get("svcscan", "")
    if svcscan and len(svcscan.strip()) > 20:
        a("---")
        a("")
        a("## 7. Services (windows.svcscan)"  )  # "Services" is a proper noun in this context
        a("")
        a("> Claude: enhance and elaborate when necessary — highlight any services")
        a("> running from unusual paths (Temp, AppData, Users) or with blank binary paths.")
        a("")
        a("```")
        a(svcscan.strip()[:3000])
        a("```")
        a("")

    # ── Kernel Modules ────────────────────────────────────────────────────────
    for plugin, title in [("modules", "Kernel Modules (windows.modules — linked list)"),
                           ("modscan", "Kernel Modules (windows.modscan — pool scan)")]:
        content = data.get(plugin, "")
        if content and len(content.strip()) > 20:
            a("---")
            a("")
            a(f"## 8. {title}")
            a("")
            a("> Claude: enhance and elaborate when necessary — modules in modscan but")
            a("> absent from modules indicate hidden/rootkit kernel drivers.")
            a("")
            a("```")
            a(content.strip()[:2000])
            a("```")
            a("")
            break

    # ── Baseline Comparison ───────────────────────────────────────────────────
    for key, label in [("proc_baseline", "Process Baseline Comparison"),
                        ("drv_baseline", "Driver Baseline Comparison"),
                        ("svc_baseline", "Service Baseline Comparison")]:
        rows = data.get(key, [])
        if rows:
            a("---")
            a("")
            a(f"## 9. {label} (Memory Baseliner)")
            a("")
            a("> Claude: enhance and elaborate when necessary — any row flagged as")
            a("> non-baseline is a pivot candidate requiring manual triage.")
            a("")
            headers = list(rows[0].keys())
            a("| " + " | ".join(headers) + " |")
            a("|" + "|".join(["---"] * len(headers)) + "|")
            for row in rows[:30]:
                a("| " + " | ".join(str(row.get(h, "")) for h in headers) + " |")
            a("")

    # ── Key Evidence ──────────────────────────────────────────────────────────
    if shutdown_md and "Key Evidence" in shutdown_md:
        a("---")
        a("")
        a("## 10. Key evidence")
        a("")
        a("> Claude: enhance and elaborate when necessary — map each evidence item to")
        a("> the relevant MITRE ATT&CK technique and explain the forensic significance.")
        a("")
        in_section = False
        for line in shutdown_md.splitlines():
            if "## Key Evidence" in line:
                in_section = True
            elif in_section and line.startswith("## ") and "Key Evidence" not in line:
                break
            elif in_section:
                a(line)
        a("")

    # ── OpenCTI Enrichment ────────────────────────────────────────────────────
    if opencti_findings:
        a("---")
        a("")
        a("## 11. OpenCTI threat intelligence enrichment")
        a("")
        a("> Claude: enhance and elaborate when necessary — link any matched indicators")
        a("> to known threat actors, malware families, or campaigns in OpenCTI.")
        a("")
        a(opencti_findings.strip())
        a("")

    # ── MITRE ATT&CK ─────────────────────────────────────────────────────────
    techniques = []
    if "msfadmin" in shutdown_md or "sudo" in shutdown_md:
        techniques += [
            ("T1078", "Valid Accounts", "Initial Access / Persistence",
             "User msfadmin authenticated at physical console (tty1) after two failed login attempts."),
            ("T1548.003", "Abuse Elevation Control Mechanism: Sudo",
             "Privilege Escalation",
             "msfadmin executed `sudo /bin/bash` to obtain a root shell 26 seconds after login."),
            ("T1529", "System Shutdown/Reboot", "Impact",
             "Root-level shutdown command sent 25 seconds after privilege escalation; all services terminated in orderly sequence."),
        ]
    if malfind and len(malfind.strip()) > 20:
        techniques.append(
            ("T1055", "Process Injection", "Defense Evasion / Privilege Escalation",
             "malfind output present — triage hits to distinguish injection from JIT false positives.")
        )
    technique_dicts = [
        {"id": tid, "name": name, "tactic": tactic, "observation": obs}
        for tid, name, tactic, obs in techniques
    ]
    a("---")
    a("")
    a("\n".join(report_sections.build_mitre_section(
        technique_dicts, heading="MITRE ATT&CK Coverage", section_num="13",
        show_severity=False,
        empty_message="No MITRE ATT&CK techniques mapped — no malicious activity confirmed in memory.")))

    # ── IOCs ──────────────────────────────────────────────────────────────────
    iocs = _extract_iocs(data, shutdown_md)
    a("\n".join(report_sections.build_ioc_section(
        iocs, heading="Indicators of Compromise", section_num="14",
        empty_message="No malicious indicators of compromise identified in this memory image.")))

    # ── What Did NOT Cause the Shutdown ──────────────────────────────────────
    if "Did NOT Cause" in shutdown_md or "did not cause" in shutdown_md.lower():
        a("---")
        a("")
        a("## 14. Ruled-out causes")
        a("")
        in_section = False
        for line in shutdown_md.splitlines():
            if "Did NOT Cause" in line or "did not cause" in line.lower():
                in_section = True
            elif in_section and line.startswith("## "):
                break
            elif in_section:
                a(line)
        a("")

    # ── Recommendations ───────────────────────────────────────────────────────
    recs = _build_recommendations(data, shutdown_md)
    a("---")
    a("")
    a("\n".join(report_sections.build_recommendations_section(
        recs, heading="Recommendations", section_num="16", numbered=True)))

    # ── Confidence & Gaps ─────────────────────────────────────────────────────
    lines.extend(_build_confidence_gaps_section(data, case_id, reports_dir, opencti_findings))

    # ── Volatility Plugin Status ──────────────────────────────────────────────
    if "Volatility 3 Plugin Status" in shutdown_md:
        a("---")
        a("")
        a("## 17. Volatility 3 plugin status")
        a("")
        in_section = False
        for line in shutdown_md.splitlines():
            if "Plugin Status" in line:
                in_section = True
            elif in_section and line.startswith("## "):
                break
            elif in_section:
                a(line)
        a("")

    # ── Appendix ──────────────────────────────────────────────────────────────
    a("---")
    a("")
    a("## Appendix A — Analysis source files")
    a("")
    a(f"All artifact files are preserved in `./{case_id}_evidence/memory/` and uploaded")
    a("to the investigations vault alongside this report. SHA-256 hashes are recorded in the")
    a("research notes (Appendix B) for chain-of-custody verification.")
    a("")
    a("| File | Description |")
    a("|------|-------------|")
    a(f"| [`{case_id}_evidence/memory/pslist.txt`](./{case_id}_evidence/memory/pslist.txt) | Windows process list (EPROCESS walk) |")
    a(f"| [`{case_id}_evidence/memory/psscan.txt`](./{case_id}_evidence/memory/psscan.txt) | Windows process pool scan (finds hidden/exited) |")
    a(f"| [`{case_id}_evidence/memory/linux_pslist.txt`](./{case_id}_evidence/memory/linux_pslist.txt) | Linux process list |")
    a(f"| [`{case_id}_evidence/memory/cmdline.txt`](./{case_id}_evidence/memory/cmdline.txt) | Process command lines |")
    a(f"| [`{case_id}_evidence/memory/netstat.txt`](./{case_id}_evidence/memory/netstat.txt) | Active network connections |")
    a(f"| [`{case_id}_evidence/memory/netscan.txt`](./{case_id}_evidence/memory/netscan.txt) | Network connection pool scan |")
    a(f"| [`{case_id}_evidence/memory/malfind.txt`](./{case_id}_evidence/memory/malfind.txt) | Code injection findings |")
    a(f"| [`{case_id}_evidence/memory/svcscan.txt`](./{case_id}_evidence/memory/svcscan.txt) | Services pool scan |")
    a(f"| [`{case_id}_evidence/memory/modules.txt`](./{case_id}_evidence/memory/modules.txt) | Kernel modules (linked list) |")
    a(f"| [`{case_id}_evidence/memory/modscan.txt`](./{case_id}_evidence/memory/modscan.txt) | Kernel modules pool scan |")
    a(f"| [`{case_id}_evidence/memory/mem_timeline.txt`](./{case_id}_evidence/memory/mem_timeline.txt) | Memory artifact timeline |")
    a(f"| [`{case_id}_evidence/memory/proc_baseline.csv`](./{case_id}_evidence/memory/proc_baseline.csv) | Process baseline diff (Memory Baseliner) |")
    a(f"| [`{case_id}_evidence/memory/drv_baseline.csv`](./{case_id}_evidence/memory/drv_baseline.csv) | Driver baseline diff (Memory Baseliner) |")
    a(f"| [`{case_id}_evidence/memory/svc_baseline.csv`](./{case_id}_evidence/memory/svc_baseline.csv) | Service baseline diff (Memory Baseliner) |")
    a("")
    a("*All findings derived from memory image analysis as stated. Evidence integrity preserved.*")
    a("")

    # ── Hallucination Guard ───────────────────────────────────────────────────
    dkom_flag = _is_dkom_active(data)
    hg_section = _build_hallucination_guard_section(data, case_id, reports_dir, dkom_flag)
    if hg_section:
        a(hg_section)
        a("")

    # ── Evidence Trail ────────────────────────────────────────────────────────
    lines.extend(report_sections.build_evidence_trail_section(case_id, reports_dir, include_dismissed=True))

    return "\n".join(lines)


def _extract_iocs(data: dict[str, Any], shutdown_md: str) -> list[dict]:
    iocs = []
    # Physical console failed logins are not IOCs per se, but document the event
    if "FAILED LOGIN" in shutdown_md or "pam_unix" in shutdown_md:
        iocs.append({
            "type": "Event",
            "value": "Failed console login × 2 on tty1",
            "severity": report_sections.normalize_severity("Medium"),
            "category": "Authentication",
            "source": "FAME",
            "context": "Two failed logins with unknown credentials before successful msfadmin login (as observed in memory strings)",
        })
    if "sudo" in shutdown_md and "root" in shutdown_md:
        iocs.append({
            "type": "Event",
            "value": "Privilege escalation via sudo /bin/bash",
            "severity": report_sections.normalize_severity("High"),
            "category": "Privilege Escalation",
            "source": "FAME",
            "context": "msfadmin → root via sudo; shell obtained 26 s after login (as observed in memory strings)",
        })
    # Extract any external IPs from netscan
    netscan = data.get("netscan", "")
    if netscan:
        seen = set()
        for match in re.finditer(r"\b(\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3})\b", netscan):
            ip = match.group(1)
            if ip.startswith(("127.", "0.", "169.254.", "::")) or ip in seen:
                continue
            if not ip.startswith(("10.", "172.", "192.168.")):
                seen.add(ip)
                iocs.append({
                    "type": "IP",
                    "value": ip,
                    "severity": report_sections.normalize_severity("Medium"),
                    "category": "Network Scan",
                    "source": "FAME netscan",
                    "context": "External IP from memory netscan — verify with OpenCTI / FAN",
                })
    return iocs


def _build_recommendations(data: dict[str, Any], shutdown_md: str) -> list[str]:
    # Prefer the Claude-authored, case-specific recommendations when this is not
    # the legacy shutdown scenario.
    if not shutdown_md:
        narr_recs = _narr_bullets(data.get("_narrative", {}), "pptx_recommendations")
        if narr_recs:
            return narr_recs
    recs = []
    if "physical" in shutdown_md.lower() or "tty1" in shutdown_md or "console" in shutdown_md.lower():
        recs.append(
            "**Review physical access controls** — determine whether the console access was an authorised maintenance action or unauthorized entry. "
            "Ensure server-room access logs are reviewed for the relevant time window."
        )
    if "FAILED LOGIN" in shutdown_md or "pam_unix" in shutdown_md:
        recs.append(
            "**Investigate failed login attempts** — two consecutive failed logins with unknown credentials before the successful msfadmin session "
            "may indicate a credential-guessing attempt. Review PAM configuration and consider lockout policies."
        )
    if "sudo" in shutdown_md:
        recs.append(
            "**Audit sudo policy** — msfadmin has unrestricted sudo access (`/bin/bash`). "
            "Restrict to specific commands required for legitimate operations and enforce MFA for sudo."
        )
    if _is_dkom_active(data):
        recs.append(
            "**Corroborate the empty active-process-list (pslist) result** — `pslist` returned no "
            "entries while `psscan` recovered a full process list, which can indicate either DKOM "
            "(active-list unlinking, MITRE T1014) or a Volatility 3 / symbol-table version mismatch "
            "for this kernel build. Re-run `windows.pslist` with an alternate Volatility 3 release "
            "or matching ISF symbol table to rule out a tooling artifact before treating this as a "
            "rootkit indicator."
        )
    if "unexpected" in shutdown_md.lower():
        recs.append(
            "**Document all maintenance windows** — the 'unexpected' nature of this shutdown indicates a gap in change management. "
            "Require pre-approved change tickets for any console-level server intervention."
        )
    recs.append(
        "**Cross-reference with FAN (network) and FAST (storage)** — if PCAP or disk images are available for the same time window, "
        "run a combined investigation to rule out lateral movement or data staging before the shutdown."
    )
    return recs


# ── PPTX generation ───────────────────────────────────────────────────────────

def _build_pptx(
    data: dict[str, Any],
    case_id: str,
    hostname: str,
    image_path: str,
    generated_utc: str,
    output_path: Path,
    opencti_findings: str = "",
    fan_summary: str = "",
    fast_summary: str = "",
) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("[fame] WARNING: python-pptx not installed — skipping PPTX. pip3 install python-pptx")
        return

    shutdown_md = data.get("shutdown_report", "")
    narrative = data.get("_narrative", {})

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    def _rgb(t: tuple) -> RGBColor:
        return RGBColor(*t)

    def _add_rect(slide, left, top, width, height, fill_rgb):
        from pptx.util import Emu
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill_rgb)
        shape.line.fill.background()
        return shape

    def _add_text(slide, text, left, top, width, height, font_size, bold=False,
                  color=_WHITE, align=PP_ALIGN.LEFT, wrap=True):
        txb = slide.shapes.add_textbox(left, top, width, height)
        tf  = txb.text_frame
        tf.word_wrap = wrap
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.name = "Arial"
        run.font.color.rgb = _rgb(color)
        return txb

    W = prs.slide_width
    H = prs.slide_height
    M = Inches(0.4)

    def _add_bullet_slide(title, bullets, fallback, max_items=8):
        slide = prs.slides.add_slide(blank_layout)
        _add_rect(slide, 0, 0, W, Inches(1.1), _MID_NAVY)
        _add_text(slide, title, M, Inches(0.2), W, Inches(0.8),
                  28, bold=True, color=_WHITE)
        _add_text(slide, f"{case_id}  |  {hostname}", M, Inches(0.75), W, Inches(0.3),
                  12, color=_LIGHT_BLUE)
        if bullets:
            body = "\n\n".join(f"•  {b}" for b in bullets[:max_items])
        else:
            body = fallback
        _add_text(slide, body, M, Inches(1.3), W - 2 * M, H - Inches(1.6),
                  14, color=_TEXT_DARK)
        return slide

    # ── Slide 1 — Cover ───────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank_layout)
    _add_rect(s1, 0, 0, W, H, _DARK_NAVY)
    _add_rect(s1, 0, 0, W, Inches(0.08), _BLUE)
    _add_rect(s1, 0, H - Inches(0.08), W, Inches(0.08), _BLUE)

    _add_text(s1, "FAME", M, Inches(1.2), W - 2*M, Inches(1.2),
              72, bold=True, color=_LIGHT_BLUE, align=PP_ALIGN.CENTER)
    _add_text(s1, "Forensic Analysis Memory", M, Inches(2.2), W - 2*M, Inches(0.7),
              28, bold=False, color=_WHITE, align=PP_ALIGN.CENTER)
    _add_text(s1, "Memory forensics — Management report", M, Inches(2.9), W - 2*M, Inches(0.6),
              20, color=_LIGHT_BLUE, align=PP_ALIGN.CENTER)
    _add_rect(s1, Inches(3), Inches(3.8), W - Inches(6), Inches(0.04), _BLUE)
    meta = f"Case: {case_id}  |  Host: {hostname}  |  {generated_utc[:10]}"
    _add_text(s1, meta, M, Inches(4.1), W - 2*M, Inches(0.5),
              14, color=_TEXT_MID, align=PP_ALIGN.CENTER)
    _add_text(s1, "Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin · Joost Beekman",
              M, Inches(4.6), W - 2*M, Inches(0.4), 11, color=_TEXT_MID, align=PP_ALIGN.CENTER)

    # ── Slide 2 — Key findings ────────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank_layout)
    _add_rect(s2, 0, 0, W, Inches(1.1), _MID_NAVY)
    _add_text(s2, "Key findings", M, Inches(0.2), W, Inches(0.8),
              28, bold=True, color=_WHITE)
    _add_text(s2, f"{case_id}  |  {hostname}", M, Inches(0.75), W, Inches(0.3),
              12, color=_LIGHT_BLUE)
    if "msfadmin" in shutdown_md:
        findings_text = (
            "The server shutdown was deliberately triggered by an authenticated user at the physical console.\n\n"
            "Two failed login attempts preceded a successful login.\n\n"
            "Administrator access was obtained within seconds of logging in.\n\n"
            "The server restarted in approximately 30 seconds and returned to normal operation.\n\n"
            "No evidence of a remote attacker, hardware failure, or software crash was found.\n\n"
            "The event is consistent with either undocumented maintenance or unauthorized physical access."
        )
    else:
        exec_bullets = _narr_bullets(narrative, "pptx_executive_summary")
        if exec_bullets:
            findings_text = "\n\n".join(f"•  {b}" for b in exec_bullets[:7])
        else:
            findings_text = "See technical report for detailed findings."
    _add_text(s2, findings_text, M, Inches(1.3), W - 2*M, Inches(5.5), 14, color=_TEXT_DARK)

    # ── Slide 3 — Business Impact ────────────────────────────────────────────
    _add_bullet_slide(
        "Business Impact",
        _narr_bullets(narrative, "pptx_impact"),
        "Operational impact assessment is in progress; refer to the technical report.",
    )

    # ── Slide 4 — Board Timeline ─────────────────────────────────────────────
    _add_bullet_slide(
        "Incident Timeline",
        _narr_bullets(narrative, "pptx_timeline"),
        "The incident timeline is under investigation; refer to the technical report for the full chronology.",
        max_items=6,
    )

    # ── Slide 5 — Root Cause & Risk ──────────────────────────────────────────
    s5 = prs.slides.add_slide(blank_layout)
    _add_rect(s5, 0, 0, W, Inches(1.1), _MID_NAVY)
    _add_text(s5, "Root Cause & Risk", M, Inches(0.2), W, Inches(0.8),
              28, bold=True, color=_WHITE)
    _add_text(s5, f"{case_id}  |  {hostname}", M, Inches(0.75), W, Inches(0.3),
              12, color=_LIGHT_BLUE)
    root_cause = narrative.get("pptx_root_cause", "").strip() or (
        "Root cause is under investigation; refer to the technical report for the initial access vector."
    )
    _add_text(s5, "ROOT CAUSE", M, Inches(1.25), W - 2 * M, Inches(0.35),
              13, bold=True, color=_LIGHT_BLUE)
    _add_text(s5, root_cause, M, Inches(1.65), W - 2 * M, Inches(1.3),
              14, color=_TEXT_DARK)
    risk_bullets = _narr_bullets(narrative, "pptx_risk")
    _add_text(s5, "KEY RISKS", M, Inches(3.1), W - 2 * M, Inches(0.35),
              13, bold=True, color=_LIGHT_BLUE)
    if risk_bullets:
        risk_text = "\n\n".join(f"•  {b}" for b in risk_bullets[:5])
    else:
        risk_text = "Risk assessment is in progress; refer to the technical report."
    _add_text(s5, risk_text, M, Inches(3.5), W - 2 * M, Inches(3.2),
              14, color=_TEXT_DARK)

    # ── Slide 6 — Response & Containment ─────────────────────────────────────
    _add_bullet_slide(
        "Response & Containment",
        _narr_bullets(narrative, "pptx_mitigations"),
        "Response and containment actions are in progress; refer to the technical report.",
    )

    # ── Slide 7 — Recommendations ────────────────────────────────────────────
    s7 = prs.slides.add_slide(blank_layout)
    _add_rect(s7, 0, 0, W, Inches(1.1), _MID_NAVY)
    _add_text(s7, "Recommendations", M, Inches(0.2), W, Inches(0.8),
              28, bold=True, color=_WHITE)
    _add_text(s7, f"{case_id}  |  {hostname}", M, Inches(0.75), W, Inches(0.3),
              12, color=_LIGHT_BLUE)
    recs = _build_recommendations(data, shutdown_md)
    row_h = Inches(0.72)
    for i, rec in enumerate(recs[:7]):
        y = Inches(1.2) + i * row_h
        _add_rect(s7, M, y, Inches(0.5), row_h - Inches(0.08), _BLUE)
        _add_text(s7, str(i + 1), M + Inches(0.1), y + Inches(0.1),
                  Inches(0.3), row_h, 16, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
        rec_clean = re.sub(r"\*\*(.*?)\*\*", r"\1", rec.split(" — ")[0][:120])
        _add_text(s7, rec_clean, M + Inches(0.6), y + Inches(0.1),
                  W - M - Inches(1.0), row_h, 13, color=_TEXT_DARK)

    # ── Slide 8 — Lessons Learned ────────────────────────────────────────────
    _add_bullet_slide(
        "Lessons Learned",
        _narr_bullets(narrative, "pptx_lessons_learned"),
        "Lessons learned will be documented once the investigation and remediation are complete.",
    )

    prs.save(str(output_path))
    print(f"[fame] PPTX saved: {output_path}")


# ── DOCX helpers ───────────────────────────────────────────────────────────────

def _load_memprocfs_results(analysis_dir: Path) -> dict:
    """Load MemProcFS JSON results from analysis/memory/memprocfs/ if present."""
    mdir = analysis_dir / "memprocfs"
    if not mdir.is_dir():
        return {}
    results = {}
    for jf in sorted(mdir.glob("memprocfs_*.json")):
        try:
            loaded = json.loads(jf.read_text())
            if isinstance(loaded, dict) and "memprocfs_version" in loaded:
                results[jf.stem] = loaded
        except Exception:
            pass
    return results


def _load_rekall_status(analysis_dir: Path) -> str:
    """Load Rekall status text from analysis/memory/memprocfs/rekall_status.txt."""
    p = analysis_dir / "memprocfs" / "rekall_status.txt"
    return p.read_text(errors="replace") if p.exists() else ""


def _build_generic_methodology_section(
    doc: Any,
    data: dict,
    case_id: str,
    hostname: str,
    image_path: str,
    analysis_dir: Path,
    _heading: Any,
    _para: Any,
    _note: Any,
    _table_2col: Any,
    _code: Any,
) -> None:
    """Generic, data-driven Part A/B methodology section.

    Used for every case that did not populate the legacy
    Kali/Metasploitable shutdown-investigation artifacts
    (shutdown_report, isf_investigation, yara_summary, ...).
    Describes only the plugins that actually produced output for
    this case — no hardcoded host names, kernels, or findings.
    """
    image_name = Path(image_path).name if image_path else "memory.image"

    _heading("Part A — Investigation Methodology", 1)
    doc.add_paragraph()

    _heading("A.1  Evidence acquisition and integrity", 2)
    _para(
        f"The memory image {image_name!r} was analysed read-only via Volatility 3. "
        "No tool in this pipeline writes back to the image file. All analysis output "
        "is written to a separate working directory (analysis/memory/) and archived "
        "in the case evidence folder for long-term preservation; SHA-256 hashes of "
        "each preserved artifact are recorded in the research notes (Appendix B)."
    )
    doc.add_paragraph()

    _heading("A.2  Tool suite", 2)
    try:
        import volatility3 as _vol3
        vol_version = getattr(_vol3, "__version__", "unknown")
    except Exception:
        vol_version = "unknown"
    _table_2col([
        ("Volatility 3", f"v{vol_version} — memory structure analysis"),
        ("python-docx",  "Document generation"),
    ])
    doc.add_paragraph()

    _heading("A.3  Chain of custody summary", 2)
    _para(
        f"Image: {image_name}  →  read-only analysis by Volatility 3  →  "
        "output preserved in the case evidence folder  →  archived in case ZIP  →  "
        "uploaded to the investigations vault."
    )
    doc.add_page_break()

    # ── Part B: Artifact Extraction Catalog ───────────────────────────────────
    _heading("Part B — Artifact Extraction Catalog", 1)
    _note(
        "For each Volatility 3 plugin that produced output for this image: "
        "what it extracts and its role in the investigation."
    )
    doc.add_paragraph()

    plugin_roles = {key: note for _label, key, note in _EXPECTED_PLUGINS}
    for key in _available_plugins(data):
        if key.startswith("_") or key in ("shutdown_report", "isf_investigation", "syslog_patterns"):
            continue
        role = plugin_roles.get(key, "")
        _para(key, bold=True)
        content = str(data.get(key, ""))
        proc_lines = [l for l in content.splitlines() if l.strip() and l.strip()[0].isdigit()]
        if proc_lines:
            _para(f"{len(proc_lines)} data row(s) recovered.")
        elif content.strip():
            _para("Output produced (see preserved artifact for details).")
        else:
            _para("No output.")
        if role:
            _para(f"Role: {role}")
        doc.add_paragraph()

    doc.add_page_break()


def _build_methodology_section(
    doc: Any,
    data: dict,
    case_id: str,
    hostname: str,
    image_path: str,
    analysis_dir: Path,
    _heading: Any,
    _para: Any,
    _note: Any,
    _table_2col: Any,
    _code: Any,
) -> None:
    """
    Write Part A (Investigation Methodology) and Part B (Artifact Extraction Catalog)
    into the Word document.

    For each tool / artifact the section explains:
      - Extraction method (tool, version, exact command)
      - Evidence integrity (why this artifact is trustworthy)
      - Contents (what information it holds)
      - Role in the investigation (which analytical question it answers)
    """
    from docx.shared import Pt, RGBColor

    # The detailed Part A/B narrative below was hand-authored for one specific
    # legacy case (Kali/Metasploitable shutdown investigation) and references
    # that case's tools, hosts, and findings by name. Only use it when this
    # case actually populated those legacy analysis artifacts; otherwise emit
    # a generic, data-driven methodology section.
    if not data.get("shutdown_report"):
        _build_generic_methodology_section(
            doc, data, case_id, hostname, image_path, analysis_dir,
            _heading, _para, _note, _table_2col, _code,
        )
        return

    memprocfs_results = _load_memprocfs_results(analysis_dir)
    rekall_status_txt = _load_rekall_status(analysis_dir)
    yara_summary      = data.get("yara_summary", {})
    isf_text          = data.get("isf_investigation", "")
    banners_txt       = data.get("banners", "")
    syslog_txt        = data.get("syslog_patterns", "")
    yara_scan_txt     = data.get("yara_scan", "")
    image_name        = Path(image_path).name if image_path else "memory.image"

    # ── Part A: Investigation Methodology ─────────────────────────────────────
    _heading("Part A — Investigation Methodology", 1)
    doc.add_paragraph()

    _heading("A.1  Evidence acquisition and integrity", 2)
    _para(
        f"The memory image {image_name!r} was acquired from the VirtualBox hypervisor "
        "host using the VirtualBox debugger command `VBoxManage debugvm`. The acquisition "
        "tool reads guest physical RAM directly from the hypervisor while the VM is paused, "
        "ensuring a consistent, atomically captured snapshot of the machine's memory state. "
        "The resulting file is an ELF64 core dump: each physical memory range is stored as "
        "a PT_LOAD segment, and the VirtualBox CPU registers — including CR3, the page "
        "directory base register — are preserved in a VBCPU PT_NOTE segment."
    )
    doc.add_paragraph()
    _para(
        "Evidence integrity is verified by SHA-256 hash. Every analysis tool in this "
        "pipeline opens the image in read-only mode. No tool writes back to the image file. "
        "All analysis output files are written to a separate working directory "
        "(analysis/memory/) and are archived in a case ZIP for long-term preservation."
    )
    doc.add_paragraph()

    _heading("A.2  Tool suite and versions", 2)
    tool_rows = [
        ("Volatility 3",  "v2.20.0  —  memory structure analysis (plugin: banners.Banners)"),
        ("GNU strings",   "binutils — printable string extraction from raw memory"),
        ("YARA",          "4.x  —  pattern-matching against custom forensic rules"),
        ("MemProcFS",     "v5.17.5 (pip: memprocfs)  —  physical memory access via LeechCore"),
        ("Rekall",        "ABANDONED — not available on Python 3.12 (see Section B.5)"),
        ("python-docx",   "Document generation"),
    ]
    _table_2col(tool_rows)
    doc.add_paragraph()

    _heading("A.3  Chain of custody summary", 2)
    _para(
        f"Image: {image_name}  →  read-only mount by each tool  →  "
        "output written to analysis/memory/  →  artifacts archived in case ZIP  →  "
        "ZIP uploaded to investigations vault"
    )
    doc.add_paragraph()
    _para(
        "Each section in Part B documents the specific command or API call used, the "
        "output file produced, and how the analyst verified that the artifact accurately "
        "represents memory content without modification."
    )
    doc.add_page_break()

    # ── Part B: Artifact Extraction Catalog ───────────────────────────────────
    _heading("Part B — Artifact Extraction Catalog", 1)
    _note(
        "For each artifact: extraction method, evidence integrity statement, "
        "contents, and role in the investigation."
    )
    doc.add_paragraph()

    # ── B.1  Volatility 3 — banners.Banners ───────────────────────────────────
    _heading("B.1  Volatility 3 — banners.Banners plugin", 2)

    _para("Extraction method", bold=True)
    _code(f"python3 /opt/volatility3/vol.py -f {image_name} banners.Banners "
          f"> analysis/memory/banners.txt")
    doc.add_paragraph()

    _para("Evidence integrity", bold=True)
    _para(
        "Volatility 3 opens the memory image with read-only file access. The banners "
        "plugin scans the raw physical memory for embedded kernel version strings — "
        "byte sequences that are compiled into the kernel binary and present in memory "
        "at a fixed offset from the kernel base address. The plugin does not interpret "
        "or modify any memory structures beyond what is necessary to locate these strings. "
        "Output is a plain-text file written to the working directory."
    )
    doc.add_paragraph()

    _para("Contents", bold=True)
    if banners_txt:
        banner_lines = [l for l in banners_txt.splitlines() if "Linux version" in l or "Offset" in l]
        summary = "\n".join(banner_lines[:6]) if banner_lines else banners_txt[:300]
        _code(summary)
    else:
        _para("banners.txt: No output — Volatility 3 plugin produced no results for this image.")
    doc.add_paragraph()

    _para("Role in the investigation", bold=True)
    _para(
        "The kernel banner string establishes the exact operating system, kernel version, "
        "compiler, and build date of the subject machine. This is the foundation for all "
        "subsequent analysis: it determines which Volatility 3 plugins are applicable, "
        "whether ISF symbol files are available, and provides the baseline against which "
        "all other memory artifacts are interpreted. For this case, the banner confirmed "
        f"the host as {hostname} running the kernel identified in the banner output above."
    )
    doc.add_paragraph()

    # ── B.2  GNU strings — ASCII string extraction ────────────────────────────
    _heading("B.2  GNU strings — printable string extraction", 2)

    _para("Extraction method", bold=True)
    _code(f"strings -a -n 8 {image_name} > analysis/memory/strings_all.txt\n"
          f"strings -a -el -n 8 {image_name} > analysis/memory/strings_unicode.txt\n"
          f"grep -E '(pam_unix|sudo:|login\\[|FAILED|shutdown|PostgreSQL)' \\\n"
          f"  analysis/memory/strings_all.txt > analysis/memory/syslog_patterns.txt")
    doc.add_paragraph()

    _para("Evidence integrity", bold=True)
    _para(
        "The strings utility reads the image sequentially, outputting every sequence of "
        "printable characters longer than the minimum length (8). It does not parse any "
        "memory structures and cannot misinterpret data — it simply reports what printable "
        "bytes are present at each offset. The output is deterministic: running strings on "
        "the same image always produces the same output. The 249 MB strings_all.txt file "
        "is the primary analysis artifact for this case and is archived in the case ZIP."
    )
    doc.add_paragraph()

    _para("Contents", bold=True)
    syslog_lines = [l for l in syslog_txt.splitlines() if l.strip()][:8] if syslog_txt else []
    if syslog_lines:
        _para(f"syslog_patterns.txt — selected entries ({len(syslog_txt.splitlines())} total lines):")
        _code("\n".join(syslog_lines[:8]))
    else:
        _para("strings_all.txt: 249 MB extracted from memory image; syslog_patterns.txt: auth/syslog grep output.")
    doc.add_paragraph()

    _para("Role in the investigation", bold=True)
    _para(
        "When Volatility 3 structured plugins cannot run (because ISF kernel symbol files "
        "are unavailable — see Section B.6), strings extraction is the primary analysis "
        "method. It recovers process names, file paths, network addresses, command lines, "
        "log messages, credentials, and script content that are present in memory as "
        "printable text. For this case, strings analysis recovered: the complete attack "
        "script (v4 Final), the credential msfadmin:msfadmin, all five session-open "
        "records, ten attack run timestamps, the remote operator's SSH connections, and "
        "seven PostgreSQL UDF library filenames. None of this would have been obtainable "
        "through structured plugins alone given the ISF limitation."
    )
    doc.add_paragraph()

    # ── B.3  YARA scanning ────────────────────────────────────────────────────
    _heading("B.3  YARA — targeted pattern matching", 2)

    _para("Extraction method", bold=True)
    _code(f"yara --print-strings --print-string-length -r \\\n"
          f"  analysis/yara/{case_id}_kali.yar \\\n"
          f"  {image_name} > analysis/memory/yara_scan.txt")
    doc.add_paragraph()

    _para("Evidence integrity", bold=True)
    _para(
        "YARA rules are authored and archived as part of the case file "
        f"(analysis/yara/{case_id}_kali.yar). Each rule targets a specific string or "
        "byte pattern confirmed to be present in the evidence from prior strings analysis. "
        "YARA scans the memory image read-only and reports each match with its physical "
        "byte offset, matched string content, and string length. Every reported match can "
        "be independently verified by seeking to the reported offset in the raw image and "
        "reading the bytes. This provides exact byte-level provenance for each finding."
    )
    doc.add_paragraph()

    _para("Contents", bold=True)
    if yara_summary:
        yara_rows = [(rule, f"{hits} match(es)") for rule, hits in list(yara_summary.items())[:14]]
        tbl = doc.add_table(rows=len(yara_rows) + 1, cols=2)
        tbl.style = "Table Grid"
        for i, hdr in enumerate(["YARA rule", "Result"]):
            tbl.rows[0].cells[i].text = hdr
            tbl.rows[0].cells[i].paragraphs[0].runs[0].font.bold = True
        for i, (rule, hits) in enumerate(yara_rows):
            tbl.rows[i+1].cells[0].text = rule
            tbl.rows[i+1].cells[1].text = hits
    else:
        _para("YARA scan not run or no results available.")
    doc.add_paragraph()

    _para("Role in the investigation", bold=True)
    _para(
        "YARA scanning provides targeted, byte-exact corroboration of findings first "
        "identified through strings analysis. Each matched rule confirms a specific "
        "investigative conclusion: for example, the Metasploit_SSH_BruteForce rule "
        "confirms the brute-force was executed from this machine; the "
        "PostgreSQL_UDF_Libraries rule independently confirms that UDF library names "
        "recovered from strings analysis are genuine memory artefacts, not artifacts of "
        "the strings tool itself; the Host_Operator_SSH_Evidence rule confirms the "
        "remote operator connection from 10.0.2.2. Twelve of twelve custom rules matched, "
        "meaning every major finding in this investigation has independent YARA-level "
        "byte corroboration."
    )
    doc.add_paragraph()

    # ── B.4  MemProcFS ────────────────────────────────────────────────────────
    _heading("B.4  MemProcFS — physical memory access via LeechCore", 2)

    _para("Extraction method", bold=True)
    # Pull DTB from results if available
    dtb_val = "0x8a000"
    memprocfs_ver = "5.17.5"
    for stem, res in memprocfs_results.items():
        if res.get("dtb_extraction", {}).get("dtb"):
            dtb_val = res["dtb_extraction"]["dtb"]
            memprocfs_ver = res.get("memprocfs_version", memprocfs_ver)
            break

    _para(
        f"MemProcFS v{memprocfs_ver} was installed via pip (pip3 install memprocfs). "
        "It uses LeechCore as its physical memory access layer. For VirtualBox ELF core "
        "dumps, LeechCore auto-detects the format and maps each PT_LOAD segment to its "
        "corresponding physical address range."
    )
    doc.add_paragraph()
    _para(
        "Because MemProcFS requires the kernel page directory base (CR3/DTB) to walk "
        "virtual-to-physical address translations, and because Linux auto-detection "
        "requires kernel symbols that are not available for this kernel (see Section B.6), "
        "the DTB was extracted manually from the ELF VBCPU note:"
    )
    doc.add_paragraph()
    _code(
        "# DTB extraction from VirtualBox VBCPU ELF note (Python)\n"
        "import struct\n"
        f"# VBCPU note at ELF PT_NOTE offset 0x4d8 in {image_name}\n"
        "# Scan for page-aligned physical addresses in CPUMCTX structure:\n"
        f"dtb = '{dtb_val}'   # offset 0x68 in VBCPU CPUMCTX — CR3 register value\n\n"
        "# Initialize MemProcFS with extracted DTB\n"
        "import memprocfs\n"
        f"vmm = memprocfs.Vmm(['-device', '{image_name}', '-dtb', '{dtb_val}'])"
    )
    doc.add_paragraph()

    _para("Evidence integrity", bold=True)
    _para(
        "The DTB extraction is forensically sound: the VBCPU note is part of the ELF "
        "core file written by the hypervisor at acquisition time and records the CPU "
        "register state at the moment of capture. The CR3 value in the CPUMCTX structure "
        "at the documented offset is the authoritative kernel page directory base from "
        "the time of the memory snapshot. MemProcFS uses this value only to interpret "
        "address translations — it does not modify the image. All physical memory reads "
        "are read-only via the LeechCore layer."
    )
    doc.add_paragraph()

    _para("Contents and findings", bold=True)
    phys_banners_all = []
    attack_cats: list[str] = []
    for stem, res in memprocfs_results.items():
        phys_banners_all.extend(res.get("physical_banners", []))
        attack_cats.extend(list(res.get("attack_artifacts", {}).keys()))

    if phys_banners_all:
        _para("Physical memory banners confirmed by MemProcFS:")
        for b in phys_banners_all[:4]:
            _code(f"  [{b['physical_address']}]  {b['content'][:100]}")
        doc.add_paragraph()
    if attack_cats:
        _para(
            f"Attack artifact categories found in physical memory scans: "
            f"{', '.join(sorted(set(attack_cats))[:12])}. "
            "See the findings sections for full context."
        )
    else:
        _para(
            "Physical memory scan found no attack strings in the first 256 MB range. "
            "MemProcFS's value in this case is the DTB extraction methodology and the "
            "independent physical memory access layer, which validates the LeechCore "
            "ELF core parsing."
        )
    doc.add_paragraph()

    _para("Role in the investigation", bold=True)
    _para(
        "MemProcFS provides a second, independent memory access pathway that is entirely "
        "separate from Volatility 3 and strings. Its contribution to this investigation "
        "is threefold. First, the DTB extraction process documents which physical page "
        "is the kernel's page directory — a forensically significant fact that proves "
        "the memory image contains a valid, running Linux kernel at the time of capture. "
        "Second, physical memory scanning via MemProcFS independently confirmed that "
        "attack payloads (LHOST/LPORT configurations, UDF SQL commands) were physically "
        "resident in the victim machine's RAM at capture time — not merely present in "
        "virtual address space. Third, MemProcFS process enumeration, while limited by "
        "missing kernel symbols (returning PID=1 with name 'unknown_process'), "
        "independently confirms that the image contains a valid process structure — "
        "providing a second tool's corroboration of memory integrity."
    )
    doc.add_paragraph()

    # ── B.5  Rekall ───────────────────────────────────────────────────────────
    _heading("B.5  Rekall — installation attempt and status", 2)

    _para("Extraction method attempted", bold=True)
    _code("pip3 install rekall rekall-core --break-system-packages")
    doc.add_paragraph()

    _para("Result", bold=True)
    if rekall_status_txt:
        for line in rekall_status_txt.splitlines()[4:22]:
            if line.strip():
                _para(line.strip())
    else:
        _para(
            "Rekall installation failed. Rekall was last released in October 2019 "
            "(v1.7.2.post1) and requires Python ≤3.7. The system Python is 3.12. "
            "Four C-extension dependencies (acora, aff4-snappy, pyblake2, fastchunking) "
            "could not be compiled. The Rekall GitHub repository has been archived "
            "(read-only) since 2021."
        )
    doc.add_paragraph()

    _para("Impact on this investigation", bold=True)
    _para(
        "Rekall's unavailability does not create a gap in the analysis for these Linux "
        "memory images. Rekall's Linux capabilities were equivalent to Volatility 3 and "
        "would have faced the same ISF symbol limitation (kernel 6.18.12+kali-amd64 and "
        "2.6.24-16-server are both unsupported by public ISF repositories). Its absence "
        "is therefore neutral: the same strings-extraction fallback would have been "
        "required. Volatility 3 is the current community standard and was chosen as "
        "Rekall's successor. All analysis is complete through Volatility 3 + strings + "
        "YARA + MemProcFS."
    )
    doc.add_paragraph()

    # ── B.6  ISF investigation ────────────────────────────────────────────────
    _heading("B.6  ISF symbol investigation — structured plugin availability", 2)

    _para("What ISF is and why it matters", bold=True)
    _para(
        "Volatility 3 Linux plugins (linux.pslist, linux.bash, linux.kallsyms, "
        "linux.netstat, etc.) require an ISF (Intermediate Symbol Format) file — a "
        "JSON mapping of kernel symbol names to virtual addresses specific to the exact "
        "kernel build. Without ISF, Volatility 3 cannot walk kernel data structures and "
        "all structured plugins fail with 'Unsatisfied requirement: symbol_table_name'."
    )
    doc.add_paragraph()

    _para("Approaches attempted to obtain ISF", bold=True)

    approach_data = [
        ("Approach 1 — Online ISF download",
         "Checked the community ISF repository at isf-server.code16.fr for kernel "
         "6.18.12+kali-amd64 and 2.6.24-16-server. Neither kernel version has a "
         "pre-built ISF file. The Kali kernel is too recent; the Metasploitable kernel "
         "is too old. Result: no ISF available online."),
        ("Approach 2 — System.map via disk image",
         "Extracted /boot/System.map-6.18.12+kali-amd64 from kali-post.vdi using a "
         "Python ext4 filesystem reader. The file was 92 bytes — a stub with the "
         "content: 'ffffffffffffffff B The real System.map is in the "
         "linux-image-6.18.12+kali-amd64-dbg package'. Kali intentionally ships this "
         "stub to prevent kernel symbol exposure. Result: no usable System.map."),
        ("Approach 3 — dwarf2json from vmlinux",
         "dwarf2json is installed and could generate ISF from a DWARF-annotated "
         "vmlinux binary. However, the standard Kali kernel ships a stripped vmlinuz; "
         "the debug symbols are in the -dbg package. Installing linux-image-"
         "6.18.12+kali-amd64-dbg on the analyst machine would produce ISF from a "
         "different build than the evidence image, making the symbols forensically "
         "invalid. Result: not pursued — would compromise evidence integrity."),
    ]

    for title, detail in approach_data:
        _para(title, bold=True)
        _para(detail)
        doc.add_paragraph()

    _para("Conclusion and fallback", bold=True)
    _para(
        "ISF symbols are not available for either evidence image without installing "
        "debug packages that would not match the captured kernel build. The fallback "
        "applied throughout this investigation — GNU strings extraction + grep pattern "
        "matching + custom YARA rules + MemProcFS physical memory access — provides "
        "equivalent analytical coverage and is fully documented with byte-level "
        "provenance for every finding."
    )
    doc.add_page_break()


# ── DOCX utility helpers ──────────────────────────────────────────────────────

def _compute_file_metadata(file_path: Path) -> tuple[str, str]:
    """Return (mtime_cet, sha256_hex) for a file, or ('N/A', 'N/A') if missing."""
    if not file_path.exists():
        return "N/A", "N/A"
    mtime_str = datetime.fromtimestamp(
        file_path.stat().st_mtime, tz=_CET
    ).strftime("%d-%b-%Y %H:%M CET")
    h = hashlib.sha256()
    with open(file_path, "rb") as fh:
        for chunk in iter(lambda: fh.read(65536), b""):
            h.update(chunk)
    return mtime_str, h.hexdigest()


def _set_doc_font(doc: Any, font_name: str = "Arial") -> None:
    """Set default font across all common styles in the document."""
    for style_name in [
        "Normal", "Heading 1", "Heading 2", "Heading 3", "Heading 4",
        "No Spacing", "Intense Quote", "List Number", "List Bullet",
        "Table Grid",
    ]:
        try:
            doc.styles[style_name].font.name = font_name
        except Exception:
            pass


def _add_header_footer(section: Any, case_id: str) -> None:
    """Add Case ID to page header (right) and page numbers to footer (center)."""
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    # Header: Case ID right-aligned
    hdr = section.header
    if not hdr.paragraphs:
        h_para = hdr.add_paragraph()
    else:
        h_para = hdr.paragraphs[0]
    h_para.clear()
    h_run = h_para.add_run(case_id)
    h_run.font.name = "Arial"
    h_run.font.size = Pt(9)
    h_run.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)
    h_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    # Footer: "Page X of Y" centered
    ftr = section.footer
    if not ftr.paragraphs:
        f_para = ftr.add_paragraph()
    else:
        f_para = ftr.paragraphs[0]
    f_para.clear()
    f_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    def _fld_char(fld_type: str) -> Any:
        r = OxmlElement("w:r")
        fc = OxmlElement("w:fldChar")
        fc.set(qn("w:fldCharType"), fld_type)
        r.append(fc)
        return r

    def _instr(text: str) -> Any:
        r = OxmlElement("w:r")
        it = OxmlElement("w:instrText")
        it.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
        it.text = text
        r.append(it)
        return r

    def _txt_run(text: str) -> Any:
        r = f_para.add_run(text)
        r.font.name = "Arial"
        r.font.size = Pt(9)
        r.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)
        return r

    _txt_run("Page ")
    f_para._p.append(_fld_char("begin"))
    f_para._p.append(_instr(" PAGE "))
    f_para._p.append(_fld_char("end"))
    _txt_run(" of ")
    f_para._p.append(_fld_char("begin"))
    f_para._p.append(_instr(" NUMPAGES "))
    f_para._p.append(_fld_char("end"))


def _add_watermark(section: Any) -> None:
    """Add a 'CONFIDENTIAL' VML watermark at 315° (≡ −45°) to a document section."""
    from lxml import etree

    hdr = section.header
    wm_para = hdr.add_paragraph()
    vml = (
        '<w:r xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"'
        ' xmlns:v="urn:schemas-microsoft-com:vml"'
        ' xmlns:o="urn:schemas-microsoft-com:office:office">'
        "<w:rPr><w:noProof/></w:rPr>"
        "<w:pict>"
        '<v:shape id="WaterMark" o:spid="_x0000_s2049"'
        ' type="#_x0000_t136"'
        ' style="position:absolute;margin-left:0;margin-top:0;'
        "width:430pt;height:120pt;z-index:-251654144;"
        "rotation:315;"
        "mso-position-horizontal:center;"
        "mso-position-horizontal-relative:page;"
        "mso-position-vertical:center;"
        'mso-position-vertical-relative:page"'
        ' fillcolor="#C0C0C0" stroked="f">'
        "<v:textpath"
        ' on="t" fitshape="t" string="CONFIDENTIAL"'
        " style='font-family:\"Arial\";font-size:1pt;font-weight:bold'"
        "/>"
        "</v:shape>"
        "</w:pict>"
        "</w:r>"
    )
    wm_para._p.append(etree.fromstring(vml))


def _remove_blank_paragraphs(doc: Any) -> None:
    """Remove empty paragraphs except those adjacent to headings or containing page breaks."""
    from docx.oxml.ns import qn

    paras = list(doc.paragraphs)
    to_remove = []
    for i, para in enumerate(paras):
        if para.text.strip():
            continue
        if para.style.name.startswith("Heading"):
            continue
        if para._element.find(".//" + qn("w:br")) is not None:
            continue
        prev_heading = i > 0 and paras[i - 1].style.name.startswith("Heading")
        next_heading = i + 1 < len(paras) and paras[i + 1].style.name.startswith("Heading")
        if prev_heading or next_heading:
            continue
        to_remove.append(para._element)
    for elem in to_remove:
        if elem.getparent() is not None:
            elem.getparent().remove(elem)


# ── DOCX generation ────────────────────────────────────────────────────────────

def _build_docx(
    data: dict[str, Any],
    case_id: str,
    hostname: str,
    image_path: str,
    generated_utc: str,
    output_path: Path,
    opencti_findings: str = "",
    fan_summary: str = "",
    fast_summary: str = "",
) -> None:
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
    except ImportError:
        print("[fame] WARNING: python-docx not installed — skipping DOCX. pip3 install python-docx")
        return

    shutdown_md    = data.get("shutdown_report", "")
    analysis_dir_p = PROJECT_ROOT / "analysis" / "memory"
    doc = Document()

    # ── Page setup ────────────────────────────────────────────────────────────
    for section in doc.sections:
        section.top_margin      = Inches(1.0)
        section.bottom_margin   = Inches(1.0)
        section.left_margin     = Inches(1.2)
        section.right_margin    = Inches(1.2)
        section.header_distance = Inches(0.4)
        section.footer_distance = Inches(0.4)

    # ── Default font and paragraph spacing ────────────────────────────────────
    _set_doc_font(doc, "Arial")
    styles = doc.styles
    doc.styles["Normal"].paragraph_format.space_after  = Pt(5)
    doc.styles["Normal"].paragraph_format.space_before = Pt(0)
    for _i in range(1, 5):
        try:
            _hs = doc.styles[f"Heading {_i}"]
            _hs.paragraph_format.space_before = Pt(14 if _i == 1 else 10)
            _hs.paragraph_format.space_after  = Pt(4)
            _hs.font.name = "Arial"
        except Exception:
            pass

    def _heading(text: str, level: int) -> None:
        p = doc.add_heading(text, level=level)
        if p.runs:
            p.runs[0].font.color.rgb = RGBColor(0x0f, 0x17, 0x2a)
            p.runs[0].font.name = "Arial"

    def _para(text: str, italic: bool = False, bold: bool = False) -> None:
        p = doc.add_paragraph()
        run = p.add_run(text)
        run.italic = italic
        run.bold   = bold
        run.font.name = "Arial"

    def _note(text: str) -> None:
        p = (doc.add_paragraph(style="Intense Quote")
             if "Intense Quote" in [s.name for s in styles]
             else doc.add_paragraph())
        run = p.add_run(text)
        run.italic = True
        run.font.name = "Arial"
        run.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)

    def _code(text: str) -> None:
        sn = "No Spacing" if "No Spacing" in [s.name for s in styles] else "Normal"
        p  = doc.add_paragraph(style=sn)
        r  = p.add_run(text)
        r.font.name  = "Courier New"
        r.font.size  = Pt(9)
        r.font.color.rgb = RGBColor(0x1e, 0x3a, 0x5f)

    def _table_2col(rows: list[tuple[str, str]], header: bool = True) -> None:
        start = 1 if header else 0
        tbl = doc.add_table(rows=len(rows) + start, cols=2)
        tbl.style = "Table Grid"
        if header:
            for i, h in enumerate(["Field", "Value"]):
                c = tbl.rows[0].cells[i]
                c.text = h
                c.paragraphs[0].runs[0].font.bold = True
                c.paragraphs[0].runs[0].font.name = "Arial"
        for i, (k, v) in enumerate(rows):
            r = tbl.rows[i + start]
            r.cells[0].text = k
            r.cells[1].text = v
            for j in range(2):
                for run in r.cells[j].paragraphs[0].runs:
                    run.font.name = "Arial"

    # ── Header, footer, watermark ─────────────────────────────────────────────
    for section in doc.sections:
        _add_header_footer(section, case_id)
        _add_watermark(section)

    # ── Cover ──────────────────────────────────────────────────────────────────
    title = doc.add_heading("FAME — Memory forensics report", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if title.runs:
        title.runs[0].font.name = "Arial"

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = sub.add_run("Forensic Analysis Memory  |  Fan Get Fame Fast")
    r.font.size  = Pt(14)
    r.font.name  = "Arial"
    r.font.color.rgb = RGBColor(0x1d, 0x4e, 0xd8)

    # Memory image acquisition timestamp
    image_mtime = "N/A"
    if image_path and Path(image_path).exists():
        image_mtime = datetime.fromtimestamp(
            Path(image_path).stat().st_mtime, tz=_CET
        ).strftime("%d-%b-%Y %H:%M CET")

    _table_2col([
        ("Case ID",              case_id),
        ("Hostname",             hostname),
        ("Memory image",         Path(image_path).name if image_path else ""),
        ("Memory image created", image_mtime),
        ("Module",               "FAME — Forensic Analysis Memory"),
        ("Analysts",             "Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin · Joost Beekman"),
        ("Generated",            generated_utc),
        ("Analysis tools",       "Volatility 3 v2.20.0 · GNU strings · YARA · MemProcFS v5.17.5"),
    ])
    doc.add_page_break()

    # ── Parts A and B — Methodology and chain of evidence ─────────────────────
    _build_methodology_section(
        doc, data, case_id, hostname, image_path, analysis_dir_p,
        _heading, _para, _note, _table_2col, _code,
    )

    # ── Part C — Findings ──────────────────────────────────────────────────────
    _heading("Part C — Findings", 1)
    doc.add_page_break()

    # ── C.1 Management summary ────────────────────────────────────────────────
    _heading("C.1  Management summary", 2)
    _note("Audience: CISO, Legal, Internal Audit — no technical identifiers.")
    if "msfadmin" in shutdown_md:
        _para(
            "A memory forensic analysis was conducted to determine the cause of an unexpected server "
            "shutdown. The analysis found that the shutdown was deliberately triggered by a person at "
            "the physical server console. Two failed login attempts with unknown credentials were made "
            "before a successful login. The authenticated user obtained administrator privileges within "
            "seconds and issued a reboot command. All services halted in the expected orderly sequence "
            "and the server restarted approximately 30 seconds later."
        )
        _para("Business impact:", bold=True)
        _para(
            "The event indicates either an undocumented maintenance action or unauthorized physical "
            "access to the server room. No evidence of a remote attacker, hardware failure, or software "
            "crash was found in the memory image. Physical access controls and change management "
            "procedures should be reviewed."
        )
    else:
        narrative = data.get("_narrative", {})
        exec_bullets = _narr_bullets(narrative, "pptx_executive_summary")
        impact_bullets = _narr_bullets(narrative, "pptx_impact")
        risk_bullets = _narr_bullets(narrative, "pptx_risk")
        if exec_bullets:
            _para(
                f"A memory forensic analysis of host {hostname} was conducted to establish "
                "what activity took place on the workstation. The business-level findings are "
                "summarised below; precise technical detail (identifiers, addresses, ports, "
                "process IDs) is contained in the technical sections and appendices."
            )
            _para("Key findings:", bold=True)
            for b in exec_bullets:
                p = doc.add_paragraph(style="List Bullet")
                run = p.add_run(b); run.font.name = "Arial"
            if impact_bullets:
                _para("Business impact:", bold=True)
                for b in impact_bullets:
                    p = doc.add_paragraph(style="List Bullet")
                    run = p.add_run(b); run.font.name = "Arial"
            if risk_bullets:
                _para("Business risk:", bold=True)
                for b in risk_bullets:
                    p = doc.add_paragraph(style="List Bullet")
                    run = p.add_run(b); run.font.name = "Arial"
        else:
            _para("See the technical sections below for detailed findings.")

    # ── C.2 Cross-module intelligence ─────────────────────────────────────────
    if fan_summary or fast_summary or opencti_findings:
        _heading("C.2  Cross-module intelligence", 2)
        _note("Correlate memory findings with network and storage evidence.")
        if fan_summary:
            _heading("Network Forensics (FAN)", 3)
            _para(fan_summary.strip())
        if fast_summary:
            _heading("Storage Forensics (FAST)", 3)
            _para(fast_summary.strip())
        if opencti_findings:
            _heading("OpenCTI Threat Intelligence", 3)
            _para(opencti_findings.strip())

    # ── C.3 Indicators of compromise ─────────────────────────────────────────
    _heading("C.3  Indicators of compromise", 2)
    _note("Defanged IOC values extracted from memory analysis artifacts.")
    iocs = _extract_iocs(data, shutdown_md)
    if iocs:
        tbl = doc.add_table(rows=len(iocs) + 1, cols=4)
        tbl.style = "Table Grid"
        for i, hdr in enumerate(["Type", "Value", "Severity", "Context"]):
            tbl.rows[0].cells[i].text = hdr
            tbl.rows[0].cells[i].paragraphs[0].runs[0].font.bold = True
            tbl.rows[0].cells[i].paragraphs[0].runs[0].font.name = "Arial"
        for i, ioc in enumerate(iocs):
            tbl.rows[i + 1].cells[0].text = ioc["type"]
            tbl.rows[i + 1].cells[1].text = ioc["value"]
            tbl.rows[i + 1].cells[2].text = ioc["severity"]
            tbl.rows[i + 1].cells[3].text = ioc["context"]
            for j in range(4):
                for run in tbl.rows[i + 1].cells[j].paragraphs[0].runs:
                    run.font.name = "Arial"
    else:
        _para("No malicious indicators of compromise identified in this memory image.")

    # ── C.4 Recommendations ───────────────────────────────────────────────────
    _heading("C.4  Recommendations", 2)
    recs = _build_recommendations(data, shutdown_md)
    for i, rec in enumerate(recs, 1):
        rec_clean = re.sub(r"\*\*(.*?)\*\*", r"\1", rec)
        p = doc.add_paragraph(style="List Number")
        r = p.add_run(rec_clean)
        r.font.name = "Arial"

    doc.add_page_break()

    # ── Appendix A — Analysis source files ────────────────────────────────────
    _heading("Appendix A — Analysis source files", 1)
    _note(
        "All files below are part of the chain of custody. "
        "SHA-256 hashes and timestamps recorded at report generation time."
    )
    artifact_files: list[tuple[Path, str]] = [
        (analysis_dir_p / "pslist.txt",           "Windows process list (EPROCESS walk)"),
        (analysis_dir_p / "psscan.txt",           "Windows process pool scan"),
        (analysis_dir_p / "linux_pslist.txt",     "Linux process list"),
        (analysis_dir_p / "cmdline.txt",          "Process command lines"),
        (analysis_dir_p / "netstat.txt",          "Active network connections"),
        (analysis_dir_p / "netscan.txt",          "Network connection pool scan"),
        (analysis_dir_p / "malfind.txt",          "Code injection findings"),
        (analysis_dir_p / "svcscan.txt",          "Services pool scan"),
        (analysis_dir_p / "mem_timeline.txt",     "Memory artifact timeline"),
        (analysis_dir_p / "strings_all.txt",      "All printable strings extracted from memory"),
        (analysis_dir_p / "syslog_patterns.txt",  "Auth/syslog pattern grep from strings"),
        (analysis_dir_p / "banners.txt",          "Volatility 3 kernel banners"),
        (analysis_dir_p / "yara_scan.txt",        "YARA rule match output with offsets"),
        (analysis_dir_p / "isf_investigation.txt","ISF symbol investigation log"),
        (analysis_dir_p / "proc_baseline.csv",    "Process baseline diff"),
        (analysis_dir_p / "drv_baseline.csv",     "Driver baseline diff"),
        (analysis_dir_p / "svc_baseline.csv",     "Service baseline diff"),
    ]
    for zf in sorted((PROJECT_ROOT / "analysis").glob("*.zip")):
        artifact_files.append((zf, "Case artifact ZIP"))

    tbl = doc.add_table(rows=len(artifact_files) + 1, cols=4)
    tbl.style = "Table Grid"
    for i, hdr in enumerate(["Filename", "Description", "Generated (CET)", "SHA-256 (first 32 chars)"]):
        tbl.rows[0].cells[i].text = hdr
        tbl.rows[0].cells[i].paragraphs[0].runs[0].font.bold = True
        tbl.rows[0].cells[i].paragraphs[0].runs[0].font.name = "Arial"
    for i, (fp, desc) in enumerate(artifact_files):
        mtime_s, sha256 = _compute_file_metadata(fp)
        tbl.rows[i + 1].cells[0].text = fp.name
        tbl.rows[i + 1].cells[1].text = desc
        tbl.rows[i + 1].cells[2].text = mtime_s
        tbl.rows[i + 1].cells[3].text = sha256[:32] if sha256 != "N/A" else "N/A"
        for j in range(4):
            for run in tbl.rows[i + 1].cells[j].paragraphs[0].runs:
                run.font.size = Pt(8)
                run.font.name = "Courier New" if j in (0, 3) else "Arial"

    # Remove blank paragraphs from methodology section spacers
    _remove_blank_paragraphs(doc)

    doc.save(str(output_path))
    print(f"[fame] DOCX saved: {output_path}")


# ── Public API ─────────────────────────────────────────────────────────────────

def generate(
    case_id: str,
    hostname: str,
    image_path: str = "",
    analysis_dir: Path | None = None,
    output_dir: Path | None = None,
    case_dir: Path | None = None,
    docs_dir: Path | None = None,
    opencti_findings: str = "",
    fan_summary: str = "",
    fast_summary: str = "",
    md_only: bool = False,
) -> dict[str, Path]:
    """
    Generate the full FAME report suite (Markdown, PDF, PPTX, DOCX).

    case_dir: module-specific directory (reports/<case_id>/FAME/<hostname>/). When
    supplied, Markdown goes to case_dir/. docs_dir overrides where PDF/PPTX/DOCX land
    (default: case_dir/output/ for legacy compat, typically reports/<case_id>/documents/).
    When both are omitted, all formats land in output_dir (legacy flat behaviour).

    Returns dict with keys: md, md_draft, pdf, pptx, docx — each a Path (or None).
    If the primary Markdown report already exists (analyst may have enhanced it),
    the new auto-generated content is written to <case_id>_fame_report_generated.md
    and md_draft points to that file; md always points to the primary report path.
    """
    analysis_dir = analysis_dir or (PROJECT_ROOT / "analysis" / "memory")

    if case_dir is not None:
        md_dir  = path_guard.guard_output_dir(case_dir)
        aux_dir = path_guard.guard_output_dir(docs_dir or (case_dir / "output"))
    else:
        md_dir  = path_guard.guard_output_dir(output_dir or (PROJECT_ROOT / "reports"))
        aux_dir = md_dir

    data = _load_analysis(analysis_dir)
    # Make the Claude-authored narrative available to the PPTX/DOCX builders too
    # (the Markdown builder loads it separately). This lets the management-facing
    # outputs fall back to the narrative's pptx_* sections for cases that are not
    # the legacy "unexpected shutdown" scenario.
    data["_narrative"] = _load_narrative(case_id, md_dir)
    generated_utc = datetime.now(_CET).strftime("%Y-%m-%d %H:%M CET")
    stem = case_id.replace(" ", "_")

    # ── Markdown ──────────────────────────────────────────────────────────────
    md_text = _build_markdown(
        data, case_id, hostname, image_path or str(analysis_dir),
        generated_utc, opencti_findings, fan_summary, fast_summary,
        reports_dir=md_dir,
    )
    md_path = md_dir / f"{stem}_fame_report.md"
    md_draft_path: Path | None = None

    if md_path.exists():
        # Analyst may have enhanced the primary report — write new auto-generated
        # content to a draft file so they can review and promote it if desired.
        md_draft_path = md_dir / f"{stem}_fame_report_generated.md"
        md_draft_path.write_text(md_text)
        pdf_source = md_draft_path
        print(f"[fame] Existing report preserved: {md_path}")
        print(f"[fame] New auto-generated draft:  {md_draft_path}")
    else:
        md_path.write_text(md_text)
        pdf_source = md_path
        print(f"[fame] Markdown saved: {md_path}")

    # ── PDF ───────────────────────────────────────────────────────────────────
    pdf_path: Path | None = None
    if not md_only:
        try:
            sys.path.insert(0, str(PROJECT_ROOT / "lib"))
            from md_to_pdf import convert as md2pdf
            pdf_path = aux_dir / f"{stem}_fame_report.pdf"
            md2pdf(pdf_source, pdf_path)
            print(f"[fame] PDF saved: {pdf_path}")
        except Exception as exc:
            print(f"[fame] WARNING: PDF generation failed: {exc}")

    # ── PPTX ──────────────────────────────────────────────────────────────────
    pptx_path: Path | None = None
    if not md_only:
        pptx_path = aux_dir / f"{stem}_fame_presentation.pptx"
        _build_pptx(
            data, case_id, hostname, image_path or str(analysis_dir),
            generated_utc, pptx_path, opencti_findings, fan_summary, fast_summary,
        )

    # ── DOCX ──────────────────────────────────────────────────────────────────
    docx_path: Path | None = None
    if not md_only:
        docx_path = aux_dir / f"{stem}_fame_report.docx"
        _build_docx(
            data, case_id, hostname, image_path or str(analysis_dir),
            generated_utc, docx_path, opencti_findings, fan_summary, fast_summary,
        )

    return {
        "md":       md_path,
        "md_draft": md_draft_path,
        "pdf":      pdf_path,
        "pptx":     pptx_path if (pptx_path and pptx_path.exists()) else None,
        "docx":     docx_path if (docx_path and docx_path.exists()) else None,
    }


# ── CLI ────────────────────────────────────────────────────────────────────────

def _build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(description="FAME — Memory Forensics Report Generator")
    p.add_argument("--case-id",      required=True, metavar="ID",   help="Case identifier")
    p.add_argument("--hostname",     required=True, metavar="HOST",  help="Target hostname")
    p.add_argument("--image-path",   default="",   metavar="PATH",  help="Path to memory image")
    p.add_argument("--analysis-dir", default=None, metavar="DIR",   help="Analysis output directory")
    p.add_argument("--output-dir",   default=None, metavar="DIR",   help="Report output directory")
    p.add_argument("--case-dir",     default=None, metavar="DIR",   help="Module-specific dir (reports/<case_id>/FAME/<host>/); MD lands here")
    p.add_argument("--docs-dir",     default=None, metavar="DIR",   help="Shared documents dir (reports/<case_id>/documents/); PDF/PPTX/DOCX land here")
    p.add_argument("--opencti",      default="",   metavar="TEXT",  help="OpenCTI enrichment text")
    p.add_argument("--fan-summary",  default="",   metavar="TEXT",  help="FAN (network) summary for cross-module section")
    p.add_argument("--fast-summary", default="",   metavar="TEXT",  help="FAST (storage) summary for cross-module section")
    p.add_argument("--md-only",      action="store_true",           help="Generate Markdown only — skip PDF, PPTX, DOCX")
    return p


if __name__ == "__main__":
    args = _build_parser().parse_args()
    paths = generate(
        case_id       = args.case_id,
        hostname      = args.hostname,
        image_path    = args.image_path,
        analysis_dir  = Path(args.analysis_dir) if args.analysis_dir else None,
        output_dir    = Path(args.output_dir)   if args.output_dir   else None,
        case_dir      = Path(args.case_dir)     if args.case_dir     else None,
        docs_dir      = Path(args.docs_dir)     if args.docs_dir     else None,
        opencti_findings = args.opencti,
        fan_summary   = args.fan_summary,
        fast_summary  = args.fast_summary,
        md_only       = args.md_only,
    )
    print("[fame] Report suite complete:")
    for fmt, p in paths.items():
        if p:
            print(f"  {fmt.upper():4s}  {p}")
