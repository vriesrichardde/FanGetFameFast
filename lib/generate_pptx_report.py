#!/usr/bin/env python3
# SPDX-License-Identifier: MIT OR Apache-2.0
# SPDX-FileCopyrightText: 2026 Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin · Joost Beekman
"""
generate_pptx_report.py — Management PowerPoint briefing for PCAP investigations.

Audience: CISO, Legal, IT leadership, Internal Audit.
Language: plain English — no raw IPs, ports, or file sizes.

Slides:
  1  Cover
  2  Executive Summary
  3  Threat Landscape
  4  Security Alerts (IDS / YARA)
  5  Indicators of Compromise
  6  Recommended Actions
  7  Investigation Coverage (22 modules)

Usage (CLI):
  python3 lib/generate_pptx_report.py \
      --stem capture --case-id FAN-2026-001 \
      [--output-dir ./analysis/_reports/capture] \
      [--base-dir ./analysis] \
      [--description "Suspected C2 beacon on DESKTOP-42"]

Python API:
  from lib.generate_pptx_report import generate
  pptx_path = generate(data, stem, case_id, output_dir, description)
"""
from __future__ import annotations

import argparse
import sys
from datetime import datetime, timezone
from pathlib import Path
sys.path.insert(0, str(Path(__file__).resolve().parent))
import path_guard  # noqa: E402  write-path policy enforcement

PROJECT_ROOT = Path(__file__).parent.parent

# ── Colour palette (RGB tuples) ────────────────────────────────────────────────
# Board-level cybersecurity dark theme

_DARK_NAVY  = (0x0a, 0x16, 0x28)   # slide background
_MID_NAVY   = (0x13, 0x2a, 0x4a)   # slightly lighter panel/card backgrounds
_BLUE       = (0x15, 0x65, 0xc0)   # title bars, accent shapes
_ELECTRIC   = (0x29, 0x79, 0xff)   # highlights, key numbers, ticks
_LIGHT_BLUE = (0x90, 0xca, 0xf9)   # secondary highlights
_WHITE      = (0xff, 0xff, 0xff)   # primary body text
_LIGHT      = (0xb0, 0xbe, 0xc5)   # secondary / de-emphasised text
_ALERT      = (0xff, 0x52, 0x52)   # risk / warning indicators
_AMBER      = (0xfb, 0xbf, 0x24)   # classification / caution

# Legacy aliases (keep for compatibility with severity tables)
_LIGHT_BG   = _MID_NAVY
_ROW_ALT    = (0x0f, 0x1e, 0x36)
_TEXT_DARK  = _WHITE
_TEXT_MID   = _LIGHT

_SEV_RGB = {
    "critical": (0xef, 0x44, 0x44),
    "high":     (0xf9, 0x73, 0x16),
    "medium":   (0xea, 0xb3, 0x08),
    "low":      (0x22, 0xc5, 0x5e),
    "info":     (0x6b, 0x72, 0x80),
}
_SEV_LABEL = {
    "critical": "CRITICAL",
    "high":     "HIGH",
    "medium":   "MEDIUM",
    "low":      "LOW",
    "info":     "INFO",
}

# Technical category names → management-friendly descriptions
_MGMT_DESC: dict[str, str] = {
    "DGA Detection":              "Malware domain patterns (DGA) — automated covert channel creation",
    "DNS Beaconing":              "Automated DNS beaconing — malware check-in behaviour",
    "DNS Exfiltration":           "Data exfiltration via DNS tunnelling",
    "DNS Fast Flux":              "Rapid DNS rotation — used to hide attacker infrastructure",
    "NXDomain Flood":             "DNS denial-of-service flooding",
    "DNS Amplification":          "DNS amplification — traffic abuse for large-scale DDoS",
    "DNS Typosquatting":          "Domain impersonation detected (typosquatting)",
    "HTTP Beaconing":             "Regular automated web requests — possible malware C2 channel",
    "HTTP Large Upload":          "Unusually large data upload — potential data exfiltration",
    "Suspicious User Agent":      "Offensive-tool signature in web traffic",
    "Deprecated TLS":             "Outdated encryption in use — downgrade attack risk",
    "Suspicious URI":             "Suspicious web request paths — possible exploitation attempt",
    "HTTP Scanning":              "Automated web scanning activity",
    "SYN Flood":                  "TCP denial-of-service attack (SYN flood)",
    "Port Scan":                  "Network reconnaissance — port scanning activity",
    "RST Flood":                  "TCP denial-of-service attack (RST flood)",
    "Stealth Scan":               "Covert network reconnaissance",
    "Session Hijacking":          "Network session takeover attempt",
    "UDP Flood":                  "UDP denial-of-service flooding",
    "UDP Amplification":          "Reflection-based denial-of-service attack",
    "ICMP Flood":                 "Ping denial-of-service attack",
    "ICMP Tunneling":             "Covert data channel via ICMP — possible exfiltration",
    "ICMP Exfiltration":          "Data exfiltration via ICMP",
    "ICMP Sweep":                 "Network host discovery sweep",
    "ARP Cache Poisoning":        "ARP spoofing — man-in-the-middle attack risk",
    "ARP Flood":                  "ARP denial-of-service flooding",
    "ARP Scan":                   "Network reconnaissance via ARP",
    "Self-Signed Certificate":    "Unverified certificate — possible attacker infrastructure",
    "Expired Certificate":        "Expired certificate in use",
    "SNI Mismatch":               "Certificate/domain mismatch — possible traffic interception",
    "Weak Cipher Suite":          "Weak encryption cipher in use",
    "Suspicious JA4/JA3":        "Suspicious TLS fingerprint — possible C2 or attack tool",
    "NBNS Spoofing":              "Windows name service spoofing — credential theft risk",
    "LLMNR Spoofing":             "Name resolution poisoning — credential theft risk",
    "NetBIOS Poisoning":          "Windows name service poisoning — credential theft risk",
    "NTLM Hash Theft":            "Windows credential theft (NTLM hash) detected",
    "SMB Relay":                  "Pass-the-hash relay attack detected",
    "DHCP Starvation":            "DHCP denial-of-service attack (address pool exhaustion)",
    "DHCP Rogue Server":          "Unauthorized DHCP server — network redirection risk",
    "NTP Amplification":          "NTP-based denial-of-service amplification",
    "SNMP Default Credentials":   "Default SNMP credentials in use — unauthorized access risk",
    "QUIC Amplification":         "QUIC protocol denial-of-service amplification",
    "mDNS Amplification":         "mDNS denial-of-service amplification",
    "SSDP Amplification":         "UPnP/SSDP denial-of-service amplification",
    "UPnP Device Exposure":       "Network device exposed via UPnP — remote access risk",
    "STUN Amplification":         "STUN-based denial-of-service amplification",
}

# Module coverage table (label, data key, user-friendly name)
_MODULES = [
    ("PCAP netflow",        "has_pcap",     "Network Flow Extraction"),
    ("DNS threats",         "has_dns",      "DNS Threat Detection"),
    ("HTTP/S threats",      "has_http",     "HTTP/S Threat Detection"),
    ("TLS session",         "has_tls",      "TLS Session Inspector"),
    ("TLS certificate",     "has_cert",     "TLS Certificate Inspector"),
    ("ICMP threats",        "has_icmp",     "ICMP Threat Detection"),
    ("TCP threats",         "has_tcp",      "TCP Threat Detection"),
    ("UDP threats",         "has_udp",      "UDP Threat Detection"),
    ("ARP threats",         "has_arp",      "ARP Threat Detection"),
    ("DHCP threats",        "has_dhcp",     "DHCP Threat Detection"),
    ("mDNS threats",        "has_mdns",     "mDNS Threat Detection"),
    ("QUIC threats",        "has_quic",     "QUIC Threat Detection"),
    ("NTP threats",         "has_ntp",      "NTP Threat Detection"),
    ("SNMP threats",        "has_snmp",     "SNMP Threat Detection"),
    ("NBNS threats",        "has_nbns",     "NBNS Threat Detection"),
    ("LLMNR threats",       "has_llmnr",    "LLMNR Threat Detection"),
    ("STUN threats",        "has_stun",     "STUN Threat Detection"),
    ("SSDP threats",        "has_ssdp",     "SSDP/UPnP Threat Detection"),
    ("NetBIOS threats",     "has_netbios",  "NetBIOS Threat Detection"),
    ("File hashes",         "has_fh",       "File Extraction & OSINT"),
    ("Suricata IDS",        "has_suricata", "Suricata IDS"),
    ("YARA rules",          "has_yara",     "YARA Signature Scan"),
    ("IP/FQDN enrichment",  "has_fan_ip",   "IP/Domain CTI Enrichment"),
]


# ── pptx helpers ───────────────────────────────────────────────────────────────

def _rgb(r_g_b: tuple[int, int, int]):
    from pptx.dml.color import RGBColor
    return RGBColor(*r_g_b)


def _set_bg(slide, color: tuple[int, int, int]) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _rect(slide, left, top, width, height,
          fill: tuple[int, int, int],
          line: tuple[int, int, int] | None = None,
          line_width_pt: float = 0.75):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def _text(slide, text: str, left, top, width, height,
          size: float = 12, bold: bool = False, italic: bool = False,
          color: tuple[int, int, int] = _TEXT_DARK,
          align=None, word_wrap: bool = True) -> None:
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    if align is None:
        align = PP_ALIGN.LEFT
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)


def _text_lines(slide, lines: list[tuple], left, top, width, height,
                default_size: float = 11, default_color: tuple = _TEXT_DARK,
                word_wrap: bool = True) -> None:
    """Add a text box with multiple paragraphs. Each entry is (text, size, bold, color, italic)."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    first = True
    for item in lines:
        text  = item[0]
        size  = item[1] if len(item) > 1 else default_size
        bold  = item[2] if len(item) > 2 else False
        color = item[3] if len(item) > 3 else default_color
        ital  = item[4] if len(item) > 4 else False
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = ital
        run.font.color.rgb = _rgb(color)


def _header_bar(slide, title: str, case_id: str, W, H) -> None:
    """Blue title bar on dark-navy slide (content slides)."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    BAR_H = Inches(0.65)
    _rect(slide, 0, 0, W, BAR_H, fill=_BLUE)
    # Thin electric-blue left accent
    _rect(slide, 0, 0, Inches(0.06), BAR_H, fill=_ELECTRIC)
    # Title
    _text(slide, title, Inches(0.18), Inches(0.07), Inches(10), Inches(0.52),
          size=18, bold=True, color=_WHITE)
    # Case ID right-aligned
    _text(slide, case_id, Inches(10.0), Inches(0.1), Inches(3.1), Inches(0.45),
          size=9, color=_LIGHT_BLUE,
          align=__import__("pptx.enum.text", fromlist=["PP_ALIGN"]).PP_ALIGN.RIGHT)


def _sev_badge(slide, severity: str, left, top, width=None, height=None) -> None:
    """Colored severity rectangle label."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    w = width  or Inches(1.1)
    h = height or Inches(0.28)
    color = _SEV_RGB.get(severity.lower(), _SEV_RGB["info"])
    _rect(slide, left, top, w, h, fill=color)
    _text(slide, _SEV_LABEL.get(severity.lower(), severity.upper()),
          left + Inches(0.04), top + Inches(0.03), w - Inches(0.08), h - Inches(0.04),
          size=9, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)


# ── Slide builders ─────────────────────────────────────────────────────────────

def _slide_cover(prs, title: str, case_id: str, description: str,
                 date_str: str, W, H) -> None:
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, _DARK_NAVY)

    # Left blue accent bar
    _rect(slide, 0, 0, Inches(0.12), H, fill=_BLUE)

    # Top label
    _text(slide, "DIGITAL FORENSICS & INCIDENT RESPONSE",
          Inches(0.3), Inches(0.6), Inches(9), Inches(0.4),
          size=9, bold=True, color=_LIGHT_BLUE)

    # Main title
    _text(slide, title,
          Inches(0.3), Inches(1.2), Inches(10.5), Inches(2.2),
          size=38, bold=True, color=_WHITE)

    # Description / subtitle
    if description:
        _text(slide, description,
              Inches(0.3), Inches(3.45), Inches(10.5), Inches(0.6),
              size=15, color=(0xbf, 0xdb, 0xfe))

    # Horizontal divider
    _rect(slide, Inches(0.3), Inches(4.15), Inches(1.2), Inches(0.05), fill=_BLUE)

    # Meta block
    meta_top = Inches(4.4)
    meta_gap = Inches(0.42)
    meta_items = [
        ("CASE ID",         case_id if case_id else "—"),
        ("DATE",            date_str + " UTC"),
        ("CLASSIFICATION",  "CONFIDENTIAL — RESTRICTED DISTRIBUTION"),
        ("PREPARED BY",     "Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin · Joost Beekman"),
    ]
    for label, value in meta_items:
        _text(slide, label,
              Inches(0.3), meta_top, Inches(2.5), Inches(0.38),
              size=8, bold=True, color=_LIGHT_BLUE)
        c = _AMBER if label == "CLASSIFICATION" else (0xe0, 0xf2, 0xfe)
        _text(slide, value,
              Inches(2.9), meta_top, Inches(8.0), Inches(0.38),
              size=11, color=c, bold=(label == "CLASSIFICATION"))
        meta_top += meta_gap

    # Bottom bar
    _rect(slide, 0, H - Inches(0.55), W, Inches(0.55),
          fill=(0x08, 0x0f, 0x1e))
    _text(slide, "FanGetFameFast  ·  Forensics Agent Network",
          Inches(0.3), H - Inches(0.48), Inches(7), Inches(0.38),
          size=9, color=_TEXT_MID)
    _text(slide, f"Generated {date_str} UTC",
          Inches(9.5), H - Inches(0.48), Inches(3.6), Inches(0.38),
          size=9, color=_TEXT_MID, align=PP_ALIGN.RIGHT)


def _slide_exec_summary(prs, data: dict, overall_sev: str,
                        case_id: str, first_ts: str, last_ts: str,
                        duration: float, W, H,
                        narrative: dict | None = None) -> None:
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    sys.path.insert(0, str(PROJECT_ROOT / "lib"))
    from generate_pcap_report import _triggered, build_recommendations, _format_duration

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, _DARK_NAVY)
    _header_bar(slide, "Executive Summary", case_id, W, H)

    CONTENT_TOP = Inches(0.8)
    CONTENT_L   = Inches(0.4)
    CONTENT_W   = Inches(12.5)

    # Severity badge
    _sev_badge(slide, overall_sev, CONTENT_L, CONTENT_TOP, width=Inches(1.3), height=Inches(0.32))

    # Opening sentence
    sev_sentences = {
        "critical": "CRITICAL threats identified — immediate containment action required.",
        "high":     "HIGH-severity threats identified — prompt investigation required.",
        "medium":   "MEDIUM-severity anomalies identified — further examination recommended.",
        "low":      "LOW-severity findings — routine follow-up recommended.",
        "info":     "No significant threat indicators identified at this time.",
    }
    _text(slide, sev_sentences.get(overall_sev, ""),
          CONTENT_L + Inches(1.45), CONTENT_TOP - Inches(0.01),
          CONTENT_W - Inches(1.5), Inches(0.35),
          size=13, bold=True, color=_WHITE)

    # Traffic scope box
    netflow_count = len(data["netflow"])
    ip_count      = len(data["unique_ips"])
    fqdn_count    = len(data["unique_fqdns"])
    dur_str       = _format_duration(duration) if duration else "—"

    BOX_TOP = CONTENT_TOP + Inches(0.45)
    BOX_H   = Inches(0.8)
    box_items = [
        (f"{netflow_count:,}", "Network flows"),
        (f"{ip_count:,}",      "Unique IPs"),
        (f"{fqdn_count:,}",   "Unique domains"),
        (dur_str,              "Capture duration"),
    ]
    box_w = Inches(2.9)
    box_gap = Inches(0.2)
    bx = CONTENT_L
    for val, lbl in box_items:
        _rect(slide, bx, BOX_TOP, box_w, BOX_H, fill=_LIGHT_BG, line=_MID_NAVY, line_width_pt=0.5)
        _text(slide, val, bx + Inches(0.12), BOX_TOP + Inches(0.05),
              box_w - Inches(0.24), Inches(0.4),
              size=20, bold=True, color=_rgb(_ELECTRIC))
        _text(slide, lbl, bx + Inches(0.12), BOX_TOP + Inches(0.44),
              box_w - Inches(0.24), Inches(0.3),
              size=10, color=_LIGHT)
        bx += box_w + box_gap

    # Key findings section
    KF_TOP = BOX_TOP + BOX_H + Inches(0.2)
    _text(slide, "KEY FINDINGS",
          CONTENT_L, KF_TOP, CONTENT_W, Inches(0.28),
          size=9, bold=True, color=_LIGHT_BLUE if False else _rgb(_BLUE))
    _rect(slide, CONTENT_L, KF_TOP + Inches(0.28), CONTENT_W, Inches(0.02), fill=_BLUE)

    all_trig = []
    for key in ("icmp_results", "dns_results", "ntp_results", "http_results",
                "cert_results", "tls_results", "arp_results", "tcp_results",
                "udp_results", "dhcp_results", "mdns_results", "quic_results",
                "snmp_results", "nbns_results", "llmnr_results",
                "stun_results", "ssdp_results", "netbios_results"):
        all_trig += _triggered(data.get(key, {}))

    from generate_pcap_report import _sev_rank
    all_trig.sort(key=lambda c: _sev_rank(c.get("severity", "info")))

    suricata = data.get("suricata_data", {})
    yara     = data.get("yara_data", {})
    fh       = data.get("fh_data", {})

    finding_lines: list[tuple] = []
    for cat in all_trig[:6]:
        sev  = cat.get("severity", "info")
        name = cat.get("name", "")
        desc = _MGMT_DESC.get(name, cat.get("description", "").split(".")[0])
        cnt  = cat.get("count", 0)
        badge = _SEV_LABEL.get(sev, "INFO").ljust(8)
        finding_lines.append(
            (f"[{badge}]  {name} — {desc}  ({cnt} finding{'s' if cnt != 1 else ''})",
             10.5, False, _SEV_RGB.get(sev, _TEXT_DARK) if sev in ("critical", "high") else _TEXT_DARK)
        )

    if suricata.get("total_alerts", 0):
        finding_lines.append(
            (f"[IDS    ]  Intrusion detection system flagged {suricata['total_alerts']} alert(s) "
             f"across {suricata.get('unique_signatures', 0)} unique rule(s)",
             10.5, False, _TEXT_DARK)
        )
    if yara.get("total_matches", 0):
        finding_lines.append(
            (f"[YARA   ]  Malware signature scan matched {yara['total_matches']} pattern(s) in network traffic",
             10.5, False, _TEXT_DARK)
        )
    if fh.get("malicious_count", 0):
        finding_lines.append(
            (f"[MALWARE]  {fh['malicious_count']} malicious file(s) extracted from network traffic",
             10.5, True, _SEV_RGB["critical"])
        )

    if not finding_lines:
        finding_lines.append(("No significant threat indicators detected.", 11, False, _LIGHT))

    # If narrative is available, replace auto-generated findings with board language
    if narrative and narrative.get("pptx_executive_summary"):
        nar_lines = [(ln, 11, False, _WHITE)
                     for ln in narrative["pptx_executive_summary"].splitlines() if ln.strip()]
        _text_lines(slide, nar_lines,
                    CONTENT_L, KF_TOP + Inches(0.35), CONTENT_W, Inches(3.2))
    else:
        _text_lines(slide, finding_lines,
                    CONTENT_L, KF_TOP + Inches(0.35), CONTENT_W, Inches(3.2))

    # Bottom CTI note
    mal_ips   = [r for r in data.get("fan_ip", []) if r.get("reputation") == "malicious"]
    mal_fqdns = [r for r in data.get("fan_correlation", []) if r.get("reputation") == "malicious"]
    if mal_ips or mal_fqdns:
        total_mal = len(mal_ips) + len(mal_fqdns)
        cti_note = (f"Threat intelligence confirmed {total_mal} malicious indicator(s) "
                    f"in this traffic capture.")
        _rect(slide, CONTENT_L, H - Inches(0.75), CONTENT_W, Inches(0.55),
              fill=(0xff, 0xf1, 0xf0), line=_SEV_RGB["critical"], line_width_pt=0.5)
        _text(slide, "⚠  " + cti_note,
              CONTENT_L + Inches(0.1), H - Inches(0.72), CONTENT_W - Inches(0.2), Inches(0.5),
              size=10, bold=True, color=_SEV_RGB["critical"])


def _slide_threats(prs, data: dict, case_id: str, W, H) -> None:
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    sys.path.insert(0, str(PROJECT_ROOT / "lib"))
    from generate_pcap_report import _triggered, _sev_rank

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, _WHITE)
    _header_bar(slide, "Threat Landscape", case_id, W, H)

    all_trig = []
    for key in ("dns_results", "http_results", "tls_results", "cert_results",
                "icmp_results", "tcp_results", "udp_results", "arp_results",
                "dhcp_results", "mdns_results", "quic_results", "ntp_results",
                "snmp_results", "nbns_results", "llmnr_results",
                "stun_results", "ssdp_results", "netbios_results"):
        for cat in _triggered(data.get(key, {})):
            cat["_proto"] = key.replace("_results", "").upper()
            all_trig.append(cat)

    all_trig.sort(key=lambda c: _sev_rank(c.get("severity", "info")))
    all_trig = all_trig[:14]  # cap for slide space

    if not all_trig:
        _text(slide, "No protocol-level threats detected in this capture.",
              Inches(0.4), Inches(1.0), Inches(12), Inches(1.0),
              size=14, color=_TEXT_MID)
        return

    # Table
    TABLE_TOP = Inches(0.72)
    TABLE_L   = Inches(0.35)
    TABLE_W   = Inches(12.6)
    TABLE_H   = H - TABLE_TOP - Inches(0.2)

    rows = len(all_trig) + 1
    tbl = slide.shapes.add_table(rows, 4, TABLE_L, TABLE_TOP, TABLE_W, TABLE_H).table

    # Column widths
    tbl.columns[0].width = Inches(1.3)   # Severity
    tbl.columns[1].width = Inches(3.1)   # Threat name
    tbl.columns[2].width = Inches(7.0)   # Description
    tbl.columns[3].width = Inches(1.2)   # Findings

    def _cell(row, col, text, bold=False, size=10,
               fg=_TEXT_DARK, bg=None, align=PP_ALIGN.LEFT):
        cell = tbl.cell(row, col)
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(bg)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(fg)

    # Header row
    for col, (hdr, algn) in enumerate([
        ("Severity",    PP_ALIGN.CENTER),
        ("Threat",      PP_ALIGN.LEFT),
        ("Description", PP_ALIGN.LEFT),
        ("Events",      PP_ALIGN.CENTER),
    ]):
        _cell(0, col, hdr, bold=True, size=9, fg=_WHITE, bg=_MID_NAVY, align=algn)

    for i, cat in enumerate(all_trig):
        row = i + 1
        sev  = cat.get("severity", "info")
        name = cat.get("name", "")
        desc = _MGMT_DESC.get(name, cat.get("description", "").split(".")[0])
        cnt  = cat.get("count", 0)
        bg = _ROW_ALT if i % 2 == 0 else _WHITE
        sev_rgb = _SEV_RGB.get(sev, _SEV_RGB["info"])

        _cell(row, 0, _SEV_LABEL.get(sev, sev.upper()),
              bold=True, size=9, fg=sev_rgb, bg=bg, align=PP_ALIGN.CENTER)
        _cell(row, 1, name, bold=False, size=9.5, fg=_TEXT_DARK, bg=bg)
        _cell(row, 2, desc or "—", bold=False, size=9, fg=_TEXT_MID, bg=bg)
        _cell(row, 3, str(cnt), bold=True, size=9.5, fg=_TEXT_DARK, bg=bg, align=PP_ALIGN.CENTER)


def _slide_alerts(prs, data: dict, case_id: str, W, H) -> None:
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, _WHITE)
    _header_bar(slide, "Security Alerts — IDS & Malware Signatures", case_id, W, H)

    CONTENT_TOP = Inches(0.82)
    CONTENT_L   = Inches(0.4)
    CONTENT_W   = Inches(12.5)

    suricata = data.get("suricata_data", {})
    yara     = data.get("yara_data", {})
    fh       = data.get("fh_data", {})

    # ── Suricata panel ──────────────────────────────────────────────────────────
    PANEL_W = Inches(6.0)
    PANEL_H = Inches(5.5)
    _rect(slide, CONTENT_L, CONTENT_TOP, PANEL_W, PANEL_H,
          fill=_LIGHT_BG, line=_MID_NAVY, line_width_pt=0.5)
    _text(slide, "INTRUSION DETECTION SYSTEM (Suricata)",
          CONTENT_L + Inches(0.15), CONTENT_TOP + Inches(0.12),
          PANEL_W - Inches(0.3), Inches(0.3),
          size=10, bold=True, color=_rgb(_MID_NAVY))

    sur_total = suricata.get("total_alerts", 0)
    if sur_total == 0:
        _text(slide, "No IDS alerts triggered.",
              CONTENT_L + Inches(0.2), CONTENT_TOP + Inches(0.55),
              PANEL_W - Inches(0.4), Inches(0.4),
              size=12, color=_TEXT_MID)
    else:
        stats = [
            (str(sur_total),                          "Total alerts"),
            (str(suricata.get("unique_signatures", 0)), "Unique rules triggered"),
            (str(suricata.get("critical_count", 0)),  "Critical severity"),
            (str(suricata.get("high_count", 0)),      "High severity"),
            (str(suricata.get("medium_count", 0)),    "Medium severity"),
        ]
        sy = CONTENT_TOP + Inches(0.52)
        for val, lbl in stats:
            _text(slide, val, CONTENT_L + Inches(0.2), sy,
                  Inches(1.0), Inches(0.32), size=16, bold=True, color=_rgb(_BLUE))
            _text(slide, lbl, CONTENT_L + Inches(1.35), sy + Inches(0.04),
                  PANEL_W - Inches(1.55), Inches(0.28), size=10.5, color=_TEXT_DARK)
            sy += Inches(0.48)

        # Top signatures
        sigs = suricata.get("top_signatures", [])[:5]
        if sigs:
            _text(slide, "Top triggered rules:",
                  CONTENT_L + Inches(0.2), sy + Inches(0.1),
                  PANEL_W - Inches(0.4), Inches(0.3),
                  size=9.5, bold=True, color=_TEXT_DARK)
            sy += Inches(0.4)
            for sig in sigs:
                name = str(sig.get("signature", sig.get("name", "—")))[:70]
                cnt  = sig.get("count", "")
                _text(slide, f"• {name}  ({cnt})",
                      CONTENT_L + Inches(0.25), sy,
                      PANEL_W - Inches(0.45), Inches(0.32),
                      size=8.5, color=_TEXT_MID)
                sy += Inches(0.3)

    # ── YARA + File panel ───────────────────────────────────────────────────────
    PANEL2_L = CONTENT_L + PANEL_W + Inches(0.3)
    PANEL2_W = W - PANEL2_L - Inches(0.35)

    # YARA sub-panel
    YARA_H = Inches(2.5)
    _rect(slide, PANEL2_L, CONTENT_TOP, PANEL2_W, YARA_H,
          fill=_LIGHT_BG, line=_MID_NAVY, line_width_pt=0.5)
    _text(slide, "YARA MALWARE SIGNATURES",
          PANEL2_L + Inches(0.15), CONTENT_TOP + Inches(0.12),
          PANEL2_W - Inches(0.3), Inches(0.3),
          size=10, bold=True, color=_rgb(_MID_NAVY))

    yara_total = yara.get("total_matches", 0)
    if yara_total == 0:
        _text(slide, "No YARA matches.",
              PANEL2_L + Inches(0.2), CONTENT_TOP + Inches(0.55),
              PANEL2_W - Inches(0.4), Inches(0.4),
              size=12, color=_TEXT_MID)
    else:
        _text(slide, str(yara_total),
              PANEL2_L + Inches(0.2), CONTENT_TOP + Inches(0.52),
              Inches(1.2), Inches(0.55), size=30, bold=True, color=_rgb(_SEV_RGB["high"]))
        _text(slide, "malware pattern match(es)",
              PANEL2_L + Inches(1.55), CONTENT_TOP + Inches(0.63),
              PANEL2_W - Inches(1.7), Inches(0.35), size=11, color=_TEXT_DARK)
        cats = yara.get("matched_categories", [])[:4]
        if cats:
            yy = CONTENT_TOP + Inches(1.2)
            for c in cats:
                _text(slide, f"• {c}",
                      PANEL2_L + Inches(0.2), yy, PANEL2_W - Inches(0.4), Inches(0.28),
                      size=9, color=_TEXT_MID)
                yy += Inches(0.28)

    # File hashes sub-panel
    FH_TOP = CONTENT_TOP + YARA_H + Inches(0.25)
    FH_H   = PANEL_H - YARA_H - Inches(0.25)
    _rect(slide, PANEL2_L, FH_TOP, PANEL2_W, FH_H,
          fill=_LIGHT_BG, line=_MID_NAVY, line_width_pt=0.5)
    _text(slide, "FILE EXTRACTION & ANALYSIS",
          PANEL2_L + Inches(0.15), FH_TOP + Inches(0.12),
          PANEL2_W - Inches(0.3), Inches(0.3),
          size=10, bold=True, color=_rgb(_MID_NAVY))

    total_files = fh.get("total_files", 0)
    mal_files   = fh.get("malicious_count", 0)
    sus_files   = fh.get("suspicious_count", 0)
    if total_files == 0:
        _text(slide, "No files extracted.",
              PANEL2_L + Inches(0.2), FH_TOP + Inches(0.55),
              PANEL2_W - Inches(0.4), Inches(0.4),
              size=12, color=_TEXT_MID)
    else:
        fh_stats = [
            (str(total_files),  "Files extracted",    _TEXT_DARK),
            (str(mal_files),    "Confirmed malicious", _SEV_RGB["critical"] if mal_files else _TEXT_DARK),
            (str(sus_files),    "Suspicious",          _SEV_RGB["high"] if sus_files else _TEXT_DARK),
        ]
        fy = FH_TOP + Inches(0.52)
        for val, lbl, col in fh_stats:
            _text(slide, val, PANEL2_L + Inches(0.2), fy,
                  Inches(1.0), Inches(0.32), size=18, bold=True, color=_rgb(col))
            _text(slide, lbl, PANEL2_L + Inches(1.35), fy + Inches(0.05),
                  PANEL2_W - Inches(1.55), Inches(0.26), size=10, color=_TEXT_DARK)
            fy += Inches(0.42)


def _slide_iocs(prs, iocs: list[dict], case_id: str, W, H) -> None:
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, _WHITE)
    _header_bar(slide, "Indicators of Compromise", case_id, W, H)

    CONTENT_TOP = Inches(0.82)
    CONTENT_L   = Inches(0.35)
    CONTENT_W   = Inches(12.6)

    # Filter to high/critical or top-N
    display = [i for i in iocs if i.get("severity") in ("critical", "high")][:20]
    if not display:
        display = iocs[:20]

    if not display:
        _text(slide, "No indicators of compromise extracted from this capture.",
              CONTENT_L, CONTENT_TOP + Inches(0.3), CONTENT_W, Inches(0.6),
              size=14, color=_TEXT_MID)
        return

    # Note about defanging
    _text(slide, "IOC values are defanged for safe handling. Refang before blocking.",
          CONTENT_L, CONTENT_TOP, CONTENT_W, Inches(0.3),
          size=9, italic=True, color=_TEXT_MID)

    TABLE_TOP = CONTENT_TOP + Inches(0.35)
    TABLE_H   = H - TABLE_TOP - Inches(0.2)
    rows = len(display) + 1
    tbl = slide.shapes.add_table(rows, 4, CONTENT_L, TABLE_TOP, CONTENT_W, TABLE_H).table

    tbl.columns[0].width = Inches(1.2)   # Severity
    tbl.columns[1].width = Inches(1.4)   # Type
    tbl.columns[2].width = Inches(7.0)   # Value (defanged)
    tbl.columns[3].width = Inches(3.0)   # Source / Category

    def _cell(row, col, text, bold=False, size=9.5,
              fg=_TEXT_DARK, bg=None, align=PP_ALIGN.LEFT):
        cell = tbl.cell(row, col)
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(bg)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(fg)

    for col, hdr in enumerate(["Severity", "Type", "Indicator (defanged)", "Source"]):
        _cell(0, col, hdr, bold=True, size=9, fg=_WHITE, bg=_MID_NAVY)

    def _defang(ioc_type: str, value: str) -> str:
        if ioc_type == "ip":
            return value.replace(".", "[.]")
        if ioc_type == "domain":
            return value.replace(".", "[.]")
        if ioc_type == "url":
            return value.replace("https://", "hxxps://").replace("http://", "hxxp://").replace(".", "[.]")
        return value

    for i, ioc in enumerate(display):
        row   = i + 1
        sev   = ioc.get("severity", "info")
        itype = ioc.get("type", "—")
        val   = _defang(itype, ioc.get("value", "—"))
        src   = ioc.get("source", "—")
        bg    = _ROW_ALT if i % 2 == 0 else _WHITE
        sev_c = _SEV_RGB.get(sev, _SEV_RGB["info"])

        _cell(row, 0, _SEV_LABEL.get(sev, sev.upper()),
              bold=True, size=9, fg=sev_c, bg=bg, align=PP_ALIGN.CENTER)
        _cell(row, 1, itype.upper(), bold=False, size=9, fg=_TEXT_MID, bg=bg)
        _cell(row, 2, val, bold=False, size=9, fg=_TEXT_DARK, bg=bg)
        _cell(row, 3, src, bold=False, size=9, fg=_TEXT_MID, bg=bg)


def _slide_recommendations(prs, recs: list[str], case_id: str, W, H,
                           narrative: dict | None = None) -> None:
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, _DARK_NAVY)
    _header_bar(slide, "Recommendations & Next Steps", case_id, W, H)

    CONTENT_TOP = Inches(0.82)
    CONTENT_L   = Inches(0.4)
    CONTENT_W   = Inches(12.5)

    def _strip_md(s: str) -> str:
        return s.replace("**IMMEDIATE**: ", "IMMEDIATE ACTION: ").replace("**", "")

    # Prefer narrative board language over auto-generated recs
    if narrative and narrative.get("pptx_recommendations"):
        display = [ln for ln in narrative["pptx_recommendations"].splitlines() if ln.strip()][:8]
    else:
        if not recs:
            recs = ["Continue routine monitoring. No critical actions required at this time."]
        display = [_strip_md(r) for r in recs[:8]]

    item_h = (H - CONTENT_TOP - Inches(0.3)) / max(len(display), 1)
    item_h = min(item_h, Inches(0.75))

    for idx, rec in enumerate(display):
        row_top = CONTENT_TOP + idx * item_h
        is_immediate = rec.upper().startswith("IMMEDIATE")
        num_color = _ALERT if is_immediate else _BLUE
        _rect(slide, CONTENT_L, row_top + Inches(0.04),
              Inches(0.36), Inches(0.36), fill=num_color)
        _text(slide, str(idx + 1),
              CONTENT_L + Inches(0.04), row_top + Inches(0.05),
              Inches(0.28), Inches(0.3),
              size=12, bold=True, color=_WHITE,
              align=__import__("pptx.enum.text", fromlist=["PP_ALIGN"]).PP_ALIGN.CENTER)
        text_color = _ALERT if is_immediate else _WHITE
        _text(slide, rec, CONTENT_L + Inches(0.48), row_top,
              CONTENT_W - Inches(0.55), item_h - Inches(0.06),
              size=11, bold=is_immediate, color=text_color)
        if idx < len(display) - 1:
            _rect(slide, CONTENT_L, row_top + item_h - Inches(0.02),
                  CONTENT_W, Inches(0.01), fill=_MID_NAVY)


def _slide_timelines(prs, case_id: str, reports_dir: Path, W, H) -> None:
    """Slides 3+ — Attack Timelines: one PPTX slide per PNG page for each timeline type."""
    from pptx.util import Inches

    CONTENT_TOP = Inches(0.75)
    CONTENT_L   = Inches(0.3)
    IMG_W       = W - Inches(0.6)

    timeline_groups = [
        ("Attacker Perspective",               f"{case_id}_timeline_attacker_p*.png"),
        ("Attacker Perspective — Unconfirmed", f"{case_id}_timeline_attacker_unconfirmed.png"),
        ("Defender Perspective",               f"{case_id}_timeline_defender_p*.png"),
        ("Combined Key Events",                f"{case_id}_timeline_combined_p*.png"),
    ]

    any_slide_added = False
    for group_title, glob_pattern in timeline_groups:
        png_files = sorted(reports_dir.glob(glob_pattern))
        if not png_files:
            continue
        for i, png_path in enumerate(png_files):
            if not png_path.exists():
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _set_bg(slide, _DARK_NAVY)
            page_label = f"Page {i + 1} of {len(png_files)}" if len(png_files) > 1 else ""
            slide_title = f"{group_title}  {page_label}".strip()
            _header_bar(slide, slide_title, case_id, W, H)
            slide.shapes.add_picture(str(png_path), CONTENT_L, CONTENT_TOP, IMG_W)
            any_slide_added = True

    if not any_slide_added:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, _DARK_NAVY)
        _header_bar(slide, "Attack Timelines", case_id, W, H)
        _text(slide, "Timeline images not available — run investigation to generate.",
              CONTENT_L, CONTENT_TOP + Inches(1.0), IMG_W, Inches(0.5),
              size=12, color=_LIGHT)


def _slide_risk_impact(prs, case_id: str, W, H, narrative: dict | None = None) -> None:
    """Slide 4 — Risk & Impact: board-level, no technical identifiers."""
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, _DARK_NAVY)
    _header_bar(slide, "Risk & Impact", case_id, W, H)

    CONTENT_TOP = Inches(0.82)
    CONTENT_L   = Inches(0.4)
    HALF_W      = Inches(6.0)
    GAP         = Inches(0.3)

    # Risk panel (left)
    _rect(slide, CONTENT_L, CONTENT_TOP, HALF_W, Inches(5.8),
          fill=_MID_NAVY, line=_BLUE, line_width_pt=0.5)
    _text(slide, "BUSINESS RISK",
          CONTENT_L + Inches(0.2), CONTENT_TOP + Inches(0.15),
          HALF_W - Inches(0.4), Inches(0.3),
          size=11, bold=True, color=_ELECTRIC)

    risk_text = (narrative or {}).get("pptx_risk",
        "Risk assessment not yet generated.\n"
        "Run the forensic skill and write the narrative file to populate this slide.")
    _text(slide, risk_text,
          CONTENT_L + Inches(0.2), CONTENT_TOP + Inches(0.55),
          HALF_W - Inches(0.4), Inches(5.0),
          size=11, color=_WHITE)

    # Impact panel (right)
    IMP_L = CONTENT_L + HALF_W + GAP
    IMP_W = W - IMP_L - Inches(0.35)
    _rect(slide, IMP_L, CONTENT_TOP, IMP_W, Inches(5.8),
          fill=_MID_NAVY, line=_BLUE, line_width_pt=0.5)
    _text(slide, "OPERATIONAL IMPACT",
          IMP_L + Inches(0.2), CONTENT_TOP + Inches(0.15),
          IMP_W - Inches(0.4), Inches(0.3),
          size=11, bold=True, color=_ELECTRIC)

    impact_text = (narrative or {}).get("pptx_impact",
        "Impact assessment not yet generated.\n"
        "Run the forensic skill and write the narrative file to populate this slide.")
    _text(slide, impact_text,
          IMP_L + Inches(0.2), CONTENT_TOP + Inches(0.55),
          IMP_W - Inches(0.4), Inches(5.0),
          size=11, color=_WHITE)


def _slide_mitigations(prs, case_id: str, W, H, narrative: dict | None = None) -> None:
    """Slide 5 — Mitigations: what has been done and what is in progress."""
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, _DARK_NAVY)
    _header_bar(slide, "Mitigations", case_id, W, H)

    CONTENT_TOP = Inches(0.85)
    CONTENT_L   = Inches(0.4)
    CONTENT_W   = Inches(12.5)

    mit_text = (narrative or {}).get("pptx_mitigations",
        "Mitigations not yet documented.\n"
        "Run the forensic skill and write the narrative file to populate this slide.")

    lines = [(ln, 12, False, _WHITE) for ln in mit_text.splitlines() if ln.strip()]
    if not lines:
        lines = [("No mitigations recorded yet.", 12, False, _LIGHT)]

    item_h = (H - CONTENT_TOP - Inches(0.4)) / max(len(lines), 1)
    item_h = min(item_h, Inches(0.65))

    for idx, (txt, size, bold, color) in enumerate(lines[:9]):
        row_top = CONTENT_TOP + idx * item_h
        _rect(slide, CONTENT_L, row_top + Inches(0.12),
              Inches(0.12), Inches(0.12), fill=_ELECTRIC)
        _text(slide, txt,
              CONTENT_L + Inches(0.28), row_top,
              CONTENT_W - Inches(0.35), item_h - Inches(0.06),
              size=size, bold=bold, color=color)


def _slide_coverage(prs, data: dict, case_id: str, W, H) -> None:
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, _DARK_NAVY)
    _header_bar(slide, "Investigation Coverage — 23 Detection Modules", case_id, W, H)

    CONTENT_TOP = Inches(0.72)
    CONTENT_L   = Inches(0.35)

    COLS     = 3
    COL_W    = Inches(4.1)
    COL_GAP  = Inches(0.12)
    ROW_H    = Inches(0.36)
    items_per_col = (len(_MODULES) + COLS - 1) // COLS

    for idx, (_, has_key, friendly) in enumerate(_MODULES):
        col = idx // items_per_col
        row = idx % items_per_col
        lx  = CONTENT_L + col * (COL_W + COL_GAP)
        ly  = CONTENT_TOP + Inches(0.55) + row * ROW_H

        ran     = data.get(has_key, False)
        dot_col = (0x22, 0xc5, 0x5e) if ran else (0xd1, 0xd5, 0xdb)
        _rect(slide, lx, ly + Inches(0.09), Inches(0.14), Inches(0.14), fill=dot_col)
        status = "Complete" if ran else "No data"
        _text(slide, friendly,
              lx + Inches(0.22), ly, COL_W - Inches(0.25), ROW_H - Inches(0.04),
              size=10, color=_WHITE if ran else _LIGHT)
        _text(slide, status,
              lx + COL_W - Inches(1.05), ly + Inches(0.06), Inches(0.95), Inches(0.24),
              size=8.5, bold=False,
              color=(0x22, 0xc5, 0x5e) if ran else _LIGHT,
              align=PP_ALIGN.RIGHT)

    # Footer count
    run_count = sum(1 for _, k, _ in _MODULES if data.get(k))
    _text(slide,
          f"{run_count} of {len(_MODULES)} modules returned results for this capture.",
          CONTENT_L, H - Inches(0.5), Inches(10), Inches(0.38),
          size=10, italic=True, color=_LIGHT)


# ── Public API ─────────────────────────────────────────────────────────────────

def _load_pptx_narrative(case_id: str, reports_dir: Path) -> dict[str, str]:
    """Load Claude-generated narrative for PPTX from {case_id}_narrative.md."""
    if not case_id:
        return {}
    path = reports_dir / f"{case_id}_narrative.md"
    if not path.exists():
        return {}
    sections: dict[str, str] = {}
    current: str | None = None
    for line in path.read_text(encoding="utf-8").splitlines():
        if line.startswith("## "):
            current = line[3:].strip()
            sections[current] = ""
        elif line.startswith("<!--"):
            continue
        elif current is not None:
            sections[current] += line + "\n"
    return {k: v.strip() for k, v in sections.items()}


def generate(
    data: dict,
    stem: str,
    case_id: str = "",
    output_dir: Path | None = None,
    description: str = "",
    reports_dir: Path | None = None,
) -> Path:
    """Build the 7-slide board management PowerPoint and return its path.

    Slide order:
      1  Cover
      2  Executive Summary
      3  Attack Timeline (PNG)
      4  Risk & Impact
      5  Mitigations
      6  Recommendations & Next Steps
      7  Investigation Coverage
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise SystemExit(
            "[pptx] python-pptx not installed.\n"
            "Install: pip3 install python-pptx"
        )

    sys.path.insert(0, str(PROJECT_ROOT / "lib"))
    from generate_pcap_report import (
        _overall_severity, _capture_window, extract_iocs, build_recommendations,
    )

    out_dir = output_dir or (PROJECT_ROOT / "analysis" / "_reports" / stem)
    path_guard.guard_output_dir(out_dir)
    out_path = out_dir / f"{stem}_management_briefing.pptx"

    rpts_dir = reports_dir or (PROJECT_ROOT / "reports")
    narrative = _load_pptx_narrative(case_id, rpts_dir)

    overall_sev = _overall_severity(data)
    first_ts, last_ts, duration = _capture_window(data)
    iocs = extract_iocs(data)
    recs = build_recommendations(data, overall_sev)

    date_str = datetime.now(timezone.utc).strftime("%Y-%m-%d")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W = prs.slide_width
    H = prs.slide_height

    # Slide 1 — Cover
    _slide_cover(prs, title="Security Incident Briefing",
                 case_id=case_id, description=description,
                 date_str=date_str, W=W, H=H)

    # Slide 2 — Executive Summary
    _slide_exec_summary(prs, data, overall_sev, case_id,
                        first_ts, last_ts, duration, W, H, narrative=narrative)

    # Slide — Risk & Impact
    _slide_risk_impact(prs, case_id, W, H, narrative=narrative)

    # Slide 5 — Mitigations
    _slide_mitigations(prs, case_id, W, H, narrative=narrative)

    # Slide 6 — Recommendations & Next Steps
    _slide_recommendations(prs, recs, case_id, W, H, narrative=narrative)

    # Slide 7 — Investigation Coverage
    _slide_coverage(prs, data, case_id, W, H)

    prs.save(str(out_path))
    print(f"[pptx] PowerPoint written: {out_path}")
    return out_path


# ── CLI ────────────────────────────────────────────────────────────────────────

def generate_board_deck(
    case_id: str,
    module: str = "fame",
    hostname: str = "",
    description: str = "",
    output_path: Path | None = None,
    reports_dir: Path | None = None,
) -> Path:
    """Build the 7-slide board deck from narrative + research notes only.

    No FAN pcap data required. Works for FAME, FAST, and FAN cases.
    Reads {case_id}_narrative.md and {case_id}_research_notes.md from reports_dir.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise SystemExit("[pptx] python-pptx not installed. Install: pip3 install python-pptx")

    sys.path.insert(0, str(PROJECT_ROOT / "lib"))

    rpts_dir = reports_dir or (PROJECT_ROOT / "reports")
    narrative = _load_pptx_narrative(case_id, rpts_dir)

    # Load research note steps and attacker events for timelines
    from research_notes import parse_steps
    steps  = parse_steps(case_id, str(rpts_dir))
    # Determine output path
    if output_path is None:
        suffix = f"_{module}_board_deck.pptx"
        output_path = rpts_dir / f"{case_id}{suffix}"

    # Severity from narrative (scan for keywords)
    nar_text = " ".join(narrative.values()).upper()
    overall_sev = "info"
    for sev in ("CRITICAL", "HIGH", "MEDIUM", "LOW"):
        if sev in nar_text:
            overall_sev = sev.lower()
            break

    # Module coverage: derive from step titles
    skip_prefixes = ("evidence preserved:", "sha256")
    analysis_steps = [
        s["title"] for s in steps
        if not any(s["title"].lower().startswith(p) for p in skip_prefixes)
    ]

    module_label = module.upper()
    cover_title  = f"Security Incident Briefing — {module_label}"
    date_str     = datetime.now(timezone.utc).strftime("%Y-%m-%d")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W = prs.slide_width
    H = prs.slide_height

    # Slide 1 — Cover
    _slide_cover(prs,
                 title=cover_title,
                 case_id=case_id,
                 description=description or (f"Host: {hostname}" if hostname else ""),
                 date_str=date_str,
                 W=W, H=H)

    # Slide 2 — Executive Summary (narrative-only variant)
    _slide_exec_summary_narrative(prs, case_id, overall_sev, hostname, narrative, W, H)

    # Slide 4 — Risk & Impact
    _slide_risk_impact(prs, case_id, W, H, narrative=narrative)

    # Slide 5 — Mitigations
    _slide_mitigations(prs, case_id, W, H, narrative=narrative)

    # Slide 6 — Recommendations
    _slide_recommendations(prs, [], case_id, W, H, narrative=narrative)

    # Slide 7 — Module Coverage (analysis steps run)
    _slide_coverage_steps(prs, case_id, module_label, analysis_steps, W, H)

    prs.save(str(output_path))
    print(f"[pptx] Board deck written: {output_path}")
    return output_path


def _slide_exec_summary_narrative(prs, case_id: str, overall_sev: str,
                                   hostname: str, narrative: dict, W, H) -> None:
    """Slide 2 variant for FAME/FAST: narrative bullets + key stat placeholders."""
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, _DARK_NAVY)
    _header_bar(slide, "Executive Summary", case_id, W, H)

    CONTENT_TOP = Inches(0.82)
    CONTENT_L   = Inches(0.4)
    CONTENT_W   = Inches(12.5)

    # Severity badge
    _sev_badge(slide, overall_sev, CONTENT_L, CONTENT_TOP, width=Inches(1.3), height=Inches(0.32))

    sev_sentences = {
        "critical": "CRITICAL threats identified — immediate action required.",
        "high":     "HIGH-severity threats identified — prompt investigation required.",
        "medium":   "MEDIUM-severity anomalies — further examination recommended.",
        "low":      "LOW-severity findings — routine follow-up recommended.",
        "info":     "No significant threat indicators at this time.",
    }
    _text(slide, sev_sentences.get(overall_sev, ""),
          CONTENT_L + Inches(1.45), CONTENT_TOP - Inches(0.01),
          CONTENT_W - Inches(1.5), Inches(0.35),
          size=13, bold=True, color=_WHITE)

    if hostname:
        _text(slide, f"Host: {hostname}",
              CONTENT_L + Inches(1.45), CONTENT_TOP + Inches(0.3),
              CONTENT_W - Inches(1.5), Inches(0.28),
              size=10, color=_LIGHT)

    # Narrative bullets
    KF_TOP = CONTENT_TOP + Inches(0.72)
    _text(slide, "KEY FINDINGS",
          CONTENT_L, KF_TOP, CONTENT_W, Inches(0.28),
          size=9, bold=True, color=_ELECTRIC)
    _rect(slide, CONTENT_L, KF_TOP + Inches(0.28), CONTENT_W, Inches(0.02), fill=_ELECTRIC)

    nar = narrative.get("pptx_executive_summary", "")
    if nar:
        bullets = [(ln, 11, False, _WHITE) for ln in nar.splitlines() if ln.strip()]
    else:
        bullets = [("Refer to the full technical report for detailed findings.", 11, False, _LIGHT)]

    _text_lines(slide, bullets[:6],
                CONTENT_L, KF_TOP + Inches(0.35), CONTENT_W, Inches(5.0))


def _slide_coverage_steps(prs, case_id: str, module_label: str,
                           analysis_steps: list[str], W, H) -> None:
    """Slide 7 for board deck: show which analysis steps were run."""
    from pptx.util import Inches
    from pptx.enum.text import PP_ALIGN

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, _DARK_NAVY)
    _header_bar(slide, f"Investigation Coverage — {module_label} Analysis", case_id, W, H)

    CONTENT_TOP = Inches(0.82)
    CONTENT_L   = Inches(0.35)

    if not analysis_steps:
        _text(slide, "No analysis steps recorded in research notes.",
              CONTENT_L, CONTENT_TOP + Inches(0.5), Inches(12), Inches(0.5),
              size=12, color=_LIGHT)
        return

    COLS = 2
    items_per_col = (len(analysis_steps) + COLS - 1) // COLS
    COL_W = Inches(6.1)
    ROW_H = Inches(0.34)

    for idx, step_title in enumerate(analysis_steps[:20]):
        col = idx // items_per_col
        row = idx % items_per_col
        lx  = CONTENT_L + col * (COL_W + Inches(0.2))
        ly  = CONTENT_TOP + Inches(0.1) + row * ROW_H

        _rect(slide, lx, ly + Inches(0.1), Inches(0.12), Inches(0.12), fill=_ELECTRIC)
        _text(slide, step_title,
              lx + Inches(0.2), ly, COL_W - Inches(0.25), ROW_H - Inches(0.04),
              size=9.5, color=_WHITE)

    _text(slide, f"{len(analysis_steps)} investigation steps completed. Full details in research notes.",
          CONTENT_L, H - Inches(0.5), Inches(12), Inches(0.35),
          size=9, italic=True, color=_LIGHT)


def _build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        description="Generate a management PowerPoint briefing"
    )
    sub = p.add_subparsers(dest="command")

    # Legacy FAN subcommand (default when no subcommand given)
    pfan = sub.add_parser("fan", help="FAN PCAP-based board deck")
    pfan.add_argument("--stem",        required=True, metavar="STEM")
    pfan.add_argument("--case-id",     default="",    metavar="ID")
    pfan.add_argument("--description", default="",    metavar="DESC")
    pfan.add_argument("--output-dir",  default="",    metavar="DIR")
    pfan.add_argument("--base-dir",    default="",    metavar="DIR")

    # Narrative-only board deck (FAME/FAST/FAN — no analysis data required)
    pboard = sub.add_parser("board-deck", help="7-slide board deck from narrative file")
    pboard.add_argument("--case-id",     required=True, metavar="ID")
    pboard.add_argument("--module",      default="fame", choices=["fame", "fast", "fan"])
    pboard.add_argument("--hostname",    default="",    metavar="NAME")
    pboard.add_argument("--description", default="",    metavar="DESC")
    pboard.add_argument("--output",      default="",    metavar="PATH",
                        help="Output PPTX path (default: reports/{case_id}_{module}_board_deck.pptx)")
    pboard.add_argument("--reports-dir", default=str(PROJECT_ROOT / "reports"), metavar="DIR")

    return p


if __name__ == "__main__":
    args = _build_parser().parse_args()

    if args.command == "board-deck":
        out = Path(args.output) if args.output else None
        generate_board_deck(
            case_id     = args.case_id,
            module      = args.module,
            hostname    = args.hostname,
            description = args.description,
            output_path = out,
            reports_dir = Path(args.reports_dir),
        )
    else:
        # Legacy FAN path (or no subcommand: --stem required)
        sys.path.insert(0, str(PROJECT_ROOT / "lib"))
        from generate_pcap_report import load_all_data, ANALYSIS_DIR
        import generate_pcap_report as _rpt

        if not hasattr(args, "stem"):
            _build_parser().print_help()
            sys.exit(1)

        if args.base_dir:
            _rpt.ANALYSIS_DIR = Path(args.base_dir)

        data = load_all_data(args.stem)
        out_dir = Path(args.output_dir) if args.output_dir else None
        generate(data, args.stem, args.case_id, out_dir, args.description)
