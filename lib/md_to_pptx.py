#!/usr/bin/env python3
# SPDX-License-Identifier: MIT OR Apache-2.0
# SPDX-FileCopyrightText: 2026 Richard de Vries · Jeffrey Everling · Malin Janssen · Suzanne Maquelin · Joost Beekman
"""
md_to_pptx.py — Generic Markdown -> PowerPoint renderer.

Renders one slide per top-level (`##`) or sub-level (`###`) heading, in the
style of the FanGetFameFast campaign-report deck: a dark-navy theme cover
slide, a "Case Overview" slide for the document's header table, then one
slide per section/subsection. Pipe tables become grids, bullet/numbered
lists become bulleted text, plain paragraphs become wrapped text blocks.
`> Claude:` blockquote authoring directives are dropped — they must never
reach the rendered deck.

This is intentionally generic: it has no knowledge of FAN/FAME/FAST report
structure beyond "one slide per heading". It is meant for hand-authored
campaign reports written against docs/campaign_report_template.md.

Usage (CLI):
    python3 lib/md_to_pptx.py /path/to/campaign_report.md \\
        --output /path/to/campaign_presentation.pptx \\
        --case-id CASE-2026-001

Python API:
    from lib.md_to_pptx import convert
    pptx_path = convert(md_path, output_path, case_id="CASE-2026-001")
"""
from __future__ import annotations

import argparse
import re
import sys
from datetime import datetime, timezone
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import path_guard  # noqa: E402  write-path policy enforcement
from generate_pptx_report import (  # noqa: E402
    _set_bg, _rect, _text, _text_lines, _header_bar, _slide_cover,
    _DARK_NAVY, _MID_NAVY, _BLUE, _ELECTRIC, _WHITE, _LIGHT, _ALERT,
    _LIGHT_BLUE, _AMBER,
)

_TABLE_ROW_RE = re.compile(r"^\|(.+)\|\s*$")
_BULLET_RE = re.compile(r"^[-*]\s+(.*)$")
_NUMBERED_RE = re.compile(r"^\d+[.)]\s+(.*)$")
_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)$")

_MAX_TABLE_ROWS_PER_SLIDE = 12
_MAX_LIST_ITEMS_PER_SLIDE = 12


# ── Markdown helpers (self-contained — avoid importing the deprecated
#    generate_combined_report module for this rendering path) ──────────────

def _strip_md(text: str) -> str:
    """Strip markdown links/emphasis/code spans for plain-text PPTX cells."""
    text = re.sub(r"\[([^\]]*)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"\*(.*?)\*", r"\1", text)
    text = re.sub(r"`(.*?)`", r"\1", text)
    return text.strip()


def _is_separator_row(cells: list[str]) -> bool:
    return all(re.fullmatch(r":?-{2,}:?", c) for c in cells)


def _parse_table_rows(lines: list[str]) -> list[list[str]]:
    rows: list[list[str]] = []
    for line in lines:
        m = _TABLE_ROW_RE.match(line.strip())
        if not m:
            continue
        cells = [c.strip() for c in m.group(1).split("|")]
        if _is_separator_row(cells):
            continue
        rows.append(cells)
    return rows


# ── Document model ───────────────────────────────────────────────────────

class _Block:
    """A heading (or the implicit pre-heading preamble) plus its body lines."""

    def __init__(self, level: int, heading: str):
        self.level = level
        self.heading = heading
        self.lines: list[str] = []


def _split_blocks(text: str) -> list[_Block]:
    """Split markdown into blocks at every ``#``/``##``/``###`` heading.

    The implicit preamble (everything before the first heading) becomes a
    level-0 block with an empty heading.
    """
    blocks = [_Block(0, "")]
    in_code_fence = False
    for line in text.splitlines():
        if line.strip().startswith("```"):
            in_code_fence = not in_code_fence
            blocks[-1].lines.append(line)
            continue
        m = _HEADING_RE.match(line) if not in_code_fence else None
        if m and len(m.group(1)) <= 3:
            blocks.append(_Block(len(m.group(1)), _strip_md(m.group(2).strip())))
        else:
            blocks[-1].lines.append(line)
    return blocks


def _parse_elements(lines: list[str]) -> list[dict]:
    """Parse block body lines into a sequence of renderable elements."""
    elements: list[dict] = []
    i = 0
    n = len(lines)
    while i < n:
        line = lines[i]
        stripped = line.strip()

        if not stripped:
            i += 1
            continue

        # Blockquote (Claude authoring directive) — drop entirely, including
        # any continuation lines that are also blockquotes or blank.
        if stripped.startswith(">"):
            i += 1
            continue

        # Sub-heading inside a slide body (e.g. #### inside a section)
        m = _HEADING_RE.match(stripped)
        if m:
            elements.append({"type": "heading", "text": _strip_md(m.group(2).strip())})
            i += 1
            continue

        # Table
        if _TABLE_ROW_RE.match(stripped):
            table_lines = []
            while i < n and _TABLE_ROW_RE.match(lines[i].strip()):
                table_lines.append(lines[i])
                i += 1
            rows = _parse_table_rows(table_lines)
            if rows:
                elements.append({"type": "table", "header": rows[0], "rows": rows[1:]})
            continue

        # Bullet / numbered list
        if _BULLET_RE.match(stripped) or _NUMBERED_RE.match(stripped):
            items: list[str] = []
            while i < n:
                s = lines[i].strip()
                bm = _BULLET_RE.match(s)
                nm = _NUMBERED_RE.match(s)
                if bm:
                    items.append(_strip_md(bm.group(1)))
                    i += 1
                elif nm:
                    items.append(_strip_md(nm.group(1)))
                    i += 1
                elif s and not s.startswith(("#", "|", ">")) and items:
                    # continuation line of the previous bullet
                    items[-1] = items[-1] + " " + _strip_md(s)
                    i += 1
                else:
                    break
            elements.append({"type": "bullets", "items": items})
            continue

        # Horizontal rule
        if stripped == "---":
            i += 1
            continue

        # Plain paragraph (collect contiguous non-empty, non-special lines)
        para_lines = []
        while i < n:
            s = lines[i].strip()
            if not s or s.startswith(("#", "|", ">", "-", "*")) or _NUMBERED_RE.match(s) or s == "---":
                break
            para_lines.append(_strip_md(s))
            i += 1
        if para_lines:
            elements.append({"type": "text", "text": " ".join(para_lines)})
        else:
            i += 1

    return elements


# ── Slide rendering ──────────────────────────────────────────────────────

def _new_content_slide(prs, title: str, case_id: str, W, H):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, _DARK_NAVY)
    _header_bar(slide, title[:90], case_id, W, H)
    return slide


def _render_table(prs, title: str, case_id: str, header: list[str],
                   rows: list[list[str]], W, H, content_top) -> None:
    from pptx.util import Inches
    M = Inches(0.4)
    avail_w = W - 2 * M
    n_cols = max(len(header), 1)
    col_w = avail_w // n_cols
    row_h = Inches(0.42)
    chars_per_cell = max(8, int(col_w / Inches(1) * 14))

    chunks = [rows[i:i + _MAX_TABLE_ROWS_PER_SLIDE]
              for i in range(0, len(rows), _MAX_TABLE_ROWS_PER_SLIDE)] or [[]]

    for chunk_idx, chunk in enumerate(chunks):
        suffix = "" if chunk_idx == 0 else f" (continued {chunk_idx + 1}/{len(chunks)})"
        slide = _new_content_slide(prs, title + suffix, case_id, W, H)
        y = content_top
        x = M
        for h, cw in zip(header, [col_w] * n_cols):
            _rect(slide, x, y, cw - Inches(0.04), row_h - Inches(0.04), fill=_MID_NAVY)
            _text(slide, _strip_md(h)[:chars_per_cell], x + Inches(0.06), y + Inches(0.04),
                  cw - Inches(0.1), row_h, size=11, bold=True, color=_WHITE)
            x += cw
        for r_idx, row in enumerate(chunk):
            y2 = y + (r_idx + 1) * row_h
            bg = _MID_NAVY if r_idx % 2 == 0 else _DARK_NAVY
            x = M
            for cell, cw in zip(row + [""] * n_cols, [col_w] * n_cols):
                _rect(slide, x, y2, cw - Inches(0.04), row_h - Inches(0.04), fill=bg)
                _text(slide, _strip_md(cell)[:chars_per_cell], x + Inches(0.06), y2 + Inches(0.04),
                      cw - Inches(0.1), row_h, size=10, color=_LIGHT)
                x += cw


def _render_bullets(prs, title: str, case_id: str, items: list[str], W, H, content_top) -> None:
    from pptx.util import Inches
    M = Inches(0.4)
    chunks = [items[i:i + _MAX_LIST_ITEMS_PER_SLIDE]
              for i in range(0, len(items), _MAX_LIST_ITEMS_PER_SLIDE)] or [[]]
    for chunk_idx, chunk in enumerate(chunks):
        suffix = "" if chunk_idx == 0 else f" (continued {chunk_idx + 1}/{len(chunks)})"
        slide = _new_content_slide(prs, title + suffix, case_id, W, H)
        lines = [(f"•  {item[:280]}", 13, False, _LIGHT, False) for item in chunk]
        _text_lines(slide, lines, M, content_top, W - 2 * M, H - content_top - Inches(0.3),
                     default_size=13, default_color=_LIGHT)


def _render_text(prs, title: str, case_id: str, paragraphs: list[str], W, H, content_top) -> None:
    from pptx.util import Inches
    M = Inches(0.4)
    slide = _new_content_slide(prs, title, case_id, W, H)
    lines = [(p[:600], 13, False, _LIGHT, False) for p in paragraphs]
    _text_lines(slide, lines, M, content_top, W - 2 * M, H - content_top - Inches(0.3),
                 default_size=13, default_color=_LIGHT)


def _render_block(prs, block: _Block, case_id: str, W, H) -> None:
    from pptx.util import Inches
    elements = _parse_elements(block.lines)
    if not elements:
        slide = _new_content_slide(prs, block.heading, case_id, W, H)
        return

    content_top = Inches(0.85)

    # Group consecutive "text"/"heading" elements into one slide; render
    # each table or bullet list on its own slide(s).
    text_buf: list[str] = []
    pending_title = block.heading
    first_slide_used = False

    def _flush_text():
        nonlocal text_buf, pending_title, first_slide_used
        if text_buf:
            _render_text(prs, pending_title, case_id, text_buf, W, H, content_top)
            text_buf = []
            first_slide_used = True
            pending_title = block.heading + " (cont.)"

    for el in elements:
        if el["type"] == "table":
            _flush_text()
            title = pending_title if not first_slide_used else block.heading + " (cont.)"
            _render_table(prs, title, case_id, el["header"], el["rows"], W, H, content_top)
            first_slide_used = True
            pending_title = block.heading + " (cont.)"
        elif el["type"] == "bullets":
            _flush_text()
            title = pending_title if not first_slide_used else block.heading + " (cont.)"
            _render_bullets(prs, title, case_id, el["items"], W, H, content_top)
            first_slide_used = True
            pending_title = block.heading + " (cont.)"
        elif el["type"] == "heading":
            text_buf.append(el["text"].upper())
        else:  # text
            text_buf.append(el["text"])

    _flush_text()


def _render_overview_slide(prs, preamble: _Block, h1_title: str, case_id: str, W, H) -> None:
    """Render the document's pre-heading preamble (typically the case header
    table) as a "Case Overview" slide."""
    from pptx.util import Inches
    elements = _parse_elements(preamble.lines)
    table = next((e for e in elements if e["type"] == "table"), None)
    slide = _new_content_slide(prs, "Case Overview", case_id, W, H)
    M = Inches(0.4)
    content_top = Inches(0.85)
    if table:
        row_h = Inches(0.45)
        col0_w = Inches(2.6)
        col1_w = W - 2 * M - col0_w
        for i, row in enumerate(table["rows"]):
            y = content_top + i * row_h
            label = _strip_md(row[0]) if row else ""
            value = _strip_md(row[1]) if len(row) > 1 else ""
            _rect(slide, M, y, col0_w - Inches(0.04), row_h - Inches(0.04), fill=_MID_NAVY)
            _text(slide, label, M + Inches(0.08), y + Inches(0.04), col0_w - Inches(0.12), row_h,
                  size=11, bold=True, color=_LIGHT_BLUE)
            _rect(slide, M + col0_w, y, col1_w - Inches(0.04), row_h - Inches(0.04), fill=_DARK_NAVY)
            _text(slide, value[:140], M + col0_w + Inches(0.08), y + Inches(0.04), col1_w - Inches(0.12), row_h,
                  size=11, color=_WHITE)
    else:
        _text(slide, h1_title, M, content_top, W - 2 * M, Inches(1.0), size=18, bold=True, color=_WHITE)


# ── Public API ───────────────────────────────────────────────────────────

_NUMBERING_RE = re.compile(r"^\d+(?:\.\d+)*\.?\s+")


def _normalize_heading(heading: str) -> str:
    """Strip a leading numbering prefix (e.g. '1. ', '4.2 ') for matching."""
    return _NUMBERING_RE.sub("", heading).strip().lower()


def convert(md_path: Path, output_path: Path, case_id: str = "",
            title: str = "", date_str: str = "",
            board_sections: list[str] | None = None) -> Path:
    """Render *md_path* as a PPTX deck — one slide per heading. Returns the
    output path.

    If *board_sections* is given, only ``##``/``###`` blocks whose heading
    (after stripping any leading numbering, case-insensitive) matches an
    entry in *board_sections* are rendered as slides. The cover (and the
    optional H1 overview slide) are always rendered. When *board_sections*
    is ``None`` (default), all ``##``/``###`` blocks are rendered."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise SystemExit(
            "[md_to_pptx] 'python-pptx' package not found.\n"
            "Install: pip3 install python-pptx"
        )

    md_path = Path(md_path)
    if not md_path.exists():
        raise FileNotFoundError(f"Markdown file not found: {md_path}")

    out = path_guard.assert_writable(Path(output_path))
    path_guard.guard_output_dir(out.parent)

    if not date_str:
        date_str = datetime.now(timezone.utc).strftime("%Y-%m-%d")

    text = md_path.read_text(encoding="utf-8")
    blocks = _split_blocks(text)

    h1_block = next((b for b in blocks if b.level == 1), None)
    h1_title = h1_block.heading if h1_block else (title or md_path.stem.replace("_", " ").title())
    cover_title = title or h1_title

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W, H = prs.slide_width, prs.slide_height

    _slide_cover(prs, title=cover_title, case_id=case_id,
                  description="Campaign Forensics Report", date_str=date_str, W=W, H=H)

    # Preamble before the H1 (rare) — skip; preamble *after* H1 but before the
    # first ## becomes the "Case Overview" slide.
    h1_index = next((i for i, b in enumerate(blocks) if b.level == 1), None)
    if h1_index is not None and h1_index + 1 < len(blocks) and blocks[h1_index + 1].level == 0:
        _render_overview_slide(prs, blocks[h1_index + 1], h1_title, case_id, W, H)

    allowlist = (
        {_normalize_heading(s) for s in board_sections}
        if board_sections is not None else None
    )
    for block in blocks:
        if block.level in (2, 3):
            if allowlist is not None and _normalize_heading(block.heading) not in allowlist:
                continue
            _render_block(prs, block, case_id, W, H)

    prs.save(str(out))
    print(f"[md_to_pptx] PPTX saved: {out}")
    return out


# ── CLI ───────────────────────────────────────────────────────────────────

def _build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(description="Render a Markdown campaign report as a PPTX deck")
    p.add_argument("markdown", help="Path to input .md file")
    p.add_argument("--output", "-o", required=True, metavar="PPTX", help="Output PPTX path")
    p.add_argument("--case-id", metavar="ID", default="", help="Case ID shown on cover/header")
    p.add_argument("--title", metavar="TITLE", default="", help="Cover title (default: H1 of the document)")
    p.add_argument("--date", metavar="YYYY-MM-DD", default="", help="Report date (default: today UTC)")
    return p


if __name__ == "__main__":
    args = _build_parser().parse_args()
    convert(
        md_path=Path(args.markdown),
        output_path=Path(args.output),
        case_id=args.case_id,
        title=args.title,
        date_str=args.date,
    )
